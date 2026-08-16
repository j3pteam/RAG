#!/usr/bin/env python3
"""
exports.py — turn an assistant reply into a downloadable document.

Supports four formats:
    docx  → Word            (python-docx)
    pptx  → PowerPoint      (python-pptx)
    xlsx  → Excel           (openpyxl)
    pdf   → PDF             (reportlab)

Everything is rendered in the J3P brand palette. All builders take the raw
reply text (lightweight markdown) and return a BytesIO ready to stream.

Add to requirements.txt:
    python-docx
    python-pptx
    openpyxl
    reportlab
"""
import io
import os
import re
from datetime import datetime

# --- Brand palette -----------------------------------------------------------
NAVY = "27334A"
GOLD = "D2BC8D"
RUST = "9D432C"
PAPER = "FAF6F0"
CHARCOAL = "3F3F44"

# Exported documents carry NO advisor branding. They are the client's own
# letters, decks and memos, so the masthead is a logo drop-zone (or the
# client's own logo) and nothing identifies the tool that produced them.
BRAND_NAME = ""
DOC_SUBTITLE = ""

# Any of these appearing in generated content is scrubbed before rendering.
_BRAND_TOKENS = [
    r"J3P\s+Healthcare\s+Solutions",
    r"J3P\s+Health",
    r"J3P\s+Advisor",
    r"J3Personica",
    r"Residency\s+Select",
    r"\bJ3P\b",
]
_BRAND_RE = re.compile("|".join(_BRAND_TOKENS), re.IGNORECASE)


def scrub_brand(text: str) -> str:
    """Remove advisor branding from document content."""
    if not text:
        return text
    out_lines = []
    for line in str(text).split("\n"):
        if not _BRAND_RE.search(line):
            out_lines.append(line)          # untouched — never reformat it
            continue
        # A short line that is essentially just branding gets dropped
        residue = _BRAND_RE.sub("", line)
        if len(re.sub(r"[^A-Za-z0-9]", "", residue)) < 12:
            continue
        # Tidy only the lines we actually edited, so list markers and other
        # formatting elsewhere in the document are left alone.
        residue = re.sub(r"[ \t]{2,}", " ", residue)
        residue = re.sub(r"\s+([,.;:!?])", r"\1", residue)
        residue = re.sub(r"^([ \t]*)[·|]\s*", r"\1", residue)
        residue = re.sub(r"^([ \t]*(?:[-*+]|\d+[.)])\s+)[·|—–]\s*", r"\1", residue)
        out_lines.append(residue.rstrip())
    out = "\n".join(out_lines)
    out = re.sub(r"\n{3,}", "\n\n", out)
    return out.strip()

# --- Logo -------------------------------------------------------------------
# Embedded at the top of every generated file. Falls back to the wordmark text
# if the image can't be found, so exports never fail over a missing asset.
# Deliberately does NOT include the app's own full_logo.png: exported documents
# are usually the client's own letters and decks, so they get a replaceable
# placeholder unless a logo is explicitly configured for exports.
LOGO_CANDIDATES = [
    os.environ.get("BRAND_LOGO_FILE", ""),
    "export_logo.png",
]


def find_logo():
    """Return a filesystem path to a brand logo the user has configured.

    Returns None unless an explicit logo file is present. When None, exports
    fall back to a replaceable "insert logo" placeholder instead, so client
    documents don't ship with someone else's branding baked in.
    """
    here = os.path.dirname(os.path.abspath(__file__))
    for candidate in LOGO_CANDIDATES:
        if not candidate:
            continue
        for base in (here, os.getcwd()):
            path = candidate if os.path.isabs(candidate) else os.path.join(base, candidate)
            if os.path.isfile(path):
                return path
    return None


# --- Logo placeholder --------------------------------------------------------
# When no logo is configured, every export gets a clearly-marked drop zone in
# the same spot the logo would occupy. It's a real image, so in Word and
# PowerPoint the user can right-click it and choose "Change Picture" to swap in
# their own logo without touching layout.

PLACEHOLDER_ALT = "Insert your logo here — right-click and choose Change Picture"
_placeholder_cache = {}


def make_placeholder(width=520, height=180, on_dark=False):
    """Generate a dashed 'INSERT LOGO HERE' PNG. Returns a path, or None.

    on_dark selects a lighter label for placement over the navy title slide.
    """
    key = (width, height, on_dark)
    if key in _placeholder_cache and os.path.isfile(_placeholder_cache[key]):
        return _placeholder_cache[key]
    try:
        from PIL import Image, ImageDraw
        import tempfile

        img = Image.new("RGBA", (width, height), (255, 255, 255, 0))
        d = ImageDraw.Draw(img)
        gold = (210, 188, 141, 255)
        muted = (222, 220, 214, 255) if on_dark else (128, 128, 136, 255)

        # Dashed border
        dash, gap, inset = 14, 10, 4
        x0, y0, x1, y1 = inset, inset, width - inset, height - inset
        for x in range(x0, x1, dash + gap):
            d.line([(x, y0), (min(x + dash, x1), y0)], fill=gold, width=3)
            d.line([(x, y1), (min(x + dash, x1), y1)], fill=gold, width=3)
        for y in range(y0, y1, dash + gap):
            d.line([(x0, y), (x0, min(y + dash, y1))], fill=gold, width=3)
            d.line([(x1, y), (x1, min(y + dash, y1))], fill=gold, width=3)

        label = "INSERT LOGO HERE"
        try:
            from PIL import ImageFont
            font = None
            for candidate in (
                "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
                "/System/Library/Fonts/Supplemental/Arial Bold.ttf",
            ):
                if os.path.isfile(candidate):
                    font = ImageFont.truetype(candidate, 26)
                    break
            if font is None:
                font = ImageFont.load_default()
            box = d.textbbox((0, 0), label, font=font)
            d.text(((width - (box[2] - box[0])) / 2,
                    (height - (box[3] - box[1])) / 2 - 4),
                   label, fill=muted, font=font)
        except Exception:
            d.text((width / 2 - 60, height / 2 - 6), label, fill=muted)

        suffix = "_dark" if on_dark else ""
        # Neutral filename: it becomes the picture's name inside the document
        path = os.path.join(
            tempfile.gettempdir(), f"logo_placeholder_{width}x{height}{suffix}.png")
        img.save(path)
        _placeholder_cache[key] = path
        return path
    except Exception:
        return None


def logo_or_placeholder(on_dark=False):
    """Return (path, is_placeholder). Path may be None if neither is available."""
    real = find_logo()
    if real:
        return real, False
    return make_placeholder(on_dark=on_dark), True


def _set_alt_text(element, text):
    """Best-effort alt text so the drop zone is self-explaining."""
    try:
        element.set("descr", text)
        element.set("title", "Insert logo")
    except Exception:
        pass


# --- Meta-language stripping -------------------------------------------------
# The assistant closes its reply with a line about the SAVE button / the file
# downloading. That's useful in the chat window but must never appear inside
# the document itself, so it's removed before rendering.
_META_LINE_PATTERNS = [
    r"\bsave\s+button\b",
    r"\bclick\s+save\b",
    r"\bre-?download\b",
    r"\bthe file is (?:downloading|being (?:downloaded|generated|created))\b",
    r"\bdownloading (?:now|automatically)\b",
    r"\bwill download it\b",
    r"\bdownload(?:ed|ing)? (?:it )?(?:in|as) (?:powerpoint|word|excel|pdf)\b",
    r"\bbelow will (?:download|save|export)\b",
    r"\bsave\b[^.]{0,40}\bdownload\b",
    r"\b(?:download|save) (?:either|both|it) (?:or both )?in (?:your|any) (?:preferred )?format\b",
    r"\bin your preferred format\b",
    r"\b(?:two|both) documents below\b",
    r"\buse the save\b",
    r"\btap save\b",
    r"\b(?:deck|document|letter|file) is (?:written and ready|ready to download)\b",
    r"\bexported? (?:to|as) an? (?:word|powerpoint|excel|pdf)\b",
    r"\bis (?:attached|downloading) (?:now|below)\b",
]
_META_RE = re.compile("|".join(_META_LINE_PATTERNS), re.IGNORECASE)


def strip_meta(text: str) -> str:
    """Remove chat-only lines about saving/downloading from document content."""
    if not text:
        return text
    kept = []
    for line in str(text).replace("\r\n", "\n").split("\n"):
        stripped = line.strip()
        # Only drop short standalone lines — never gut a real paragraph that
        # happens to mention one of these words in passing.
        if stripped and len(stripped) < 320 and _META_RE.search(stripped):
            continue
        kept.append(line)
    out = "\n".join(kept)
    # Collapse blank runs left behind and trailing separators
    out = re.sub(r"\n{3,}", "\n\n", out)
    out = re.sub(r"(?:\n\s*[-*_]{3,}\s*)+$", "", out)
    return out.strip()



# ---------------------------------------------------------------------------
# Lightweight markdown parser
# ---------------------------------------------------------------------------

def parse_blocks(text: str) -> list:
    """Parse reply text into a flat list of blocks.

    Returns a list of dicts, each one of:
        {"type": "heading", "level": 1|2|3, "text": str}
        {"type": "bullet",  "text": str}
        {"type": "number",  "text": str, "num": int}
        {"type": "para",    "text": str}

    Inline **bold** markers are left in place — each renderer strips or
    converts them as appropriate for its format.
    """
    blocks = []
    if not text:
        return blocks

    # Drop fenced code blocks — they don't translate well to these formats
    text = re.sub(r"```[\s\S]*?```", "", text)

    lines = text.replace("\r\n", "\n").split("\n")
    para_buffer = []

    def flush_para():
        if para_buffer:
            joined = " ".join(l.strip() for l in para_buffer).strip()
            if joined:
                blocks.append({"type": "para", "text": joined})
            para_buffer.clear()

    for raw in lines:
        line = raw.rstrip()
        stripped = line.strip()

        if not stripped:
            flush_para()
            continue

        # Horizontal rule
        if re.fullmatch(r"[-*_]{3,}", stripped):
            flush_para()
            continue

        # Headings: # / ## / ###
        m = re.match(r"^(#{1,6})\s+(.*)$", stripped)
        if m:
            flush_para()
            level = min(len(m.group(1)), 3)
            blocks.append({"type": "heading", "level": level, "text": m.group(2).strip()})
            continue

        # A standalone all-bold line reads as a heading
        m = re.fullmatch(r"\*\*(.+?)\*\*:?", stripped)
        if m:
            flush_para()
            blocks.append({"type": "heading", "level": 2, "text": m.group(1).strip()})
            continue

        # Bulleted list
        m = re.match(r"^[-*+]\s+(.*)$", stripped)
        if m:
            flush_para()
            blocks.append({"type": "bullet", "text": m.group(1).strip()})
            continue

        # Numbered list
        m = re.match(r"^(\d+)[.)]\s+(.*)$", stripped)
        if m:
            flush_para()
            blocks.append({
                "type": "number",
                "num": int(m.group(1)),
                "text": m.group(2).strip(),
            })
            continue

        para_buffer.append(stripped)

    flush_para()
    return blocks


def strip_inline(s: str) -> str:
    """Remove inline markdown emphasis, links, and inline code markers."""
    s = re.sub(r"\*\*(.+?)\*\*", r"\1", s)
    s = re.sub(r"__(.+?)__", r"\1", s)
    s = re.sub(r"(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)", r"\1", s)
    s = re.sub(r"`([^`]+?)`", r"\1", s)
    s = re.sub(r"\[([^\]]+?)\]\([^)]+?\)", r"\1", s)
    return s.strip()


def split_bold_runs(s: str) -> list:
    """Split a string into [(text, is_bold), ...] on **bold** boundaries."""
    parts = []
    for chunk in re.split(r"(\*\*.+?\*\*)", s):
        if not chunk:
            continue
        if chunk.startswith("**") and chunk.endswith("**") and len(chunk) > 4:
            parts.append((strip_inline(chunk[2:-2]), True))
        else:
            parts.append((strip_inline(chunk), False))
    return [(t, b) for t, b in parts if t]


def split_documents(text: str) -> list:
    """Split a reply that contains several deliverables into separate documents.

    Returns a list of {"title": str, "body": str}. A reply with one document
    returns a single entry, so callers can treat the result uniformly.

    Boundaries, in order of preference:
      1. Top-level '# ' headings (the model is instructed to use these)
      2. Horizontal rules ('---') between substantial chunks
    A normal single document full of '## ' section headings is NOT split.
    """
    body = scrub_brand(strip_meta(text or ""))
    if not body.strip():
        return [{"title": "", "body": ""}]

    lines = body.replace("\r\n", "\n").split("\n")

    # --- 1. Split on level-1 headings ---
    h1_idx = [i for i, l in enumerate(lines)
              if re.match(r"^#\s+\S", l.strip()) and not l.strip().startswith("##")]
    MAX_DOCUMENTS = 4          # more than this means they're section headings
    MIN_DOC_CHARS = 400        # a real deliverable has substance

    if 2 <= len(h1_idx) <= MAX_DOCUMENTS:
        parts = []
        for n, start in enumerate(h1_idx):
            end = h1_idx[n + 1] if n + 1 < len(h1_idx) else len(lines)
            chunk = "\n".join(lines[start:end]).strip()
            title = re.sub(r"^#\s+", "", lines[start].strip())
            parts.append({"title": strip_inline(title), "body": chunk})
        parts = [p for p in parts if p["body"]]
        # Every part must be substantial; otherwise treat the whole reply as one
        if len(parts) >= 2 and all(len(p["body"]) >= MIN_DOC_CHARS for p in parts):
            return parts

    # NOTE: horizontal rules ('---') are deliberately NOT a boundary. Models
    # use them to separate sections *within* one document — between slides of a
    # deck, or between sections of a summary — so splitting on them turned a
    # single executive summary into one file per section. A deliverable
    # boundary requires an explicit '# ' title.

    return [{"title": derive_title(body), "body": body}]


# Words that suggest a document is meant to be slides rather than prose
_DECK_HINTS = re.compile(
    r"\b(deck|slide|slides|presentation|powerpoint|talk track|pitch)\b", re.I)
_SHEET_HINTS = re.compile(
    r"\b(spreadsheet|excel|budget|tracker|matrix|roster|workbook|line items?)\b", re.I)


def suggest_format(title: str, body: str = "") -> str:
    """Best-guess file format for a document, based on its title and shape."""
    head = f"{title or ''}\n{(body or '')[:600]}"
    if _SHEET_HINTS.search(head):
        return "xlsx"
    if _DECK_HINTS.search(head):
        return "pptx"

    # Slide-shaped content: several headings, each mostly short bullets
    blocks = parse_blocks(body or "")
    headings = [b for b in blocks if b["type"] == "heading"]
    bullets = [b for b in blocks if b["type"] == "bullet"]
    paras = [b for b in blocks if b["type"] == "para"]
    if len(headings) >= 3 and len(bullets) >= 6 and len(bullets) > len(paras) * 2:
        return "pptx"
    return "docx"


# Openings that are part of the letter, never its title
_SALUTATION_RE = re.compile(
    r"^\s*(dear\b|to whom\b|hello\b|hi\b|greetings\b|attn\b|attention:)", re.I)
# Subject lines make good titles
_SUBJECT_RE = re.compile(r"^\s*(?:re|subject)\s*:\s*(.+)$", re.I)


def derive_title(text: str, fallback: str = "") -> str:
    """Pick a document title, or return '' when there isn't a real one.

    A letter that opens with "Dear Dr. Harris" has no title. Inventing one from
    the first sentence produced headings like "Dear Dr." sitting above the
    actual salutation, so an empty string is returned instead and the renderers
    simply omit the title line.
    """
    body = (text or "").strip()
    if not body:
        return fallback

    # 1. A real markdown heading is the best title
    for block in parse_blocks(body):
        if block["type"] == "heading":
            heading = strip_inline(block["text"])
            if heading and not _SALUTATION_RE.match(heading):
                return heading[:90]

    # 2. A subject line ("Re: Application for ...") is the next best thing
    for line in body.split("\n")[:12]:
        m = _SUBJECT_RE.match(line.strip())
        if m:
            subject = strip_inline(m.group(1)).strip()
            if subject:
                return subject[:90]

    # 3. No heading and no subject: don't fabricate one from the prose
    return fallback


def drop_leading_title(text: str, title: str) -> str:
    """Remove the opening heading when it's already the document's title.

    The renderers print the title in the masthead, so leaving the same heading
    at the top of the body prints it twice.
    """
    if not text or not title:
        return text
    lines = str(text).replace("\r\n", "\n").split("\n")
    norm = lambda v: re.sub(r"[^a-z0-9]", "", strip_inline(v or "").lower())
    target = norm(title)
    if not target:
        return text

    for i, line in enumerate(lines):
        stripped = line.strip()
        if not stripped:
            continue                      # skip blank lines above the heading
        m = re.match(r"^#{1,6}\s+(.*)$", stripped)
        if not m:
            m2 = re.fullmatch(r"\*\*(.+?)\*\*:?", stripped)
            if not m2:
                return text               # body doesn't start with a heading
            heading = m2.group(1)
        else:
            heading = m.group(1)
        if norm(heading) == target:
            rest = lines[i + 1:]
            # Also drop a rule or blank lines immediately after it
            while rest and (not rest[0].strip()
                            or re.fullmatch(r"\s*[-*_]{3,}\s*", rest[0])):
                rest.pop(0)
            return "\n".join(rest).strip()
        return text                       # different heading — leave it alone
    return text


def safe_filename(title: str, ext: str) -> str:
    slug = re.sub(r"[^A-Za-z0-9]+", "_", strip_inline(title or "")).strip("_")[:50]
    if not slug:
        slug = "document"
    return f"{slug}_{datetime.now().strftime('%Y%m%d_%H%M')}.{ext}"


# ---------------------------------------------------------------------------
# Word (.docx)
# ---------------------------------------------------------------------------

def build_docx(text: str, title: str = None) -> io.BytesIO:
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    title = title or derive_title(text)
    doc = Document()

    # Base style
    normal = doc.styles["Normal"]
    normal.font.name = "Calibri"
    normal.font.size = Pt(11)
    normal.paragraph_format.space_after = Pt(8)
    normal.paragraph_format.line_spacing = 1.15

    for section in doc.sections:
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1)
        section.right_margin = Inches(1)

    # Masthead — logo image when available, wordmark text otherwise
    logo_path, is_placeholder = logo_or_placeholder()
    brand = doc.add_paragraph()
    brand.paragraph_format.space_after = Pt(2)
    if logo_path:
        try:
            brand.add_run().add_picture(
                logo_path, height=Inches(0.55 if is_placeholder else 0.42))
            if is_placeholder and doc.inline_shapes:
                _set_alt_text(doc.inline_shapes[-1]._inline.docPr, PLACEHOLDER_ALT)
        except Exception:
            logo_path = None
    if not logo_path:
        # No logo and no placeholder available — leave the masthead empty
        # rather than stamping the advisor's name on a client document.
        brand.paragraph_format.space_after = Pt(0)

    if title:
        head = doc.add_paragraph()
        head_run = head.add_run(title)
        head_run.bold = True
        head_run.font.size = Pt(20)
        head_run.font.color.rgb = RGBColor.from_string(NAVY)
        head.paragraph_format.space_after = Pt(4)

    meta = doc.add_paragraph()
    meta_run = meta.add_run(
        datetime.now().strftime('%B %d, %Y')
    )
    meta_run.font.size = Pt(9)
    meta_run.font.color.rgb = RGBColor.from_string(CHARCOAL)
    meta.paragraph_format.space_after = Pt(14)

    # Body
    for block in parse_blocks(text):
        btype = block["type"]

        if btype == "heading":
            level = block["level"]
            p = doc.add_paragraph()
            run = p.add_run(strip_inline(block["text"]))
            run.bold = True
            run.font.size = Pt({1: 15, 2: 13, 3: 11.5}[level])
            run.font.color.rgb = RGBColor.from_string(NAVY)
            p.paragraph_format.space_before = Pt(14)
            p.paragraph_format.space_after = Pt(5)

        elif btype == "bullet":
            p = doc.add_paragraph(style="List Bullet")
            for chunk, is_bold in split_bold_runs(block["text"]):
                run = p.add_run(chunk)
                run.bold = is_bold
            p.paragraph_format.space_after = Pt(4)

        elif btype == "number":
            p = doc.add_paragraph(style="List Number")
            for chunk, is_bold in split_bold_runs(block["text"]):
                run = p.add_run(chunk)
                run.bold = is_bold
            p.paragraph_format.space_after = Pt(4)

        else:
            p = doc.add_paragraph()
            for chunk, is_bold in split_bold_runs(block["text"]):
                run = p.add_run(chunk)
                run.bold = is_bold

    # No footer disclaimer: exported documents are the client's own work
    # product and carry nothing from the advisor tool.

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


# ---------------------------------------------------------------------------
# PowerPoint (.pptx)
# ---------------------------------------------------------------------------

def build_pptx(text: str, title: str = None) -> io.BytesIO:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    title = title or derive_title(text)
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)

    navy = RGBColor.from_string(NAVY)
    gold = RGBColor.from_string(GOLD)
    paper = RGBColor.from_string(PAPER)
    MAX_BULLETS = 6

    def set_bg(slide, color):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

    # --- Title slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank
    set_bg(slide, navy)

    logo_path, is_placeholder = logo_or_placeholder(on_dark=True)
    # Content slides use the paper background, so they need the dark-label version
    content_logo_path, _ = logo_or_placeholder(on_dark=False)

    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.5), Inches(11.5), Inches(2.2))
    tf = box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    if logo_path:
        try:
            pic = slide.shapes.add_picture(
                logo_path, Inches(0.9), Inches(1.05),
                height=Inches(0.8 if is_placeholder else 0.95))
            if is_placeholder:
                pic.name = "INSERT LOGO HERE"
                _set_alt_text(pic._element.nvPicPr.cNvPr, PLACEHOLDER_ALT)
            p.text = ""                    # image replaces the wordmark text
        except Exception:
            p.text = ""
    else:
        p.text = ""
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = gold

    if title:
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(38)
        p.font.bold = True
        p.font.color.rgb = paper
        p.space_before = Pt(10)

    p = tf.add_paragraph()
    p.text = datetime.now().strftime('%B %d, %Y')
    p.font.size = Pt(13)
    p.font.color.rgb = gold
    p.space_before = Pt(12)

    # --- Group blocks into sections keyed by heading ---
    sections = []
    current = {"title": title, "items": []}
    for block in parse_blocks(text):
        if block["type"] == "heading":
            if current["items"]:
                sections.append(current)
            current = {"title": strip_inline(block["text"]), "items": []}
        else:
            current["items"].append(strip_inline(block["text"]))
    if current["items"]:
        sections.append(current)
    if not sections:
        sections = [{"title": title, "items": [strip_inline(text or "")]}]

    # --- Content slides — split any section over MAX_BULLETS items ---
    for section in sections:
        items = [i for i in section["items"] if i]
        pages = [items[i:i + MAX_BULLETS] for i in range(0, len(items), MAX_BULLETS)] or [[]]
        for idx, page_items in enumerate(pages):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_bg(slide, paper)

            # Gold accent rule under the heading
            head_text = section["title"] + ("" if idx == 0 else f" (cont. {idx + 1})")
            head_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(11.7), Inches(1.0))
            htf = head_box.text_frame
            htf.word_wrap = True
            hp = htf.paragraphs[0]
            hp.text = head_text
            hp.font.size = Pt(28)
            hp.font.bold = True
            hp.font.color.rgb = navy

            line = slide.shapes.add_shape(1, Inches(0.85), Inches(1.62), Inches(2.0), Inches(0.045))
            line.fill.solid()
            line.fill.fore_color.rgb = gold
            line.line.fill.background()
            line.shadow.inherit = False

            body_box = slide.shapes.add_textbox(Inches(0.85), Inches(2.0), Inches(11.6), Inches(4.7))
            btf = body_box.text_frame
            btf.word_wrap = True
            for j, item in enumerate(page_items):
                bp = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
                bp.text = "•  " + item
                bp.font.size = Pt(17)
                bp.font.color.rgb = RGBColor.from_string(CHARCOAL)
                bp.space_after = Pt(13)

            # Footer — small logo, falling back to the wordmark
            placed_logo = False
            if content_logo_path or logo_path:
                try:
                    fpic = slide.shapes.add_picture(
                        content_logo_path or logo_path, Inches(0.85), Inches(6.72),
                        height=Inches(0.42 if is_placeholder else 0.34))
                    if is_placeholder:
                        fpic.name = "INSERT LOGO HERE"
                        _set_alt_text(fpic._element.nvPicPr.cNvPr, PLACEHOLDER_ALT)
                    placed_logo = True
                except Exception:
                    placed_logo = False
            # No text footer — client decks carry no advisor branding

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ---------------------------------------------------------------------------
# Excel (.xlsx)
# ---------------------------------------------------------------------------

def build_xlsx(text: str, title: str = None) -> io.BytesIO:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    title = title or derive_title(text)
    wb = Workbook()
    ws = wb.active
    ws.title = "Response"

    navy_fill = PatternFill("solid", fgColor=NAVY)
    gold_fill = PatternFill("solid", fgColor=GOLD)
    paper_fill = PatternFill("solid", fgColor=PAPER)
    wrap = Alignment(vertical="top", wrap_text=True)
    bottom_rule = Border(bottom=Side(style="thin", color=GOLD))

    # Masthead
    ws["A1"] = ""
    ws["A1"].fill = navy_fill
    ws["B1"].fill = navy_fill
    ws["C1"].fill = navy_fill

    logo_path, is_placeholder = logo_or_placeholder()
    if logo_path:
        try:
            from openpyxl.drawing.image import Image as XLImage
            img = XLImage(logo_path)
            ratio = (img.height or 1) / float(img.width or 1)
            img.width = 120
            img.height = max(int(120 * ratio), 1)
            ws.row_dimensions[1].height = max(img.height * 0.78, 18)
            ws["A1"] = ""                       # logo replaces the text wordmark
            ws.add_image(img, "A1")
        except Exception:
            pass                                 # keep the text masthead

    ws["A2"] = title or ""
    ws["A2"].font = Font(bold=True, size=16, color=NAVY)
    ws.merge_cells("A2:C2")
    ws["A2"].alignment = Alignment(vertical="center", wrap_text=True)
    ws.row_dimensions[2].height = 30

    ws["A3"] = datetime.now().strftime('%B %d, %Y')
    ws["A3"].font = Font(size=9, italic=True, color=CHARCOAL)
    ws.merge_cells("A3:C3")

    # Column headers
    header_row = 5
    for col, label in enumerate(["Section", "Type", "Content"], start=1):
        cell = ws.cell(row=header_row, column=col, value=label)
        cell.font = Font(bold=True, size=10, color=GOLD)
        cell.fill = navy_fill
        cell.alignment = Alignment(horizontal="left", vertical="center")

    TYPE_LABEL = {"heading": "Heading", "bullet": "Bullet",
                  "number": "Numbered", "para": "Paragraph"}

    row = header_row + 1
    current_section = "—"
    for block in parse_blocks(text):
        content = strip_inline(block["text"])
        if not content:
            continue
        if block["type"] == "heading":
            current_section = content
            ws.cell(row=row, column=1, value=content).font = Font(bold=True, size=11, color=NAVY)
            ws.cell(row=row, column=2, value="Heading").font = Font(size=9, color=CHARCOAL)
            ws.cell(row=row, column=3, value=content).font = Font(bold=True, size=11, color=NAVY)
            for col in range(1, 4):
                ws.cell(row=row, column=col).fill = paper_fill
                ws.cell(row=row, column=col).border = bottom_rule
                ws.cell(row=row, column=col).alignment = wrap
        else:
            prefix = f"{block['num']}. " if block["type"] == "number" else ""
            ws.cell(row=row, column=1, value=current_section).font = Font(size=9, color=CHARCOAL)
            ws.cell(row=row, column=2, value=TYPE_LABEL.get(block["type"], "Text")).font = Font(size=9, color=CHARCOAL)
            ws.cell(row=row, column=3, value=prefix + content).font = Font(size=10)
            for col in range(1, 4):
                ws.cell(row=row, column=col).alignment = wrap
        row += 1

    ws.column_dimensions["A"].width = 28
    ws.column_dimensions["B"].width = 13
    ws.column_dimensions["C"].width = 95
    ws.freeze_panes = f"A{header_row + 1}"

    # Plain-text sheet for anyone who just wants to copy the whole thing
    ws2 = wb.create_sheet("Full text")
    ws2["A1"] = "Full response"
    ws2["A1"].font = Font(bold=True, size=11, color=NAVY)
    ws2["A1"].fill = gold_fill
    ws2["A2"] = strip_inline(text or "")
    ws2["A2"].alignment = wrap
    ws2.column_dimensions["A"].width = 120
    ws2.row_dimensions[2].height = 400

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

def build_pdf(text: str, title: str = None) -> io.BytesIO:
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, HRFlowable, ListFlowable, ListItem,
    )

    title = title or derive_title(text)
    navy = colors.HexColor("#" + NAVY)
    gold = colors.HexColor("#" + GOLD)
    charcoal = colors.HexColor("#" + CHARCOAL)

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=LETTER,
        leftMargin=0.9 * inch, rightMargin=0.9 * inch,
        topMargin=0.85 * inch, bottomMargin=0.85 * inch,
        title=title,
    )

    base = getSampleStyleSheet()
    styles = {
        "brand": ParagraphStyle("brand", parent=base["Normal"], fontName="Helvetica-Bold",
                                fontSize=8.5, textColor=gold, spaceAfter=3, leading=11),
        "title": ParagraphStyle("titleS", parent=base["Normal"], fontName="Helvetica-Bold",
                                fontSize=20, textColor=navy, leading=25, spaceAfter=4),
        "meta": ParagraphStyle("meta", parent=base["Normal"], fontName="Helvetica-Oblique",
                               fontSize=8.5, textColor=charcoal, spaceAfter=12),
        "h1": ParagraphStyle("h1S", parent=base["Normal"], fontName="Helvetica-Bold",
                             fontSize=14.5, textColor=navy, spaceBefore=14, spaceAfter=5, leading=18),
        "h2": ParagraphStyle("h2S", parent=base["Normal"], fontName="Helvetica-Bold",
                             fontSize=12.5, textColor=navy, spaceBefore=12, spaceAfter=4, leading=16),
        "h3": ParagraphStyle("h3S", parent=base["Normal"], fontName="Helvetica-Bold",
                             fontSize=11, textColor=navy, spaceBefore=10, spaceAfter=3, leading=14),
        "body": ParagraphStyle("bodyS", parent=base["Normal"], fontName="Helvetica",
                               fontSize=10.5, textColor=colors.HexColor("#2B2B2B"),
                               leading=15.5, spaceAfter=8),
        "item": ParagraphStyle("itemS", parent=base["Normal"], fontName="Helvetica",
                               fontSize=10.5, textColor=colors.HexColor("#2B2B2B"),
                               leading=15, spaceAfter=4),
        "foot": ParagraphStyle("footS", parent=base["Normal"], fontName="Helvetica-Oblique",
                               fontSize=7.5, textColor=charcoal, alignment=TA_CENTER),
    }

    def inline_html(s: str) -> str:
        """Convert **bold** to <b> and escape XML-unsafe characters."""
        s = s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        s = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", s)
        s = re.sub(r"`([^`]+?)`", r"\1", s)
        s = re.sub(r"\[([^\]]+?)\]\([^)]+?\)", r"\1", s)
        return s

    from reportlab.platypus import Image as RLImage
    logo_path, is_placeholder = logo_or_placeholder()
    masthead = None
    if logo_path:
        try:
            from reportlab.lib.utils import ImageReader
            iw, ih = ImageReader(logo_path).getSize()
            h = (0.55 if is_placeholder else 0.42) * inch
            masthead = RLImage(logo_path, width=h * (iw / float(ih)), height=h)
            masthead.hAlign = "LEFT"
        except Exception:
            masthead = None

    story = [
        masthead if masthead else Spacer(1, 0),
        Spacer(1, 6) if masthead else Spacer(1, 0),
        Paragraph(inline_html(title), styles["title"]) if title else Spacer(1, 0),
        HRFlowable(width="100%", thickness=1.5, color=gold,
                   spaceBefore=2, spaceAfter=6),
        Paragraph(datetime.now().strftime('%B %d, %Y'),
                  styles["meta"]),
    ]

    # Consecutive list items get grouped into a single ListFlowable
    pending = []
    pending_kind = None

    def flush_list():
        nonlocal pending, pending_kind
        if not pending:
            return
        items = [ListItem(Paragraph(inline_html(t), styles["item"]), leftIndent=16)
                 for t in pending]
        if pending_kind == "bullet":
            story.append(ListFlowable(
                items, bulletType="bullet", start="\u2022",
                bulletColor=gold, bulletFontSize=11,
                bulletFontName="Helvetica", leftIndent=20,
            ))
        else:
            story.append(ListFlowable(
                items, bulletType="1", start=1, bulletFormat="%s.",
                bulletColor=navy, bulletFontSize=10.5,
                bulletFontName="Helvetica-Bold", leftIndent=20,
            ))
        story.append(Spacer(1, 6))
        pending = []
        pending_kind = None

    for block in parse_blocks(text):
        btype = block["type"]
        if btype in ("bullet", "number"):
            if pending_kind and pending_kind != btype:
                flush_list()
            pending_kind = btype
            pending.append(block["text"])
            continue
        flush_list()
        if btype == "heading":
            story.append(Paragraph(inline_html(block["text"]),
                                   styles[f"h{block['level']}"]))
        else:
            story.append(Paragraph(inline_html(block["text"]), styles["body"]))
    flush_list()


    doc.build(story)
    buf.seek(0)
    return buf


# ---------------------------------------------------------------------------
# Dispatch
# ---------------------------------------------------------------------------

BUILDERS = {
    "docx": build_docx,
    "pptx": build_pptx,
    "xlsx": build_xlsx,
    "pdf": build_pdf,
}

MIME_TYPES = {
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pdf": "application/pdf",
}


def build(fmt: str, text: str, title: str = None):
    """Return (BytesIO, filename, mimetype) for the requested format."""
    fmt = (fmt or "").lower().strip()
    if fmt not in BUILDERS:
        raise ValueError(f"Unsupported format: {fmt}")
    # Chat-only instructions and advisor branding must never reach the document
    text = scrub_brand(strip_meta(text))
    resolved_title = (title or "").strip() or derive_title(text)
    body = drop_leading_title(text, resolved_title)
    buffer = BUILDERS[fmt](body, resolved_title)
    return buffer, safe_filename(resolved_title, fmt), MIME_TYPES[fmt]
