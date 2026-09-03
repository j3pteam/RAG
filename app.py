#!/usr/bin/env python3
"""
J3P Persona Bot — with RAG knowledge base and admin panel.

NEW since persona template:
  - Document upload / chunking / embedding pipeline
  - Semantic search retrieves relevant chunks before each response
  - Admin page (password-protected) for uploads, doc management, feedback review
  - Feedback persisted to Postgres instead of just logs

NEW environment variables:
  DATABASE_URL          Auto-set by Railway when you add a Postgres plugin
  OPENAI_API_KEY        For embeddings (~$0.02/1M tokens, very cheap)
  ADMIN_PASSWORD        Password for /admin page

All other env vars from the persona template still apply.
"""
import os
import tempfile
import re
import threading
from datetime import datetime
from pathlib import Path
from functools import wraps
from flask import (
    Flask, request, jsonify, session, render_template_string,
    send_from_directory, redirect, url_for, flash,
)
import anthropic

import database as db
import embeddings as emb
import paywall
import exports


# Ensure paywall schema exists (no-op if DB unavailable)
try:
    paywall.ensure_schema()
except Exception as _e:
    print(f"[paywall] Schema init warning: {_e}")


# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

def load_system_prompt():
    inline = os.environ.get("PERSONA_SYSTEM_PROMPT")
    if inline:
        return inline
    filename = os.environ.get("PERSONA_SYSTEM_PROMPT_FILE")
    if filename and Path(filename).exists():
        return Path(filename).read_text(encoding="utf-8")
    try:
        from system_prompt import SYSTEM_PROMPT
        return SYSTEM_PROMPT
    except ImportError:
        pass
    return "You are a helpful assistant."


# Per-file upload ceiling. Slide decks with embedded media routinely run past
# 25 MB, so the default is 100 MB and it's tunable without a code change.
# Bump this whenever the file changes so it's obvious which build is live.
# Visible at /health and in the admin header.
APP_VERSION = "2026-09-03-c"
APP_BUILD_NOTES = "admin field for the name shown with the photo"

MAX_UPLOAD_MB = int(os.environ.get("MAX_UPLOAD_MB", "100"))
MAX_UPLOAD_BYTES = MAX_UPLOAD_MB * 1024 * 1024

# Anthropic caps a single image at roughly 5 MB once base64-encoded, so images
# get their own tighter limit than documents.
MAX_IMAGE_MB = int(os.environ.get("MAX_IMAGE_MB", "5"))
MAX_IMAGE_BYTES = MAX_IMAGE_MB * 1024 * 1024

CONFIG = {
    "persona_name": os.environ.get("PERSONA_NAME", "J3P Advisor"),
    "opening": os.environ.get(
        "PERSONA_OPENING",
        "Hello, welcome to your session with the J3P Advisor.",
    ),
    "placeholder": os.environ.get("PERSONA_PLACEHOLDER", "How can I help you?"),
    "system_prompt": load_system_prompt(),

    "logo_url": os.environ.get("BRAND_LOGO_URL", "/full_logo.png"),
    "favicon_url": os.environ.get("BRAND_FAVICON_URL", "/monogram.jpg"),
    # Avatar shown beside advisor replies. Set ADVISOR_AVATAR_URL="" to hide it.
    "avatar_url": os.environ.get("ADVISOR_AVATAR_URL", "/advisor_avatar.jpg"),
    # Looping clip for the avatar's resting state. Empty string = still photo.
    "avatar_loop_url": os.environ.get("ADVISOR_LOOP_URL", "/advisor_idle.mp4"),
    "talking_avatar": os.environ.get("TALKING_AVATAR", "off").lower(),
    "navy": os.environ.get("BRAND_NAVY", "#27334A"),
    "gold": os.environ.get("BRAND_GOLD", "#D2BC8D"),
    "paper": os.environ.get("BRAND_PAPER", "#FAF6F0"),

    "footer_disclaimer": os.environ.get(
        "FOOTER_DISCLAIMER",
        "For informational purposes only. Not official advice.",
    ),
    "footer_cta_text": os.environ.get(
        "FOOTER_CTA_TEXT",
        "To schedule time with a J3P Advisor, please",
    ),
    "contact_email": os.environ.get("CONTACT_EMAIL", "clientservices@j3p.health"),
    "max_upload_mb": MAX_UPLOAD_MB,
    "max_image_mb": MAX_IMAGE_MB,
    "footer_ai_note": os.environ.get(
        "FOOTER_AI_NOTE",
        "The J3P Advisor is AI and can make mistakes. Please double-check responses.",
    ),
    "footer_cta_label": os.environ.get(
        "FOOTER_CTA_LABEL",
        "Schedule Time With a J3P Advisor",
    ),
    "footer_cta_url": os.environ.get(
        "FOOTER_CTA_URL",
        "https://app.acuityscheduling.com/catalog.php"
        "?owner=29987697&action=addCart&clear=1&id=2262965",
    ),

    "model": os.environ.get("ANTHROPIC_MODEL", "claude-sonnet-4-6"),
    # 1024 truncated two-page documents mid-sentence ("...built with this
    # faculty, not *"). 4096 covers a cover letter plus a strategic plan.
    "max_tokens": int(os.environ.get("MAX_TOKENS", "4096")),
    "rag_top_k": int(os.environ.get("RAG_TOP_K", "4")),
    "rag_min_similarity": float(os.environ.get("RAG_MIN_SIMILARITY", "0.3")),
    "admin_password": os.environ.get("ADMIN_PASSWORD", ""),
}


# ---------------------------------------------------------------------------
# App setup — initialize DB schema on startup
# ---------------------------------------------------------------------------

app = Flask(__name__)
app.secret_key = os.environ.get("FLASK_SECRET_KEY", os.urandom(24).hex())
# The whole request has to fit several files plus multipart overhead
app.config["MAX_CONTENT_LENGTH"] = int(MAX_UPLOAD_BYTES * 3.5)
client = anthropic.Anthropic()

# Initialize DB schema once at startup (idempotent)
try:
    if db.is_enabled():
        db.init_schema()
        app.logger.info("Database schema initialized")
    else:
        app.logger.warning("Database not configured — RAG and feedback persistence disabled")
except Exception as e:
    app.logger.error(f"DB init failed: {e}")


# ---------------------------------------------------------------------------
# Auth helper for admin routes
# ---------------------------------------------------------------------------

# ---------------------------------------------------------------------------
# Runtime settings (toggles the admin can flip without a redeploy)
# ---------------------------------------------------------------------------
# Kept in its own tiny table so database.py is untouched. If Postgres isn't
# available the values live in a JSON file instead — that survives restarts but
# not redeploys, which is fine for a display preference.

import json as _json

_SETTINGS_DEFAULTS = {
    "show_scheduling_button": True,
    # Seeded from REQUIRE_LOGIN, then owned by the admin panel
    "require_login": os.environ.get("REQUIRE_LOGIN", "off").lower() in ("on", "1", "true"),
    # Auto-approve eligible feedback as lessons on a schedule
    "auto_learning": False,
    # Show the advisor photo beside replies
    "show_avatar": True,
    # Let participants add their own documents and writing
    "allow_materials": True,
    # Name shown with the default advisor photo. Blank uses the app name.
    "avatar_name": "",
}
_settings_cache = None
_SETTINGS_FILE = os.path.join(tempfile.gettempdir(), "j3p_settings.json")


def _settings_db_conn():
    """Connection for the settings table, or None when Postgres isn't set up."""
    url = os.environ.get("DATABASE_URL")
    if not url:
        return None
    try:
        import psycopg
        return psycopg.connect(url)
    except Exception as e:
        app.logger.warning(f"[settings] Postgres unavailable: {e}")
        return None


def _settings_ensure_table(conn):
    with conn.cursor() as cur:
        cur.execute("""
            CREATE TABLE IF NOT EXISTS app_settings (
                key         TEXT PRIMARY KEY,
                value       TEXT NOT NULL,
                updated_at  TIMESTAMPTZ NOT NULL DEFAULT NOW()
            )
        """)
    conn.commit()


def load_settings(force: bool = False) -> dict:
    """Current settings, defaults filled in for anything unset."""
    global _settings_cache
    if _settings_cache is not None and not force:
        return _settings_cache

    values = dict(_SETTINGS_DEFAULTS)
    conn = _settings_db_conn()
    if conn:
        try:
            _settings_ensure_table(conn)
            with conn.cursor() as cur:
                cur.execute("SELECT key, value FROM app_settings")
                for key, raw in cur.fetchall():
                    if key not in values:
                        continue
                    # Booleans were the only setting type originally; text
                    # values (like the name shown with the photo) need to
                    # survive the round trip intact.
                    if isinstance(values[key], bool):
                        values[key] = (raw == "true")
                    else:
                        values[key] = raw or ""
        except Exception as e:
            app.logger.error(f"[settings] read failed: {e}")
        finally:
            conn.close()
    else:
        try:
            if os.path.isfile(_SETTINGS_FILE):
                with open(_SETTINGS_FILE, encoding="utf-8") as fh:
                    stored = _json.load(fh)
                for key in values:
                    if key not in stored:
                        continue
                    if isinstance(values[key], bool):
                        values[key] = bool(stored[key])
                    else:
                        values[key] = stored[key] or ""
        except Exception as e:
            app.logger.error(f"[settings] file read failed: {e}")

    _settings_cache = values
    return values


def save_setting(key, value) -> bool:
    """Persist one setting — boolean or text. Returns True when written."""
    global _settings_cache
    if key not in _SETTINGS_DEFAULTS:
        return False

    written = False
    conn = _settings_db_conn()
    if conn:
        try:
            _settings_ensure_table(conn)
            with conn.cursor() as cur:
                cur.execute("""
                    INSERT INTO app_settings (key, value, updated_at)
                    VALUES (%s, %s, NOW())
                    ON CONFLICT (key)
                    DO UPDATE SET value = EXCLUDED.value, updated_at = NOW()
                """, (key, ("true" if value else "false")
                        if isinstance(value, bool) else str(value)))
            conn.commit()
            written = True
        except Exception as e:
            app.logger.error(f"[settings] write failed: {e}")
        finally:
            conn.close()
    else:
        try:
            current = load_settings(force=True)
            # Text settings must keep their text: bool() turned a name like
            # "Alan Friedman" into True on any deployment without Postgres.
            current[key] = (bool(value)
                            if isinstance(_SETTINGS_DEFAULTS[key], bool)
                            else str(value))
            with open(_SETTINGS_FILE, "w", encoding="utf-8") as fh:
                _json.dump(current, fh)
            written = True
        except Exception as e:
            app.logger.error(f"[settings] file write failed: {e}")

    _settings_cache = None      # force a reload on next read
    return written


# ---------------------------------------------------------------------------
# Approximate location for the conversation log
# ---------------------------------------------------------------------------
# Resolves the caller's IP to city / region / country and stores only that.
# The IP itself is never written down. Lookups run on a background thread so
# they never slow a reply, and results are cached per IP for the process life.
# Set GEO_LOOKUP=off to disable entirely.

GEO_LOOKUP_ENABLED = os.environ.get("GEO_LOOKUP", "on").lower() not in ("off", "0", "false")
_geo_cache = {}
_geo_lock = threading.Lock()


def client_ip() -> str:
    """Caller's IP, accounting for Railway's proxy."""
    forwarded = request.headers.get("X-Forwarded-For", "")
    if forwarded:
        return forwarded.split(",")[0].strip()
    return request.headers.get("X-Real-IP") or request.remote_addr or ""


def _is_private_ip(ip: str) -> bool:
    try:
        import ipaddress
        addr = ipaddress.ip_address(ip)
        return addr.is_private or addr.is_loopback or addr.is_reserved
    except Exception:
        return True


def resolve_location(ip: str):
    """Return {'city','region','country'} for an IP, or None."""
    if not ip or _is_private_ip(ip):
        return None
    with _geo_lock:
        if ip in _geo_cache:
            return _geo_cache[ip]
    result = None
    try:
        import urllib.request as _url
        import json as _j
        req = _url.Request(
            f"http://ip-api.com/json/{ip}?fields=status,city,regionName,country",
            headers={"User-Agent": "j3p-advisor"},
        )
        with _url.urlopen(req, timeout=4) as resp:
            data = _j.loads(resp.read().decode("utf-8"))
        if data.get("status") == "success":
            result = {
                "city": data.get("city") or "",
                "region": data.get("regionName") or "",
                "country": data.get("country") or "",
            }
    except Exception as e:
        app.logger.warning(f"[geo] lookup failed: {e}")
    with _geo_lock:
        _geo_cache[ip] = result
    return result


def _ack_ensure_table(conn):
    with conn.cursor() as cur:
        cur.execute("""
            CREATE TABLE IF NOT EXISTS interaction_ack (
                interaction_id  BIGINT PRIMARY KEY,
                acknowledged_at TIMESTAMPTZ,
                ack_version     TEXT
            )
        """)
    conn.commit()


def record_acknowledgement(interaction_id: int):
    """Note that this session had accepted the release, against one interaction."""
    when = session.get("release_ack_at")
    if not (interaction_id and when):
        return
    conn = _settings_db_conn()
    if not conn:
        return
    try:
        _ack_ensure_table(conn)
        with conn.cursor() as cur:
            cur.execute("""
                INSERT INTO interaction_ack (interaction_id, acknowledged_at, ack_version)
                VALUES (%s, %s, %s)
                ON CONFLICT (interaction_id) DO NOTHING
            """, (int(interaction_id), when, session.get("release_ack_version", "")))
        conn.commit()
    except Exception as e:
        app.logger.error(f"[ack] store failed: {e}")
    finally:
        conn.close()


def acknowledgements_for(interaction_ids):
    """Map interaction_id -> acknowledgement timestamp for the admin table."""
    out = {}
    ids = [int(i) for i in interaction_ids if i]
    if not ids:
        return out
    conn = _settings_db_conn()
    if not conn:
        return out
    try:
        _ack_ensure_table(conn)
        with conn.cursor() as cur:
            cur.execute("""
                SELECT interaction_id, acknowledged_at FROM interaction_ack
                WHERE interaction_id = ANY(%s)
            """, (ids,))
            for iid, when in cur.fetchall():
                if when is None:
                    continue
                try:
                    out[iid] = when.strftime("%m/%d %H:%M")
                except AttributeError:
                    # Some drivers hand back an ISO string rather than a datetime
                    try:
                        out[iid] = datetime.fromisoformat(
                            str(when).replace("Z", "")).strftime("%m/%d %H:%M")
                    except Exception:
                        out[iid] = str(when)[:16]
    except Exception as e:
        app.logger.error(f"[ack] read failed: {e}")
    finally:
        conn.close()
    return out


def _geo_ensure_table(conn):
    with conn.cursor() as cur:
        cur.execute("""
            CREATE TABLE IF NOT EXISTS interaction_geo (
                interaction_id BIGINT PRIMARY KEY,
                city           TEXT,
                region         TEXT,
                country        TEXT,
                captured_at    TIMESTAMPTZ NOT NULL DEFAULT NOW()
            )
        """)
    conn.commit()


def record_location(interaction_id: int, ip: str):
    """Look up and store the location for one logged interaction."""
    if not (GEO_LOOKUP_ENABLED and interaction_id and ip):
        return
    loc = resolve_location(ip)
    if not loc:
        return
    conn = _settings_db_conn()
    if not conn:
        return
    try:
        _geo_ensure_table(conn)
        with conn.cursor() as cur:
            cur.execute("""
                INSERT INTO interaction_geo (interaction_id, city, region, country)
                VALUES (%s, %s, %s, %s)
                ON CONFLICT (interaction_id) DO NOTHING
            """, (int(interaction_id), loc["city"], loc["region"], loc["country"]))
        conn.commit()
    except Exception as e:
        app.logger.error(f"[geo] store failed: {e}")
    finally:
        conn.close()


def record_location_async(interaction_id: int, ip: str):
    """Fire the lookup in the background so the reply isn't delayed."""
    if not (GEO_LOOKUP_ENABLED and interaction_id and ip):
        return
    threading.Thread(
        target=record_location, args=(interaction_id, ip), daemon=True
    ).start()


def locations_for(interaction_ids):
    """Map of interaction_id -> 'City, Region, Country' for the admin table."""
    out = {}
    ids = [int(i) for i in interaction_ids if i]
    if not ids:
        return out
    conn = _settings_db_conn()
    if not conn:
        return out
    try:
        _geo_ensure_table(conn)
        with conn.cursor() as cur:
            cur.execute("""
                SELECT interaction_id, city, region, country
                FROM interaction_geo WHERE interaction_id = ANY(%s)
            """, (ids,))
            for iid, city, region, country in cur.fetchall():
                parts = [p for p in (city, region, country) if p]
                out[iid] = ", ".join(parts)
    except Exception as e:
        app.logger.error(f"[geo] read failed: {e}")
    finally:
        conn.close()
    return out


# ---------------------------------------------------------------------------
# Safety: risk detection and escalation
# ---------------------------------------------------------------------------
# If someone signals suicidal intent, self-harm, or intent to harm another
# person, the advisor stops coaching, points them to emergency services and a
# qualified professional, and an alert is emailed. This runs before any model
# call so the response is deterministic rather than left to the model.

SAFETY_ALERT_EMAIL = os.environ.get("SAFETY_ALERT_EMAIL", "afriedman@j3p.health")
SAFETY_ALERT_FROM = os.environ.get("SAFETY_ALERT_FROM", "")

# Explicit statements of intent or ideation. Deliberately specific — vague
# distress ("I'm exhausted", "this is killing me") must not trip this.
_RISK_PATTERNS = [
    # Suicidal ideation / intent
    r"\b(?:want|going|plan|planning|about)\s+to\s+(?:kill|end)\s+(?:myself|my\s+life)\b",
    r"\bkill(?:ing)?\s+myself\b",
    r"\bend(?:ing)?\s+(?:my\s+life|it\s+all)\b",
    r"\btake\s+my\s+own\s+life\b",
    r"\bcommit(?:ting)?\s+suicide\b",
    r"\bsuicidal\b",
    r"\bwant\s+to\s+die\b",
    r"\bbetter\s+off\s+dead\b",
    r"\bdon'?t\s+want\s+to\s+(?:be\s+here|live)\s+(?:anymore|any\s+more)\b",
    r"\bno\s+reason\s+to\s+(?:live|go\s+on)\b",
    r"\bthinking\s+about\s+(?:killing|ending)\b",
    # Self-harm
    r"\b(?:hurt|hurting|harm|harming|cut|cutting|injur\w*|burn|burning)\s+myself\b",
    r"\bself[-\s]?harm(?:ing)?\b",
    # Harm to others
    r"\b(?:want|going|plan|planning)\s+to\s+(?:kill|hurt|harm|attack|shoot)\s+(?:him|her|them|someone|somebody|my\s+\w+)\b",
    r"\bkill\s+(?:him|her|them)\b",
    r"\bhurt\s+(?:someone|somebody|him|her|them)\b",
    # Crisis phrasing
    r"\bin\s+crisis\b.*\b(?:hurt|harm|die|end)\b",
]
_RISK_RE = re.compile("|".join(_RISK_PATTERNS), re.IGNORECASE)

# Figures of speech that use the same words harmlessly
_RISK_FALSE_POSITIVES = re.compile(
    r"(this|it|that|deadline|schedule|commute|meeting|workload|job|"
    r"paperwork|emr|charting)\s+is\s+killing\s+me|"
    r"\bkilling\s+it\b|\bdying\s+to\b|\bdead\s+line\b|"
    r"\bshoot\s+(?:me\s+)?an?\s+(?:email|message|note)\b|"
    r"\bkill(?:ing)?\s+the\s+(?:project|initiative|program|idea)\b",
    re.IGNORECASE)


def detect_risk(text: str) -> bool:
    """True when a message signals risk of harm to self or others."""
    if not text:
        return False
    body = str(text)
    if _RISK_FALSE_POSITIVES.search(body) and not _RISK_RE.search(
            _RISK_FALSE_POSITIVES.sub(" ", body)):
        return False
    return bool(_RISK_RE.search(body))


SAFETY_RESPONSE = (
    "I want to stop and address what you just said, because it matters more "
    "than anything else we were working on.\n\n"
    "I'm not able to help with this, and you deserve support from someone who "
    "can. Please reach out right now:\n\n"
    "- **If you are in immediate danger, call 911.**\n"
    "- **Call or text 988** (Suicide & Crisis Lifeline, US) to reach a trained "
    "counsellor any time, day or night.\n"
    "- If you're outside the US, contact your local emergency number or crisis line.\n\n"
    "Please also contact a qualified professional — your physician, a licensed "
    "mental health clinician, or your employee assistance program — as soon as "
    "you can. If there is someone you trust nearby, tell them what's going on "
    "and let them stay with you.\n\n"
    "You don't have to work through this alone, and reaching out is the right "
    "next step."
)


def send_email(to_addr: str, subject: str, body: str) -> bool:
    """Send one plain-text email. Tries Postmark, then SMTP, then logs."""
    token = (os.environ.get("POSTMARK_SERVER_TOKEN")
             or os.environ.get("POSTMARK_TOKEN") or "")
    from_addr = (SAFETY_ALERT_FROM or os.environ.get("POSTMARK_FROM")
                 or os.environ.get("SMTP_FROM") or SAFETY_ALERT_EMAIL)

    if token:
        try:
            import urllib.request as _url
            import json as _j
            payload = _j.dumps({
                "From": from_addr, "To": to_addr,
                "Subject": subject, "TextBody": body,
                "MessageStream": os.environ.get("POSTMARK_STREAM", "outbound"),
            }).encode("utf-8")
            req = _url.Request(
                "https://api.postmarkapp.com/email", data=payload,
                headers={"Accept": "application/json",
                         "Content-Type": "application/json",
                         "X-Postmark-Server-Token": token})
            with _url.urlopen(req, timeout=8) as resp:
                resp.read()
            return True
        except Exception as e:
            app.logger.error(f"[safety] Postmark alert failed: {e}")

    host = os.environ.get("SMTP_HOST")
    if host:
        try:
            import smtplib
            from email.message import EmailMessage
            msg = EmailMessage()
            msg["Subject"] = subject
            msg["From"] = from_addr
            msg["To"] = to_addr
            msg.set_content(body)
            port = int(os.environ.get("SMTP_PORT", "587"))
            with smtplib.SMTP(host, port, timeout=10) as smtp:
                smtp.starttls()
                user = os.environ.get("SMTP_USER")
                if user:
                    smtp.login(user, os.environ.get("SMTP_PASSWORD", ""))
                smtp.send_message(msg)
            return True
        except Exception as e:
            app.logger.error(f"[safety] SMTP alert failed: {e}")

    # Nothing configured — make sure it is at least visible in the logs
    app.logger.error(f"[email] NOT SENT (no mail transport configured). "
                     f"To: {to_addr} Subject: {subject}\n{body}")
    return False


def send_alert_email(subject: str, body: str) -> bool:
    """Operational alert to the safety address."""
    return send_email(SAFETY_ALERT_EMAIL, subject, body)


def raise_safety_alert(user_message: str, ip: str = ""):
    """Email the alert on a background thread so the reply isn't delayed."""
    when = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    excerpt = (user_message or "")[:1500]
    location = ""
    try:
        loc = resolve_location(ip) if ip else None
        if loc:
            location = ", ".join(p for p in (loc["city"], loc["region"], loc["country"]) if p)
    except Exception:
        pass
    body = (
        "A J3P Advisor session contained language indicating possible risk of "
        "harm to self or others.\n\n"
        f"Time: {when}\n"
        f"Approximate location: {location or 'unknown'}\n\n"
        "The participant was shown emergency guidance (911 and 988) and "
        "directed to a qualified professional. The advisor did not continue "
        "the coaching conversation.\n\n"
        "--- Participant message ---\n"
        f"{excerpt}\n"
        "---------------------------\n\n"
        "This alert is generated automatically. Review in the admin "
        "conversation log for full context."
    )
    threading.Thread(
        target=send_alert_email,
        args=("J3P Advisor — participant safety alert", body),
        daemon=True,
    ).start()


# ---------------------------------------------------------------------------
# Setting ratings from the admin panel
# ---------------------------------------------------------------------------
# database.py owns the feedback table, so rather than assume its name the
# table is discovered once from information_schema by looking for the columns
# this app knows it has. Clearing a rating back to unrated needs direct SQL,
# which db.update_feedback_rating doesn't cover.

_feedback_table_cache = None


def _feedback_table(conn):
    """Find the table holding feedback rows. Cached after the first lookup."""
    global _feedback_table_cache
    if _feedback_table_cache:
        return _feedback_table_cache
    with conn.cursor() as cur:
        cur.execute("""
            SELECT table_name
            FROM information_schema.columns
            WHERE table_schema = 'public'
              AND column_name IN ('rating', 'bot_reply', 'user_message')
            GROUP BY table_name
            HAVING COUNT(DISTINCT column_name) = 3
            LIMIT 1
        """)
        row = cur.fetchone()
    if row:
        _feedback_table_cache = row[0]
    return _feedback_table_cache


def set_feedback_rating(feedback_id: int, rating) -> bool:
    """Set a row's rating to 'up', 'down', or None to clear it."""
    if rating not in ("up", "down", None):
        return False

    # Prefer the database module's own path where it applies
    if rating in ("up", "down"):
        try:
            if db.update_feedback_rating(int(feedback_id), rating, ""):
                return True
        except Exception as e:
            app.logger.warning(f"[rating] db.update_feedback_rating failed: {e}")

    conn = _settings_db_conn()
    if not conn:
        return False
    try:
        table = _feedback_table(conn)
        if not table:
            app.logger.error("[rating] could not locate the feedback table")
            return False
        with conn.cursor() as cur:
            cur.execute(
                f'UPDATE "{table}" SET rating = %s WHERE id = %s',
                (rating, int(feedback_id)))
            changed = cur.rowcount
        conn.commit()
        return changed > 0
    except Exception as e:
        app.logger.error(f"[rating] update failed: {e}")
        return False
    finally:
        conn.close()


# ---------------------------------------------------------------------------
# Sign-in by magic link
# ---------------------------------------------------------------------------
# Independent of paywall.py (which ties sign-in to Stripe). A signed, expiring
# token is emailed; following it establishes identity. Identity is what makes
# memory persist across visits — history is keyed to the signed-in address
# rather than a throwaway browser token.

REQUIRE_LOGIN = os.environ.get("REQUIRE_LOGIN", "off").lower() in ("on", "1", "true")
LOGIN_LINK_TTL = int(os.environ.get("LOGIN_LINK_TTL_MIN", "30")) * 60
# Optional allow-list: comma-separated addresses or @domain entries.
LOGIN_ALLOWED = [e.strip().lower() for e in
                 os.environ.get("LOGIN_ALLOWED", "").split(",") if e.strip()]


def mail_transport_configured() -> bool:
    """True when a sign-in link could actually be delivered."""
    return bool(os.environ.get("POSTMARK_SERVER_TOKEN")
                or os.environ.get("POSTMARK_TOKEN")
                or os.environ.get("SMTP_HOST"))


def login_is_required() -> bool:
    """Admin setting wins; falls back to the environment variable."""
    try:
        return bool(load_settings()["require_login"])
    except Exception:
        return REQUIRE_LOGIN


def _login_serializer():
    from itsdangerous import URLSafeTimedSerializer
    return URLSafeTimedSerializer(app.secret_key, salt="j3p-login")


def email_allowed(email: str) -> bool:
    if not LOGIN_ALLOWED:
        return True
    email = (email or "").lower()
    domain = "@" + email.split("@")[-1]
    return email in LOGIN_ALLOWED or domain in LOGIN_ALLOWED


def make_login_link(email: str, base_url: str) -> str:
    token = _login_serializer().dumps(email.strip().lower())
    return f"{base_url.rstrip('/')}/login/verify?token={token}"


def read_login_token(token: str):
    from itsdangerous import BadSignature, SignatureExpired
    try:
        return _login_serializer().loads(token, max_age=LOGIN_LINK_TTL)
    except (BadSignature, SignatureExpired):
        return None


def current_user() -> str:
    """Signed-in email, or '' when anonymous."""
    return session.get("user_email", "")


def login_required(f):
    """Require a signed-in participant when REQUIRE_LOGIN is on."""
    @wraps(f)
    def wrapper(*args, **kwargs):
        if login_is_required() and not current_user():
            return redirect(url_for("login_page", next=request.path))
        return f(*args, **kwargs)
    return wrapper


LOGIN_HTML = """<!DOCTYPE html>
<html lang="en"><head>
<meta charset="UTF-8" /><meta name="viewport" content="width=device-width, initial-scale=1.0" />
<title>Sign in — {{ cfg.persona_name }}</title>
<link rel="icon" href="{{ cfg.favicon_url }}" />
<link rel="preconnect" href="https://fonts.googleapis.com">
<link href="https://fonts.googleapis.com/css2?family=Jost:wght@300;400;500;600&display=swap" rel="stylesheet">
<style>
:root { --navy:#27334A; --gold:#D2BC8D; --paper:#FAF6F0; --line:rgba(39,51,74,0.12); --muted:#6B7280; }
body { margin:0; font-family:'Jost',-apple-system,sans-serif; background:var(--paper);
       color:var(--navy); display:flex; align-items:center; justify-content:center;
       min-height:100vh; padding:1.25rem; }
.box { background:#fff; border-radius:4px; box-shadow:0 18px 50px rgba(39,51,74,.18);
       max-width:460px; width:100%; overflow:hidden; }
.head { background:var(--navy); border-bottom:2px solid var(--gold); padding:1rem 1.75rem;
        display:flex; align-items:center; gap:1rem; }
.head img { height:46px; }
.head span { color:var(--gold); font-size:.8rem; letter-spacing:.22em; text-transform:uppercase; }
.content { padding:1.75rem 2rem 1.6rem; }
h1 { font-size:.85rem; letter-spacing:.16em; text-transform:uppercase; margin:0 0 1rem;
     padding-bottom:.65rem; border-bottom:1px solid var(--line); font-weight:500; }
p { font-size:.92rem; line-height:1.6; margin:0 0 1rem; }
input[type=email] { width:100%; padding:.8rem .9rem; border:1px solid var(--line);
    border-radius:2px; font-family:inherit; font-size:.95rem; background:var(--paper); }
input[type=email]:focus { outline:none; border-color:var(--gold); }
button { width:100%; margin-top:.9rem; padding:.85rem; background:var(--navy); color:var(--gold);
    border:1px solid var(--navy); border-radius:2px; font-family:inherit; font-size:.78rem;
    letter-spacing:.18em; text-transform:uppercase; cursor:pointer; }
button:hover { background:var(--gold); color:var(--navy); }
.note { margin-top:1rem; font-size:.72rem; color:var(--muted); line-height:1.6; }
.msg { padding:.7rem .9rem; border-radius:2px; font-size:.86rem; margin-bottom:1rem; }
.ok { background:var(--gold); color:var(--navy); }
.err { background:#FEE; color:#9D432C; border:1px solid #E7C3BA; }
</style></head><body>
<div class="box">
  <div class="head">
    <img src="{{ cfg.logo_url }}" alt="{{ cfg.persona_name }}" />
    <span>{{ cfg.persona_name }}</span>
  </div>
  <div class="content">
    <h1>Sign in</h1>
    {% if notice %}<div class="msg {{ 'ok' if notice_ok else 'err' }}">{{ notice }}</div>{% endif %}
    {% if not sent %}
    <p>Enter your email and we'll send you a sign-in link. Signing in lets the
       advisor pick up where you left off next time.</p>
    <form method="POST" action="/login/send">
      <input type="hidden" name="next" value="{{ next_path }}" />
      <input type="email" name="email" placeholder="you@example.com" required autofocus />
      <button type="submit">Email me a sign-in link</button>
    </form>
    <p class="note">The link works once and expires in {{ ttl_minutes }} minutes.</p>
    {% endif %}
  </div>
</div></body></html>"""


@app.route("/login", methods=["GET"])
def login_page():
    return render_template_string(
        LOGIN_HTML, cfg=CONFIG, notice=request.args.get("notice"),
        notice_ok=request.args.get("ok") == "1", sent=False,
        next_path=request.args.get("next", "/"),
        ttl_minutes=LOGIN_LINK_TTL // 60)


@app.route("/login/send", methods=["POST"])
def login_send():
    email = (request.form.get("email") or "").strip().lower()
    next_path = request.form.get("next") or "/"
    if not email or "@" not in email:
        return render_template_string(
            LOGIN_HTML, cfg=CONFIG, notice="Please enter a valid email address.",
            notice_ok=False, sent=False, next_path=next_path,
            ttl_minutes=LOGIN_LINK_TTL // 60)
    if not email_allowed(email):
        # Deliberately identical to the success message — don't reveal the list
        return render_template_string(
            LOGIN_HTML, cfg=CONFIG,
            notice=f"If {email} is registered, a sign-in link is on its way.",
            notice_ok=True, sent=True, next_path=next_path,
            ttl_minutes=LOGIN_LINK_TTL // 60)

    base = paywall.PUBLIC_BASE_URL or request.host_url.rstrip("/")
    link = make_login_link(email, base)
    body = (
        f"Here is your sign-in link for the {CONFIG['persona_name']}:\n\n"
        f"{link}\n\n"
        f"It works once and expires in {LOGIN_LINK_TTL // 60} minutes. "
        "If you didn't request it, you can ignore this message."
    )
    sent_ok = send_email(email, f"Your {CONFIG['persona_name']} sign-in link", body)
    if not sent_ok:
        app.logger.error("[login] could not send sign-in link")
    return render_template_string(
        LOGIN_HTML, cfg=CONFIG,
        notice=f"Sign-in link sent to {email}. Check your inbox and spam folder.",
        notice_ok=True, sent=True, next_path=next_path,
        ttl_minutes=LOGIN_LINK_TTL // 60)


@app.route("/login/verify", methods=["GET"])
def login_verify():
    email = read_login_token(request.args.get("token", ""))
    if not email:
        return redirect(url_for(
            "login_page",
            notice="That link has expired or is invalid. Please request a new one."))
    session["user_email"] = email
    session.pop("chat_token", None)      # switch to this person's own history
    app.logger.info(f"[login] signed in: {email}")
    nxt = request.args.get("next") or "/"
    return redirect(nxt if nxt.startswith("/") else "/")


@app.route("/logout", methods=["GET", "POST"])
def logout():
    session.pop("user_email", None)
    session.pop("chat_token", None)
    session["messages"] = []
    return redirect(url_for("login_page", notice="You've been signed out.", ok="1"))


# ---------------------------------------------------------------------------
# Continuous learning from feedback
# ---------------------------------------------------------------------------
# Thumbs-down rows with a comment are the raw material for lessons. Approving
# them one at a time doesn't scale, so this finds eligible rows, embeds the
# question, and approves them in a batch — either on demand from the admin
# panel or on a schedule.
#
# Deliberately NOT automatic for everything: a lesson is injected into other
# participants' prompts, so anything auto-approved needs a comment explaining
# what went wrong, and needs its text checked for confidential detail first.

LEARNING_MIN_COMMENT_CHARS = 8
LEARNING_BATCH_CAP = 25


def _learning_ensure_table(conn):
    with conn.cursor() as cur:
        cur.execute("""
            CREATE TABLE IF NOT EXISTS learning_runs (
                id          BIGSERIAL PRIMARY KEY,
                ran_at      TIMESTAMPTZ NOT NULL DEFAULT NOW(),
                considered  INTEGER NOT NULL DEFAULT 0,
                approved    INTEGER NOT NULL DEFAULT 0,
                skipped     INTEGER NOT NULL DEFAULT 0,
                detail      TEXT,
                trigger     TEXT
            )
        """)
        # Added later — runs that changed nothing are archived rather than
        # deleted, so the visible history stays readable but nothing is lost.
        try:
            cur.execute("ALTER TABLE learning_runs ADD COLUMN IF NOT EXISTS "
                        "archived BOOLEAN NOT NULL DEFAULT FALSE")
        except Exception:
            pass
    conn.commit()


def _record_learning_run(considered, approved, skipped, detail, trigger):
    conn = _settings_db_conn()
    if not conn:
        return
    try:
        _learning_ensure_table(conn)
        # A run that learned nothing is archived immediately — those repeat
        # every day and would otherwise bury the runs that mattered.
        archived = (approved == 0)
        with conn.cursor() as cur:
            cur.execute("""
                INSERT INTO learning_runs
                    (considered, approved, skipped, detail, trigger, archived)
                VALUES (%s, %s, %s, %s, %s, %s)
            """, (considered, approved, skipped, detail[:4000], trigger, archived))
            # Keep only the most recent meaningful runs on the main view
            cur.execute("""
                UPDATE learning_runs SET archived = TRUE
                WHERE archived = FALSE AND id NOT IN (
                    SELECT id FROM learning_runs
                    WHERE archived = FALSE ORDER BY id DESC LIMIT 10
                )
            """)
        conn.commit()
    except Exception as e:
        app.logger.error(f"[learning] could not record run: {e}")
    finally:
        conn.close()


def recent_learning_runs(limit=8, archived=False):
    conn = _settings_db_conn()
    if not conn:
        return []
    out = []
    try:
        _learning_ensure_table(conn)
        with conn.cursor() as cur:
            cur.execute("""
                SELECT ran_at, considered, approved, skipped, trigger, detail
                FROM learning_runs WHERE archived = %s
                ORDER BY id DESC LIMIT %s
            """, (archived, limit))
            for ran_at, considered, approved, skipped, trigger, detail in cur.fetchall():
                out.append({
                    "when": _fmt_ts(ran_at), "considered": considered,
                    "approved": approved, "skipped": skipped,
                    "trigger": trigger or "", "detail": detail or "",
                })
    except Exception as e:
        app.logger.error(f"[learning] could not read runs: {e}")
    finally:
        conn.close()
    return out


# Detail that must never be carried into another participant's prompt
_CONFIDENTIAL_RE = re.compile(
    r"\b(salary|compensation|rvu target|severance|terminat\w+|fired|"
    r"lawsuit|litigation|malpractice|hipaa|patient name|mrn|"
    r"diagnos\w+|medication|prescri\w+|probation|grievance|"
    r"disciplinary|investigation)\b", re.IGNORECASE)


def lesson_is_safe_to_share(row) -> tuple:
    """Should this feedback row become a lesson visible in others' sessions?"""
    text = " ".join([
        row.get("user_message") or "", row.get("bot_reply") or "",
        row.get("comment") or "",
    ])
    if row.get("attachment_info") == "SAFETY ALERT":
        return False, "safety escalation"
    hit = _CONFIDENTIAL_RE.search(text)
    if hit:
        return False, f"mentions '{hit.group(0)}' — needs manual review"
    if len(row.get("user_message") or "") > 4000:
        return False, "question too long to generalise"
    return True, ""


def run_learning_cycle(trigger="manual", dry_run=False) -> dict:
    """Approve eligible thumbs-down feedback as lessons. Returns a summary."""
    if not (db.is_enabled() and emb.is_enabled()):
        return {"ok": False, "error": "Database or embeddings not configured."}

    considered = approved = skipped = 0
    up_count = down_count = 0
    notes = []
    try:
        # Both ratings teach something: thumbs-down shows what to avoid,
        # thumbs-up shows the shape of a reply that landed well.
        rows = list(db.list_feedback(limit=500, rating="down"))
        rows += list(db.list_feedback(limit=500, rating="up"))
    except Exception as e:
        app.logger.error(f"[learning] could not list feedback: {e}")
        return {"ok": False, "error": str(e)[:200]}

    for row in rows:
        if row.get("approved_for_learning"):
            continue
        rating = (row.get("rating") or "").lower()
        comment = (row.get("comment") or "").strip()

        # A thumbs-down needs an explanation — without one there's nothing to
        # learn from. A thumbs-up is self-explanatory: the reply worked.
        if rating == "down" and len(comment) < LEARNING_MIN_COMMENT_CHARS:
            skipped += 1
            notes.append(f"#{row.get('id')}: thumbs-down with no comment — nothing to learn")
            continue
        if rating == "up" and len(row.get("bot_reply") or "") < 120:
            skipped += 1
            notes.append(f"#{row.get('id')}: reply too short to be a useful example")
            continue

        considered += 1
        safe, why = lesson_is_safe_to_share(row)
        if not safe:
            skipped += 1
            notes.append(f"#{row.get('id')}: held back — {why}")
            continue
        if dry_run:
            approved += 1
            notes.append(f"#{row.get('id')}: would be approved")
            continue
        if approved >= LEARNING_BATCH_CAP:
            notes.append(f"batch cap of {LEARNING_BATCH_CAP} reached — run again for the rest")
            break
        try:
            question_embedding = emb.embed_text(row.get("user_message") or "")
            if db.approve_feedback_as_lesson(row["id"], question_embedding):
                approved += 1
                if rating == "up":
                    up_count += 1
                else:
                    down_count += 1
                notes.append(f"#{row['id']}: learned ({'what worked' if rating == 'up' else 'what to avoid'})")
            else:
                skipped += 1
                notes.append(f"#{row['id']}: not eligible")
        except Exception as e:
            skipped += 1
            notes.append(f"#{row.get('id')}: failed — {str(e)[:80]}")

    detail = "\n".join(notes[:120])
    if not dry_run:
        _record_learning_run(considered, approved, skipped, detail, trigger)
    app.logger.info(f"[learning] {trigger}: considered={considered} "
                    f"approved={approved} skipped={skipped}")
    return {"ok": True, "considered": considered, "approved": approved,
            "skipped": skipped, "detail": detail, "dry_run": dry_run,
            "up": up_count, "down": down_count}


def archived_run_count() -> int:
    conn = _settings_db_conn()
    if not conn:
        return 0
    try:
        _learning_ensure_table(conn)
        with conn.cursor() as cur:
            cur.execute("SELECT COUNT(*) FROM learning_runs WHERE archived = TRUE")
            row = cur.fetchone()
        return int(row[0]) if row else 0
    except Exception as e:
        app.logger.error(f"[learning] archive count failed: {e}")
        return 0
    finally:
        conn.close()


def archive_all_learning_runs() -> bool:
    """Clear the visible history; nothing is deleted."""
    conn = _settings_db_conn()
    if not conn:
        return False
    try:
        _learning_ensure_table(conn)
        with conn.cursor() as cur:
            cur.execute("UPDATE learning_runs SET archived = TRUE WHERE archived = FALSE")
        conn.commit()
        return True
    except Exception as e:
        app.logger.error(f"[learning] archive failed: {e}")
        return False
    finally:
        conn.close()


# --- Scheduler -------------------------------------------------------------
# A background thread rather than an external cron, so there's nothing extra
# to configure on Railway. Runs only when the admin has switched it on.

LEARNING_INTERVAL_HOURS = int(os.environ.get("LEARNING_INTERVAL_HOURS", "24"))
_learning_thread_started = False


def _claim_scheduled_run(conn) -> bool:
    """Only one worker should run the cycle.

    Gunicorn starts several worker processes and each was launching its own
    scheduler thread, so the same pass ran two or three times — visible as
    duplicate rows with identical timestamps. A worker now has to claim the
    slot in the database before running, and the claim only succeeds if enough
    time has passed since the last run.
    """
    try:
        _learning_ensure_table(conn)
        with conn.cursor() as cur:
            cur.execute("""
                SELECT ran_at FROM learning_runs
                WHERE trigger = 'scheduled' ORDER BY id DESC LIMIT 1
            """)
            row = cur.fetchone()
        if not row or not row[0]:
            return True
        last = row[0]
        try:
            from datetime import timezone
            now = datetime.now(last.tzinfo) if last.tzinfo else datetime.now()
            elapsed_h = (now - last).total_seconds() / 3600.0
        except Exception:
            return True
        # A margin below the interval avoids a worker being locked out by
        # clock drift, while still blocking near-simultaneous duplicates.
        return elapsed_h >= max(1, LEARNING_INTERVAL_HOURS) * 0.9
    except Exception as e:
        app.logger.error(f"[learning] claim check failed: {e}")
        return False


def _learning_loop():
    import time
    # Let the app finish starting before the first pass, and stagger workers
    # so two don't reach the claim check in the same instant.
    time.sleep(120 + (os.getpid() % 47))
    while True:
        try:
            if load_settings().get("auto_learning"):
                conn = _settings_db_conn()
                may_run = True
                if conn:
                    try:
                        may_run = _claim_scheduled_run(conn)
                    finally:
                        conn.close()
                if may_run:
                    with app.app_context():
                        run_learning_cycle(trigger="scheduled")
                else:
                    app.logger.info("[learning] another worker ran recently — skipping")
        except Exception as e:
            app.logger.error(f"[learning] scheduled run failed: {e}")
        time.sleep(max(1, LEARNING_INTERVAL_HOURS) * 3600)


def start_learning_scheduler():
    global _learning_thread_started
    if _learning_thread_started:
        return
    _learning_thread_started = True
    threading.Thread(target=_learning_loop, daemon=True).start()
    app.logger.info(f"[learning] scheduler started "
                    f"(every {LEARNING_INTERVAL_HOURS}h when enabled)")


def admin_required(f):
    @wraps(f)
    def wrapper(*args, **kwargs):
        if not CONFIG["admin_password"]:
            return ("Admin disabled. Set ADMIN_PASSWORD environment variable.", 503)
        if not session.get("is_admin"):
            return redirect(url_for("admin_login"))
        return f(*args, **kwargs)
    return wrapper


# ---------------------------------------------------------------------------
# Main chat HTML
# ---------------------------------------------------------------------------

# ---------------------------------------------------------------------------
# Release & acknowledgment language
# ---------------------------------------------------------------------------
# Defined once and rendered into BOTH the session-entry gate and the
# scheduling gate, so the two can never drift apart. Edit here to change both.

RELEASE_HEADING = "Release &amp; Acknowledgment"

RELEASE_BODY_HTML = """
  <p>
    By checking the box below, I acknowledge that I am voluntarily using
    the J3P Advisor and understand that the content, coaching and guidance
    provided are for personal and professional development purposes only.
    I understand that these activities are not medical, psychological,
    legal, or other professional advice, and I am responsible for my own
    decisions and actions.
  </p>
  <p>
    To the extent permitted by law, I release Residency Select LLC dba
    J3P Health, its coaches, employees, and representatives from liability
    arising from my voluntary use of the J3P Advisor.
  </p>
"""

RELEASE_CHECKBOX_LABEL = "I have read, understood, and agree to the above."


INDEX_HTML = r"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1.0" />
  <title>{{ cfg.persona_name }}</title>
  <link rel="icon" href="{{ cfg.favicon_url }}" />
  <link rel="preconnect" href="https://fonts.googleapis.com">
  <link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>
  <link href="https://fonts.googleapis.com/css2?family=Jost:wght@300;400;500;600&display=swap" rel="stylesheet">
  <style>
    :root {
      --navy: {{ cfg.navy }};
      --gold: {{ cfg.gold }};
      --paper: {{ cfg.paper }};
      --paper-2: #FFFFFF;
      --line: rgba(39, 51, 74, 0.12);
      --muted: #6B7280;
      --text: #27334A;
      --rust: #9D432C;
      --shadow: 0 1px 2px rgba(39, 51, 74, 0.05), 0 8px 28px rgba(39, 51, 74, 0.07);
    }
    * { box-sizing: border-box; }
    html, body { height: 100%; }
    body {
      margin: 0;
      font-family: 'Jost', -apple-system, BlinkMacSystemFont, sans-serif;
      background: var(--paper); color: var(--text);
      display: flex; flex-direction: column;
      font-size: 15px; line-height: 1.6;
      -webkit-font-smoothing: antialiased;
    }
    header {
      background: var(--navy); color: var(--paper-2);
      padding: 1rem 1.75rem;
      display: flex; justify-content: space-between; align-items: center;
      border-bottom: 2px solid var(--gold);
      gap: 0.75rem;
    }
    .brand { display: flex; align-items: center; gap: 1rem; min-width: 0; flex: 1; }
    .brand-logo { height: 60px; width: auto; display: block; flex-shrink: 0; }
    .brand-divider { width: 1px; height: 38px; background: rgba(210, 188, 141, 0.35); flex-shrink: 0; }
    .brand-tag {
      font-size: 0.92rem; letter-spacing: 0.22em;
      text-transform: uppercase; color: var(--gold);
      overflow: hidden; text-overflow: ellipsis; white-space: nowrap;
    }
    header button {
      background: transparent; color: var(--paper-2);
      border: 1px solid rgba(210, 188, 141, 0.35);
      padding: 0.5rem 1rem; border-radius: 2px;
      cursor: pointer; font-size: 0.75rem;
      font-family: inherit; letter-spacing: 0.14em;
      text-transform: uppercase; transition: all 0.2s ease;
      flex-shrink: 0; white-space: nowrap;
      display: inline-flex; align-items: center; gap: 0.4rem;
    }
    header button:hover {
      background: rgba(210, 188, 141, 0.08);
      border-color: var(--gold); color: var(--gold);
    }
    /* Show full label on desktop, icon-only label on small screens */
    .reset-icon { width: 16px; height: 16px; display: none; }
    .reset-label { display: inline; }
    .autospeak-icon { width: 16px; height: 16px; }
    .autospeak-label { display: inline; }
    .reset-short { display: none; }
    /* Highlight the auto-speak button when it's ON */
    #autospeak-btn.on {
      background: var(--gold);
      color: var(--navy);
      border-color: var(--gold);
    }
    #autospeak-btn.on:hover {
      background: var(--gold);
      color: var(--navy);
    }
    /* Voice picker — uses the device's own voices, so no service cost */
    .voice-wrap { position: relative; display: inline-flex; align-items: center; flex-shrink: 0; }
    .voice-menu {
      position: absolute; top: calc(100% + 10px); right: 0;
      background: var(--paper-2); color: var(--text);
      border: 1px solid var(--line); border-top: 2px solid var(--gold);
      border-radius: 4px; box-shadow: 0 12px 34px rgba(39, 51, 74, 0.26);
      padding: 0.9rem 1rem 0.85rem; min-width: 268px;
      display: none; z-index: 60; text-align: left;
    }
    .voice-menu.open { display: block; }
    .voice-row { margin-bottom: 0.8rem; }
    .voice-row label {
      display: block; margin-bottom: 0.35rem;
      font-size: 0.64rem; letter-spacing: 0.14em;
      text-transform: uppercase; color: var(--muted);
    }
    .voice-menu select {
      width: 100%; padding: 0.5rem 0.6rem;
      border: 1px solid var(--line); border-radius: 2px;
      background: var(--paper); color: var(--text);
      font-family: inherit; font-size: 0.84rem; cursor: pointer;
    }
    .voice-menu select:focus { outline: none; border-color: var(--gold); }
    .voice-menu input[type="range"] {
      width: 100%; accent-color: var(--navy); cursor: pointer; margin: 0;
    }
    header .voice-preview {
      width: 100%; margin: 0.1rem 0 0 0;
      background: var(--navy); color: var(--gold);
      border: 1px solid var(--navy); border-radius: 2px;
      padding: 0.55rem 0.8rem; font-size: 0.68rem;
      letter-spacing: 0.14em; text-transform: uppercase;
      display: block; text-align: center;
    }
    header .voice-preview:hover { background: var(--gold); color: var(--navy); border-color: var(--gold); }
    .voice-note {
      margin: 0.6rem 0 0 0; font-size: 0.66rem; line-height: 1.45;
      color: var(--muted); letter-spacing: 0; text-transform: none;
    }
    @media (max-width: 640px) {
      .voice-menu { min-width: 240px; right: -40px; }
    }

    #chat-wrap { flex: 1; overflow-y: auto; }
    #chat { max-width: 760px; margin: 0 auto; padding: 2.25rem 1.5rem 1rem; }
    /* Participant materials */
    .materials-link {
      background: none; border: none; cursor: pointer;
      font-family: inherit; font-size: 0.62rem; letter-spacing: 0.1em;
      text-transform: uppercase; color: var(--muted);
      padding: 0.2rem 0; border-bottom: 1px solid var(--gold);
    }
    .materials-link:hover { color: var(--navy); }
    .mat-overlay {
      position: fixed; inset: 0; background: rgba(39,51,74,0.55);
      display: flex; align-items: center; justify-content: center;
      z-index: 900; padding: 1.25rem;
    }
    .mat-overlay[hidden] { display: none; }
    .mat-box {
      background: #fff; border-radius: 4px; width: 100%; max-width: 560px;
      max-height: 88vh; overflow-y: auto;
      box-shadow: 0 20px 60px rgba(39,51,74,0.3);
    }
    .mat-head {
      background: var(--navy); border-bottom: 2px solid var(--gold);
      padding: 0.9rem 1.3rem; display: flex; justify-content: space-between;
      align-items: center; position: sticky; top: 0;
    }
    .mat-head h3 {
      margin: 0; color: var(--gold); font-size: 0.78rem; font-weight: 500;
      letter-spacing: 0.16em; text-transform: uppercase;
    }
    .mat-close {
      background: none; border: none; color: var(--paper);
      font-size: 1.3rem; cursor: pointer; line-height: 1; padding: 0 0.2rem;
    }
    .mat-body { padding: 1.2rem 1.3rem 1.4rem; }
    .mat-intro { font-size: 0.86rem; line-height: 1.6; margin: 0 0 1.1rem; }
    .mat-tabs { display: flex; gap: 0.4rem; margin-bottom: 1rem; }
    .mat-tab {
      flex: 1; padding: 0.5rem; font-family: inherit; cursor: pointer;
      font-size: 0.68rem; letter-spacing: 0.1em; text-transform: uppercase;
      background: var(--paper); border: 1px solid var(--line);
      border-radius: 2px; color: var(--navy);
    }
    .mat-tab.active { background: var(--navy); color: var(--gold); border-color: var(--navy); }
    .mat-pane[hidden] { display: none; }
    .mat-field {
      width: 100%; padding: 0.6rem 0.7rem; margin-bottom: 0.6rem;
      border: 1px solid var(--line); border-radius: 2px;
      font-family: inherit; font-size: 0.9rem; box-sizing: border-box;
    }
    textarea.mat-field { min-height: 170px; resize: vertical; line-height: 1.5; }
    .mat-share {
      display: flex; align-items: flex-start; gap: 0.5rem;
      font-size: 0.78rem; color: var(--muted); margin: 0.4rem 0 0.9rem;
      line-height: 1.5;
    }
    .mat-list { margin-top: 1.3rem; padding-top: 1rem; border-top: 1px dashed var(--line); }
    .mat-item {
      display: flex; align-items: center; gap: 0.6rem;
      padding: 0.5rem 0; border-bottom: 1px solid var(--line);
      font-size: 0.85rem;
    }
    .mat-item-title { flex: 1; min-width: 0; word-break: break-word; }
    .mat-item-meta { font-size: 0.7rem; color: var(--muted); white-space: nowrap; }
    .mat-remove {
      background: none; border: none; color: var(--rust); cursor: pointer;
      font-size: 0.68rem; letter-spacing: 0.08em; text-transform: uppercase;
    }
    .mat-box .btn {
      background: var(--navy); color: var(--gold);
      border: 1px solid var(--navy); border-radius: 2px;
      padding: 0.7rem 1.2rem; cursor: pointer; font-family: inherit;
      font-size: 0.72rem; letter-spacing: 0.16em; text-transform: uppercase;
      transition: background 0.2s ease, color 0.2s ease;
    }
    .mat-box .btn:hover { background: var(--gold); color: var(--navy); }
    .mat-box .btn:disabled { opacity: 0.5; cursor: not-allowed; }
    .mat-profile {
      display: flex; align-items: center; gap: 0.6rem; flex-wrap: wrap;
      background: var(--paper); border-left: 3px solid var(--gold);
      padding: 0.6rem 0.8rem; margin-bottom: 1rem; font-size: 0.84rem;
    }
    .mat-profile[hidden] { display: none; }
    .mat-profile-label {
      font-size: 0.62rem; letter-spacing: 0.1em; text-transform: uppercase;
      color: var(--muted);
    }
    #mat-profile-text { flex: 1; min-width: 0; }
    .mat-msg { font-size: 0.82rem; margin: 0.6rem 0 0; }
    .mat-msg.ok { color: #2F6B4F; }
    .mat-msg.err { color: var(--rust); }

    /* Avatar beside advisor replies — animated states */
    .avatar-wrap {
      position: relative; flex: 0 0 auto; width: 38px; height: 38px;
      margin-top: 0.2rem; cursor: pointer; border-radius: 50%;
      background: none; border: none; padding: 0;
    }
    .avatar-wrap:focus-visible { outline: 2px solid var(--gold); outline-offset: 3px; }
    /* Idle: a slow breath so it doesn't look like a dead image */
    @keyframes av-breathe {
      0%, 100% { transform: scale(1); }
      50%      { transform: scale(1.035); }
    }
    /* Thinking: gold ring travelling round the edge */
    @keyframes av-think {
      0%   { transform: rotate(0deg); }
      100% { transform: rotate(360deg); }
    }
    /* Speaking: concentric pulses radiating outward */
    @keyframes av-pulse {
      0%   { transform: scale(1); opacity: 0.55; }
      70%  { transform: scale(1.75); opacity: 0; }
      100% { transform: scale(1.75); opacity: 0; }
    }
    .avatar-wrap .avatar { animation: av-breathe 5.5s ease-in-out infinite; }
    .avatar-wrap.is-thinking .avatar,
    .avatar-wrap.is-speaking .avatar { animation: none; }
    .avatar-wrap.is-thinking .avatar { opacity: 0.85; }

    .avatar-ring {
      position: absolute; inset: -3px; border-radius: 50%;
      pointer-events: none; display: none;
    }
    .avatar-wrap.is-thinking .avatar-ring {
      display: block;
      border: 2px solid transparent;
      border-top-color: var(--gold);
      border-right-color: rgba(210,188,141,0.4);
      animation: av-think 0.9s linear infinite;
    }
    .avatar-pulse {
      position: absolute; inset: 0; border-radius: 50%;
      border: 2px solid var(--gold); pointer-events: none;
      opacity: 0; 
    }
    /* Responding: a gentle nod plus a soft glow while the reply lands */
    @keyframes av-respond {
      0%, 100% { transform: scale(1) translateY(0); }
      25%      { transform: scale(1.06) translateY(-1px); }
      60%      { transform: scale(1.02) translateY(1px); }
    }
    @keyframes av-glow {
      0%, 100% { box-shadow: 0 0 0 0 rgba(210,188,141,0.0); }
      50%      { box-shadow: 0 0 0 5px rgba(210,188,141,0.45); }
    }
    .avatar-wrap.is-responding .avatar {
      animation: av-respond 1.15s ease-in-out infinite,
                 av-glow 1.15s ease-in-out infinite;
    }
    .avatar-wrap.is-speaking .avatar-pulse { animation: av-pulse 1.6s ease-out infinite; }
    .avatar-wrap.is-speaking .avatar-pulse:nth-of-type(2) { animation-delay: 0.55s; }
    .avatar-wrap.is-speaking .avatar-pulse:nth-of-type(3) { animation-delay: 1.1s; }
    .avatar-wrap.is-speaking .avatar { box-shadow: 0 0 0 2px var(--gold); }

    /* Talking-head video replaces the photo in the same circle */
    .avatar-video {
      position: absolute; inset: 0; width: 100%; height: 100%;
      border-radius: 50%; object-fit: cover; display: none;
      border: 1.5px solid var(--gold);
    }
    .avatar-wrap.is-video .avatar-video { display: block; }
    .avatar-wrap.is-video .avatar { visibility: hidden; }
    .avatar-wrap.is-loading .avatar { opacity: 0.55; }
    .avatar-wrap.is-loading .avatar-ring {
      display: block; border: 2px solid transparent;
      border-top-color: var(--gold);
      animation: av-think 0.9s linear infinite;
    }

    /* Small hint on hover, matching the mic tooltip */
    .avatar-hint {
      position: absolute; left: 50%; transform: translateX(-50%);
      bottom: calc(100% + 8px); white-space: nowrap;
      background: var(--navy); color: var(--paper);
      font-size: 0.66rem; letter-spacing: 0.02em;
      padding: 0.3rem 0.55rem; border-radius: 3px;
      opacity: 0; visibility: hidden; transition: opacity 0.15s ease;
      pointer-events: none; z-index: 20;
    }
    .avatar-wrap:hover .avatar-hint,
    .avatar-wrap:focus-visible .avatar-hint { opacity: 1; visibility: visible; }

    @media (prefers-reduced-motion: reduce) {
      .avatar-wrap .avatar,
      .avatar-wrap.is-speaking .avatar-pulse,
      .avatar-wrap.is-thinking .avatar-ring { animation: none !important; }
      .avatar-wrap.is-speaking .avatar-ring,
      .avatar-wrap.is-responding .avatar-ring { display: block; border: 2px solid var(--gold); }
    }

    /* Avatar beside advisor replies */
    .msg-row { display: flex; align-items: flex-start; gap: 0.7rem; }
    .avatar {
      flex: 0 0 auto; width: 38px; height: 38px; border-radius: 50%;
      object-fit: cover; margin-top: 0.2rem;
      border: 1.5px solid var(--gold);
      box-shadow: 0 2px 6px rgba(39,51,74,0.14);
    }
    .msg-row > .msg { flex: 1 1 auto; min-width: 0; }
    @media (max-width: 640px) {
      .avatar { width: 30px; height: 30px; }
      .msg-row { gap: 0.5rem; }
    }

    .msg {
      margin-bottom: 1.25rem; padding: 1rem 1.2rem; border-radius: 4px;
      white-space: pre-wrap; word-wrap: break-word; font-size: 0.95rem;
      animation: fadeIn 0.3s ease-out;
    }
    @keyframes fadeIn { from { opacity: 0; transform: translateY(6px); } to { opacity: 1; transform: translateY(0); } }
    .user { background: var(--navy); color: var(--paper); margin-left: 18%; box-shadow: var(--shadow); }
    .assistant {
      background: var(--paper-2); border: 1px solid var(--line);
      margin-right: 12%; box-shadow: var(--shadow); position: relative;
    }
    .assistant::before {
      content: ""; position: absolute; left: 0; top: 0; bottom: 0;
      width: 3px; background: var(--gold);
    }
    /* Acknowledgment gate — shown at the start of every session */
    .ack-overlay {
      position: fixed; inset: 0; z-index: 100;
      background: rgba(39, 51, 74, 0.78);
      backdrop-filter: blur(3px);
      display: flex; align-items: center; justify-content: center;
      padding: 1.25rem;
    }
    .ack-overlay[hidden] { display: none; }
    .ack-box {
      background: var(--paper-2); color: var(--text);
      border-radius: 4px; box-shadow: 0 18px 50px rgba(39, 51, 74, 0.34);
      max-width: 560px; width: 100%;
      max-height: 90vh; overflow-y: auto;
      animation: ackIn 0.28s ease-out;
    }
    @keyframes ackIn {
      from { opacity: 0; transform: translateY(10px); }
      to   { opacity: 1; transform: translateY(0); }
    }
    /* Header lockup mirrors the main app header */
    .ack-head {
      background: var(--navy);
      border-bottom: 2px solid var(--gold);
      border-radius: 4px 4px 0 0;
      padding: 1rem 1.75rem;
      display: flex; align-items: center; gap: 1rem;
    }
    .ack-logo { height: 46px; width: auto; display: block; flex-shrink: 0; }
    .ack-divider {
      width: 1px; height: 30px; flex-shrink: 0;
      background: rgba(210, 188, 141, 0.35);
    }
    .ack-tag {
      font-size: 0.8rem; letter-spacing: 0.22em;
      text-transform: uppercase; color: var(--gold);
      white-space: nowrap; overflow: hidden; text-overflow: ellipsis;
    }
    .ack-content { padding: 1.75rem 2rem 1.5rem; }
    .ack-content h2 {
      margin: 0 0 1.1rem 0;
      font-size: 0.85rem; font-weight: 500;
      letter-spacing: 0.14em; text-transform: uppercase;
      color: var(--navy); line-height: 1.5;
      padding-bottom: 0.65rem;
      border-bottom: 1px solid var(--line);
    }
    .ack-text p { margin: 0 0 0.85rem 0; font-size: 0.92rem; line-height: 1.65; }
    .ack-text p:last-child { margin-bottom: 0; }
    .ack-check {
      display: flex; align-items: flex-start; gap: 0.7rem;
      background: var(--paper);
      border: 1px solid var(--line);
      border-left: 3px solid var(--gold);
      border-radius: 2px;
      padding: 0.9rem 1rem;
      margin: 1.25rem 0; cursor: pointer;
      font-size: 0.88rem; line-height: 1.55;
      transition: border-color 0.18s ease;
    }
    .ack-check:hover { border-color: var(--gold); }
    .ack-check input {
      margin: 0.18rem 0 0 0; width: 17px; height: 17px;
      accent-color: var(--navy); flex-shrink: 0; cursor: pointer;
    }
    #ack-continue, #sched-continue {
      width: 100%; padding: 0.85rem 1rem;
      background: var(--navy); color: var(--gold);
      border: 1px solid var(--navy); border-radius: 2px;
      font-family: inherit; font-size: 0.78rem;
      letter-spacing: 0.18em; text-transform: uppercase;
      cursor: pointer; transition: all 0.2s ease;
    }
    #ack-continue:hover:not(:disabled), #sched-continue:hover:not(:disabled) { background: var(--gold); color: var(--navy); }
    #ack-continue:disabled, #sched-continue:disabled { opacity: 0.4; cursor: not-allowed; }
    .ack-secondary {
      width: 100%; margin-top: 0.55rem;
      padding: 0.7rem 1rem;
      background: transparent; color: var(--muted);
      border: 1px solid var(--line); border-radius: 2px;
      font-family: inherit; font-size: 0.72rem;
      letter-spacing: 0.14em; text-transform: uppercase;
      cursor: pointer; transition: all 0.18s ease;
    }
    .ack-secondary:hover { border-color: var(--gold); color: var(--navy); }

    .ack-foot {
      margin: 1.1rem 0 0 0; text-align: center;
      font-size: 0.62rem; color: var(--muted);
      letter-spacing: 0.14em; text-transform: uppercase; line-height: 1.7;
    }
    @media (max-width: 640px) {
      .ack-head { padding: 0.8rem 1.1rem; gap: 0.7rem; }
      .ack-logo { height: 36px; }
      .ack-divider { height: 24px; }
      .ack-tag { font-size: 0.68rem; letter-spacing: 0.18em; }
      .ack-content { padding: 1.35rem 1.25rem 1.25rem; }
      .ack-text p { font-size: 0.88rem; }
    }

    .typing { color: var(--muted); font-style: italic; }

    .typing { color: var(--muted); font-style: italic; }

    /* Rendered markdown inside assistant replies */
    .msg-body { white-space: normal; }
    .msg-body > *:first-child { margin-top: 0; }
    .msg-body > *:last-child { margin-bottom: 0; }
    .msg-body p { margin: 0 0 0.75rem 0; }
    .msg-body h3, .msg-body h4, .msg-body h5, .msg-body h6 {
      margin: 1.15rem 0 0.5rem 0; color: var(--navy);
      font-weight: 500; line-height: 1.35;
    }
    .msg-body h3 { font-size: 1.08rem; letter-spacing: 0.01em; }
    .msg-body h4 { font-size: 0.98rem; }
    .msg-body h5, .msg-body h6 { font-size: 0.92rem; color: var(--muted); }
    .msg-body ul, .msg-body ol { margin: 0 0 0.75rem 0; padding-left: 1.35rem; }
    .msg-body li { margin-bottom: 0.35rem; }
    .msg-body li::marker { color: var(--gold); }
    .msg-body strong { font-weight: 600; color: var(--navy); }
    .msg-body code {
      background: var(--paper); border: 1px solid var(--line);
      border-radius: 3px; padding: 0.08rem 0.32rem;
      font-family: ui-monospace, SFMono-Regular, Menlo, monospace;
      font-size: 0.88em;
    }
    .msg-body a { color: var(--navy); border-bottom: 1px solid var(--gold); text-decoration: none; }
    .msg-body a:hover { color: var(--rust); }
    .msg-body hr { border: none; border-top: 1px solid var(--line); margin: 1rem 0; }
    .feedback {
      display: flex; align-items: center;
      gap: 0.5rem; margin-top: 0.6rem;
      padding-top: 0.6rem; border-top: 1px solid var(--line);
      flex-wrap: wrap;
    }
    .feedback-label {
      font-size: 0.7rem; letter-spacing: 0.12em; text-transform: uppercase;
      color: var(--muted); margin-right: 0.3rem;
    }
    .feedback-btn {
      background: transparent; border: 1px solid var(--line);
      color: var(--muted); width: 30px; height: 30px; border-radius: 50%;
      cursor: pointer; display: flex; align-items: center; justify-content: center;
      padding: 0; transition: all 0.18s ease;
    }
    .feedback-btn svg { width: 14px; height: 14px; }
    .feedback-btn:hover { border-color: var(--gold); color: var(--navy); background: var(--paper); }
    .feedback-btn.selected-up { background: var(--navy); border-color: var(--navy); color: var(--gold); }
    .feedback-btn.selected-down { background: var(--rust); border-color: var(--rust); color: #fff; }
    .feedback-btn:disabled { cursor: default; }
    .feedback-thanks { font-size: 0.7rem; color: var(--muted); margin-left: 0.4rem; font-style: italic; }

    /* Action buttons (copy + share) — labeled pill style */
    .action-sep {
      width: 1px; height: 22px; background: var(--line);
      margin: 0 0.3rem;
    }
    .action-btn {
      background: transparent;
      border: 1px solid var(--line);
      color: var(--muted);
      padding: 0.4rem 0.85rem;
      border-radius: 4px;
      cursor: pointer;
      display: inline-flex;
      align-items: center;
      gap: 0.4rem;
      font-family: inherit;
      font-size: 0.7rem;
      letter-spacing: 0.14em;
      text-transform: uppercase;
      transition: all 0.18s ease;
      white-space: nowrap;
    }
    .action-btn svg { width: 14px; height: 14px; }
    .action-btn:hover {
      border-color: var(--gold);
      color: var(--navy);
      background: var(--paper);
    }
    .action-btn.copied {
      background: var(--navy);
      border-color: var(--navy);
      color: var(--gold);
    }
    .action-btn.copied:hover {
      background: var(--navy); color: var(--gold);
    }
    .action-btn.speaking {
      background: var(--navy);
      border-color: var(--navy);
      color: var(--gold);
      animation: speak-pulse 1.4s ease-in-out infinite;
    }
    .action-btn.speaking:hover { background: var(--navy); color: var(--gold); }
    .action-btn.paused {
      background: var(--paper);
      border-color: var(--gold);
      color: var(--navy);
    }
    @keyframes speak-pulse {
      0%, 100% { box-shadow: 0 0 0 0 rgba(210, 188, 141, 0); }
      50%      { box-shadow: 0 0 0 6px rgba(210, 188, 141, 0.28); }
    }

    /* Share menu popover */
    .share-wrap, .download-wrap {
      position: relative; display: inline-flex;
      align-items: center; line-height: 0;
    }
    .share-menu {
      position: absolute; bottom: calc(100% + 8px); right: 0;
      background: var(--paper-2); border: 1px solid var(--line);
      border-radius: 4px; box-shadow: var(--shadow);
      padding: 0.4rem; min-width: 208px;
      display: none; flex-direction: column; gap: 0.1rem;
      z-index: 10;
      /* Enough entries now that it can exceed the viewport — scroll instead */
      max-height: min(62vh, 420px); overflow-y: auto;
    }
    .share-menu.open { display: flex; }
    /* Flipped when there isn't room above the button */
    .share-menu.drop-down { top: calc(100% + 8px); bottom: auto; }
    .share-menu a, .share-menu button {
      display: flex; align-items: center; gap: 0.6rem;
      padding: 0.5rem 0.7rem; border-radius: 2px;
      background: transparent; border: none; cursor: pointer;
      color: var(--text); font-size: 0.82rem;
      font-family: inherit; text-decoration: none;
      text-align: left; width: 100%;
      letter-spacing: 0; text-transform: none;
      white-space: nowrap;
    }
    .share-menu a:hover, .share-menu button:hover { background: var(--paper); color: var(--navy); }
    .share-menu svg { width: 16px; height: 16px; flex-shrink: 0; color: var(--muted); }

    /* Persistent avatar presence — always on screen, always animating */
    .presence {
      position: fixed; right: 1.5rem; bottom: 8.5rem; z-index: 400;
      width: 104px; text-align: center; user-select: none;
    }
    .presence-frame {
      position: relative; width: 104px; height: 104px; cursor: pointer;
      border: none; background: none; padding: 0; border-radius: 50%;
    }
    .presence-frame:focus-visible { outline: 2px solid var(--gold); outline-offset: 4px; }
    .presence-photo {
      width: 100%; height: 100%; border-radius: 50%; object-fit: cover;
      border: 2px solid var(--gold);
      box-shadow: 0 6px 20px rgba(39,51,74,0.28);
      animation: presence-breathe 6s ease-in-out infinite;
      display: block;
    }
    @keyframes presence-breathe {
      0%, 100% { transform: scale(1); }
      50%      { transform: scale(1.028); }
    }
    /* Thinking: ring travelling round the frame */
    .presence-loop {
      position: absolute; inset: 0; width: 100%; height: 100%;
      border-radius: 50%; object-fit: cover; border: 2px solid var(--gold);
      box-shadow: 0 6px 20px rgba(39,51,74,0.28);
      display: block;
    }
    /* A generated talking-head video replaces the loop while it plays */
    .presence.video .presence-loop { display: none; }
    /* The looping portrait already moves, so drop the CSS breath */
    .presence-loop ~ .presence-photo,
    .presence:has(.presence-loop) .presence-photo { animation: none; }
    .presence.speaking .presence-loop {
      box-shadow: 0 0 0 3px rgba(210,188,141,0.5), 0 6px 20px rgba(39,51,74,0.28);
      animation: presence-talk 0.62s ease-in-out infinite;
    }
    .presence.responding .presence-loop { animation: presence-talk 1.1s ease-in-out infinite; }
    .presence.thinking .presence-loop { opacity: 0.9; }
    /* When the loop can't play, the still photo breathes instead */
    .presence.no-loop .presence-photo {
      animation: presence-breathe 6s ease-in-out infinite !important;
    }

    .presence-ring {
      position: absolute; inset: -5px; border-radius: 50%;
      border: 2px solid transparent; display: none; pointer-events: none;
    }
    .presence.thinking .presence-ring {
      display: block;
      border-top-color: var(--gold);
      border-right-color: rgba(210,188,141,0.35);
      animation: av-think 0.9s linear infinite;
    }
    .presence.thinking .presence-photo { animation: none; opacity: 0.9; }
    /* Speaking: concentric pulses */
    .presence-pulse {
      position: absolute; inset: 0; border-radius: 50%;
      border: 2px solid var(--gold); opacity: 0; pointer-events: none;
    }
    .presence.speaking .presence-pulse { animation: av-pulse 1.7s ease-out infinite; }
    .presence.speaking .presence-pulse:nth-of-type(3) { animation-delay: 0.55s; }
    .presence.speaking .presence-pulse:nth-of-type(4) { animation-delay: 1.1s; }
    .presence.speaking .presence-photo {
      animation: presence-talk 0.62s ease-in-out infinite;
      box-shadow: 0 0 0 3px rgba(210,188,141,0.5), 0 6px 20px rgba(39,51,74,0.28);
    }
    @keyframes presence-talk {
      0%, 100% { transform: scale(1) translateY(0); }
      50%      { transform: scale(1.035) translateY(-1px); }
    }
    /* Responding: nod while the reply lands */
    .presence.responding .presence-photo {
      animation: presence-talk 1.1s ease-in-out infinite;
    }
    .presence-name {
      margin-top: 0.5rem; font-size: 0.72rem; font-weight: 500;
      line-height: 1.25; color: var(--navy);
    }
    .presence-status {
      margin-top: 0.15rem; font-size: 0.56rem; letter-spacing: 0.12em;
      text-transform: uppercase; color: var(--muted); min-height: 1em;
    }
    .presence.speaking .presence-status,
    .presence.thinking .presence-status { color: var(--navy); }
    /* A video takes over the frame when a talking head is generated */
    .presence-video {
      position: absolute; inset: 0; width: 100%; height: 100%;
      border-radius: 50%; object-fit: cover; display: none;
      border: 2px solid var(--gold);
    }
    .presence.video .presence-video { display: block; }
    .presence.video .presence-photo { visibility: hidden; }

    @media (max-width: 900px) {
      /* Smaller, and above the composer — the send button owns the bottom
         right at these widths. */
      /* Clear of the composer row, which reaches ~13rem from the bottom here */
      .presence { right: 0.7rem; bottom: 14rem; width: 76px; }
      .presence-frame { width: 60px; height: 60px; margin: 0 auto; }
      .presence-name { font-size: 0.62rem; }
      .presence-status { display: none; }
    }
    @media (max-width: 640px) {
      /* Hidden on phones: there isn't room without covering the conversation,
         and every reply already carries its own avatar there. */
      .presence { display: none; }
    }
    @media (prefers-reduced-motion: reduce) {
      .presence-photo, .presence-pulse, .presence-ring { animation: none !important; }
      .presence.speaking .presence-ring,
      .presence.thinking .presence-ring { display: block; border: 2px solid var(--gold); }
    }

    /* End-of-session offer of a follow-up plan */
    .plan-offer {
      display: flex; align-items: center; gap: 0.6rem; flex-wrap: wrap;
      margin-top: 0.8rem; padding: 0.7rem 0.85rem;
      background: var(--paper); border: 1px solid var(--line);
      border-left: 3px solid var(--gold); border-radius: 3px;
      font-size: 0.85rem; color: var(--navy);
      animation: nudge-in 0.35s ease;
    }
    .plan-offer span { flex: 1 1 240px; }
    .plan-offer button {
      background: var(--navy); color: var(--gold); border: 1px solid var(--navy);
      border-radius: 2px; padding: 0.45rem 0.8rem; cursor: pointer;
      font-family: inherit; font-size: 0.66rem; letter-spacing: 0.1em;
      text-transform: uppercase;
    }
    .plan-offer button.secondary {
      background: transparent; color: var(--muted); border-color: var(--line);
    }
    .plan-offer button:hover { background: var(--gold); color: var(--navy); }
    .plan-offer button:disabled { opacity: 0.55; cursor: not-allowed; }

    /* Gentle nudge to rate the last reply when a session looks finished */
    .rate-nudge {
      display: flex; align-items: center; gap: 0.5rem; flex-wrap: wrap;
      margin-top: 0.8rem; padding: 0.55rem 0.75rem;
      background: var(--paper); border: 1px solid var(--line);
      border-left: 3px solid var(--gold); border-radius: 3px;
      font-size: 0.82rem; color: var(--navy);
      animation: nudge-in 0.35s ease;
    }
    @keyframes nudge-in { from { opacity: 0; transform: translateY(4px); } to { opacity: 1; } }
    .rate-nudge .nudge-dismiss {
      margin-left: auto; background: none; border: none; cursor: pointer;
      color: var(--muted); font-size: 1rem; line-height: 1; padding: 0 0.2rem;
    }
    .rate-nudge .nudge-dismiss:hover { color: var(--navy); }
    /* Draw the eye to the thumbs on the row above */
    .feedback.nudged .feedback-btn {
      border-color: var(--gold);
      box-shadow: 0 0 0 3px rgba(210,188,141,0.28);
    }

    /* Opt-in download offer under a reply — nothing saves without a click */
    .export-offer {
      display: flex; align-items: center; flex-wrap: wrap; gap: 0.5rem;
      margin-top: 0.85rem; padding: 0.6rem 0.75rem;
      background: var(--paper); border: 1px solid var(--line);
      border-left: 3px solid var(--gold); border-radius: 3px;
    }
    .export-offer-label {
      font-size: 0.82rem; color: var(--navy); margin-right: 0.2rem;
    }
    .export-offer-chips { display: flex; flex-wrap: wrap; gap: 0.35rem; }
    .export-chip {
      background: var(--paper-2); border: 1px solid var(--line);
      color: var(--muted); border-radius: 2px;
      padding: 0.3rem 0.62rem; cursor: pointer;
      font-family: inherit; font-size: 0.72rem; letter-spacing: 0.02em;
      transition: all 0.15s ease;
    }
    .export-chip:hover:not(:disabled) {
      border-color: var(--gold); color: var(--navy);
    }
    .export-chip.primary {
      background: var(--navy); border-color: var(--navy); color: var(--gold);
    }
    .export-chip.primary:hover:not(:disabled) {
      background: var(--gold); color: var(--navy); border-color: var(--gold);
    }
    .export-chip.dismiss { border-style: dashed; }
    .export-chip:disabled { opacity: 0.55; cursor: default; }
    @media (max-width: 640px) {
      .export-offer { flex-direction: column; align-items: flex-start; }
    }

    /* Multi-document SAVE menu: one row per deliverable */
    .share-menu.multi-doc { min-width: 268px; padding: 0.55rem; }
    .doc-group { padding: 0.5rem 0.55rem 0.6rem; }
    .doc-group + .doc-group { border-top: 1px solid var(--line); }
    .doc-title {
      font-size: 0.78rem; font-weight: 500; color: var(--navy);
      margin-bottom: 0.45rem; letter-spacing: 0;
      text-transform: none; line-height: 1.35;
      overflow: hidden; text-overflow: ellipsis; white-space: nowrap;
    }
    .doc-formats { display: flex; flex-wrap: wrap; gap: 0.3rem; }
    .share-menu .fmt-chip {
      display: inline-flex; align-items: center; justify-content: center;
      width: auto; padding: 0.32rem 0.6rem;
      background: var(--paper); border: 1px solid var(--line);
      border-radius: 2px; cursor: pointer;
      font-family: inherit; font-size: 0.7rem; color: var(--muted);
      letter-spacing: 0.04em; text-transform: none;
    }
    .share-menu .fmt-chip:hover { border-color: var(--gold); color: var(--navy); background: var(--paper-2); }
    .share-menu .fmt-chip.primary {
      background: var(--navy); border-color: var(--navy); color: var(--gold);
    }
    .share-menu .fmt-chip.primary:hover { background: var(--gold); color: var(--navy); border-color: var(--gold); }

    /* On the feedback row, push actions to the right */
    .feedback-actions {
      display: flex;
      align-items: center;
      gap: 0.5rem;
    }
    .feedback-comment {
      margin-top: 0.7rem;
      padding-top: 0.7rem;
      border-top: 1px dashed var(--line);
      display: flex;
      flex-direction: column;
      gap: 0.5rem;
    }
    .feedback-comment label {
      font-size: 0.72rem;
      letter-spacing: 0.1em;
      text-transform: uppercase;
      color: var(--muted);
    }
    .feedback-comment textarea {
      width: 100%;
      padding: 0.6rem 0.75rem;
      border: 1px solid var(--line);
      border-radius: 4px;
      font-family: inherit;
      font-size: 0.9rem;
      outline: none;
      resize: vertical;
      min-height: 64px;
      background: var(--paper);
      color: var(--text);
      transition: border-color 0.15s ease, box-shadow 0.15s ease;
    }
    .feedback-comment textarea:focus {
      border-color: var(--rust);
      box-shadow: 0 0 0 3px rgba(157, 67, 44, 0.12);
    }
    .feedback-comment-actions {
      display: flex;
      gap: 0.5rem;
      align-items: center;
    }
    .feedback-comment-btn {
      background: var(--rust);
      color: #fff;
      border: 1px solid var(--rust);
      border-radius: 2px;
      padding: 0.45rem 0.95rem;
      font-size: 0.7rem;
      font-family: inherit;
      letter-spacing: 0.14em;
      text-transform: uppercase;
      cursor: pointer;
      transition: opacity 0.15s ease;
    }
    .feedback-comment-btn:hover:not(:disabled) { opacity: 0.85; }
    .feedback-comment-btn.secondary {
      background: transparent;
      color: var(--muted);
      border-color: var(--line);
    }
    .feedback-comment-btn:disabled { opacity: 0.5; cursor: default; }
    .composer-wrap { background: var(--paper-2); border-top: 1px solid var(--line); }
    form { display: flex; gap: 0.6rem; padding: 1rem 1.5rem; max-width: 760px; margin: 0 auto; }
    .input-wrap { flex: 1; position: relative; display: flex; align-items: center; }
    .attach-btn {
      position: absolute; right: 3rem; top: 50%; transform: translateY(-50%);
      background: transparent; border: none;
      color: var(--muted); cursor: pointer;
      width: 32px; height: 32px; border-radius: 50%;
      display: flex; align-items: center; justify-content: center;
      transition: all 0.2s ease; padding: 0;
    }
    .attach-btn:hover { color: var(--navy); background: var(--paper); }
    .attach-btn svg { width: 18px; height: 18px; }
    .folder-btn {
      position: absolute; right: 5rem; top: 50%; transform: translateY(-50%);
      background: transparent; border: none;
      color: var(--muted); cursor: pointer;
      width: 32px; height: 32px; border-radius: 50%;
      display: flex; align-items: center; justify-content: center;
      transition: all 0.2s ease; padding: 0;
    }
    .folder-btn:hover { color: var(--navy); background: var(--paper); }
    .folder-btn svg { width: 18px; height: 18px; }
    #file-input, #folder-input-chat { display: none; }
    .attached-file {
      display: none; align-items: center; gap: 0.5rem;
      background: var(--paper); border: 1px solid var(--line);
      border-radius: 4px; padding: 0.4rem 0.6rem;
      margin: 0 1.75rem 0.5rem; font-size: 0.82rem;
      color: var(--navy);
    }
    .attached-file.visible { display: inline-flex; }
    .attached-file svg { width: 14px; height: 14px; color: var(--muted); flex-shrink: 0; }
    .attached-file .filename { max-width: 240px; overflow: hidden; text-overflow: ellipsis; white-space: nowrap; }
    .attached-file .filesize { color: var(--muted); font-size: 0.72rem; }
    .attached-file button {
      background: transparent; border: none; cursor: pointer;
      color: var(--muted); padding: 2px; display: flex;
      align-items: center; justify-content: center;
    }
    .attached-file button:hover { color: var(--rust); }
    .attached-file button svg { width: 14px; height: 14px; }
    input[type="text"] {
      flex: 1; padding: 0.85rem 7.6rem 0.85rem 1.1rem;
      border: 1px solid var(--line); border-radius: 2px;
      font-size: 0.95rem; font-family: inherit; outline: none;
      background: var(--paper); color: var(--text); width: 100%;
      transition: border-color 0.15s ease, box-shadow 0.15s ease;
    }
    input[type="text"]:focus { border-color: var(--gold); box-shadow: 0 0 0 3px rgba(210, 188, 141, 0.18); }
    .mic-btn {
      position: absolute; right: 6px; top: 50%; transform: translateY(-50%);
      width: 38px; height: 38px; display: flex; align-items: center;
      justify-content: center; background: var(--navy); color: var(--gold);
      border: none; border-radius: 50%; cursor: pointer; padding: 0;
      transition: all 0.2s ease;
    }
    .mic-btn:hover { background: var(--gold); color: var(--navy); }
    .mic-btn svg { width: 18px; height: 18px; }
    .mic-btn.recording { background: var(--rust); color: #fff; animation: pulse 1.2s ease-in-out infinite; }
    .mic-btn.unsupported { display: none; }
    /* Hover tooltips on the attach and folder buttons */
    .input-tip {
      position: absolute; bottom: calc(100% + 10px);
      background: var(--navy); color: var(--paper);
      padding: 0.4rem 0.7rem; border-radius: 3px;
      font-size: 0.72rem; letter-spacing: 0.02em; white-space: nowrap;
      box-shadow: var(--shadow);
      opacity: 0; visibility: hidden;
      transform: translateY(4px);
      transition: opacity 0.15s ease, transform 0.15s ease, visibility 0.15s;
      pointer-events: none; z-index: 20;
    }
    .input-tip::after {
      content: ""; position: absolute; top: 100%;
      border: 5px solid transparent; border-top-color: var(--navy);
    }
    /* Anchored under each button, which sits at a fixed offset from the right */
    .folder-tip { right: 4.2rem; }
    .folder-tip::after { right: 18px; }
    .attach-tip { right: 2.2rem; }
    .attach-tip::after { right: 18px; }
    .folder-btn:hover ~ .folder-tip,
    .folder-btn:focus-visible ~ .folder-tip,
    .attach-btn:hover ~ .attach-tip,
    .attach-btn:focus-visible ~ .attach-tip {
      opacity: 1; visibility: visible; transform: translateY(0);
    }
    /* The recording bar takes priority once a hold starts */
    .mic-btn.recording ~ .input-tip,
    .mic-btn.cancelling ~ .input-tip { opacity: 0; visibility: hidden; }
    @media (max-width: 640px) {
      .input-tip { font-size: 0.68rem; padding: 0.35rem 0.6rem; white-space: normal; max-width: 60vw; }
    }

    /* Hover tooltip on the mic — appears instantly, unlike the native title */
    .mic-tip {
      position: absolute; right: 0; bottom: calc(100% + 10px);
      background: var(--navy); color: var(--paper);
      padding: 0.4rem 0.7rem; border-radius: 3px;
      font-size: 0.72rem; letter-spacing: 0.02em; white-space: nowrap;
      box-shadow: var(--shadow);
      opacity: 0; visibility: hidden;
      transform: translateY(4px);
      transition: opacity 0.15s ease, transform 0.15s ease, visibility 0.15s;
      pointer-events: none; z-index: 20;
    }
    .mic-tip::after {
      content: ""; position: absolute; top: 100%; right: 14px;
      border: 5px solid transparent; border-top-color: var(--navy);
    }
    .mic-btn:hover ~ .mic-tip,
    .mic-btn:focus-visible ~ .mic-tip {
      opacity: 1; visibility: visible; transform: translateY(0);
    }
    /* Once recording starts the hint bar takes over */
    .mic-btn.recording ~ .mic-tip,
    .mic-btn.cancelling ~ .mic-tip { opacity: 0; visibility: hidden; }
    @media (max-width: 640px) {
      .mic-tip { font-size: 0.68rem; padding: 0.35rem 0.6rem; }
    }

    /* Press-and-hold voice recording */
    .mic-btn { touch-action: none; -webkit-user-select: none; user-select: none; }
    .mic-btn.cancelling { background: var(--muted); color: #fff; animation: none; }
    .voice-hint {
      display: none; align-items: center; gap: 0.55rem;
      margin: 0 1.75rem 0.5rem; padding: 0.5rem 0.8rem;
      background: var(--navy); color: var(--paper);
      border-radius: 4px; font-size: 0.8rem;
      box-shadow: var(--shadow);
    }
    .voice-hint.visible { display: inline-flex; }
    .voice-hint.cancel { background: var(--rust); }
    .voice-dot {
      width: 9px; height: 9px; border-radius: 50%;
      background: var(--gold); flex-shrink: 0;
      animation: voicePulse 1.1s ease-in-out infinite;
    }
    .voice-hint.cancel .voice-dot { background: #fff; animation: none; }
    @keyframes voicePulse {
      0%, 100% { opacity: 1; transform: scale(1); }
      50%      { opacity: 0.35; transform: scale(0.75); }
    }
    .voice-timer {
      margin-left: auto; font-variant-numeric: tabular-nums;
      font-size: 0.75rem; color: var(--gold); letter-spacing: 0.06em;
    }
    .voice-hint.cancel .voice-timer { color: #fff; }
    @media (max-width: 640px) {
      .voice-hint { margin: 0 1rem 0.45rem; font-size: 0.76rem; }
    }

    @keyframes pulse {
      0%, 100% { box-shadow: 0 0 0 0 rgba(157, 67, 44, 0.6); }
      50% { box-shadow: 0 0 0 8px rgba(157, 67, 44, 0); }
    }
    button[type="submit"] {
      padding: 0.85rem 1.75rem; background: var(--navy); color: var(--gold);
      border: 1px solid var(--navy); border-radius: 2px;
      font-size: 0.78rem; font-family: inherit; letter-spacing: 0.18em;
      text-transform: uppercase; cursor: pointer; transition: all 0.2s ease;
    }
    button[type="submit"]:hover:not(:disabled) { background: var(--gold); color: var(--navy); }
    button[type="submit"]:disabled { opacity: 0.5; cursor: not-allowed; }
    /* Scheduling call-to-action — the main conversion point, so it reads as
       a button rather than fine print inside the disclaimer line. */
    .footer-cta {
      display: flex; justify-content: center;
      padding: 0.4rem 1rem 0.8rem;
    }
    .cta-btn {
      display: inline-flex; align-items: center; gap: 0.6rem;
      background: var(--gold); color: var(--navy);
      border: 1px solid var(--gold); border-radius: 2px;
      padding: 0.72rem 1.6rem;
      font-family: inherit; font-size: 0.82rem; font-weight: 500;
      letter-spacing: 0.16em; text-transform: uppercase;
      text-decoration: none; white-space: nowrap;
      box-shadow: 0 2px 10px rgba(210, 188, 141, 0.45);
      transition: background 0.2s ease, color 0.2s ease,
                  box-shadow 0.2s ease, transform 0.12s ease;
    }
    .cta-btn { cursor: pointer; user-select: none; }
    .cta-btn svg { width: 17px; height: 17px; flex-shrink: 0; }
    .cta-btn:hover {
      background: var(--navy); color: var(--gold);
      border-color: var(--navy);
      box-shadow: 0 4px 16px rgba(39, 51, 74, 0.28);
    }
    .cta-btn:active { transform: translateY(1px); }
    .cta-btn:focus-visible { outline: 2px solid var(--navy); outline-offset: 3px; }

    .footer-note {
      text-align: center; font-size: 0.68rem; color: var(--muted);
      padding: 0 1rem 0.9rem; letter-spacing: 0.14em;
      text-transform: uppercase; line-height: 1.7;
    }
    .footer-ai-note { display: block; margin-top: 0.3rem; }
    .footer-note a { color: var(--navy); text-decoration: none; border-bottom: 1px solid var(--gold); }
    .footer-note a:hover { color: var(--rust); }
    @media (max-width: 640px) {
      .user { margin-left: 8%; } .assistant { margin-right: 6%; }
      header { padding: 0.75rem 0.9rem; gap: 0.5rem; }
      .brand-logo { height: 40px; }
      .brand-tag { font-size: 0.7rem; letter-spacing: 0.18em; }
      .brand { gap: 0.6rem; } .brand-divider { height: 26px; }
      header button { padding: 0.45rem 0.7rem; font-size: 0.68rem; letter-spacing: 0.1em; }
      #chat { padding: 1.5rem 1rem 0.75rem; }
      form { padding: 0.75rem 1rem; gap: 0.4rem; }
      input[type="text"] { padding: 0.75rem 7.2rem 0.75rem 0.9rem; font-size: 16px; }
      button[type="submit"] { padding: 0.75rem 1rem; font-size: 0.7rem; letter-spacing: 0.12em; }
      .footer-note { font-size: 0.62rem; letter-spacing: 0.1em; }
      .footer-cta { padding: 0.3rem 0.75rem 0.65rem; }
      .cta-btn {
        font-size: 0.72rem; letter-spacing: 0.12em;
        padding: 0.7rem 1.1rem; gap: 0.45rem;
        white-space: normal; text-align: center;
      }
      .cta-btn svg { width: 15px; height: 15px; }
    }
    /* Very narrow phones: drop the persona tag, and stack a short caption under
       each header icon. Icons alone were ambiguous — the voice button in
       particular read as an unlabelled box. */
    @media (max-width: 480px) {
      /* Keep the J3P ADVISOR wordmark beside the logo, as on desktop.
         It shrinks rather than disappearing. */
      .brand-divider { display: block; height: 22px; }
      .brand-tag { display: inline; font-size: 0.56rem; letter-spacing: 0.1em; }
      .brand { gap: 0.45rem; min-width: 0; }
      .brand-logo { height: 34px; }
      header { gap: 0.35rem; padding: 0.6rem 0.7rem; }
      .reset-icon { display: inline-block; }
      header button {
        padding: 0.35rem 0.45rem; min-width: 46px; min-height: 42px;
        display: inline-flex; flex-direction: column; align-items: center;
        justify-content: center; gap: 0.12rem; line-height: 1;
      }
      header button svg { width: 16px; height: 16px; }
      .autospeak-label, .voice-label, .reset-label {
        display: block; font-size: 0.5rem; letter-spacing: 0.04em;
        text-transform: uppercase; white-space: nowrap;
      }
      .reset-full { display: none; }
      .reset-short { display: inline; }
      /* Hide action button text labels — keep icons only */
      .action-btn span { display: none; }
      .action-btn { padding: 0.45rem 0.55rem; }
    }
  </style>
</head>
<body>
  <div id="ack-overlay" class="ack-overlay" role="dialog" aria-modal="true"
       aria-labelledby="ack-title" aria-describedby="ack-body" hidden>
    <div class="ack-box">
      <div class="ack-head">
        <img src="{{ cfg.logo_url }}" alt="{{ cfg.persona_name }}" class="ack-logo" />
        <span class="ack-divider"></span>
        <span class="ack-tag">{{ cfg.persona_name }}</span>
      </div>
      <div class="ack-content">
        <h2 id="ack-title">{{ release_heading|safe }}</h2>
        <div id="ack-body" class="ack-text">
          {{ release_body|safe }}
        </div>
        <label class="ack-check" for="ack-checkbox">
          <input type="checkbox" id="ack-checkbox" />
          <span>{{ release_checkbox_label }}</span>
        </label>
        <button type="button" id="ack-continue" disabled>Enter session</button>
        <p class="ack-foot">{{ cfg.footer_disclaimer }}</p>
      </div>
    </div>
  </div>

  <div id="sched-overlay" class="ack-overlay" role="dialog" aria-modal="true"
       aria-labelledby="sched-title" aria-describedby="sched-body" hidden>
    <div class="ack-box">
      <div class="ack-head">
        <img src="{{ cfg.logo_url }}" alt="{{ cfg.persona_name }}" class="ack-logo" />
        <span class="ack-divider"></span>
        <span class="ack-tag">Schedule a session</span>
      </div>
      <div class="ack-content">
        <h2 id="sched-title">{{ release_heading|safe }}</h2>
        <div id="sched-body" class="ack-text">
          {{ release_body|safe }}
          <p style="margin: 0.9rem 0 0;">
            <strong>So your time is well spent:</strong> when you continue, a
            short summary of this conversation is shared with the advisor you
            are booking, so they can prepare and you needn't start over.
          </p>
        </div>
        <label class="ack-check" for="sched-checkbox">
          <input type="checkbox" id="sched-checkbox" />
          <span>{{ release_checkbox_label }}</span>
        </label>
        <button type="button" id="sched-continue" disabled>Continue to scheduling</button>
        <button type="button" id="sched-cancel" class="ack-secondary">Not now</button>
      </div>
    </div>
  </div>

  <header>
    <div class="brand">
      <img src="{{ cfg.logo_url }}" alt="{{ cfg.persona_name }}" class="brand-logo" />
      <span class="brand-divider"></span>
      <span class="brand-tag">{{ cfg.persona_name }}</span>
    </div>
    <button id="autospeak-btn" aria-label="Toggle speak mode" title="Speak — read every response aloud">
      <svg class="autospeak-icon" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round" aria-hidden="true">
        <polygon points="11 5 6 9 2 9 2 15 6 15 11 19 11 5"/>
        <path d="M15.54 8.46a5 5 0 0 1 0 7.07"/>
        <path d="M19.07 4.93a10 10 0 0 1 0 14.14"/>
      </svg>
      <span class="autospeak-label">Speak</span>
    </button>
    <div class="voice-wrap">
      <button id="voice-btn" aria-label="Voice settings" aria-haspopup="true"
              title="Choose the reading voice and speed">
        <svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2"
             stroke-linecap="round" stroke-linejoin="round" aria-hidden="true">
          <line x1="4" y1="21" x2="4" y2="14"/><line x1="4" y1="10" x2="4" y2="3"/>
          <line x1="12" y1="21" x2="12" y2="12"/><line x1="12" y1="8" x2="12" y2="3"/>
          <line x1="20" y1="21" x2="20" y2="16"/><line x1="20" y1="12" x2="20" y2="3"/>
          <line x1="1" y1="14" x2="7" y2="14"/><line x1="9" y1="8" x2="15" y2="8"/>
          <line x1="17" y1="16" x2="23" y2="16"/>
        </svg>
        <span class="voice-label">Voice</span>
      </button>
      <div class="voice-menu" id="voice-menu" role="menu">
        <div class="voice-row">
          <label for="voice-select">Reading voice</label>
          <select id="voice-select"></select>
        </div>
        <div class="voice-row">
          <label for="rate-range">Speed <span id="rate-val">1.0&times;</span></label>
          <input type="range" id="rate-range" min="0.7" max="1.4" step="0.1" value="1" />
        </div>
        <button type="button" id="voice-preview" class="voice-preview">Preview voice</button>
        <p class="voice-note">Uses the voices installed on this device.</p>
      </div>
    </div>
    <button id="reset-btn" aria-label="Start a new conversation" title="New conversation">
      <svg class="reset-icon" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round" aria-hidden="true">
        <path d="M12 5v14M5 12h14"/>
      </svg>
      <span class="reset-label"><span class="reset-full">New conversation</span><span class="reset-short">New chat</span></span>
    </button>
  </header>

  <div id="chat-wrap">
    <div id="chat">
      <div class="msg assistant">{{ cfg.opening }}</div>
    </div>
  </div>

  <div class="composer-wrap">
    <div id="voice-hint" class="voice-hint" aria-live="polite">
      <span class="voice-dot"></span>
      <span id="voice-hint-text">Listening&hellip; release to send</span>
      <span id="voice-timer" class="voice-timer">0:00</span>
    </div>
    <div id="attached-file" class="attached-file" aria-live="polite">
      <svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M14 2H6a2 2 0 0 0-2 2v16a2 2 0 0 0 2 2h12a2 2 0 0 0 2-2V8z"/><polyline points="14 2 14 8 20 8"/></svg>
      <span class="filename" id="attached-file-name">document.pdf</span>
      <span class="filesize" id="attached-file-size"></span>
      <button type="button" id="remove-file-btn" aria-label="Remove attachment" title="Remove">
        <svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><line x1="18" y1="6" x2="6" y2="18"/><line x1="6" y1="6" x2="18" y2="18"/></svg>
      </button>
    </div>
    <form id="chat-form">
      <div class="input-wrap">
        <input type="text" id="message" placeholder="{{ cfg.placeholder }}" autocomplete="off" autofocus />
        <input type="file" id="file-input" accept=".pdf,.docx,.doc,.xlsx,.xlsm,.xls,.pptx,.ppt,.csv,.tsv,.txt,.md,.rtf,.jpg,.jpeg,.png,.gif,.webp" multiple />
        <input type="file" id="folder-input-chat" webkitdirectory directory multiple />
        <button type="button" id="folder-btn" class="folder-btn" aria-label="Attach folder" title="Attach a folder of documents">
          <svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
            <path d="M22 19a2 2 0 0 1-2 2H4a2 2 0 0 1-2-2V5a2 2 0 0 1 2-2h5l2 3h9a2 2 0 0 1 2 2z"/>
          </svg>
        </button>
        <span class="input-tip folder-tip" id="folder-tip" role="tooltip">Attach a whole folder of documents</span>
        <button type="button" id="attach-btn" class="attach-btn" aria-label="Attach file" title="Attach a document or image">
          <svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
            <path d="M21.44 11.05 12.25 20.24a6 6 0 0 1-8.49-8.49l9.19-9.19a4 4 0 0 1 5.66 5.66l-9.2 9.19a2 2 0 0 1-2.83-2.83l8.49-8.48"/>
          </svg>
        </button>
        <span class="input-tip attach-tip" id="attach-tip" role="tooltip">Attach a document or image</span>
        <button type="button" id="mic-btn" class="mic-btn" aria-label="Hold to record a voice message" title="Press and hold to record — release to send">
          <svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
            <path d="M12 1a3 3 0 0 0-3 3v8a3 3 0 0 0 6 0V4a3 3 0 0 0-3-3z"/>
            <path d="M19 10v2a7 7 0 0 1-14 0v-2"/>
            <line x1="12" y1="19" x2="12" y2="23"/>
            <line x1="8" y1="23" x2="16" y2="23"/>
          </svg>
        </button>
        <span class="mic-tip" id="mic-tip" role="tooltip">Press and hold to record</span>
      </div>
      <button type="submit" id="send-btn">Send</button>
    </form>
    {% if show_scheduling_button %}
    <div class="footer-cta">
      <a class="cta-btn" role="button" tabindex="0"
         data-cta-url="{{ cfg.footer_cta_url }}"
         aria-haspopup="dialog">
        <svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2"
             stroke-linecap="round" stroke-linejoin="round" aria-hidden="true">
          <rect x="3" y="4" width="18" height="18" rx="2"/>
          <line x1="16" y1="2" x2="16" y2="6"/><line x1="8" y1="2" x2="8" y2="6"/>
          <line x1="3" y1="10" x2="21" y2="10"/>
        </svg>
        <span>{{ cfg.footer_cta_label }}</span>
      </a>
    </div>
    {% endif %}
    {% if allow_materials %}
    <div id="queue-hint" hidden
         style="text-align: center; margin-top: 0.4rem; font-size: 0.72rem;
                color: var(--muted); letter-spacing: 0.04em;"></div>
    <div style="text-align: center; margin-top: 0.55rem;">
      <button type="button" class="materials-link" id="materials-open">
        Add your documents &amp; writing
      </button>
    </div>
    {% endif %}
    <div class="footer-note">
      {{ cfg.footer_disclaimer }}
      <span class="footer-ai-note">{{ cfg.footer_ai_note }}</span>
    </div>
  </div>

  {% if show_avatar and cfg.avatar_url %}
  <div class="presence" id="presence">
    <button type="button" class="presence-frame" id="presence-frame"
            aria-label="Read the latest reply aloud">
      <img class="presence-photo" src="{{ cfg.avatar_url }}" alt=""
           onerror="document.getElementById('presence').style.display='none'" />
      {% if cfg.avatar_loop_url %}
      <video class="presence-loop" id="presence-loop" muted loop playsinline
             autoplay preload="auto" aria-hidden="true"
             onerror="this.style.display='none'">
        <source src="{{ cfg.avatar_loop_url|replace('.mp4', '.webm') }}" type="video/webm" />
        <source src="{{ cfg.avatar_loop_url }}" type="video/mp4" />
      </video>
      {% endif %}
      <video class="presence-video" id="presence-video" playsinline></video>
      <span class="presence-ring"></span>
      <span class="presence-pulse"></span>
      <span class="presence-pulse"></span>
      <span class="presence-pulse"></span>
    </button>
    <div class="presence-name">{{ cfg.avatar_name or cfg.persona_name }}</div>
    <div class="presence-status" id="presence-status">Listening</div>
  </div>
  {% endif %}

  {% if allow_materials %}
  <div class="mat-overlay" id="mat-overlay" hidden>
    <div class="mat-box" role="dialog" aria-modal="true" aria-labelledby="mat-title">
      <div class="mat-head">
        <h3 id="mat-title">Your documents &amp; writing</h3>
        <button type="button" class="mat-close" id="mat-close"
                aria-label="Close">&times;</button>
      </div>
      <div class="mat-body">
        <div id="mat-profile" class="mat-profile" hidden>
          <span class="mat-profile-label">The advisor remembers</span>
          <span id="mat-profile-text"></span>
          <button type="button" class="mat-remove" id="mat-profile-forget">Clear</button>
        </div>

        <p class="mat-intro">
          Add your CV, a strategic plan, a talk, an article you've written —
          anything that helps the advisor understand your work and match your
          voice. These stay <strong>private to you</strong> and are used only in
          your own sessions.
        </p>

        <div class="mat-tabs">
          <button type="button" class="mat-tab active" data-pane="upload">Upload files</button>
          <button type="button" class="mat-tab" data-pane="paste">Paste writing</button>
        </div>

        <div class="mat-pane" id="mat-pane-upload">
          <input type="file" id="mat-files" class="mat-field" multiple
                 accept=".pdf,.docx,.doc,.pptx,.xlsx,.csv,.txt,.md,.rtf" />
          <label class="mat-share">
            <input type="checkbox" id="mat-share-upload" />
            <span>Also share with J3P so the team can review it. Leave unticked
              to keep it entirely private to your sessions.</span>
          </label>
          <button type="button" class="btn" id="mat-upload-btn">Add to my library</button>
        </div>

        <div class="mat-pane" id="mat-pane-paste" hidden>
          <input type="text" id="mat-title-input" class="mat-field"
                 placeholder="Title (e.g. My leadership philosophy)" />
          <textarea id="mat-text" class="mat-field"
                    placeholder="Paste an article, notes, a bio, a draft…"></textarea>
          <label class="mat-share">
            <input type="checkbox" id="mat-share-text" />
            <span>Also share with J3P so the team can review it.</span>
          </label>
          <button type="button" class="btn" id="mat-save-btn">Add to my library</button>
        </div>

        <p class="mat-msg" id="mat-msg"></p>

        <div class="mat-list" id="mat-list"></div>
      </div>
    </div>
  </div>
  {% endif %}

  <script>
    const chat = document.getElementById("chat");

    // -------------------------------------------------------------
    // Acknowledgment gate — required at the start of every session
    // -------------------------------------------------------------
    // sessionStorage (not localStorage) is deliberate: acknowledgment is scoped
    // to the browsing session, so closing the tab and returning requires it
    // again. A reload mid-session won't re-prompt. Bump the key if the
    // disclaimer wording changes.
    const ACK_KEY = "j3p_ack_v2";
    (function initAckGate() {
      const overlay = document.getElementById("ack-overlay");
      const checkbox = document.getElementById("ack-checkbox");
      const continueBtn = document.getElementById("ack-continue");
      if (!overlay) return;

      let alreadyAcked = false;
      try {
        alreadyAcked = !!sessionStorage.getItem(ACK_KEY);
      } catch (e) {
        // Storage blocked (private mode / strict settings) — show the gate
        // rather than silently skipping it.
        alreadyAcked = false;
      }
      if (alreadyAcked) return;

      // Block the composer while the gate is up
      const composer = document.querySelector(".composer-wrap");
      if (composer) composer.setAttribute("aria-hidden", "true");

      overlay.hidden = false;
      setTimeout(() => checkbox.focus(), 120);

      checkbox.addEventListener("change", () => {
        continueBtn.disabled = !checkbox.checked;
      });

      continueBtn.addEventListener("click", () => {
        if (!checkbox.checked) return;
        try {
          sessionStorage.setItem(ACK_KEY, new Date().toISOString());
        } catch (e) { /* not persistable — gate will show again */ }
        // Record it server-side so the admin log can show who acknowledged
        try {
          fetch("/acknowledge", {
            method: "POST",
            headers: { "Content-Type": "application/json" },
            body: JSON.stringify({ version: ACK_KEY }),
            keepalive: true,
          }).catch(() => {});
        } catch (e) { /* non-blocking — never hold up entry */ }
        overlay.hidden = true;
        if (composer) composer.removeAttribute("aria-hidden");
        const input = document.getElementById("message");
        if (input) input.focus();
      });

      // Enter on the checkbox accepts; Esc must not dismiss the gate
      overlay.addEventListener("keydown", (e) => {
        if (e.key === "Enter" && checkbox.checked) continueBtn.click();
        if (e.key === "Escape") e.preventDefault();
      });
    })();

    const form = document.getElementById("chat-form");
    const input = document.getElementById("message");
    const sendBtn = document.getElementById("send-btn");
    const resetBtn = document.getElementById("reset-btn");
    const chatWrap = document.getElementById("chat-wrap");
    const OPENING = {{ cfg.opening|tojson }};

    // -------------------------------------------------------------
    // Cross-platform speech engine (macOS, Windows, iOS, Android, Linux)
    // -------------------------------------------------------------
    // The Web Speech API behaves differently on every platform. This wrapper
    // normalizes the known problems:
    //
    //   1. Chrome (all desktop OSes) silently stops speaking after ~15 seconds
    //      on a long utterance. Fixed by splitting text into sentence-sized
    //      chunks played in sequence, plus a periodic pause/resume keepalive.
    //   2. iOS Safari refuses speak() unless it originates from a user gesture,
    //      which breaks auto-speak. Fixed by "priming" the engine with a silent
    //      utterance on the first tap anywhere on the page.
    //   3. Chrome and Android populate getVoices() asynchronously and
    //      onvoiceschanged sometimes never fires. Fixed by polling with a cap.
    //   4. pause()/resume() are unreliable on iOS and Android. Detected up
    //      front; on those platforms the control becomes stop-only.
    //   5. Calling speak() immediately after cancel() drops the utterance in
    //      some builds. Fixed with a short delay between the two.
    const J3PSpeech = (function () {
      const supported =
        typeof window !== "undefined" &&
        "speechSynthesis" in window &&
        "SpeechSynthesisUtterance" in window;

      const ua = (navigator.userAgent || "");
      const isIOS = /iPad|iPhone|iPod/.test(ua) ||
                    (navigator.platform === "MacIntel" && navigator.maxTouchPoints > 1);
      const isAndroid = /Android/.test(ua);
      const isChrome = /Chrome|CriOS|Chromium/.test(ua) && !/Edg|OPR/.test(ua);
      // pause/resume is broken or a no-op on mobile WebKit and Android
      const canPause = supported && !isIOS && !isAndroid;
      // Chrome's watchdog is only needed where the 15s cutoff exists
      const needsKeepalive = supported && isChrome && !isAndroid;

      // Preferred voices per platform, best first. Falls through to any
      // English voice, then whatever the device offers.
      const PREFERRED_VOICES = [
        "Samantha",                  // macOS / iOS
        "Karen", "Moira", "Fiona",   // macOS / iOS alternates
        "Google US English",         // Chrome desktop / Android
        "Google UK English Female",
        "Microsoft Aria",            // Windows 11
        "Microsoft Jenny",           // Windows 11
        "Microsoft Zira",            // Windows 10
        "Microsoft Hazel",
        "English United States",     // some Android builds
        "en-us-x-sfg#female_1",      // Android TTS engine ids
        "English (United States)",
      ];

      let voices = [];
      let preferred = null;
      let primed = false;
      let keepaliveTimer = null;

      // --- state for the currently playing item ---
      let chunks = [];
      let chunkIndex = 0;
      let token = 0;              // increments on every stop/new play
      let callbacks = {};
      let speaking = false;
      let paused = false;

      // User preferences (free — these are the voices already on the device)
      let savedVoiceName = null;
      let rate = 1.0;
      try {
        savedVoiceName = localStorage.getItem("j3p_voice") || null;
        const r = parseFloat(localStorage.getItem("j3p_rate"));
        if (!isNaN(r) && r >= 0.6 && r <= 1.6) rate = r;
      } catch (e) { /* storage blocked — fall back to defaults */ }

      function refreshVoices() {
        if (!supported) return;
        voices = window.speechSynthesis.getVoices() || [];
        if (!voices.length) return;
        // An explicit user choice always wins over the automatic pick
        if (savedVoiceName) {
          const chosen = voices.find(v => v.name === savedVoiceName);
          if (chosen) { preferred = chosen; return; }
        }
        for (const name of PREFERRED_VOICES) {
          const hit = voices.find(v =>
            v.name === name ||
            v.name.toLowerCase().startsWith(name.toLowerCase()));
          if (hit) { preferred = hit; return; }
        }
        preferred =
          voices.find(v => (v.lang || "").toLowerCase().replace("_", "-").startsWith("en-us")) ||
          voices.find(v => (v.lang || "").toLowerCase().startsWith("en")) ||
          voices[0] || null;
      }

      if (supported) {
        refreshVoices();
        // onvoiceschanged is the documented path...
        window.speechSynthesis.onvoiceschanged = refreshVoices;
        // ...but it doesn't always fire on Android/Chrome, so poll briefly too.
        let tries = 0;
        const poll = setInterval(() => {
          tries += 1;
          if (voices.length || tries > 20) { clearInterval(poll); return; }
          refreshVoices();
        }, 250);

        // iOS requires a gesture-originated speak() before any programmatic one
        // is allowed. Prime on first interaction with a whitespace utterance.
        // Note: volume 0 does NOT reliably unlock WebKit, so this runs at full
        // volume — a single space is inaudible either way.
        const prime = () => {
          if (primed) return;
          primed = true;
          try {
            refreshVoices();   // iOS often populates the list only after a gesture
            const u = new SpeechSynthesisUtterance(" ");
            u.rate = 1; u.volume = 1;
            if (preferred) u.voice = preferred;
            window.speechSynthesis.speak(u);
          } catch (e) { /* non-fatal */ }
        };
        ["pointerdown", "touchstart", "touchend", "click", "keydown"].forEach(evt =>
          window.addEventListener(evt, prime, { once: true, passive: true }));

        // Don't keep talking after the tab is closed or navigated away
        window.addEventListener("beforeunload", () => {
          try { window.speechSynthesis.cancel(); } catch (e) {}
        });
        // Chrome keeps queued speech alive across tab switches inconsistently;
        // stopping on hide is more predictable than half-spoken audio resuming.
        document.addEventListener("visibilitychange", () => {
          if (document.hidden && speaking) stop();
        });
      }

      // Split into chunks small enough to dodge Chrome's ~15s cutoff, breaking
      // on sentence boundaries so the pauses land naturally.
      function chunkText(text, maxLen) {
        maxLen = maxLen || 180;
        const out = [];
        const sentences = String(text || "")
          .split(/(?<=[.!?:;])\s+/)
          .filter(Boolean);
        let buf = "";
        for (let s of sentences) {
          // A single sentence longer than the cap gets broken on commas/spaces
          while (s.length > maxLen) {
            let cut = s.lastIndexOf(",", maxLen);
            if (cut < maxLen * 0.5) cut = s.lastIndexOf(" ", maxLen);
            if (cut < maxLen * 0.5) cut = maxLen;
            if (buf) { out.push(buf.trim()); buf = ""; }
            out.push(s.slice(0, cut + 1).trim());
            s = s.slice(cut + 1);
          }
          if ((buf + " " + s).trim().length > maxLen) {
            if (buf) out.push(buf.trim());
            buf = s;
          } else {
            buf = (buf ? buf + " " : "") + s;
          }
        }
        if (buf.trim()) out.push(buf.trim());
        return out.filter(c => c.length);
      }

      function startKeepalive() {
        if (!needsKeepalive || keepaliveTimer) return;
        keepaliveTimer = setInterval(() => {
          const synth = window.speechSynthesis;
          if (synth.speaking && !synth.paused) {
            // The pause/resume pair resets Chrome's internal idle timer
            try { synth.pause(); synth.resume(); } catch (e) {}
          }
        }, 8000);
      }
      function stopKeepalive() {
        if (keepaliveTimer) { clearInterval(keepaliveTimer); keepaliveTimer = null; }
      }

      function speakChunk(myToken) {
        if (myToken !== token) return;               // superseded
        if (chunkIndex >= chunks.length) {           // finished cleanly
          speaking = false; paused = false;
          stopKeepalive();
          if (callbacks.onEnd) callbacks.onEnd();
          return;
        }
        const u = new SpeechSynthesisUtterance(chunks[chunkIndex]);
        u.rate = rate; u.pitch = 1.0; u.volume = 1.0;
        if (preferred) { u.voice = preferred; u.lang = preferred.lang || "en-US"; }
        else { u.lang = "en-US"; }

        u.onend = () => {
          if (myToken !== token) return;
          chunkIndex += 1;
          speakChunk(myToken);
        };
        u.onerror = (e) => {
          if (myToken !== token) return;
          const reason = (e && e.error) || "";
          // These fire on normal cancellation — not real failures
          if (reason === "interrupted" || reason === "canceled") return;
          speaking = false; paused = false;
          stopKeepalive();
          if (callbacks.onError) callbacks.onError(reason);
        };
        try {
          window.speechSynthesis.speak(u);
        } catch (err) {
          speaking = false; stopKeepalive();
          if (callbacks.onError) callbacks.onError(err && err.message);
        }
      }

      function stop() {
        token += 1;
        chunks = []; chunkIndex = 0;
        speaking = false; paused = false;
        stopKeepalive();
        try { window.speechSynthesis.cancel(); } catch (e) {}
      }

      function play(text, cbs, opts) {
        if (!supported) return false;
        opts = opts || {};
        callbacks = cbs || {};
        stop();                                   // clears any prior playback
        const body = String(text || "").trim();
        if (!body) return false;
        chunks = chunkText(body);
        chunkIndex = 0;
        const myToken = ++token;
        speaking = true; paused = false;
        startKeepalive();
        if (callbacks.onStart) callbacks.onStart();

        // iOS Safari only permits speak() while still inside the user gesture
        // that triggered it. Any setTimeout in between puts us outside that
        // window and WebKit silently refuses. So when the call originates from
        // a tap, fire the first chunk synchronously.
        if (opts.fromGesture) {
          speakChunk(myToken);
        } else {
          // Programmatic (auto-speak): rely on the earlier priming. The small
          // delay avoids builds that drop a speak() issued right after cancel().
          setTimeout(() => speakChunk(myToken), isIOS ? 120 : 60);
        }

        // Watchdog: if nothing is actually speaking shortly after we started,
        // the platform refused the request. Surface it instead of leaving the
        // button stuck in a "speaking" state forever.
        setTimeout(() => {
          if (myToken !== token) return;
          if (!speaking) return;
          let reallySpeaking = false;
          try { reallySpeaking = window.speechSynthesis.speaking || window.speechSynthesis.pending; }
          catch (e) {}
          if (!reallySpeaking) {
            speaking = false; stopKeepalive();
            if (callbacks.onError) callbacks.onError("blocked");
          }
        }, 1200);
        return true;
      }

      function pause() {
        if (!canPause || !speaking) return false;
        try { window.speechSynthesis.pause(); paused = true; return true; }
        catch (e) { return false; }
      }
      function resume() {
        if (!canPause || !speaking) return false;
        try { window.speechSynthesis.resume(); paused = false; return true; }
        catch (e) { return false; }
      }

      // --- Curation: professional voices only, one per language ------------
      // Platforms ship a lot of novelty and low-fidelity voices alongside the
      // good ones. These are filtered out, then the best remaining voice for
      // each language is kept so the menu stays short and every option is
      // presentable to a client.

      // Novelty / joke / legacy low-quality voices (mostly macOS and iOS)
      const NOVELTY_VOICES = new RegExp("^(" + [
        "albert", "bad news", "bahh", "bells", "boing", "bubbles", "cellos",
        "deranged", "good news", "hysterical", "jester", "junior", "kathy",
        "organ", "pipe organ", "princess", "ralph", "superstar", "trinoids",
        "whisper", "wobble", "zarvox", "bruce", "fred", "agnes", "victoria",
        "grandma", "grandpa", "rocko", "sandy", "shelley", "flo", "eddy", "reed",
      ].join("|") + ")\\b", "i");

      // Apple's professional-grade voices across languages
      const APPLE_PRO = new RegExp("^(" + [
        "samantha", "alex", "ava", "allison", "susan", "tom", "nicky", "aaron",
        "serena", "kate", "oliver", "daniel", "karen", "moira", "fiona", "tessa",
        "rishi", "veena", "zoe", "evan", "nathan", "joelle", "noelle", "isha",
        "amelie", "audrey", "aurelie", "thomas", "chantal", "nicolas",
        "anna", "petra", "markus", "yannick", "helena", "martin",
        "monica", "paulina", "juan", "diego", "jorge", "marisol", "soledad",
        "alice", "luca", "federica", "emma", "paola",
        "kyoko", "otoya", "hattori", "o-ren",
        "ting-ting", "sin-ji", "li-mu", "yu-shu", "meijia",
        "yuna", "sora", "milena", "katya", "yuri",
        "zosia", "krzysztof", "ellen", "xander", "nora", "alva", "klara",
        "satu", "melina", "damayanti", "lekha", "carmit", "maged", "tarik",
        "luciana", "joana", "catarina", "ioana", "laura", "lesya", "mariska",
        "magnus", "sara", "kanya", "linh", "mai",
      ].join("|") + ")\\b", "i");

      function isProfessional(v) {
        const name = v.name || "";
        if (NOVELTY_VOICES.test(name)) return false;
        // "Compact" variants are the low-bitrate fallbacks on Apple platforms
        if (/compact|eloquence/i.test(name)) return false;
        // Google and Microsoft ship uniformly presentable voices
        if (/^(Google|Microsoft)\b/i.test(name)) return true;
        if (APPLE_PRO.test(name)) return true;
        return false;
      }

      // Higher score wins when two voices share a language
      function voiceScore(v) {
        let score = 0;
        const name = (v.name || "").toLowerCase();
        const idx = PREFERRED_VOICES.findIndex(
          p => name === p.toLowerCase() || name.startsWith(p.toLowerCase()));
        if (idx >= 0) score += 1000 - idx * 10;         // explicitly preferred
        if (/natural|premium|enhanced|neural/i.test(v.name || "")) score += 60;
        if (v.localService) score += 25;                // offline = no lag
        if (/en-us/i.test(v.lang || "")) score += 15;   // house default
        return score;
      }

      function listVoices() {
        if (!voices.length) refreshVoices();
        const pro = voices.filter(isProfessional);
        // One entry per primary language subtag: en, fr, es, de …
        const best = new Map();
        for (const v of pro) {
          const key = (v.lang || "und").toLowerCase().replace("_", "-").split("-")[0];
          const current = best.get(key);
          if (!current || voiceScore(v) > voiceScore(current)) best.set(key, v);
        }
        // Always include whatever is currently selected, even if curation
        // would otherwise have dropped it (e.g. a previously saved choice).
        if (preferred && ![...best.values()].some(v => v.name === preferred.name)) {
          const key = (preferred.lang || "und").toLowerCase().split("-")[0];
          if (!best.has(key)) best.set(key, preferred);
        }
        return [...best.values()]
          .map(v => ({ name: v.name, lang: v.lang || "", local: !!v.localService }))
          .sort((a, b) => {
            // English first, then alphabetical by language
            const ae = a.lang.toLowerCase().startsWith("en") ? 0 : 1;
            const be = b.lang.toLowerCase().startsWith("en") ? 0 : 1;
            return ae - be || a.lang.localeCompare(b.lang);
          });
      }
      function setVoice(name) {
        const hit = voices.find(v => v.name === name);
        if (!hit) return false;
        preferred = hit;
        savedVoiceName = name;
        try { localStorage.setItem("j3p_voice", name); } catch (e) {}
        return true;
      }
      function setRate(r) {
        r = parseFloat(r);
        if (isNaN(r) || r < 0.6 || r > 1.6) return false;
        rate = r;
        try { localStorage.setItem("j3p_rate", String(r)); } catch (e) {}
        return true;
      }

      return {
        supported, canPause, isIOS, isAndroid,
        play, stop, pause, resume,
        isSpeaking: () => speaking,
        isPaused: () => paused,
        voiceName: () => (preferred ? preferred.name : null),
        listVoices, setVoice, setRate,
        getRate: () => rate,
        onVoicesReady: (cb) => {
          if (voices.length) { cb(); return; }
          let n = 0;
          const t = setInterval(() => {
            n += 1;
            if (voices.length || n > 20) { clearInterval(t); cb(); }
          }, 250);
        },
      };
    })();

    // -------------------------------------------------------------
    // Voice picker — lists the voices already installed on the device.
    // No API, no per-character cost; the preference persists locally.
    // -------------------------------------------------------------
    (function initVoicePicker() {
      const wrap = document.querySelector(".voice-wrap");
      const btn = document.getElementById("voice-btn");
      const menu = document.getElementById("voice-menu");
      const select = document.getElementById("voice-select");
      const range = document.getElementById("rate-range");
      const rateVal = document.getElementById("rate-val");
      const preview = document.getElementById("voice-preview");
      if (!wrap || !J3PSpeech.supported) {
        if (wrap) wrap.style.display = "none";
        return;
      }

      // Map a language tag to a readable name, so the menu reads
      // "English — Samantha" rather than "Samantha (en-US)".
      let langNamer = null;
      try {
        if (typeof Intl !== "undefined" && Intl.DisplayNames) {
          langNamer = new Intl.DisplayNames(["en"], { type: "language" });
        }
      } catch (e) { /* older browser — fall back to the raw tag */ }

      function languageLabel(tag) {
        const base = (tag || "").replace("_", "-").split("-")[0];
        if (!base) return "Other";
        if (langNamer) {
          try {
            const n = langNamer.of(base);
            if (n && n !== base) return n.charAt(0).toUpperCase() + n.slice(1);
          } catch (e) {}
        }
        return base.toUpperCase();
      }

      function populate() {
        const all = J3PSpeech.listVoices();
        // One curated voice per language — if that leaves nothing to choose
        // between, the picker has no purpose, so hide it entirely.
        if (all.length < 2) {
          wrap.style.display = "none";
          menu.classList.remove("open");
          return;
        }
        wrap.style.display = "";
        select.disabled = false;
        preview.disabled = false;

        // Tidy platform-verbose names:
        //   "Microsoft Aria Online (Natural) - English (United States)" -> "Microsoft Aria"
        const prettyName = (n) => (n || "")
          .split(" - ")[0]
          .replace(/\s*online\s*/i, " ")
          .replace(/\s*\((natural|premium|enhanced|compact)\)\s*/i, " ")
          .replace(/\s{2,}/g, " ")
          .trim();

        select.innerHTML = all.map(v => {
          const safe = v.name.replace(/"/g, "&quot;");   // value must stay exact
          return `<option value="${safe}">${languageLabel(v.lang)} — ${prettyName(v.name)}</option>`;
        }).join("");

        const current = J3PSpeech.voiceName();
        if (current && all.some(v => v.name === current)) select.value = current;
      }

      // Voices load asynchronously on Chrome/Android, so wait for them
      J3PSpeech.onVoicesReady(populate);

      const r = J3PSpeech.getRate();
      range.value = String(r);
      rateVal.textContent = Number(r).toFixed(1) + "\u00d7";

      btn.addEventListener("click", (e) => {
        e.stopPropagation();
        const opening = !menu.classList.contains("open");
        menu.classList.toggle("open");
        if (opening) populate();     // refresh in case voices arrived late
      });

      select.addEventListener("change", () => {
        if (J3PSpeech.setVoice(select.value)) {
          // Speak a short sample so the choice is audible immediately
          J3PSpeech.play("This is how I'll sound.", {}, { fromGesture: true });
        }
      });

      range.addEventListener("input", () => {
        rateVal.textContent = Number(range.value).toFixed(1) + "\u00d7";
        J3PSpeech.setRate(range.value);
      });

      preview.addEventListener("click", () => {
        J3PSpeech.play(
          "Naming the hard thing early is usually the work. This is the voice you'll hear.",
          {}, { fromGesture: true }
        );
      });

      document.addEventListener("click", (e) => {
        if (!wrap.contains(e.target)) menu.classList.remove("open");
      });
      document.addEventListener("keydown", (e) => {
        if (e.key === "Escape") menu.classList.remove("open");
      });
    })();

    // Auto-speak state — persists across visits
    const autoSpeakBtn = document.getElementById("autospeak-btn");
    let autoSpeakEnabled = false;
    try {
      // Key deliberately versioned: speak-everything was on for anyone who
      // had enabled it once, which is rarely what they want on return.
      autoSpeakEnabled = localStorage.getItem("j3p_autospeak_v2") === "1";
    } catch (e) { /* localStorage may be blocked; fall back to session default */ }
    function refreshAutoSpeakUI() {
      if (!autoSpeakBtn) return;
      if (autoSpeakEnabled) {
        autoSpeakBtn.classList.add("on");
        autoSpeakBtn.title = "Speak is ON — every response is read aloud. Click to turn off";
      } else {
        autoSpeakBtn.classList.remove("on");
        autoSpeakBtn.title = "Speak is OFF — click to have every response read aloud";
      }
    }
    refreshAutoSpeakUI();
    // If browser doesn't support speech, hide the toggle
    if (!J3PSpeech.supported) {
      if (autoSpeakBtn) autoSpeakBtn.style.display = "none";
      autoSpeakEnabled = false;
    } else if (autoSpeakBtn) {
      autoSpeakBtn.addEventListener("click", () => {
        autoSpeakEnabled = !autoSpeakEnabled;
        try { localStorage.setItem("j3p_autospeak_v2", autoSpeakEnabled ? "1" : "0"); } catch (e) {}
        refreshAutoSpeakUI();
        // Turning OFF should stop anything currently speaking
        if (!autoSpeakEnabled) {
          J3PSpeech.stop();
          if (window.__activeSpeakBtn) {
            const prev = window.__activeSpeakBtn;
            prev.classList.remove("speaking", "paused");
            const lbl = prev.querySelector(".speak-label");
            if (lbl) lbl.textContent = "Speak";
            window.__activeSpeakBtn = null;
          }
        }
      });
    }
    // Expose to addMessage so a new bot reply can auto-play when enabled
    window.__isAutoSpeakOn = () => autoSpeakEnabled;

    // -------------------------------------------------------------
    // Scheduling gate — the release must be acknowledged before the
    // calendar opens. Deliberately NOT remembered: it is re-shown on
    // every attempt to book, so each booking carries an acknowledgment.
    // -------------------------------------------------------------
    (function initSchedGate() {
      const overlay = document.getElementById("sched-overlay");
      const checkbox = document.getElementById("sched-checkbox");
      const continueBtn = document.getElementById("sched-continue");
      const cancelBtn = document.getElementById("sched-cancel");
      if (!overlay) return;

      let pendingUrl = null;

      function closeGate() {
        overlay.hidden = true;
        checkbox.checked = false;
        continueBtn.disabled = true;
        pendingUrl = null;
      }

      function openGate(url) {
        pendingUrl = url;
        checkbox.checked = false;        // always start unchecked
        continueBtn.disabled = true;
        overlay.hidden = false;
        setTimeout(() => checkbox.focus(), 120);
      }

      // Delegated + capture phase: catches the click no matter when the button
      // was added to the page or what else is listening. The button carries no
      // href, so if this script ever fails to run the calendar simply isn't
      // reachable — it can't be opened ungated.
      document.addEventListener("click", (e) => {
        const btn = e.target.closest && e.target.closest(".cta-btn");
        if (!btn) return;
        e.preventDefault();
        e.stopPropagation();
        openGate(btn.getAttribute("data-cta-url"));
      }, true);

      // Keyboard activation, since it's no longer a real link
      document.addEventListener("keydown", (e) => {
        if (e.key !== "Enter" && e.key !== " ") return;
        const btn = e.target.closest && e.target.closest(".cta-btn");
        if (!btn) return;
        e.preventDefault();
        openGate(btn.getAttribute("data-cta-url"));
      }, true);

      checkbox.addEventListener("change", () => {
        continueBtn.disabled = !checkbox.checked;
      });

      continueBtn.addEventListener("click", () => {
        if (!checkbox.checked || !pendingUrl) return;
        const url = pendingUrl;
        // Brief the advisor on what this person has been working through.
        // Fire-and-forget: the calendar must open regardless.
        try {
          fetch("/briefing/schedule", { method: "POST", keepalive: true })
            .catch(() => {});
        } catch (e) { /* never block the booking */ }
        closeGate();
        // A synthetic anchor click is the reliable way to open a new tab here:
        // window.open() returns null whenever "noopener" is set — even on
        // success — so its return value can't be used to detect a blocked
        // pop-up, and acting on it would navigate this page away and drop the
        // user's session.
        const a = document.createElement("a");
        a.href = url;
        a.target = "_blank";
        a.rel = "noopener noreferrer";
        document.body.appendChild(a);
        a.click();
        document.body.removeChild(a);
      });

      cancelBtn.addEventListener("click", closeGate);

      // This gate is dismissible (unlike the entry gate) — Esc and backdrop close
      overlay.addEventListener("keydown", (e) => {
        if (e.key === "Escape") closeGate();
        if (e.key === "Enter" && checkbox.checked) continueBtn.click();
      });
      overlay.addEventListener("click", (e) => {
        if (e.target === overlay) closeGate();
      });
    })();

    // Generates the file server-side and hands it to the browser. Shared by
    // the SAVE menu and by automatic export when the user asked for a format.
    async function downloadExport(fmt, text, part) {
      const payload = { text: text, title: "" };
      if (part !== undefined && part !== null) payload.part = part;
      const resp = await fetch(`/export/${fmt}`, {
        method: "POST",
        headers: { "Content-Type": "application/json" },
        body: JSON.stringify(payload),
      });
      if (!resp.ok) {
        let msg = `Export failed (${resp.status})`;
        try { const j = await resp.json(); if (j.error) msg = j.error; } catch (e) {}
        throw new Error(msg);
      }
      // Prefer the server-generated filename from Content-Disposition
      let filename = `j3p_response.${fmt}`;
      const cd = resp.headers.get("Content-Disposition") || "";
      const match = cd.match(/filename="?([^"]+)"?/);
      if (match) filename = match[1];

      const blob = await resp.blob();
      const url = URL.createObjectURL(blob);
      const a = document.createElement("a");
      a.href = url;
      a.download = filename;
      document.body.appendChild(a);
      a.click();
      document.body.removeChild(a);
      setTimeout(() => URL.revokeObjectURL(url), 4000);
      return filename;
    }

    // Popover menus sit above their button by default. Near the top of the
    // page that clips them, so measure and flip downward when needed.
    function positionMenu(menu) {
      if (!menu) return;
      // The menu lives inside the scrolling chat area, so it's clipped by that
      // container's edges rather than the viewport's. Measure against both.
      const scroller = document.getElementById("chat-wrap");
      const sb = scroller ? scroller.getBoundingClientRect() : null;
      const topLimit = Math.max(8, sb ? sb.top + 4 : 8);
      const bottomLimit = Math.min(window.innerHeight - 8,
                                   sb ? sb.bottom - 4 : window.innerHeight - 8);

      menu.classList.remove("drop-down");
      const above = menu.getBoundingClientRect();
      if (above.top >= topLimit) return;            // fits as-is

      menu.classList.add("drop-down");
      const below = menu.getBoundingClientRect();
      if (below.bottom > bottomLimit) {
        // Neither direction fits cleanly — take whichever clips less
        const clipAbove = topLimit - above.top;
        const clipBelow = below.bottom - bottomLimit;
        if (clipAbove <= clipBelow) menu.classList.remove("drop-down");
      }
    }

    // Keep an open menu correctly placed if the chat scrolls behind it
    ["scroll", "resize"].forEach(evt => {
      const target = evt === "scroll" ? document.getElementById("chat-wrap") : window;
      if (target) target.addEventListener(evt, () => {
        document.querySelectorAll(".share-menu.open").forEach(positionMenu);
      }, { passive: true });
    });

    const FORMAT_NAMES = { docx: "Word", pptx: "PowerPoint", xlsx: "Excel", pdf: "PDF" };

    // ---------------------------------------------------------------
    // Rating nudge
    // ---------------------------------------------------------------
    // A session has no explicit end, so the nudge appears on the most recent
    // unrated reply when the conversation looks finished: the person has gone
    // quiet, has said something valedictory, or is starting a new conversation.
    let idleNudgeTimer = null;
    const IDLE_NUDGE_MS = 75000;

    function lastUnratedReply() {
      const msgs = Array.from(document.querySelectorAll(".msg.assistant"));
      for (let i = msgs.length - 1; i >= 0; i--) {
        const fb = msgs[i].querySelector(".feedback");
        if (!fb) continue;
        if (fb.dataset.rated === "1") return null;   // already rated: no nudge
        return msgs[i];
      }
      return null;
    }

    function showRatingNudge(reason) {
      const msgDiv = lastUnratedReply();
      if (!msgDiv) return;
      if (msgDiv.querySelector(".rate-nudge")) return;
      if (msgDiv.querySelector(".safety-reply")) return;   // never on a crisis reply

      const fb = msgDiv.querySelector(".feedback");
      if (fb) fb.classList.add("nudged");

      const bar = document.createElement("div");
      bar.className = "rate-nudge";
      const text = document.createElement("span");
      text.textContent = reason === "farewell"
        ? "Before you go — was this helpful? Use the thumbs above; it's how the advisor improves."
        : "Was this helpful? A thumbs up or down above helps improve the advisor.";
      bar.appendChild(text);

      const close = document.createElement("button");
      close.type = "button";
      close.className = "nudge-dismiss";
      close.setAttribute("aria-label", "Dismiss");
      close.innerHTML = "&times;";
      close.addEventListener("click", () => {
        bar.remove();
        if (fb) fb.classList.remove("nudged");
      });
      bar.appendChild(close);

      msgDiv.appendChild(bar);
      chatWrap.scrollTop = chatWrap.scrollHeight;
    }

    // ---------------------------------------------------------------
    // End-of-session follow-up plan
    // ---------------------------------------------------------------
    // Offered on the same signals as the rating nudge: the conversation has
    // gone quiet, they've signed off, or they're starting fresh. Offered once.
    let planOffered = false;

    function showPlanOffer() {
      if (planOffered) return;
      const msgs = Array.from(document.querySelectorAll(".msg.assistant"));
      const last = msgs[msgs.length - 1];
      if (!last || last.querySelector(".plan-offer")) return;
      if (last.querySelector(".safety-reply")) return;   // never after a crisis reply
      // Needs enough of a conversation to be worth planning from
      if (document.querySelectorAll(".msg.user").length < 2) return;
      planOffered = true;

      const bar = document.createElement("div");
      bar.className = "plan-offer";
      const text = document.createElement("span");
      text.textContent = "Before you go — want a follow-up plan? "
                       + "What you said you'd do, by when.";
      bar.appendChild(text);

      const yes = document.createElement("button");
      yes.type = "button";
      yes.textContent = "Build my plan";
      bar.appendChild(yes);

      const no = document.createElement("button");
      no.type = "button";
      no.className = "secondary";
      no.textContent = "No thanks";
      no.addEventListener("click", () => bar.remove());
      bar.appendChild(no);

      yes.addEventListener("click", async () => {
        yes.disabled = true;
        no.disabled = true;
        text.textContent = "Building your plan\u2026";
        try {
          const resp = await fetch("/plan/create", { method: "POST" });
          const data = await resp.json();
          if (!data.ok) throw new Error(data.error || "Could not build it");
          bar.remove();
          // Render as a normal reply so SAVE, COPY, SHARE and SPEAK all work
          addMessage(data.plan, "assistant", true, null, null);
        } catch (err) {
          text.textContent = String(err.message || err);
          no.disabled = false;
          no.textContent = "Close";
        }
      });

      last.appendChild(bar);
      chatWrap.scrollTop = chatWrap.scrollHeight;
    }

    function clearRatingNudge() {
      document.querySelectorAll(".rate-nudge").forEach(n => n.remove());
      document.querySelectorAll(".feedback.nudged").forEach(f => f.classList.remove("nudged"));
    }

    function armIdleNudge() {
      if (idleNudgeTimer) clearTimeout(idleNudgeTimer);
      idleNudgeTimer = setTimeout(() => {
        showPlanOffer();
        showRatingNudge("idle");
      }, IDLE_NUDGE_MS);
    }

    // "thanks, that's all" and similar — a session ending politely
    const FAREWELL_RE = new RegExp(
      "\\b(thank you|thanks|thx|that'?s all|that helps|perfect|appreciate it"
      + "|goodbye|bye|talk soon|see you|have a good|that'?s everything"
      + "|all set|nothing else|we'?re done|i'?m done|very helpful)\\b", "i");

    function looksLikeFarewell(text) {
      const t = (text || "").trim();
      if (!t || t.length > 90) return false;         // long messages aren't sign-offs
      return FAREWELL_RE.test(t);
    }

    // Offers a download beneath a reply. Nothing is written to disk unless
    // the user picks a format — replies never export themselves.
    // Shows a readable failure with a Retry button, and puts the user's text
    // back in the composer so nothing they typed is lost.
    function showRetry(message, originalText, files) {
      const div = document.createElement("div");
      div.className = "msg assistant";
      const body = document.createElement("div");
      body.className = "msg-body";
      body.textContent = message;
      div.appendChild(body);

      const bar = document.createElement("div");
      bar.className = "export-offer";
      const retry = document.createElement("button");
      retry.type = "button";
      retry.className = "export-chip primary";
      retry.textContent = "Try again";
      retry.addEventListener("click", () => {
        div.remove();
        input.value = originalText || "";
        form.dispatchEvent(new Event("submit", { cancelable: true, bubbles: true }));
      });
      bar.appendChild(retry);
      div.appendChild(bar);
      chat.appendChild(div);
      chatWrap.scrollTop = chatWrap.scrollHeight;

      // Restore what they typed so it isn't lost either way
      if (originalText && !input.value) input.value = originalText;
      if (files && files.length) {
        attachedFiles = files.slice();
        refreshAttachmentChip();
      }
    }

    function offerExport(msgDiv, text, suggested, docs) {
      if (!msgDiv || msgDiv.querySelector(".export-offer")) return;
      const FMT_ORDER = ["docx", "pptx", "xlsx", "pdf"];
      const primary = FMT_ORDER.includes(suggested) ? suggested : "docx";

      const bar = document.createElement("div");
      bar.className = "export-offer";

      const label = document.createElement("span");
      label.className = "export-offer-label";
      label.textContent = (docs && docs.length > 1)
        ? `Save these ${docs.length} documents as files?`
        : "Save this as a file?";
      bar.appendChild(label);

      const chips = document.createElement("span");
      chips.className = "export-offer-chips";

      const makeChip = (fmt, text2, isPrimary) => {
        const btn = document.createElement("button");
        btn.type = "button";
        btn.className = "export-chip" + (isPrimary ? " primary" : "");
        btn.textContent = text2;
        btn.addEventListener("click", async () => {
          const original = btn.textContent;
          chips.querySelectorAll("button").forEach(x => x.disabled = true);
          btn.textContent = "Building\u2026";
          try {
            if (docs && docs.length > 1) {
              for (const doc of docs) {
                await downloadExport(fmt, text, doc.index);
                await new Promise(r => setTimeout(r, 800));
              }
            } else {
              await downloadExport(fmt, text);
            }
            btn.textContent = "Saved";
            setTimeout(() => bar.remove(), 1500);
          } catch (err) {
            console.error("Export failed:", err);
            btn.textContent = "Failed";
            setTimeout(() => {
              btn.textContent = original;
              chips.querySelectorAll("button").forEach(x => x.disabled = false);
            }, 1800);
          }
        });
        return btn;
      };

      chips.appendChild(makeChip(primary, FORMAT_NAMES[primary], true));
      FMT_ORDER.filter(f => f !== primary)
               .forEach(f => chips.appendChild(makeChip(f, FORMAT_NAMES[f], false)));

      const dismiss = document.createElement("button");
      dismiss.type = "button";
      dismiss.className = "export-chip dismiss";
      dismiss.textContent = "No thanks";
      dismiss.addEventListener("click", () => bar.remove());
      chips.appendChild(dismiss);

      bar.appendChild(chips);
      const actions = msgDiv.querySelector(".feedback");
      if (actions) msgDiv.insertBefore(bar, actions);
      else msgDiv.appendChild(bar);
    }

    // Minimal, safe markdown renderer for assistant replies. Escapes HTML
    // first, then converts headings, lists, bold/italic and inline code so
    // "## Cover Letter" displays as a heading instead of literal characters.
    function escapeHtml(s) {
      return String(s)
        .replace(/&/g, "&amp;").replace(/</g, "&lt;").replace(/>/g, "&gt;");
    }

    function renderInline(s) {
      return escapeHtml(s)
        .replace(/\*\*(.+?)\*\*/g, "<strong>$1</strong>")
        .replace(/(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)/g, "<em>$1</em>")
        .replace(/`([^`]+?)`/g, "<code>$1</code>")
        .replace(/\[([^\]]+?)\]\((https?:\/\/[^)\s]+?)\)/g,
                 '<a href="$2" target="_blank" rel="noopener">$1</a>');
    }

    function renderMarkdown(src) {
      const lines = String(src || "").replace(/\r\n/g, "\n").split("\n");
      let html = "";
      let listType = null;      // "ul" | "ol" | null
      let paraBuffer = [];

      const flushPara = () => {
        if (paraBuffer.length) {
          html += "<p>" + renderInline(paraBuffer.join(" ")) + "</p>";
          paraBuffer = [];
        }
      };
      const closeList = () => {
        if (listType) { html += `</${listType}>`; listType = null; }
      };

      for (const raw of lines) {
        const line = raw.trim();

        if (!line) { flushPara(); closeList(); continue; }

        if (/^([-*_])\1{2,}$/.test(line)) {
          flushPara(); closeList(); html += "<hr />"; continue;
        }

        let m = line.match(/^(#{1,6})\s+(.*)$/);
        if (m) {
          flushPara(); closeList();
          const level = Math.min(m[1].length + 2, 6);   // # -> h3, keeps bubbles tidy
          html += `<h${level}>${renderInline(m[2])}</h${level}>`;
          continue;
        }

        m = line.match(/^[-*+]\s+(.*)$/);
        if (m) {
          flushPara();
          if (listType !== "ul") { closeList(); html += "<ul>"; listType = "ul"; }
          html += "<li>" + renderInline(m[1]) + "</li>";
          continue;
        }

        m = line.match(/^\d+[.)]\s+(.*)$/);
        if (m) {
          flushPara();
          if (listType !== "ol") { closeList(); html += "<ol>"; listType = "ol"; }
          html += "<li>" + renderInline(m[1]) + "</li>";
          continue;
        }

        closeList();
        paraBuffer.push(line);
      }
      flushPara();
      closeList();
      return html;
    }

    const ADVISOR_AVATAR = {% if show_avatar %}"{{ cfg.avatar_url }}"{% else %}""{% endif %};
    const TALKING_AVATAR_ON = {{ 'true' if cfg.talking_avatar in ('demo','live') else 'false' }};

    // ---------------------------------------------------------------
    // Sending while the advisor is still working
    // ---------------------------------------------------------------
    // Rather than locking the composer during generation, a message sent mid
    // reply is held and dispatched the moment the current one completes.
    let awaitingReply = false;
    let queuedMessages = [];

    function updateQueueHint() {
      const hint = document.getElementById("queue-hint");
      if (!hint) return;
      if (queuedMessages.length) {
        hint.textContent = queuedMessages.length === 1
          ? "1 message queued — it'll send when the reply lands"
          : queuedMessages.length + " messages queued";
        hint.hidden = false;
      } else {
        hint.hidden = true;
      }
    }

    let sendingQueued = false;

    function flushQueue() {
      if (awaitingReply || !queuedMessages.length) return;
      const next = queuedMessages.shift();
      sendingQueued = true;      // its bubble is already on screen
      updateQueueHint();
      // Put it back in the composer and send it through the normal path, so
      // attachments, history and every guard behave identically.
      input.value = next.text;
      setTimeout(() => {
        // Go through the form's own submit path so every guard, attachment
        // and history rule behaves exactly as for a typed message.
        if (typeof form.requestSubmit === "function") form.requestSubmit();
        else form.dispatchEvent(new Event("submit", { cancelable: true }));
      }, 60);
    }

    // ---------------------------------------------------------------
    // Participant materials
    // ---------------------------------------------------------------
    (function () {
      const overlay = document.getElementById("mat-overlay");
      const openBtn = document.getElementById("materials-open");
      if (!overlay || !openBtn) return;
      const closeBtn = document.getElementById("mat-close");
      const msg = document.getElementById("mat-msg");
      const list = document.getElementById("mat-list");

      function say(text, ok) {
        msg.textContent = text || "";
        msg.className = "mat-msg" + (text ? (ok ? " ok" : " err") : "");
      }

      async function refresh() {
        list.innerHTML = "";
        try {
          const resp = await fetch("/materials");
          const data = await resp.json();
          if (!data.ok) return;
          if (!data.docs.length) {
            const p = document.createElement("p");
            p.className = "mat-item-meta";
            p.textContent = "Nothing added yet.";
            list.appendChild(p);
            return;
          }
          const head = document.createElement("p");
          head.className = "mat-item-meta";
          head.textContent = data.docs.length + " of " + data.limit + " items";
          list.appendChild(head);

          data.docs.forEach(doc => {
            const row = document.createElement("div");
            row.className = "mat-item";

            const title = document.createElement("span");
            title.className = "mat-item-title";
            title.textContent = doc.title;
            row.appendChild(title);

            const meta = document.createElement("span");
            meta.className = "mat-item-meta";
            meta.textContent = (doc.shared ? "shared · " : "private · ")
                             + Math.max(1, Math.round(doc.chars / 1000)) + "k chars";
            row.appendChild(meta);

            const rm = document.createElement("button");
            rm.type = "button";
            rm.className = "mat-remove";
            rm.textContent = "Remove";
            rm.addEventListener("click", async () => {
              rm.disabled = true;
              try {
                const r = await fetch("/materials/delete/" + doc.id, { method: "POST" });
                const d = await r.json();
                if (d.ok) { say("Removed.", true); refresh(); }
                else { say(d.error || "Could not remove it.", false); rm.disabled = false; }
              } catch (e) { say("Could not remove it.", false); rm.disabled = false; }
            });
            row.appendChild(rm);
            list.appendChild(row);
          });
        } catch (e) { /* leave the list empty */ }
      }

      const profBox = document.getElementById("mat-profile");
      const profText = document.getElementById("mat-profile-text");
      const profForget = document.getElementById("mat-profile-forget");

      async function refreshProfile() {
        try {
          const r = await fetch("/profile");
          const d = await r.json();
          const p = (d && d.profile) || {};
          const bits = [];
          if (p.first_name) bits.push(p.first_name);
          if (p.role) bits.push(p.role);
          if (p.specialty) bits.push(p.specialty);
          if (!bits.length) { profBox.hidden = true; return; }
          profText.textContent = bits.join(" · ");
          profBox.hidden = false;
        } catch (e) { profBox.hidden = true; }
      }

      if (profForget) {
        profForget.addEventListener("click", async () => {
          profForget.disabled = true;
          try {
            await fetch("/profile/forget", { method: "POST" });
            say("Cleared. The advisor may ask again next time.", true);
            refreshProfile();
          } catch (e) { say("Could not clear it.", false); }
          finally { profForget.disabled = false; }
        });
      }

      function open() { overlay.hidden = false; say(""); refresh(); refreshProfile(); }
      function close() { overlay.hidden = true; }

      openBtn.addEventListener("click", open);
      closeBtn.addEventListener("click", close);
      overlay.addEventListener("click", e => { if (e.target === overlay) close(); });
      document.addEventListener("keydown", e => {
        if (e.key === "Escape" && !overlay.hidden) close();
      });

      // Tabs
      document.querySelectorAll(".mat-tab").forEach(tab => {
        tab.addEventListener("click", () => {
          document.querySelectorAll(".mat-tab").forEach(t => t.classList.remove("active"));
          tab.classList.add("active");
          const which = tab.dataset.pane;
          document.getElementById("mat-pane-upload").hidden = (which !== "upload");
          document.getElementById("mat-pane-paste").hidden = (which !== "paste");
          say("");
        });
      });

      // Upload files
      const upBtn = document.getElementById("mat-upload-btn");
      upBtn.addEventListener("click", async () => {
        const input = document.getElementById("mat-files");
        if (!input.files || !input.files.length) { say("Choose a file first.", false); return; }
        const fd = new FormData();
        for (const f of input.files) fd.append("files", f);
        if (document.getElementById("mat-share-upload").checked) fd.append("shared", "1");
        upBtn.disabled = true;
        say("Reading your documents\u2026", true);
        try {
          const r = await fetch("/materials/upload", { method: "POST", body: fd });
          const d = await r.json();
          if (d.ok) {
            let note = "Added " + d.added.length + " item"
                     + (d.added.length === 1 ? "" : "s") + ".";
            if (d.failed && d.failed.length) note += " Skipped: " + d.failed.join("; ");
            say(note, true);
            input.value = "";
            refresh();
          } else {
            say(d.error || "Could not add those.", false);
          }
        } catch (e) {
          say("Could not add those. Check your connection and try again.", false);
        } finally {
          upBtn.disabled = false;
        }
      });

      // Paste writing
      const saveBtn = document.getElementById("mat-save-btn");
      saveBtn.addEventListener("click", async () => {
        const title = document.getElementById("mat-title-input").value.trim();
        const content = document.getElementById("mat-text").value.trim();
        if (content.length < 20) { say("Add a little more text than that.", false); return; }
        saveBtn.disabled = true;
        try {
          const r = await fetch("/materials/text", {
            method: "POST",
            headers: { "Content-Type": "application/json" },
            body: JSON.stringify({ title: title, content: content,
                                   shared: document.getElementById("mat-share-text").checked }),
          });
          const d = await r.json();
          if (d.ok) {
            say("Added to your library.", true);
            document.getElementById("mat-title-input").value = "";
            document.getElementById("mat-text").value = "";
            refresh();
          } else {
            say(d.error || "Could not save it.", false);
          }
        } catch (e) {
          say("Could not save it. Check your connection and try again.", false);
        } finally {
          saveBtn.disabled = false;
        }
      });
    })();

    // ---------------------------------------------------------------
    // Persistent avatar presence
    // ---------------------------------------------------------------
    // One always-visible avatar that reflects what the advisor is doing:
    // listening, thinking, responding, or speaking. Clicking it reads the
    // latest reply aloud, or plays a talking-head video when that's enabled.
    const Presence = (function () {
      const root = document.getElementById("presence");
      if (!root) return { set: function () {}, video: function () {} };

      // If the looping portrait can't play, hide it so the still photo shows
      // through rather than leaving an empty frame on top of it.
      const loop = document.getElementById("presence-loop");
      if (loop) {
        const hideLoop = () => {
          loop.style.display = "none";
          // Hand the motion back to CSS so the portrait still looks alive
          root.classList.add("no-loop");
        };
        loop.addEventListener("error", hideLoop);
        Array.from(loop.querySelectorAll("source")).forEach(sc =>
          sc.addEventListener("error", () => {
            if (loop.networkState === 3 || !loop.videoWidth) {
              setTimeout(() => { if (!loop.videoWidth) hideLoop(); }, 1200);
            }
          }));
        setTimeout(() => { if (!loop.videoWidth) hideLoop(); }, 3500);
      }
      const status = document.getElementById("presence-status");
      const frame = document.getElementById("presence-frame");
      const video = document.getElementById("presence-video");
      let respondTimer = null;

      const LABELS = {
        idle: "Listening",
        thinking: "Thinking",
        responding: "Responding",
        speaking: "Speaking",
      };

      function set(state) {
        root.classList.remove("thinking", "responding", "speaking");
        if (state && state !== "idle") root.classList.add(state);
        if (status) status.textContent = LABELS[state] || LABELS.idle;
        if (frame) {
          frame.setAttribute("aria-label", state === "speaking"
            ? "Stop reading aloud" : "Read the latest reply aloud");
        }
      }

      function respondBriefly(seconds) {
        set("responding");
        if (respondTimer) clearTimeout(respondTimer);
        respondTimer = setTimeout(() => {
          if (!root.classList.contains("speaking")) set("idle");
        }, (seconds || 4) * 1000);
      }

      // Clicking the presence speaks the most recent reply
      if (frame) {
        frame.addEventListener("click", async () => {
          const msgs = Array.from(document.querySelectorAll(".msg.assistant"))
            .filter(m => !m.className.includes("typing"));
          const last = msgs[msgs.length - 1];
          if (!last) return;

          if (TALKING_AVATAR_ON) {
            if (root.classList.contains("video")) {
              video.pause();
              root.classList.remove("video");
              set("idle");
              return;
            }
            set("thinking");
            if (status) status.textContent = "Preparing";
            try {
              const text = (last.querySelector(".msg-body") || last).innerText;
              const r = await fetch("/avatar/speak", {
                method: "POST",
                headers: { "Content-Type": "application/json" },
                body: JSON.stringify({ text: text }),
              });
              const d = await r.json();
              if (!d.ok) throw new Error(d.error || "failed");
              while (video.firstChild) video.removeChild(video.firstChild);
              const webm = document.createElement("source");
              webm.src = d.video_url.replace(/\.mp4$/, ".webm");
              webm.type = "video/webm";
              const mp4 = document.createElement("source");
              mp4.src = d.video_url;
              mp4.type = "video/mp4";
              video.appendChild(webm);
              video.appendChild(mp4);
              video.load();
              root.classList.add("video");
              set("speaking");
              video.play().catch(() => {});
              video.onended = () => { root.classList.remove("video"); set("idle"); };
              return;
            } catch (e) {
              root.classList.remove("video");
              // fall through to browser speech
            }
          }

          const speakBtn = last.querySelector(".speak-btn");
          if (speakBtn) speakBtn.click();
        });
      }

      return { set: set, respondBriefly: respondBriefly };
    })();

    // ---------------------------------------------------------------
    // Animated, interactive avatar
    // ---------------------------------------------------------------
    // Three states: idle breathes gently, thinking shows a travelling ring
    // while a reply is being generated, speaking radiates pulses while the
    // reply is read aloud. Clicking it starts or stops that reply's audio,
    // which is the same action as the SPEAK button beside the message.
    // Removes a message and its avatar row together. Removing the message
    // alone used to leave an orphaned row, which rendered as a second avatar.
    function removeMessage(el) {
      if (!el) return;
      const row = el.parentElement;
      el.remove();
      if (row && row.classList.contains("msg-row") && !row.querySelector(".msg")) {
        row.remove();
      }
    }

    function buildAvatar(msgDiv) {
      const wrap = document.createElement("button");
      wrap.type = "button";
      wrap.className = "avatar-wrap";
      wrap.setAttribute("aria-label", "Read this reply aloud");

      const img = document.createElement("img");
      img.className = "avatar";
      img.src = ADVISOR_AVATAR;
      img.alt = "";
      img.addEventListener("error", () => wrap.remove());
      wrap.appendChild(img);

      const ring = document.createElement("span");
      ring.className = "avatar-ring";
      wrap.appendChild(ring);
      for (let i = 0; i < 3; i++) {
        const pulse = document.createElement("span");
        pulse.className = "avatar-pulse";
        wrap.appendChild(pulse);
      }

      const hint = document.createElement("span");
      hint.className = "avatar-hint";
      hint.textContent = "Click to listen";
      wrap.appendChild(hint);

      const video = document.createElement("video");
      video.className = "avatar-video";
      video.setAttribute("playsinline", "");
      video.muted = false;
      wrap.appendChild(video);

      // Point the element at a URL, offering WebM and MP4 where both exist so
      // Safari (H.264) and Chromium builds without proprietary codecs both play.
      function setVideoSource(url) {
        while (video.firstChild) video.removeChild(video.firstChild);
        video.removeAttribute("src");
        const alt = url.endsWith(".mp4") ? url.replace(/\.mp4$/, ".webm") : null;
        if (alt) {
          const w = document.createElement("source");
          w.src = alt; w.type = "video/webm";
          video.appendChild(w);
        }
        const m = document.createElement("source");
        m.src = url;
        m.type = url.endsWith(".webm") ? "video/webm" : "video/mp4";
        video.appendChild(m);
        video.load();
      }

      // With the talking avatar on, clicking asks the server for a video of
      // the advisor speaking this reply. Otherwise it falls back to the
      // browser's own speech, which is instant and free.
      wrap.addEventListener("click", async () => {
        if (!TALKING_AVATAR_ON) {
          const speakBtn = msgDiv ? msgDiv.querySelector(".speak-btn") : null;
          if (speakBtn) speakBtn.click();
          return;
        }
        if (wrap.classList.contains("is-video")) {
          video.pause();
          wrap.classList.remove("is-video", "is-speaking");
          return;
        }
        if (wrap.classList.contains("is-loading")) return;

        const text = msgDiv ? (msgDiv.querySelector(".msg-body") || msgDiv).innerText : "";
        wrap.classList.add("is-loading");
        const hint = wrap.querySelector(".avatar-hint");
        if (hint) hint.textContent = "Preparing video\u2026";
        try {
          const resp = await fetch("/avatar/speak", {
            method: "POST",
            headers: { "Content-Type": "application/json" },
            body: JSON.stringify({ text: text }),
          });
          const data = await resp.json();
          if (!data.ok) throw new Error(data.error || "Generation failed");
          setVideoSource(data.video_url);
          wrap.classList.remove("is-loading");
          wrap.classList.add("is-video", "is-speaking");
          if (hint) hint.textContent = data.demo ? "Demo clip — click to stop"
                                                 : "Click to stop";
          video.play().catch(() => {});
          video.onended = () => {
            wrap.classList.remove("is-video", "is-speaking");
            if (hint) hint.textContent = "Click to listen";
          };
        } catch (err) {
          // Never leave a dead frame — fall back to browser speech
          wrap.classList.remove("is-loading", "is-video");
          if (hint) hint.textContent = "Click to listen";
          const speakBtn = msgDiv ? msgDiv.querySelector(".speak-btn") : null;
          if (speakBtn) speakBtn.click();
        }
      });

      if (msgDiv) msgDiv.__avatarWrap = wrap;
      // The placeholder's avatar starts in the thinking state immediately
      if (msgDiv && msgDiv.className.includes("typing")) {
        wrap.classList.add("is-thinking");
        const h = wrap.querySelector(".avatar-hint");
        if (h) h.textContent = "Thinking\u2026";
        wrap.setAttribute("aria-label", "Preparing a reply");
      }
      return wrap;
    }

    // While a reply is being generated, the newest avatar shows the thinking ring
    function setAvatarThinking(on) {
      Presence.set(on ? "thinking" : "idle");
      const wraps = document.querySelectorAll(".avatar-wrap");
      const last = wraps[wraps.length - 1];
      if (!last) return;
      last.classList.toggle("is-thinking", !!on);
      const hint = last.querySelector(".avatar-hint");
      if (hint) hint.textContent = on ? "Thinking…" : "Click to listen";
    }

    // Plays for a few seconds as a reply appears, so the avatar is visibly
    // active while responding even when the reply isn't being read aloud.
    function setAvatarResponding(msgDiv, seconds = 4) {
      if (Presence.respondBriefly) Presence.respondBriefly(seconds);
      const wrap = msgDiv && msgDiv.__avatarWrap ? msgDiv.__avatarWrap : null;
      if (!wrap) return;
      wrap.classList.add("is-responding");
      if (wrap.__respondTimer) clearTimeout(wrap.__respondTimer);
      wrap.__respondTimer = setTimeout(() => {
        wrap.classList.remove("is-responding");
      }, seconds * 1000);
    }

    function setAvatarSpeaking(msgDiv, on) {
      const wrap = msgDiv && msgDiv.__avatarWrap
        ? msgDiv.__avatarWrap
        : (msgDiv && msgDiv.parentElement
            ? msgDiv.parentElement.querySelector(".avatar-wrap") : null);
      if (!wrap) return;
      wrap.classList.toggle("is-speaking", !!on);
      const hint = wrap.querySelector(".avatar-hint");
      if (hint) hint.textContent = on ? "Click to stop" : "Click to listen";
      wrap.setAttribute("aria-label", on ? "Stop reading aloud" : "Read this reply aloud");
    }

    // The speech engine can stop without a callback on some platforms
    // (navigation, silent switch, interrupted chunk). Poll so a pulsing ring
    // can never be left running with nothing playing.
    let quietTicks = 0;
    setInterval(() => {
      if (!J3PSpeech) return;
      if (J3PSpeech.isSpeaking()) { quietTicks = 0; return; }
      // Two consecutive quiet reads, so this can't race the start of playback
      if (++quietTicks < 2) return;
      document.querySelectorAll(".avatar-wrap.is-speaking").forEach(w => {
        w.classList.remove("is-speaking");
        const h = w.querySelector(".avatar-hint");
        if (h) h.textContent = "Click to listen";
      });
    }, 900);

    function clearAllAvatarStates() {
      document.querySelectorAll(".avatar-wrap").forEach(w => {
        w.classList.remove("is-speaking", "is-thinking");
        const h = w.querySelector(".avatar-hint");
        if (h) h.textContent = "Click to listen";
      });
    }

    function addMessage(text, role, withFeedback = false, interactionId = null, documents = null) {
      const div = document.createElement("div");
      div.className = "msg " + role;
      if (interactionId) div.dataset.interactionId = String(interactionId);
      const textNode = document.createElement("div");
      // Assistant replies get markdown formatting; user messages and the
      // "Thinking…" placeholder stay as plain text.
      if (role.startsWith("assistant") && !role.includes("typing")) {
        textNode.className = "msg-body";
        textNode.innerHTML = renderMarkdown(text);
      } else {
        textNode.textContent = text;
      }
      div.appendChild(textNode);
      if (withFeedback) attachFeedback(div, text, interactionId, documents);
      // Advisor replies are preceded by the avatar, when one is configured.
      // If the image fails to load it removes itself, so a missing file just
      // leaves the reply looking as it did before.
      // Per-reply avatars are deliberately not added: the persistent avatar
      // beside the chat already shows who is speaking, and repeating it on
      // every message was visually noisy.
      chat.appendChild(div);
      chatWrap.scrollTop = chatWrap.scrollHeight;
      // If auto-speak is enabled and this is a fresh assistant reply (with feedback
      // row, meaning it was just received from the API), start reading it aloud.
      // A tiny delay gives the DOM time to attach the speak button.
      if (withFeedback && role.startsWith("assistant") && window.__isAutoSpeakOn && window.__isAutoSpeakOn()) {
        setTimeout(() => {
          const sb = div.querySelector(".speak-btn");
          if (sb) sb.click();
        }, 60);
      }
      return div;
    }

    function attachFeedback(msgDiv, replyText, interactionId, documents) {
      const wrap = document.createElement("div");
      wrap.className = "feedback";
      wrap.innerHTML = `
        <span class="feedback-label">Helpful?</span>
        <button class="feedback-btn" data-rating="up" aria-label="Thumbs up" title="Yes, helpful">
          <svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
            <path d="M7 10v12"/><path d="M15 5.88 14 10h5.83a2 2 0 0 1 1.92 2.56l-2.33 8A2 2 0 0 1 17.5 22H7v-12L11.69 2.5a2 2 0 0 1 3.31 3.38z"/>
          </svg>
        </button>
        <button class="feedback-btn" data-rating="down" aria-label="Thumbs down" title="No, not helpful">
          <svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
            <path d="M17 14V2"/><path d="M9 18.12 10 14H4.17a2 2 0 0 1-1.92-2.56l2.33-8A2 2 0 0 1 6.5 2H17v12l-4.69 7.5a2 2 0 0 1-3.31-3.38z"/>
          </svg>
        </button>
        <button class="action-btn speak-btn" style="margin-left: auto;" aria-label="Read answer aloud" title="Read answer aloud">
          <svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
            <polygon points="11 5 6 9 2 9 2 15 6 15 11 19 11 5"/>
            <path d="M15.54 8.46a5 5 0 0 1 0 7.07"/>
            <path d="M19.07 4.93a10 10 0 0 1 0 14.14"/>
          </svg>
          <span class="speak-label">Speak</span>
        </button>
        <button class="action-btn copy-btn" aria-label="Copy answer" title="Copy answer">
          <svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
            <rect x="9" y="9" width="13" height="13" rx="2"/><path d="M5 15H4a2 2 0 0 1-2-2V4a2 2 0 0 1 2-2h9a2 2 0 0 1 2 2v1"/>
          </svg>
          <span class="copy-label">Copy</span>
        </button>
        <span class="download-wrap">
          <button class="action-btn download-btn" aria-label="Download as document" title="Download as Word, PowerPoint, Excel, or PDF">
            <svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
              <path d="M21 15v4a2 2 0 0 1-2 2H5a2 2 0 0 1-2-2v-4"/><polyline points="7 10 12 15 17 10"/><line x1="12" y1="15" x2="12" y2="3"/>
            </svg>
            <span class="download-label">Save</span>
          </button>
          <div class="share-menu download-menu" role="menu">
            <button type="button" data-fmt="docx" role="menuitem">
              <svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M14 2H6a2 2 0 0 0-2 2v16a2 2 0 0 0 2 2h12a2 2 0 0 0 2-2V8z"/><polyline points="14 2 14 8 20 8"/></svg>
              Word (.docx)
            </button>
            <button type="button" data-fmt="pptx" role="menuitem">
              <svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><rect x="2" y="3" width="20" height="14" rx="2"/><line x1="8" y1="21" x2="16" y2="21"/><line x1="12" y1="17" x2="12" y2="21"/></svg>
              PowerPoint (.pptx)
            </button>
            <button type="button" data-fmt="xlsx" role="menuitem">
              <svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><rect x="3" y="3" width="18" height="18" rx="2"/><line x1="3" y1="9" x2="21" y2="9"/><line x1="3" y1="15" x2="21" y2="15"/><line x1="9" y1="3" x2="9" y2="21"/></svg>
              Excel (.xlsx)
            </button>
            <button type="button" data-fmt="pdf" role="menuitem">
              <svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M14 2H6a2 2 0 0 0-2 2v16a2 2 0 0 0 2 2h12a2 2 0 0 0 2-2V8z"/><polyline points="14 2 14 8 20 8"/><line x1="9" y1="15" x2="15" y2="15"/></svg>
              PDF (.pdf)
            </button>
          </div>
        </span>
        <span class="share-wrap">
          <button class="action-btn share-btn" aria-label="Share answer" title="Share answer">
            <svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
              <circle cx="18" cy="5" r="3"/><circle cx="6" cy="12" r="3"/><circle cx="18" cy="19" r="3"/>
              <line x1="8.59" y1="13.51" x2="15.42" y2="17.49"/><line x1="15.41" y1="6.51" x2="8.59" y2="10.49"/>
            </svg>
            <span>Share</span>
          </button>
          <div class="share-menu" role="menu"></div>
        </span>
      `;

      async function sendFeedback(rating, comment) {
        try {
          await fetch("/feedback", {
            method: "POST",
            headers: { "Content-Type": "application/json" },
            body: JSON.stringify({
              rating, reply: replyText, comment: comment || "",
              interaction_id: interactionId || null,
            }),
          });
        } catch (err) {
          console.error("Feedback error:", err);
        }
      }

      const buttons = wrap.querySelectorAll(".feedback-btn");
      buttons.forEach(btn => {
        btn.addEventListener("click", async () => {
          if (btn.disabled) return;
          const rating = btn.dataset.rating;
          // Rated — retire the nudge for this reply
          wrap.dataset.rated = "1";
          wrap.classList.remove("nudged");
          const ownNudge = wrap.parentElement
            ? wrap.parentElement.querySelector(".rate-nudge") : null;
          if (ownNudge) ownNudge.remove();

          // A rating can always be changed — people mis-tap, and a wrong
          // rating is worse than none because it teaches the wrong lesson.
          // Clear any previous selection and its UI before applying this one.
          buttons.forEach(b => {
            b.classList.remove("selected-up", "selected-down");
            b.disabled = false;
          });
          const oldThanks = wrap.querySelector(".feedback-thanks");
          if (oldThanks) oldThanks.remove();
          const oldComment = wrap.querySelector(".feedback-comment");
          if (oldComment) oldComment.remove();

          if (rating === "up") {
            // Thumbs up: simple submit, no comment needed
            btn.classList.add("selected-up");
            await sendFeedback("up", "");
            const thanks = document.createElement("span");
            thanks.className = "feedback-thanks";
            thanks.textContent = "Thanks for the feedback \u2014 tap either thumb to change it";
            wrap.appendChild(thanks);
          } else {
            // Thumbs down: record it straight away, then invite a comment.
            // Waiting for the comment meant a participant who switched from
            // up to down and didn't type anything left "up" on the server —
            // the opposite of what they meant.
            btn.classList.add("selected-down");
            await sendFeedback("down", "");

            const commentBox = document.createElement("div");
            commentBox.className = "feedback-comment";
            commentBox.innerHTML = `
              <label>What was wrong? (optional)</label>
              <textarea placeholder="Tell us what would have been more helpful..." maxlength="2000"></textarea>
              <div class="feedback-comment-actions">
                <button type="button" class="feedback-comment-btn" data-action="submit">Submit feedback</button>
                <button type="button" class="feedback-comment-btn secondary" data-action="skip">Skip</button>
              </div>
            `;
            wrap.appendChild(commentBox);

            const textarea = commentBox.querySelector("textarea");
            textarea.focus();

            const submitBtn = commentBox.querySelector('[data-action="submit"]');
            const skipBtn = commentBox.querySelector('[data-action="skip"]');

            async function finalize(commentText) {
              submitBtn.disabled = true;
              skipBtn.disabled = true;
              textarea.disabled = true;
              await sendFeedback("down", commentText);
              commentBox.innerHTML = '<span class="feedback-thanks">Thanks for the feedback \u2014 tap either thumb to change it</span>';
            }

            submitBtn.addEventListener("click", () => finalize(textarea.value.trim()));
            skipBtn.addEventListener("click", () => finalize(""));
            textarea.addEventListener("keydown", (e) => {
              if (e.key === "Enter" && (e.metaKey || e.ctrlKey)) {
                finalize(textarea.value.trim());
              }
            });
          }
        });
      });
      // === SPEAK button (browser text-to-speech) ===
      const speakBtn = wrap.querySelector(".speak-btn");
      const speakLabel = speakBtn ? speakBtn.querySelector(".speak-label") : null;
      if (!J3PSpeech.supported && speakBtn) {
        // Browser doesn't support speech synthesis — hide the button entirely
        speakBtn.style.display = "none";
      } else if (speakBtn) {
        // Strip markdown so the reader doesn't literally say "star star bold star star"
        const stripMarkdown = (t) => (t || "")
          .replace(/```[\s\S]*?```/g, " ")             // fenced code blocks
          .replace(/\*\*(.+?)\*\*/g, "$1")             // **bold**
          .replace(/__([^_]+?)__/g, "$1")              // __bold__
          .replace(/(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)/g, "$1")  // *italic*
          .replace(/(?<!_)_(?!_)([^_]+?)(?<!_)_(?!_)/g, "$1")     // _italic_
          .replace(/^#{1,6}\s+/gm, "")                 // # headers
          .replace(/`([^`]+?)`/g, "$1")                // `inline code`
          .replace(/^[-*+]\s+/gm, "")                  // bullet markers
          .replace(/^\d+\.\s+/gm, "")                  // numbered list markers
          .replace(/\[([^\]]+?)\]\([^)]+?\)/g, "$1")   // [text](url) -> text
          .replace(/^-{3,}\s*$/gm, ". ")               // horizontal rules
          .replace(/\n{2,}/g, ". ")                    // paragraph breaks -> pause
          .replace(/\n/g, " ")                         // remaining newlines
          .replace(/\s{2,}/g, " ")
          .trim();

        function resetSpeakUI() {
          speakBtn.classList.remove("speaking", "paused");
          if (speakLabel) speakLabel.textContent = "Speak";
          setAvatarSpeaking(msgDiv, false);
          Presence.set("idle");
        }
        function markSpeaking() {
          speakBtn.classList.add("speaking");
          speakBtn.classList.remove("paused");
          // Where pause isn't available the control is stop-only, so say so
          if (speakLabel) speakLabel.textContent = J3PSpeech.canPause ? "Speaking" : "Stop";
          clearAllAvatarStates();
          setAvatarSpeaking(msgDiv, true);
          Presence.set("speaking");
        }

        // On platforms without pause support, make the intent clear up front
        if (!J3PSpeech.canPause) {
          speakBtn.title = "Read answer aloud (tap again to stop)";
        }

        speakBtn.addEventListener("click", () => {
          // This button is the active one → toggle pause/resume, or stop
          if (window.__activeSpeakBtn === speakBtn) {
            if (!J3PSpeech.canPause) {
              J3PSpeech.stop();
              resetSpeakUI();
              window.__activeSpeakBtn = null;
              return;
            }
            if (J3PSpeech.isPaused()) {
              J3PSpeech.resume();
              markSpeaking();
            } else {
              J3PSpeech.pause();
              speakBtn.classList.remove("speaking");
              speakBtn.classList.add("paused");
              if (speakLabel) speakLabel.textContent = "Paused";
            }
            return;
          }

          // A different message was playing — reset its button first
          if (window.__activeSpeakBtn && window.__activeSpeakBtn !== speakBtn) {
            const prev = window.__activeSpeakBtn;
            prev.classList.remove("speaking", "paused");
            const prevLabel = prev.querySelector(".speak-label");
            if (prevLabel) prevLabel.textContent = "Speak";
          }

          const ok = J3PSpeech.play(stripMarkdown(replyText), {
            onStart: () => {
              window.__activeSpeakBtn = speakBtn;
              markSpeaking();
            },
            onEnd: () => {
              resetSpeakUI();
              if (window.__activeSpeakBtn === speakBtn) window.__activeSpeakBtn = null;
            },
            onError: (reason) => {
              console.error("Speech error:", reason);
              resetSpeakUI();
              if (window.__activeSpeakBtn === speakBtn) window.__activeSpeakBtn = null;
              if (reason === "blocked" && speakLabel) {
                // Most common cause on iPhone/iPad is the physical silent switch
                speakLabel.textContent = J3PSpeech.isIOS ? "Check mute" : "Unavailable";
                speakBtn.title = J3PSpeech.isIOS
                  ? "No audio — check the side silent switch and volume, then try again"
                  : "Speech is unavailable in this browser";
                setTimeout(() => { if (speakLabel) speakLabel.textContent = "Speak"; }, 3200);
              }
            },
          }, { fromGesture: true });
          if (!ok) resetSpeakUI();
        });
      }

      // === COPY button ===
      const copyBtn = wrap.querySelector(".copy-btn");
      const copyLabel = copyBtn.querySelector(".copy-label");
      copyBtn.addEventListener("click", async () => {
        try {
          if (navigator.clipboard && navigator.clipboard.writeText) {
            await navigator.clipboard.writeText(replyText);
          } else {
            // Fallback for older browsers / non-HTTPS contexts
            const ta = document.createElement("textarea");
            ta.value = replyText; ta.style.position = "fixed"; ta.style.opacity = "0";
            document.body.appendChild(ta); ta.select();
            document.execCommand("copy"); document.body.removeChild(ta);
          }
          copyBtn.classList.add("copied");
          copyLabel.textContent = "Copied";
          setTimeout(() => {
            copyBtn.classList.remove("copied");
            copyLabel.textContent = "Copy";
          }, 1600);
        } catch (err) {
          console.error("Copy failed:", err);
          copyLabel.textContent = "Failed";
          setTimeout(() => { copyLabel.textContent = "Copy"; }, 1600);
        }
      });

      // === SAVE / DOWNLOAD button ===
      const downloadBtn = wrap.querySelector(".download-btn");
      const downloadMenu = wrap.querySelector(".download-menu");
      const downloadLabel = downloadBtn.querySelector(".download-label");

      downloadBtn.addEventListener("click", () => {
        // Close the share menu if it's open, then toggle this one
        const sm = wrap.querySelector(".share-wrap > .share-menu");
        if (sm) sm.classList.remove("open");
        downloadMenu.classList.toggle("open");
        if (downloadMenu.classList.contains("open")) positionMenu(downloadMenu);
      });

      // When the reply holds more than one deliverable, rebuild the menu so
      // each document can be downloaded separately in its own format.
      const multiDocs = Array.isArray(documents) && documents.length > 1;
      if (multiDocs) {
        const FMT_LABEL = { docx: "Word", pptx: "PowerPoint", xlsx: "Excel", pdf: "PDF" };
        downloadMenu.classList.add("multi-doc");
        downloadMenu.innerHTML = documents.map(doc => {
          const safeTitle = String(doc.title || "Document")
            .replace(/&/g, "&amp;").replace(/</g, "&lt;").replace(/>/g, "&gt;");
          const chips = ["docx", "pptx", "xlsx", "pdf"].map(f => {
            const primary = f === doc.suggested ? " primary" : "";
            return `<button type="button" class="fmt-chip${primary}" ` +
                   `data-fmt="${f}" data-part="${doc.index}" role="menuitem">` +
                   `${FMT_LABEL[f]}</button>`;
          }).join("");
          return `<div class="doc-group">
                    <div class="doc-title">${safeTitle}</div>
                    <div class="doc-formats">${chips}</div>
                  </div>`;
        }).join("");
      }

      downloadMenu.querySelectorAll("[data-fmt]").forEach(item => {
        item.addEventListener("click", async () => {
          const fmt = item.dataset.fmt;
          const part = item.dataset.part !== undefined ? Number(item.dataset.part) : undefined;
          downloadMenu.classList.remove("open");
          downloadBtn.classList.add("copied");
          downloadLabel.textContent = "Building\u2026";
          try {
            await downloadExport(fmt, replyText, part);
            downloadLabel.textContent = "Saved";
          } catch (err) {
            console.error("Export failed:", err);
            downloadLabel.textContent = "Failed";
            alert(err.message || "Could not generate the document.");
          } finally {
            setTimeout(() => {
              downloadBtn.classList.remove("copied");
              downloadLabel.textContent = "Save";
            }, 1800);
          }
        });
      });

      // === SHARE button ===
      const shareBtn = wrap.querySelector(".share-btn");
      const shareMenu = wrap.querySelector(".share-wrap > .share-menu");
      const shareTitle = "From J3P Advisor";
      // Truncate share text to keep social/SMS messages under sane limits
      const shareText = replyText.length > 600
        ? replyText.slice(0, 600).trim() + "…"
        : replyText;
      const shareUrl = window.location.origin;

      shareBtn.addEventListener("click", async () => {
        // Try native share sheet first (mobile + modern desktop browsers)
        // The OS share sheet is the better experience on phones and tablets,
        // where it lists the mail apps actually installed. On desktop it only
        // offers registered share extensions — macOS shows AirDrop, Mail and
        // Notes but never Outlook — so desktop always gets our own menu.
        const ua = navigator.userAgent || "";
        const isMobileDevice =
          /Android|iPhone|iPod/.test(ua) ||
          (/iPad|Macintosh/.test(ua) && navigator.maxTouchPoints > 1);

        if (navigator.share && isMobileDevice) {
          try {
            await navigator.share({ title: shareTitle, text: shareText, url: shareUrl });
            return;
          } catch (err) {
            // User cancelled or share failed — fall through to menu
            if (err.name === "AbortError") return;
          }
        }
        // Our own menu — the only place Outlook can be offered
        if (shareMenu.classList.contains("open")) {
          shareMenu.classList.remove("open");
          return;
        }
        const emailSubject = encodeURIComponent(shareTitle);
        const emailBody = encodeURIComponent(shareText + "\\n\\n" + shareUrl);
        const smsBody = encodeURIComponent(shareText + " " + shareUrl);
        const twText = encodeURIComponent(shareText.slice(0, 240) + " " + shareUrl);
        // Outlook on the web compose deeplink. Work/school accounts use
        // outlook.office.com; personal accounts use outlook.live.com.
        const outlookWork = "https://outlook.office.com/mail/deeplink/compose"
                          + `?subject=${emailSubject}&body=${emailBody}`;
        const liUrl = encodeURIComponent(shareUrl);
        const fbUrl = encodeURIComponent(shareUrl);

        shareMenu.innerHTML = `
          <a href="mailto:?subject=${emailSubject}&body=${emailBody}" role="menuitem">
            <svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M4 4h16c1.1 0 2 .9 2 2v12c0 1.1-.9 2-2 2H4c-1.1 0-2-.9-2-2V6c0-1.1.9-2 2-2z"/><polyline points="22,6 12,13 2,6"/></svg>
            Email
          </a>
          <a href="${outlookWork}" target="_blank" rel="noopener" role="menuitem">
            <svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><rect x="2" y="5" width="12" height="14" rx="2"/><path d="M14 8h6a2 2 0 0 1 2 2v7a2 2 0 0 1-2 2h-6"/><path d="M5 9.5 8 12l3-2.5"/></svg>
            Outlook
          </a>
          <a href="sms:?body=${smsBody}" role="menuitem">
            <svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M21 15a2 2 0 0 1-2 2H7l-4 4V5a2 2 0 0 1 2-2h14a2 2 0 0 1 2 2z"/></svg>
            Text message
          </a>
          <a href="https://twitter.com/intent/tweet?text=${twText}" target="_blank" rel="noopener" role="menuitem">
            <svg viewBox="0 0 24 24" fill="currentColor"><path d="M18.244 2.25h3.308l-7.227 8.26 8.502 11.24H16.17l-5.214-6.817L4.99 21.75H1.68l7.73-8.835L1.254 2.25H8.08l4.713 6.231zm-1.161 17.52h1.833L7.084 4.126H5.117z"/></svg>
            X / Twitter
          </a>
          <a href="https://www.linkedin.com/sharing/share-offsite/?url=${liUrl}" target="_blank" rel="noopener" role="menuitem">
            <svg viewBox="0 0 24 24" fill="currentColor"><path d="M20.447 20.452h-3.554v-5.569c0-1.328-.027-3.037-1.852-3.037-1.853 0-2.136 1.445-2.136 2.939v5.667H9.351V9h3.414v1.561h.046c.477-.9 1.637-1.85 3.37-1.85 3.601 0 4.267 2.37 4.267 5.455v6.286zM5.337 7.433a2.062 2.062 0 01-2.063-2.065 2.063 2.063 0 112.063 2.065zm1.782 13.019H3.555V9h3.564v11.452zM22.225 0H1.771C.792 0 0 .774 0 1.729v20.542C0 23.227.792 24 1.771 24h20.451C23.2 24 24 23.227 24 22.271V1.729C24 .774 23.2 0 22.222 0h.003z"/></svg>
            LinkedIn
          </a>
          <a href="https://www.facebook.com/sharer/sharer.php?u=${fbUrl}" target="_blank" rel="noopener" role="menuitem">
            <svg viewBox="0 0 24 24" fill="currentColor"><path d="M24 12.073c0-6.627-5.373-12-12-12s-12 5.373-12 12c0 5.99 4.388 10.954 10.125 11.854v-8.385H7.078v-3.47h3.047V9.43c0-3.007 1.792-4.669 4.533-4.669 1.312 0 2.686.235 2.686.235v2.953H15.83c-1.491 0-1.956.925-1.956 1.874v2.25h3.328l-.532 3.47h-2.796v8.385C19.612 23.027 24 18.062 24 12.073z"/></svg>
            Facebook
          </a>
          <button type="button" data-action="copy-link" role="menuitem">
            <svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M10 13a5 5 0 0 0 7.54.54l3-3a5 5 0 0 0-7.07-7.07l-1.72 1.71"/><path d="M14 11a5 5 0 0 0-7.54-.54l-3 3a5 5 0 0 0 7.07 7.07l1.71-1.71"/></svg>
            Copy link
          </button>
        `;
        shareMenu.classList.add("open");
        positionMenu(shareMenu);

        const copyLinkBtn = shareMenu.querySelector('[data-action="copy-link"]');
        if (copyLinkBtn) {
          copyLinkBtn.addEventListener("click", async () => {
            try {
              await navigator.clipboard.writeText(shareUrl);
              copyLinkBtn.textContent = "Link copied";
            } catch (err) { console.error(err); }
            setTimeout(() => shareMenu.classList.remove("open"), 700);
          });
        }
      });

      // Close both popover menus when clicking outside
      document.addEventListener("click", (e) => {
        if (!wrap.contains(e.target)) {
          shareMenu.classList.remove("open");
          downloadMenu.classList.remove("open");
        }
      });

      msgDiv.appendChild(wrap);
    }

    // Attach file / folder handling
    const fileInput = document.getElementById("file-input");
    const folderInput = document.getElementById("folder-input-chat");
    const attachBtn = document.getElementById("attach-btn");
    const folderBtn = document.getElementById("folder-btn");
    const attachedFileDiv = document.getElementById("attached-file");
    const attachedFileName = document.getElementById("attached-file-name");
    const attachedFileSize = document.getElementById("attached-file-size");
    const removeFileBtn = document.getElementById("remove-file-btn");
    const MAX_FILE_MB = {{ cfg.max_upload_mb }};
    const MAX_IMAGE_MB = {{ cfg.max_image_mb }};
    const MAX_FOLDER_FILES = 20;
    const DOC_RE = /\.(pdf|docx|xlsx|xlsm|pptx|csv|tsv|txt|md|rtf)$/i;
    // Track state: either a single file OR a list of folder files, never both
    let attachedFolderFiles = [];   // list of File objects when a folder is picked
    let attachedFolderName = "";

    function formatFileSize(bytes) {
      if (bytes < 1024) return bytes + " B";
      if (bytes < 1024 * 1024) return (bytes / 1024).toFixed(1) + " KB";
      return (bytes / (1024 * 1024)).toFixed(1) + " MB";
    }
    // Files picked with the paperclip accumulate here across multiple
    // selections. Reading fileInput.files directly would mean each new pick
    // silently replaced the previous one.
    let attachedFiles = [];
    const MAX_ATTACHMENTS = 20;

    function clearAttachment() {
      fileInput.value = "";
      folderInput.value = "";
      attachedFiles = [];
      attachedFolderFiles = [];
      attachedFolderName = "";
      attachedFileDiv.classList.remove("visible");
    }

    // Refresh the pill to describe everything currently attached
    function refreshAttachmentChip() {
      if (attachedFiles.length === 0) {
        attachedFileDiv.classList.remove("visible");
        return;
      }
      if (attachedFiles.length === 1) {
        attachedFileName.textContent = attachedFiles[0].name;
        attachedFileSize.textContent = "· " + formatFileSize(attachedFiles[0].size);
      } else {
        const total = attachedFiles.reduce((sum, f) => sum + f.size, 0);
        const names = attachedFiles.slice(0, 3).map(f => f.name).join(", ");
        attachedFileName.textContent =
          `${attachedFiles.length} files: ${names}${attachedFiles.length > 3 ? "…" : ""}`;
        attachedFileSize.textContent = "· " + formatFileSize(total);
      }
      attachedFileDiv.classList.add("visible");
    }

    attachBtn.addEventListener("click", () => {
      // Folder mode and paperclip mode are mutually exclusive, but previously
      // attached individual files are kept so more can be added to the batch.
      folderInput.value = "";
      attachedFolderFiles = [];
      attachedFolderName = "";
      fileInput.click();
    });
    folderBtn.addEventListener("click", () => {
      clearAttachment();          // clear any prior single-file selection
      folderInput.click();
    });
    removeFileBtn.addEventListener("click", clearAttachment);

    fileInput.addEventListener("change", () => {
      const picked = Array.from(fileInput.files || []);
      if (picked.length === 0) return;
      // Merge with anything already attached, ignoring exact duplicates
      const seen = new Set(attachedFiles.map(f => f.name + ":" + f.size));
      const fresh = picked.filter(f => !seen.has(f.name + ":" + f.size));
      const all = fresh;
      // Validate every file
      const okRe = /\.(pdf|docx|xlsx|xlsm|pptx|csv|tsv|txt|md|rtf|jpe?g|png|gif|webp)$/i;
      // Legacy Office formats (.ppt/.doc/.xls) can't be read by the server, so
      // catch them here and say exactly how to convert, before any upload.
      const legacyRe = /\.(ppt|doc|xls)$/i;
      const legacy = all.filter(f => legacyRe.test(f.name));
      if (legacy.length) {
        const LEGACY_HINT = {
          ppt: ["PowerPoint", ".pptx"], doc: ["Word", ".docx"], xls: ["Excel", ".xlsx"],
        };
        const lines = legacy.map(f => {
          const ext = f.name.split(".").pop().toLowerCase();
          const [app, target] = LEGACY_HINT[ext] || ["Office", ".docx"];
          return `• ${f.name}\n   Open it in ${app}, choose File \u203a Save As, ` +
                 `pick ${target}, then attach the new file.`;
        });
        alert("These are in an older Office format that can't be read:\n\n" +
              lines.join("\n\n"));
      }
      const nonLegacy = all.filter(f => !legacyRe.test(f.name));

      // Reject only the offending files — anything already attached stays put
      const bad = nonLegacy.filter(f => !okRe.test(f.name));
      const imageRe = /\.(jpe?g|png|gif|webp)$/i;
      const oversized = nonLegacy.filter(f => okRe.test(f.name)
                                       && !imageRe.test(f.name)
                                       && f.size > MAX_FILE_MB * 1024 * 1024);
      const bigImages = nonLegacy.filter(f => imageRe.test(f.name)
                                       && f.size > MAX_IMAGE_MB * 1024 * 1024);
      const problems = [];
      if (bad.length) {
        problems.push("Unsupported type: " + bad.map(f => f.name).join(", ") +
          "\nSupported: PDF, Word, Excel, PowerPoint, CSV, TXT, MD, RTF, and images.");
      }
      if (oversized.length) {
        problems.push(`Over ${MAX_FILE_MB} MB: ` + oversized.map(f => f.name).join(", "));
      }
      if (bigImages.length) {
        problems.push(`Images must be under ${MAX_IMAGE_MB} MB: `
                      + bigImages.map(f => f.name).join(", "));
      }
      const rejected = new Set([...bad, ...oversized, ...bigImages]);
      const accepted = nonLegacy.filter(f => !rejected.has(f));
      if (problems.length) {
        alert(problems.join("\n\n") +
              (accepted.length ? `\n\nThe other ${accepted.length} file(s) were attached.` : ""));
      }
      if (accepted.length === 0) {
        fileInput.value = "";
        refreshAttachmentChip();      // keep whatever was already there
        return;
      }
      attachedFiles = attachedFiles.concat(accepted);
      if (attachedFiles.length > MAX_ATTACHMENTS) {
        alert(`You can attach up to ${MAX_ATTACHMENTS} files at once. ` +
              `Keeping the first ${MAX_ATTACHMENTS}.`);
        attachedFiles = attachedFiles.slice(0, MAX_ATTACHMENTS);
      }
      // Reset the input so re-picking the same file later still fires a change
      fileInput.value = "";
      refreshAttachmentChip();
    });

    folderInput.addEventListener("change", () => {
      const all = Array.from(folderInput.files || []);
      // Filter to documents only (images not supported in folder-attach mode)
      const docs = all.filter(f => DOC_RE.test(f.name));
      if (docs.length === 0) {
        const legacyInFolder = all.filter(f => /\.(ppt|doc|xls)$/i.test(f.name));
        alert(legacyInFolder.length
          ? "This folder only contains older Office files (.ppt/.doc/.xls), which " +
            "can't be read. Re-save them as .pptx/.docx/.xlsx and try again."
          : "No supported documents found in this folder (PDF, Word, Excel, " +
            "PowerPoint, CSV, TXT, MD, RTF).");
        clearAttachment(); return;
      }
      const capped = docs.slice(0, MAX_FOLDER_FILES);
      // webkitRelativePath looks like "FolderName/sub/file.pdf" — extract the top-level folder name
      const first = capped[0].webkitRelativePath || capped[0].name;
      const folderName = first.split("/")[0] || "Folder";
      attachedFolderFiles = capped;
      attachedFolderName = folderName;
      let msg = `${folderName} · ${capped.length} file${capped.length === 1 ? "" : "s"}`;
      if (docs.length > MAX_FOLDER_FILES) msg += ` (first ${MAX_FOLDER_FILES} only)`;
      if (docs.length < all.length) msg += ` · ${all.length - docs.length} skipped`;
      attachedFileName.textContent = msg;
      attachedFileSize.textContent = "";
      attachedFileDiv.classList.add("visible");
    });

    form.addEventListener("submit", async (e) => {
      e.preventDefault();
      const text = input.value.trim();
      const wasFarewell = looksLikeFarewell(text);
      const paperclipFiles = attachedFiles.slice();
      const folderFiles = attachedFolderFiles.slice();
      const folderName = attachedFolderName;
      // Require at least one of: text, single/multi file(s), or folder
      if (!text && paperclipFiles.length === 0 && folderFiles.length === 0) return;

      // Already waiting on a reply? Hold this one rather than running two
      // requests at once, and show it in the transcript so the participant
      // can see it was received.
      if (awaitingReply) {
        addMessage(text, "user");
        queuedMessages.push({ text: text });
        input.value = "";
        clearAttachment();
        updateQueueHint();
        input.focus();
        return;
      }

      // Build the message shown in chat
      let displayText;
      if (folderFiles.length > 0) {
        displayText = (text || "[No message]") +
          `\n\n📂 Attached folder: ${folderName} (${folderFiles.length} file${folderFiles.length === 1 ? "" : "s"})`;
      } else if (paperclipFiles.length === 1) {
        displayText = (text || "[No message]") + `\n\n📎 Attached: ${paperclipFiles[0].name}`;
      } else if (paperclipFiles.length > 1) {
        const names = paperclipFiles.slice(0, 3).map(f => f.name).join(", ") + (paperclipFiles.length > 3 ? "…" : "");
        displayText = (text || "[No message]") + `\n\n📎 Attached ${paperclipFiles.length} files: ${names}`;
      } else {
        displayText = text;
      }
      if (sendingQueued) sendingQueued = false;   // already shown when queued
      else addMessage(displayText, "user");

      input.value = "";
      const paperclipFilesForRequest = paperclipFiles.slice();
      const folderFilesForRequest = folderFiles;
      const folderNameForRequest = folderName;
      clearAttachment();
      // The send button stays live: a follow-up typed while the advisor is
      // still working is queued and sent as soon as the reply lands, so the
      // participant never has to wait to add a thought.
      awaitingReply = true;
      updateQueueHint();
      const thinking = addMessage("Thinking…", "assistant typing");
      setAvatarThinking(true);

      try {
        let res;
        if (folderFilesForRequest.length > 0) {
          // Folder upload — send multiple files under the "files" field
          const fd = new FormData();
          fd.append("message", text || "Please review these attached documents.");
          fd.append("folder_name", folderNameForRequest);
          folderFilesForRequest.forEach(f => fd.append("files", f, f.name));
          res = await fetch("/chat", { method: "POST", body: fd });
        } else if (paperclipFilesForRequest.length === 1) {
          // Single-file path (preserves image vision handling for one image)
          const fd = new FormData();
          fd.append("message", text || "Please review this attached file.");
          fd.append("file", paperclipFilesForRequest[0]);
          res = await fetch("/chat", { method: "POST", body: fd });
        } else if (paperclipFilesForRequest.length > 1) {
          // Multi-file paperclip — use the same multi-file field as folder,
          // so the backend concatenates docs and adds each image as a vision block.
          const fd = new FormData();
          fd.append("message", text || "Please review these attached files.");
          fd.append("attachment_label", "Attachments");   // shown in chat/logs
          paperclipFilesForRequest.forEach(f => fd.append("files", f, f.name));
          res = await fetch("/chat", { method: "POST", body: fd });
        } else {
          res = await fetch("/chat", {
            method: "POST", headers: { "Content-Type": "application/json" },
            body: JSON.stringify({ message: text }),
          });
        }
        // Railway's proxy returns plain text like "upstream error" on a 502,
        // so parsing blindly as JSON produced "Unexpected token 'u'".
        let data;
        const raw = await res.text();
        try {
          data = JSON.parse(raw);
        } catch (parseErr) {
          removeMessage(thinking);
        setAvatarThinking(false);
          const detail = (raw || "").trim().slice(0, 120);
          const friendly = res.status === 502 || /upstream/i.test(detail)
            ? "The server didn't finish that one — usually a large attachment or a very long conversation. Try again, or start a New conversation if it keeps happening."
            : `The server returned an unexpected response (${res.status}). Please try again.`;
          showRetry(friendly, text, paperclipFilesForRequest);
          return;
        }
        removeMessage(thinking);
        setAvatarThinking(false);
        if (data.reply) {
          const msgDiv = addMessage(data.reply, "assistant", true,
                                    data.interaction_id || null, data.documents || null);
          // Nothing downloads on its own. When the reply looks like a
          // deliverable, offer a download and let the user decide.
          setAvatarResponding(msgDiv);
          clearRatingNudge();
          if (wasFarewell) setTimeout(() => {
            showPlanOffer();
            showRatingNudge("farewell");
          }, 900);
          else armIdleNudge();
          const docs = data.documents || [];
          if (docs.length > 1 && data.separate_files) {
            offerExport(msgDiv, data.reply, data.export_format, docs);
          } else if (data.export_format) {
            offerExport(msgDiv, data.reply, data.export_format, null);
          }
        }
        else addMessage("Error: " + (data.error || "Unknown error"), "assistant");
      } catch (err) {
        removeMessage(thinking);
        setAvatarThinking(false);
        showRetry(
          "Couldn't reach the server. Check your connection and try again.",
          text, paperclipFilesForRequest);
      } finally {
        awaitingReply = false;
        updateQueueHint();
        input.focus();
        flushQueue();
      }
    });

    let resetArmed = false;
    resetBtn.addEventListener("click", async () => {
      // If the last reply is unrated, ask once before wiping the conversation.
      // Showing the nudge and then clearing the chat would be pointless, so the
      // first click surfaces it and the second proceeds.
      if (!resetArmed && lastUnratedReply()) {
        showPlanOffer();
        showRatingNudge("farewell");
        resetArmed = true;
        setTimeout(() => { resetArmed = false; }, 12000);
        return;
      }
      resetArmed = false;
      // Stop any in-progress speech before wiping the chat
      J3PSpeech.stop();
      window.__activeSpeakBtn = null;
      await fetch("/reset", { method: "POST" });
      chat.innerHTML = "";
      const div = document.createElement("div");
      div.className = "msg assistant";
      div.textContent = OPENING;
      chat.appendChild(div);
      input.focus();
    });

    // -------------------------------------------------------------
    // Voice input — press and hold to record, release to send
    // -------------------------------------------------------------
    // Hold the mic, speak, release: the transcript is placed in the composer
    // and submitted. Slide away from the button before releasing to cancel.
    const micBtn = document.getElementById("mic-btn");
    const SR = window.SpeechRecognition || window.webkitSpeechRecognition;
    const voiceHint = document.getElementById("voice-hint");
    const voiceHintText = document.getElementById("voice-hint-text");
    const voiceTimer = document.getElementById("voice-timer");

    if (!SR) {
      micBtn.classList.add("unsupported");
    } else {
      const recognition = new SR();
      recognition.continuous = true;        // keep listening for the whole hold
      recognition.interimResults = true;
      recognition.lang = "en-US";

      const MIN_HOLD_MS = 350;              // shorter than this is a stray tap
      const CANCEL_DISTANCE = 90;           // px away from the button = cancel

      let holding = false;
      let cancelled = false;
      let tooShort = false;
      let startedAt = 0;
      let baseText = "";
      let finalText = "";
      let interimText = "";
      let timerId = null;
      let originRect = null;

      function setHint(msg, danger) {
        if (!voiceHintText) return;
        voiceHintText.textContent = msg;
        voiceHint.classList.toggle("cancel", !!danger);
      }

      function showHint(show) {
        if (!voiceHint) return;
        voiceHint.classList.toggle("visible", !!show);
      }

      function tickTimer() {
        if (!voiceTimer) return;
        const secs = Math.floor((Date.now() - startedAt) / 1000);
        const m = String(Math.floor(secs / 60)).padStart(1, "0");
        const ss = String(secs % 60).padStart(2, "0");
        voiceTimer.textContent = `${m}:${ss}`;
      }

      function beginHold(e) {
        if (holding) return;
        // Don't fight text selection or page scroll while holding
        if (e.cancelable) e.preventDefault();
        holding = true;
        cancelled = false;
        tooShort = false;
        startedAt = Date.now();
        finalText = "";
        interimText = "";
        baseText = input.value.trim();
        originRect = micBtn.getBoundingClientRect();

        // Speaking while recording would feed the mic back into itself
        if (window.J3PSpeech) J3PSpeech.stop();

        micBtn.classList.add("recording");
        setHint("Listening\u2026 release to send", false);
        showHint(true);
        tickTimer();
        timerId = setInterval(tickTimer, 250);

        try { recognition.start(); }
        catch (err) { /* already started — harmless */ }
      }

      function moveDuringHold(e) {
        if (!holding || !originRect) return;
        const pt = e.touches ? e.touches[0] : e;
        const cx = originRect.left + originRect.width / 2;
        const cy = originRect.top + originRect.height / 2;
        const dist = Math.hypot(pt.clientX - cx, pt.clientY - cy);
        const wasCancelled = cancelled;
        cancelled = dist > CANCEL_DISTANCE;
        if (cancelled !== wasCancelled) {
          setHint(cancelled ? "Release to cancel" : "Listening\u2026 release to send", cancelled);
          micBtn.classList.toggle("cancelling", cancelled);
        }
      }

      function endHold(e) {
        if (!holding) return;
        holding = false;
        const heldFor = Date.now() - startedAt;
        if (timerId) { clearInterval(timerId); timerId = null; }
        micBtn.classList.remove("recording", "cancelling");

        try { recognition.stop(); } catch (err) {}

        if (cancelled) {
          setHint("Cancelled", true);
          input.value = baseText;
          setTimeout(() => showHint(false), 900);
          return;
        }
        if (heldFor < MIN_HOLD_MS) {
          tooShort = true;
          setHint("Hold the mic to record", false);
          input.value = baseText;
          setTimeout(() => showHint(false), 1400);
          return;
        }
        setHint("Transcribing\u2026", false);
        // recognition "end" fires shortly after stop(); submission happens there
      }

      micBtn.addEventListener("pointerdown", beginHold);
      window.addEventListener("pointermove", moveDuringHold, { passive: true });
      window.addEventListener("pointerup", endHold);
      window.addEventListener("pointercancel", endHold);
      // Holding the mic shouldn't open a context menu on mobile
      micBtn.addEventListener("contextmenu", (e) => e.preventDefault());

      recognition.addEventListener("result", (event) => {
        interimText = "";
        for (let i = event.resultIndex; i < event.results.length; i++) {
          const chunk = event.results[i][0].transcript;
          if (event.results[i].isFinal) finalText += chunk;
          else interimText += chunk;
        }
        // Live preview in the composer while speaking
        const preview = (baseText ? baseText + " " : "") + (finalText + interimText).trim();
        input.value = preview;
      });

      recognition.addEventListener("end", () => {
        if (holding) return;                 // stopped early; a new hold is running
        const spoken = (finalText + interimText).trim();
        const message = ((baseText ? baseText + " " : "") + spoken).trim();

        if (cancelled || tooShort || !spoken) {
          input.value = baseText;
          if (!cancelled && !tooShort && !spoken) {
            setHint("Didn\u2019t catch that \u2014 try again", true);
          }
          setTimeout(() => showHint(false), 1400);
          return;
        }

        input.value = message;
        showHint(false);
        // Submit the transcription as the message
        if (typeof form.requestSubmit === "function") form.requestSubmit();
        else form.dispatchEvent(new Event("submit", { cancelable: true, bubbles: true }));
      });

      recognition.addEventListener("error", (event) => {
        holding = false;
        if (timerId) { clearInterval(timerId); timerId = null; }
        micBtn.classList.remove("recording", "cancelling");
        if (event.error === "not-allowed" || event.error === "service-not-allowed") {
          setHint("Microphone blocked \u2014 allow access in your browser", true);
          setTimeout(() => showHint(false), 3000);
        } else if (event.error === "no-speech" && !tooShort) {
          setHint("Didn\u2019t catch that \u2014 try again", true);
          setTimeout(() => showHint(false), 1600);
        } else {
          showHint(false);
        }
      });
    }
  </script>
</body>
</html>
"""


# ---------------------------------------------------------------------------
# Chat routes
# ---------------------------------------------------------------------------

# Flask stores the session in a signed cookie. Browsers hard-cap cookies at
# about 4 KB, and an oversized cookie is silently discarded — which wipes the
# entire conversation. Signing and base64 add overhead, so we budget well under
# that for the message payload itself.
# ---------------------------------------------------------------------------
# Conversation history
# ---------------------------------------------------------------------------
# History used to live in the Flask session cookie, which browsers cap at ~4 KB.
# That meant a single cover letter filled the budget and earlier turns were
# evicted — the advisor would then say it had no cover letter to revise. When
# Postgres is available the transcript is stored server-side instead, keyed by a
# random token held in the cookie, so a full working session stays intact.

CHAT_HISTORY_BUDGET = int(os.environ.get("CHAT_HISTORY_CHARS", "60000"))


def _history_ensure_table(conn):
    with conn.cursor() as cur:
        cur.execute("""
            CREATE TABLE IF NOT EXISTS chat_history (
                id         BIGSERIAL PRIMARY KEY,
                token      TEXT NOT NULL,
                role       TEXT NOT NULL,
                content    TEXT NOT NULL,
                created_at TIMESTAMPTZ NOT NULL DEFAULT NOW()
            )
        """)
        cur.execute("""
            CREATE INDEX IF NOT EXISTS chat_history_token_idx
            ON chat_history (token, id)
        """)
    conn.commit()


def _history_token() -> str:
    """Identifies whose transcript this is.

    Signed in, it's derived from the email so the conversation continues across
    visits and devices. Anonymous, it's a random per-browser token as before.
    """
    email = session.get("user_email", "")
    if email:
        import hashlib
        digest = hashlib.sha256(
            (app.secret_key + "|" + email.lower()).encode("utf-8")).hexdigest()[:32]
        return "u_" + digest
    token = session.get("chat_token")
    if not token:
        token = os.urandom(16).hex()
        session["chat_token"] = token
    return token


def load_history() -> list:
    """Recent turns for this session, oldest first, within the character budget."""
    conn = _settings_db_conn()
    if not conn:
        return session.get("messages", [])      # cookie fallback
    token = _history_token()
    rows = []
    try:
        _history_ensure_table(conn)
        with conn.cursor() as cur:
            cur.execute("""
                SELECT role, content FROM chat_history
                WHERE token = %s ORDER BY id DESC LIMIT 200
            """, (token,))
            used = 0
            for role, content in cur.fetchall():
                used += len(content or "")
                if used > CHAT_HISTORY_BUDGET and rows:
                    break
                rows.append({"role": role, "content": content})
        rows.reverse()
    except Exception as e:
        app.logger.error(f"[history] read failed: {e}")
        return session.get("messages", [])
    finally:
        conn.close()
    return rows


def append_history(role: str, content: str):
    """Persist one turn."""
    if not content:
        return
    conn = _settings_db_conn()
    if not conn:
        msgs = session.get("messages", [])
        msgs.append({"role": role, "content": content})
        session["messages"] = _fit_history(msgs)
        return
    token = _history_token()
    try:
        _history_ensure_table(conn)
        with conn.cursor() as cur:
            cur.execute(
                "INSERT INTO chat_history (token, role, content) VALUES (%s, %s, %s)",
                (token, role, content))
        conn.commit()
    except Exception as e:
        app.logger.error(f"[history] write failed: {e}")
    finally:
        conn.close()


def clear_history():
    """Wipe this session's transcript — used by New Conversation."""
    session["messages"] = []
    session.pop("specialty", None)
    session.pop("specialty_asked", None)
    # The profile deliberately survives: New Conversation clears the
    # transcript, not who the participant is.
    # advisor_slug and force_scheduling deliberately survive a reset: the
    # visitor is still on that advisor's link.
    conn = _settings_db_conn()
    token = session.get("chat_token")
    if conn and token:
        try:
            _history_ensure_table(conn)
            with conn.cursor() as cur:
                cur.execute("DELETE FROM chat_history WHERE token = %s", (token,))
            conn.commit()
        except Exception as e:
            app.logger.error(f"[history] clear failed: {e}")
        finally:
            conn.close()
    session.pop("chat_token", None)


SESSION_HISTORY_BUDGET = 2600   # characters of message content
SESSION_MSG_CAP = 1600          # max characters kept for any single message


# Phrases that mean "give me the actual file", mapped to an export format.
# When one of these appears in the user's message the file is generated and
# downloaded automatically rather than waiting for a SAVE click.
_EXPORT_PATTERNS = [
    ("pptx", r"\b(powerpoint|power point|pptx|slide deck|slides?|deck|presentation)\b"),
    ("xlsx", r"\b(excel|xlsx|spreadsheet|workbook)\b"),
    ("pdf",  r"\bpdf\b"),
    ("docx", r"\b(word|docx|word doc(?:ument)?|\.doc)\b"),
]

# Only auto-export when the user is actually asking for something to be made
# or delivered — not when they merely mention the word in passing.
_EXPORT_INTENT = (
    r"(creat|mak|writ|draft|build|generat|produc|prepar|export|download|"
    r"save|send|give|provide|turn .* into|convert|put .* in|as an? |in a |"
    r"i need|i want|can you)"
)


def looks_like_document(text: str) -> bool:
    """True if the reply is an actual deliverable rather than clarifying questions.

    A document has structure (headings, a salutation, or real length) and does
    not read as a list of questions back to the user.
    """
    import re as _r
    if not text:
        return False
    body = text.strip()

    # Structural signals of a real document
    has_headings = bool(_r.search(r"^#{1,6}\s+\S", body, _r.M))
    has_salutation = bool(_r.search(r"^\s*(dear|to the|re:)\b", body, _r.M | _r.I))
    is_long = len(body) > 1400

    # Interrogation signals
    questions = body.count("?")
    asks_upfront = bool(_r.search(
        r"(before i (can |could )?(start|begin|draft|write|build)|"
        r"i need (a few|some|the following)|"
        r"a few (quick )?(questions|answers)|"
        r"(could|can) you (tell|share|confirm|let me know)|"
        r"to (build|write|draft|create) (you )?a|"
        r"what (is|are) (the|your)\b)", body, _r.I))
    mostly_questions = questions >= 3 and len(body) < 2500

    if has_headings or has_salutation or is_long:
        # Even a structured reply is an interrogation if it opens by asking
        if asks_upfront and questions >= 2 and len(body) < 1800:
            return False
        return True
    if asks_upfront or mostly_questions or body.rstrip().endswith("?"):
        return False
    return len(body) > 700


# Cues that a format mention is being REJECTED rather than requested:
# "not a slide deck", "instead of a deck", "rather than a pdf", "no slides".
_NEGATION_CUES = (
    r"(?:not|no|don'?t|do not|doesn'?t|dont|instead of|rather than|"
    r"as opposed to|other than|without|avoid|skip|never)"
)
_NEG_WINDOW = 45          # characters before the format word to inspect


def _is_negated(low: str, match_start: int) -> bool:
    """True when a negation cue sits just before the format word."""
    import re as _r
    window = low[max(0, match_start - _NEG_WINDOW):match_start]
    # Only look at the current clause — a comma or period ends the negation
    clause = _r.split(r"[.;!?]|,\s+(?:and|but)\b", window)[-1]
    return bool(_r.search(_NEGATION_CUES + r"\s+(?:an?\s+|the\s+)?[\w\s]{0,18}$", clause))


def detect_export_format(text: str):
    """Return 'docx' | 'pptx' | 'xlsx' | 'pdf' if the user asked for that file.

    Formats the user explicitly rules out ("an executive summary, not a slide
    deck") are excluded rather than matched.
    """
    import re as _r
    if not text:
        return None
    low = text.lower()
    if not _r.search(_EXPORT_INTENT, low):
        return None

    wanted, rejected = [], set()
    for fmt, pattern in _EXPORT_PATTERNS:
        for m in _r.finditer(pattern, low):
            if _is_negated(low, m.start()):
                rejected.add(fmt)
            else:
                wanted.append(fmt)
    for fmt, _pattern in _EXPORT_PATTERNS:
        if fmt in wanted and fmt not in rejected:
            return fmt
    return None


# A request for a written deliverable, even when no file format is named.
# Used to enforce that the model produces the document instead of asking
# questions about it.
_DELIVERABLE_NOUNS = (
    r"(cover letter|letter of intent|letter|memo|executive summary|summary|"
    r"strategic plan|plan|deck|presentation|one[- ]pager|briefing|brief|"
    r"proposal|agenda|talking points|script|outline|draft|bio|biography|cv|"
    r"resume|curriculum vitae|recommendation|statement|report|packet|"
    r"materials|document|export|file)"
)
_DELIVERABLE_VERBS = (
    r"(creat|mak|writ|draft|build|generat|produc|prepar|revis|rework|rewrite|"
    r"tighten|shorten|expand|polish|update|turn .* into|put together|"
    r"help (?:me )?with|i need|i want|can you|give|send|share|provide|get me)"
)


def rejected_formats(text: str) -> set:
    """Formats the user explicitly ruled out ('not a slide deck')."""
    import re as _r
    out = set()
    if not text:
        return out
    low = text.lower()
    for fmt, pattern in _EXPORT_PATTERNS:
        for m in _r.finditer(pattern, low):
            if _is_negated(low, m.start()):
                out.add(fmt)
    return out


# ---------------------------------------------------------------------------
# Specialty awareness
# ---------------------------------------------------------------------------
# Replies were defaulting to surgical framing. Rather than guess, the advisor
# asks once — and once the participant says it, it's remembered for the session
# so they're never asked twice.

_SPECIALTY_PATTERNS = [
    (r"\b(orthopaed?ic|orthopedic|ortho)\b", "orthopaedics"),
    (r"\bsports medicine\b", "sports medicine"),
    (r"\b(cardiolog|cardiac)\w*\b", "cardiology"),
    (r"\bcardiothoracic|CT surgery\b", "cardiothoracic surgery"),
    (r"\b(neurosurg|neurolog)\w*\b", "neurosciences"),
    (r"\b(oncolog|hematolog|haematolog)\w*\b", "oncology"),
    (r"\b(paediatric|pediatric|peds)\w*\b", "paediatrics"),
    (r"\b(psychiatr|behavioral health|behavioural health)\w*\b", "psychiatry"),
    (r"\b(radiolog|imaging)\w*\b", "radiology"),
    (r"\bpatholog\w*\b", "pathology"),
    (r"\b(anesthesiolog|anaesthesiolog|anesthesia|anaesthesia)\w*\b", "anaesthesiology"),
    (r"\b(emergency medicine|emergency department|\bED\b|\bER\b)", "emergency medicine"),
    (r"\b(internal medicine|internist|hospitalist)\b", "internal medicine"),
    (r"\b(primary care|family medicine|family practice)\b", "primary care"),
    (r"\b(obstetric|gynecolog|gynaecolog|OB-?GYN)\w*\b", "obstetrics and gynaecology"),
    (r"\b(general surgery|surgical oncology|colorectal|vascular surgery)\b", "surgery"),
    (r"\b(urolog)\w*\b", "urology"),
    (r"\b(dermatolog)\w*\b", "dermatology"),
    (r"\b(ophthalmolog|eye institute)\w*\b", "ophthalmology"),
    (r"\b(otolaryngolog|ENT|head and neck)\b", "otolaryngology"),
    (r"\b(pulmonolog|critical care|intensivist|ICU)\b", "pulmonary and critical care"),
    (r"\b(gastroenterolog|\bGI\b)\w*\b", "gastroenterology"),
    (r"\b(nephrolog)\w*\b", "nephrology"),
    (r"\b(endocrinolog)\w*\b", "endocrinology"),
    (r"\b(rheumatolog)\w*\b", "rheumatology"),
    (r"\b(infectious disease|\bID\b)\b", "infectious disease"),
    (r"\b(physiatr|rehabilitation medicine|PM&R)\b", "rehabilitation medicine"),
    (r"\b(radiation oncolog)\w*\b", "radiation oncology"),
    (r"\b(palliative)\w*\b", "palliative care"),
    (r"\b(geriatric)\w*\b", "geriatrics"),
    (r"\b(plastic surgery|plastics)\b", "plastic surgery"),
    (r"\b(transplant)\w*\b", "transplant"),
    (r"\b(trauma surgery|trauma service)\b", "trauma surgery"),
]


def detect_specialty(text: str):
    """Pick up a specialty the participant mentions, so we needn't ask again."""
    if not text:
        return None
    low = str(text).lower()
    for pattern, name in _SPECIALTY_PATTERNS:
        if re.search(pattern, low, re.IGNORECASE):
            return name
    return None


def specialty_guidance() -> str:
    """Prompt text: use a known specialty, or ask once for it."""
    known = session.get("specialty")
    if known:
        return (
            "\n\n---\n"
            f"PARTICIPANT'S FIELD: they work in {known}. Write for that context "
            "specifically — use its language, its pressures and its examples. "
            "Do not ask them what their specialty is; you already know.\n"
        )
    if session.get("specialty_asked"):
        return (
            "\n\n---\n"
            "PARTICIPANT'S FIELD: unknown, and you have already asked once. Do "
            "not ask again. Write for a physician leader generally.\n"
        )
    return (
        "\n\n---\n"
        "PARTICIPANT'S FIELD: not yet known. Where the answer would genuinely "
        "differ by specialty, ask which field they work in — but ONLY as a short "
        "closing question AFTER giving a substantive answer. Never withhold help "
        "pending the answer, never open with the question, and never ask it when "
        "they have requested a document or a revision: produce that first. One "
        "sentence, e.g. 'What's your specialty? I can make this more specific.' "
        "Ask at most once in the conversation.\n"
    )


# ---------------------------------------------------------------------------
# Talking-head avatar (D-ID)
# ---------------------------------------------------------------------------
# Sends the reply text to D-ID, which returns a video of the photo speaking it.
# Three modes:
#   off   — the animated still photo only (default; no cost, no dependency)
#   demo  — plays a bundled placeholder clip so the experience can be judged
#           before committing to an account. It is NOT lip-synced.
#   live  — real generation, requires DID_API_KEY
#
# Generation takes several seconds, so it is on-demand: the video is requested
# when the participant clicks the avatar, never automatically for every reply.

TALKING_AVATAR_MODE = os.environ.get("TALKING_AVATAR", "off").lower()
DID_API_KEY = os.environ.get("DID_API_KEY", "")
DID_SOURCE_URL = os.environ.get("DID_SOURCE_URL", "")   # public URL of the photo
DID_VOICE_ID = os.environ.get("DID_VOICE_ID", "en-US-GuyNeural")
TALKING_MAX_CHARS = int(os.environ.get("TALKING_MAX_CHARS", "800"))


def talking_avatar_enabled() -> bool:
    if TALKING_AVATAR_MODE == "demo":
        return True
    return TALKING_AVATAR_MODE == "live" and bool(DID_API_KEY and DID_SOURCE_URL)


def _did_request(text: str):
    """Ask D-ID for a talking-head clip. Returns a video URL, or raises."""
    import urllib.request as _url
    import json as _j
    import base64, time

    auth = base64.b64encode(f"{DID_API_KEY}:".encode()).decode()
    payload = _j.dumps({
        "source_url": DID_SOURCE_URL,
        "script": {
            "type": "text",
            "input": text[:TALKING_MAX_CHARS],
            "provider": {"type": "microsoft", "voice_id": DID_VOICE_ID},
        },
        "config": {"stitch": True},
    }).encode("utf-8")

    req = _url.Request("https://api.d-id.com/talks", data=payload, headers={
        "Authorization": f"Basic {auth}",
        "Content-Type": "application/json",
        "Accept": "application/json",
    })
    with _url.urlopen(req, timeout=20) as resp:
        created = _j.loads(resp.read().decode("utf-8"))
    talk_id = created.get("id")
    if not talk_id:
        raise RuntimeError("D-ID did not return a talk id")

    # Poll until the clip is rendered
    poll = _url.Request(f"https://api.d-id.com/talks/{talk_id}", headers={
        "Authorization": f"Basic {auth}", "Accept": "application/json"})
    for _ in range(30):
        time.sleep(2)
        with _url.urlopen(poll, timeout=20) as resp:
            data = _j.loads(resp.read().decode("utf-8"))
        status = data.get("status")
        if status == "done":
            return data.get("result_url")
        if status in ("error", "rejected"):
            raise RuntimeError(f"D-ID {status}: {str(data.get('error'))[:120]}")
    raise RuntimeError("D-ID timed out")


# ---------------------------------------------------------------------------
# Advisor photo storage
# ---------------------------------------------------------------------------
# Stored in Postgres, not on disk: Railway rebuilds the filesystem on every
# deploy, so an uploaded file would silently revert to the bundled one.

AVATAR_MAX_BYTES = 8 * 1024 * 1024
_AVATAR_TYPES = {"image/jpeg": "jpg", "image/png": "png",
                 "image/webp": "webp", "image/gif": "gif"}


def _avatar_ensure_table(conn):
    with conn.cursor() as cur:
        cur.execute("""
            CREATE TABLE IF NOT EXISTS advisor_avatar (
                id         INTEGER PRIMARY KEY,
                mime       TEXT NOT NULL,
                data       BYTEA NOT NULL,
                updated_at TIMESTAMPTZ NOT NULL DEFAULT NOW()
            )
        """)
    conn.commit()


def store_avatar(data: bytes, mime: str) -> bool:
    conn = _settings_db_conn()
    if not conn:
        return False
    try:
        _avatar_ensure_table(conn)
        with conn.cursor() as cur:
            cur.execute("""
                INSERT INTO advisor_avatar (id, mime, data, updated_at)
                VALUES (1, %s, %s, NOW())
                ON CONFLICT (id) DO UPDATE
                SET mime = EXCLUDED.mime, data = EXCLUDED.data, updated_at = NOW()
            """, (mime, data))
        conn.commit()
        return True
    except Exception as e:
        app.logger.error(f"[avatar] store failed: {e}")
        return False
    finally:
        conn.close()


def load_avatar():
    """Return (bytes, mime, updated_at) for the uploaded photo, or None."""
    conn = _settings_db_conn()
    if not conn:
        return None
    try:
        _avatar_ensure_table(conn)
        with conn.cursor() as cur:
            cur.execute("SELECT data, mime, updated_at FROM advisor_avatar WHERE id = 1")
            row = cur.fetchone()
        if not row:
            return None
        data = row[0]
        return (bytes(data) if not isinstance(data, bytes) else data, row[1], row[2])
    except Exception as e:
        app.logger.error(f"[avatar] read failed: {e}")
        return None
    finally:
        conn.close()


def clear_avatar() -> bool:
    conn = _settings_db_conn()
    if not conn:
        return False
    try:
        _avatar_ensure_table(conn)
        with conn.cursor() as cur:
            cur.execute("DELETE FROM advisor_avatar WHERE id = 1")
        conn.commit()
        return True
    except Exception as e:
        app.logger.error(f"[avatar] clear failed: {e}")
        return False
    finally:
        conn.close()


def prepare_avatar(raw: bytes):
    """Crop to a centred square and resize. Falls back to the original bytes."""
    try:
        from PIL import Image
        import io as _io
        im = Image.open(_io.BytesIO(raw))
        im = im.convert("RGB")
        w, h = im.size
        side = min(w, h)
        # Bias the crop upward — faces sit above centre in most headshots
        left = (w - side) // 2
        top = max(0, int((h - side) * 0.28))
        im = im.crop((left, top, left + side, top + side))
        im = im.resize((320, 320), Image.LANCZOS)
        out = _io.BytesIO()
        im.save(out, "JPEG", quality=88, optimize=True)
        return out.getvalue(), "image/jpeg"
    except Exception as e:
        app.logger.warning(f"[avatar] could not process image, storing as-is: {e}")
        return raw, "image/jpeg"


# ---------------------------------------------------------------------------
# Advisor profiles
# ---------------------------------------------------------------------------
# Several named advisors can share one deployment. Each has a slug, a display
# name and its own photo, and gets its own pair of links — with and without the
# scheduling button. Stored in Postgres so photos survive redeploys.

def _advisors_ensure_table(conn):
    with conn.cursor() as cur:
        cur.execute("""
            CREATE TABLE IF NOT EXISTS advisors (
                slug        TEXT PRIMARY KEY,
                name        TEXT NOT NULL,
                photo       BYTEA,
                photo_mime  TEXT,
                created_at  TIMESTAMPTZ NOT NULL DEFAULT NOW()
            )
        """)
    conn.commit()


def slugify_advisor(name: str) -> str:
    base = re.sub(r"[^a-z0-9]+", "-", (name or "").strip().lower()).strip("-")
    return base[:40] or "advisor"


def list_advisors():
    """All advisor profiles, alphabetical, without the photo bytes."""
    conn = _settings_db_conn()
    if not conn:
        return []
    out = []
    try:
        _advisors_ensure_table(conn)
        with conn.cursor() as cur:
            cur.execute("""
                SELECT slug, name, (photo IS NOT NULL) FROM advisors
                ORDER BY name
            """)
            for slug, name, has_photo in cur.fetchall():
                out.append({"slug": slug, "name": name, "has_photo": bool(has_photo)})
    except Exception as e:
        app.logger.error(f"[advisors] list failed: {e}")
    finally:
        conn.close()
    return out


def get_advisor(slug: str):
    if not slug:
        return None
    conn = _settings_db_conn()
    if not conn:
        return None
    try:
        _advisors_ensure_table(conn)
        with conn.cursor() as cur:
            cur.execute("SELECT slug, name FROM advisors WHERE slug = %s", (slug,))
            row = cur.fetchone()
        return {"slug": row[0], "name": row[1]} if row else None
    except Exception as e:
        app.logger.error(f"[advisors] get failed: {e}")
        return None
    finally:
        conn.close()


def get_advisor_photo(slug: str):
    conn = _settings_db_conn()
    if not conn:
        return None
    try:
        _advisors_ensure_table(conn)
        with conn.cursor() as cur:
            cur.execute("SELECT photo, photo_mime FROM advisors WHERE slug = %s", (slug,))
            row = cur.fetchone()
        if not row or row[0] is None:
            return None
        data = row[0]
        return (bytes(data) if not isinstance(data, bytes) else data,
                row[1] or "image/jpeg")
    except Exception as e:
        app.logger.error(f"[advisors] photo read failed: {e}")
        return None
    finally:
        conn.close()


def save_advisor(slug: str, name: str, photo=None, mime=None) -> bool:
    conn = _settings_db_conn()
    if not conn:
        return False
    try:
        _advisors_ensure_table(conn)
        with conn.cursor() as cur:
            if photo is not None:
                cur.execute("""
                    INSERT INTO advisors (slug, name, photo, photo_mime)
                    VALUES (%s, %s, %s, %s)
                    ON CONFLICT (slug) DO UPDATE
                    SET name = EXCLUDED.name, photo = EXCLUDED.photo,
                        photo_mime = EXCLUDED.photo_mime
                """, (slug, name, photo, mime))
            else:
                cur.execute("""
                    INSERT INTO advisors (slug, name) VALUES (%s, %s)
                    ON CONFLICT (slug) DO UPDATE SET name = EXCLUDED.name
                """, (slug, name))
        conn.commit()
        return True
    except Exception as e:
        app.logger.error(f"[advisors] save failed: {e}")
        return False
    finally:
        conn.close()


def delete_advisor(slug: str) -> bool:
    conn = _settings_db_conn()
    if not conn:
        return False
    try:
        _advisors_ensure_table(conn)
        with conn.cursor() as cur:
            cur.execute("DELETE FROM advisors WHERE slug = %s", (slug,))
        conn.commit()
        return True
    except Exception as e:
        app.logger.error(f"[advisors] delete failed: {e}")
        return False
    finally:
        conn.close()


# ---------------------------------------------------------------------------
# Participant materials
# ---------------------------------------------------------------------------
# Participants can add their own documents and writing — a CV, a strategic
# plan, a blog post, notes they've drafted. These are PRIVATE to that
# participant: they are used only in that person's own sessions and never
# enter the shared knowledge base, because these transcripts contain
# personnel and institutional detail that must not surface for anyone else.
# Sharing with J3P is a separate, explicit opt-in.

PARTICIPANT_DOC_MAX_BYTES = 20 * 1024 * 1024
PARTICIPANT_DOC_MAX_CHARS = 60000      # per item, kept in the prompt budget
PARTICIPANT_DOC_LIMIT = 25             # per participant


def _participant_docs_ensure_table(conn):
    with conn.cursor() as cur:
        cur.execute("""
            CREATE TABLE IF NOT EXISTS participant_documents (
                id          BIGSERIAL PRIMARY KEY,
                token       TEXT NOT NULL,
                title       TEXT NOT NULL,
                kind        TEXT NOT NULL DEFAULT 'upload',
                content     TEXT NOT NULL,
                shared      BOOLEAN NOT NULL DEFAULT FALSE,
                created_at  TIMESTAMPTZ NOT NULL DEFAULT NOW()
            )
        """)
        cur.execute("""CREATE INDEX IF NOT EXISTS participant_documents_token
                       ON participant_documents (token)""")
    conn.commit()


def participant_token() -> str:
    """Identifies whose materials these are.

    Deliberately NOT the chat history token: that one is cleared by New
    Conversation and regenerated for anonymous visitors on each visit, which
    would silently orphan someone's library. Signed in, this is derived from
    the email so the library follows them across devices.
    """
    email = session.get("user_email", "")
    if email:
        import hashlib
        digest = hashlib.sha256(
            (app.secret_key + "|materials|" + email.lower()).encode("utf-8")
        ).hexdigest()[:32]
        return "m_" + digest
    token = session.get("materials_token")
    if not token:
        token = "a_" + os.urandom(16).hex()
        session["materials_token"] = token
        session.permanent = True
    return token


def list_participant_docs(token=None):
    token = token or participant_token()
    conn = _settings_db_conn()
    if not conn:
        return []
    out = []
    try:
        _participant_docs_ensure_table(conn)
        with conn.cursor() as cur:
            cur.execute("""
                SELECT id, title, kind, shared, created_at, LENGTH(content)
                FROM participant_documents WHERE token = %s
                ORDER BY id DESC LIMIT %s
            """, (token, PARTICIPANT_DOC_LIMIT))
            for row in cur.fetchall():
                out.append({"id": row[0], "title": row[1], "kind": row[2],
                            "shared": bool(row[3]), "when": _fmt_ts(row[4]),
                            "chars": row[5] or 0})
    except Exception as e:
        app.logger.error(f"[materials] list failed: {e}")
    finally:
        conn.close()
    return out


def add_participant_doc(title: str, content: str, kind="upload", shared=False) -> bool:
    token = participant_token()
    conn = _settings_db_conn()
    if not conn:
        return False
    try:
        _participant_docs_ensure_table(conn)
        with conn.cursor() as cur:
            cur.execute("SELECT COUNT(*) FROM participant_documents WHERE token = %s",
                        (token,))
            row = cur.fetchone()
            if row and int(row[0]) >= PARTICIPANT_DOC_LIMIT:
                app.logger.info("[materials] limit reached for this participant")
                return False
            cur.execute("""
                INSERT INTO participant_documents (token, title, kind, content, shared)
                VALUES (%s, %s, %s, %s, %s)
            """, (token, title[:200], kind, content[:PARTICIPANT_DOC_MAX_CHARS], shared))
        conn.commit()
        return True
    except Exception as e:
        app.logger.error(f"[materials] add failed: {e}")
        return False
    finally:
        conn.close()


def delete_participant_doc(doc_id: int) -> bool:
    """Scoped to this participant's own token, so nobody can delete another's."""
    token = participant_token()
    conn = _settings_db_conn()
    if not conn:
        return False
    try:
        _participant_docs_ensure_table(conn)
        with conn.cursor() as cur:
            cur.execute("DELETE FROM participant_documents WHERE id = %s AND token = %s",
                        (doc_id, token))
        conn.commit()
        return True
    except Exception as e:
        app.logger.error(f"[materials] delete failed: {e}")
        return False
    finally:
        conn.close()


def participant_materials_block(budget=24000) -> str:
    """The participant's own materials, for their prompt only."""
    if not load_settings().get("allow_materials"):
        return ""
    docs = []
    conn = _settings_db_conn()
    if not conn:
        return ""
    try:
        _participant_docs_ensure_table(conn)
        with conn.cursor() as cur:
            cur.execute("""
                SELECT title, content FROM participant_documents
                WHERE token = %s ORDER BY id DESC LIMIT 12
            """, (participant_token(),))
            docs = cur.fetchall()
    except Exception as e:
        app.logger.error(f"[materials] read failed: {e}")
        return ""
    finally:
        conn.close()

    if not docs:
        return ""

    parts, used = [], 0
    for title, content in docs:
        text = (content or "").strip()
        if not text:
            continue
        room = budget - used
        if room < 500:
            break
        excerpt = text[:room]
        used += len(excerpt)
        parts.append(f"--- {title} ---\n{excerpt}")

    if not parts:
        return ""
    return (
        "\n\n---\n"
        "THE PARTICIPANT'S OWN MATERIALS — documents and writing this person has "
        "added to their library. Treat these as authoritative about their own "
        "situation, work and voice. When they ask you to draft something, match "
        "the voice and reuse the substance found here. Do not mention that you "
        "are reading from an uploaded library; just use it.\n\n"
        + "\n\n".join(parts) + "\n"
    )


# ---------------------------------------------------------------------------
# Participant profile
# ---------------------------------------------------------------------------
# A deliberately thin profile: first name, role, specialty. No institution, no
# surname, no email beyond what sign-in already holds, nothing that would
# identify the person or their organisation. Enough for the advisor to stop
# being generic; not enough to be a personnel record.

PROFILE_FIELDS = ("first_name", "role", "specialty")


def _profile_ensure_table(conn):
    with conn.cursor() as cur:
        cur.execute("""
            CREATE TABLE IF NOT EXISTS participant_profile (
                token       TEXT PRIMARY KEY,
                first_name  TEXT,
                role        TEXT,
                specialty   TEXT,
                asked       TEXT,
                updated_at  TIMESTAMPTZ NOT NULL DEFAULT NOW()
            )
        """)
    conn.commit()


def load_profile() -> dict:
    conn = _settings_db_conn()
    if not conn:
        return dict(session.get("profile_fallback") or {})
    try:
        _profile_ensure_table(conn)
        with conn.cursor() as cur:
            cur.execute("""SELECT first_name, role, specialty, asked
                           FROM participant_profile WHERE token = %s""",
                        (participant_token(),))
            row = cur.fetchone()
        if not row:
            return {}
        return {"first_name": row[0], "role": row[1], "specialty": row[2],
                "asked": (row[3] or "")}
    except Exception as e:
        app.logger.error(f"[profile] read failed: {e}")
        return {}
    finally:
        conn.close()


def save_profile(**fields) -> bool:
    """Write only the fields given; leave the rest alone."""
    fields = {k: v for k, v in fields.items()
              if k in PROFILE_FIELDS + ("asked",) and v}
    if not fields:
        return True
    conn = _settings_db_conn()
    if not conn:
        # No database — keep it for this session at least
        cur = dict(session.get("profile_fallback") or {})
        cur.update(fields)
        session["profile_fallback"] = cur
        return True
    try:
        _profile_ensure_table(conn)
        cols = ", ".join(fields.keys())
        marks = ", ".join(["%s"] * len(fields))
        updates = ", ".join(f"{k} = EXCLUDED.{k}" for k in fields)
        with conn.cursor() as cur:
            cur.execute(f"""
                INSERT INTO participant_profile (token, {cols})
                VALUES (%s, {marks})
                ON CONFLICT (token) DO UPDATE
                SET {updates}, updated_at = NOW()
            """, (participant_token(), *fields.values()))
        conn.commit()
        return True
    except Exception as e:
        app.logger.error(f"[profile] save failed: {e}")
        return False
    finally:
        conn.close()


def forget_profile() -> bool:
    session.pop("profile_fallback", None)
    conn = _settings_db_conn()
    if not conn:
        return True
    try:
        _profile_ensure_table(conn)
        with conn.cursor() as cur:
            cur.execute("DELETE FROM participant_profile WHERE token = %s",
                        (participant_token(),))
        conn.commit()
        return True
    except Exception as e:
        app.logger.error(f"[profile] forget failed: {e}")
        return False
    finally:
        conn.close()


# Leadership roles worth recognising, longest first so "associate program
# director" wins over "program director".
_ROLE_PATTERNS = [
    (r"\b(department chair|chair of the department|chairman)\b", "department chair"),
    (r"\b(vice chair|associate chair)\b", "vice chair"),
    (r"\b(division chief|chief of the division|division head)\b", "division chief"),
    (r"\b(section chief|section head)\b", "section chief"),
    (r"\b(service line (?:director|lead|chief))\b", "service line director"),
    (r"\b(associate program director)\b", "associate program director"),
    (r"\b(program director)\b", "program director"),
    (r"\b(fellowship director)\b", "fellowship director"),
    (r"\b(residency director)\b", "residency director"),
    (r"\b(medical director)\b", "medical director"),
    (r"\b(chief medical officer|\bCMO\b)", "chief medical officer"),
    (r"\b(chief of staff)\b", "chief of staff"),
    (r"\b(vice president|\bVP\b)", "vice president"),
    (r"\b(dean|associate dean|vice dean)\b", "dean"),
    (r"\b(clinical director)\b", "clinical director"),
    (r"\b(quality (?:director|officer|lead))\b", "quality director"),
    (r"\b(attending|faculty member|staff physician)\b", "attending"),
]


def detect_role(text: str):
    low = (text or "").lower()
    for pattern, label in _ROLE_PATTERNS:
        if re.search(pattern, low, re.IGNORECASE):
            return label
    return None


def detect_first_name(text: str):
    """Only from an explicit self-introduction — never inferred."""
    if not text:
        return None
    patterns = [
        r"\bmy name'?s? (?:is )?([A-Z][a-z]{1,18})\b",
        r"\bi'?m ([A-Z][a-z]{1,18})\b",
        r"\bthis is ([A-Z][a-z]{1,18})\b",
        r"\bcall me ([A-Z][a-z]{1,18})\b",
        r"^([A-Z][a-z]{1,18}) here\b",
    ]
    # Words that follow "I'm" but aren't names
    stop = {"a", "an", "the", "not", "in", "on", "at", "trying", "looking",
            "working", "going", "just", "still", "also", "happy", "glad",
            "sorry", "new", "here", "chief", "chair", "director", "faculty",
            "attending", "concerned", "worried", "curious", "hoping"}
    for pattern in patterns:
        # Case-insensitive so "My name is" matches as well as "my name is",
        # but the captured word must still be capitalised in the original —
        # that's what separates a name from "I'm trying".
        m = re.search(pattern, text, re.IGNORECASE)
        if m:
            candidate = m.group(1).strip()
            if not candidate[:1].isupper():
                continue
            if candidate.lower() in stop:
                continue
            return candidate
    return None


def update_profile_from_message(text: str) -> dict:
    """Pick up anything the participant volunteers. Never asks, never guesses."""
    found = {}
    name = detect_first_name(text)
    role = detect_role(text)
    specialty = detect_specialty(text)
    current = load_profile()
    if name and not current.get("first_name"):
        found["first_name"] = name
    if role and not current.get("role"):
        found["role"] = role
    if specialty and not current.get("specialty"):
        found["specialty"] = specialty
    if found:
        save_profile(**found)
        app.logger.info(f"[profile] learned: {', '.join(found.keys())}")
    return found


def profile_guidance() -> str:
    """Tell the model what it knows, and what it may ask for next."""
    prof = load_profile()
    known = [f"{k.replace('_', ' ')}: {prof[k]}"
             for k in PROFILE_FIELDS if prof.get(k)]
    asked = set((prof.get("asked") or "").split(","))
    missing = [k for k in PROFILE_FIELDS if not prof.get(k) and k not in asked]

    lines = ["\n\n---\nWHAT YOU KNOW ABOUT THIS PARTICIPANT"]
    if known:
        lines.append("  " + "; ".join(known))
        lines.append("  Use this. Address them by first name occasionally — not "
                     "every message. Write for their role and specialty "
                     "specifically. Never ask again for anything listed here.")
    else:
        lines.append("  Nothing yet.")

    if missing:
        nice = {"first_name": "their first name", "role": "their leadership role",
                "specialty": "their specialty"}
        lines.append(
            "  Not yet known: " + ", ".join(nice[m] for m in missing) + ". "
            "You may ask for ONE of these, as a short closing question AFTER a "
            "substantive answer — never as a gate, never before helping, and "
            "never when they have asked for a document or a revision. At most "
            "one such question per conversation."
        )

    lines.append(
        "  NEVER ask for, and never record, their institution, hospital, "
        "employer, city, surname, or anything else that would identify them or "
        "their organisation. If they volunteer such detail, use it in the reply "
        "if helpful but do not treat it as something to collect. Keep what you "
        "hold to first name, role and specialty."
    )
    return "\n".join(lines) + "\n"


# ---------------------------------------------------------------------------
# Follow-up plans
# ---------------------------------------------------------------------------
# At the end of a session the participant can ask for an accountability plan:
# what they said they'd do, by when, and what to watch for. It's built from
# the conversation, saved against their profile, and surfaced at the start of
# the next session so the advisor can follow up rather than start cold.

PLAN_PROMPT = (
    "Build a short accountability plan from this conversation.\n\n"
    "Rules:\n"
    "- Only include commitments THEY made or clearly accepted. Do not invent "
    "actions they never agreed to, and do not pad the list.\n"
    "- If they committed to nothing concrete, say so plainly in one line and "
    "suggest at most two things they could commit to. Do not fabricate a plan.\n"
    "- Use their own words for what they're doing where you can.\n"
    "- Give each item a realistic timeframe. If they named one, use theirs; "
    "otherwise suggest one and mark it as a suggestion.\n"
    "- End with one line on the single thing most likely to derail this, drawn "
    "from what they actually told you.\n\n"
    "Format as markdown with a '# Follow-up Plan' heading, then a short "
    "'## What you're doing' list, then '## Watch for'. Keep the whole thing "
    "under 350 words. No preamble, no closing pleasantries."
)


def _plans_ensure_table(conn):
    with conn.cursor() as cur:
        cur.execute("""
            CREATE TABLE IF NOT EXISTS participant_plans (
                id          BIGSERIAL PRIMARY KEY,
                token       TEXT NOT NULL,
                content     TEXT NOT NULL,
                created_at  TIMESTAMPTZ NOT NULL DEFAULT NOW()
            )
        """)
        cur.execute("""CREATE INDEX IF NOT EXISTS participant_plans_token
                       ON participant_plans (token)""")
    conn.commit()


def save_plan(content: str) -> bool:
    conn = _settings_db_conn()
    if not conn:
        return False
    try:
        _plans_ensure_table(conn)
        with conn.cursor() as cur:
            cur.execute("""INSERT INTO participant_plans (token, content)
                           VALUES (%s, %s)""",
                        (participant_token(), content[:20000]))
        conn.commit()
        return True
    except Exception as e:
        app.logger.error(f"[plans] save failed: {e}")
        return False
    finally:
        conn.close()


def list_plans(limit=10):
    conn = _settings_db_conn()
    if not conn:
        return []
    out = []
    try:
        _plans_ensure_table(conn)
        with conn.cursor() as cur:
            cur.execute("""SELECT id, content, created_at FROM participant_plans
                           WHERE token = %s ORDER BY id DESC LIMIT %s""",
                        (participant_token(), limit))
            for row in cur.fetchall():
                out.append({"id": row[0], "content": row[1],
                            "when": _fmt_ts(row[2])})
    except Exception as e:
        app.logger.error(f"[plans] list failed: {e}")
    finally:
        conn.close()
    return out


def delete_plan(plan_id: int) -> bool:
    conn = _settings_db_conn()
    if not conn:
        return False
    try:
        _plans_ensure_table(conn)
        with conn.cursor() as cur:
            cur.execute("""DELETE FROM participant_plans
                           WHERE id = %s AND token = %s""",
                        (plan_id, participant_token()))
        conn.commit()
        return True
    except Exception as e:
        app.logger.error(f"[plans] delete failed: {e}")
        return False
    finally:
        conn.close()


def open_commitments_block() -> str:
    """The most recent plan, so the advisor can follow up next session."""
    plans = list_plans(limit=1)
    if not plans:
        return ""
    plan = plans[0]
    return (
        "\n\n---\n"
        f"COMMITMENTS FROM THEIR LAST SESSION ({plan['when']}) — they asked for "
        "this plan themselves. Early in this conversation, ask how one of these "
        "went, naming it specifically. Ask once, briefly, and then get on with "
        "whatever they've come with; do not work through the list or nag. If "
        "they say it slipped, treat that as information, not failure.\n\n"
        + plan["content"][:4000] + "\n"
    )


# ---------------------------------------------------------------------------
# Pre-session briefings
# ---------------------------------------------------------------------------
# When a participant books time, the advisor gets a short brief on what the
# participant has been working through — so the call starts where the
# conversation left off rather than from scratch.
#
# The brief is always stored in the admin panel. It is emailed as well when a
# mail transport is configured; without one the panel is the delivery route.

BRIEFING_PROMPT = (
    "Write a pre-call brief for an advisor about to meet this person. The "
    "advisor has not seen this conversation.\n\n"
    "Rules:\n"
    "- Lead with the single thing they most want help with, in one sentence.\n"
    "- Then the specifics that matter: what they've tried, who else is "
    "involved (by role, not name), any dates or constraints they mentioned.\n"
    "- Note anything they seemed hesitant or guarded about — useful for the "
    "advisor to handle carefully.\n"
    "- End with two or three questions the advisor could open with.\n"
    "- Only what they actually said. No speculation about motives, no "
    "diagnosis, no advice to the advisor about what to recommend.\n"
    "- Under 250 words. Markdown with a '# Pre-Call Brief' heading."
)


def _briefings_ensure_table(conn):
    with conn.cursor() as cur:
        cur.execute("""
            CREATE TABLE IF NOT EXISTS session_briefings (
                id            BIGSERIAL PRIMARY KEY,
                advisor_slug  TEXT,
                advisor_name  TEXT,
                participant   TEXT,
                summary       TEXT NOT NULL,
                emailed       BOOLEAN NOT NULL DEFAULT FALSE,
                created_at    TIMESTAMPTZ NOT NULL DEFAULT NOW()
            )
        """)
    conn.commit()


def save_briefing(advisor_slug, advisor_name, participant, summary, emailed):
    conn = _settings_db_conn()
    if not conn:
        return False
    try:
        _briefings_ensure_table(conn)
        with conn.cursor() as cur:
            cur.execute("""
                INSERT INTO session_briefings
                    (advisor_slug, advisor_name, participant, summary, emailed)
                VALUES (%s, %s, %s, %s, %s)
            """, (advisor_slug, advisor_name, participant, summary[:12000], emailed))
        conn.commit()
        return True
    except Exception as e:
        app.logger.error(f"[briefing] save failed: {e}")
        return False
    finally:
        conn.close()


def list_briefings(limit=25):
    conn = _settings_db_conn()
    if not conn:
        return []
    out = []
    try:
        _briefings_ensure_table(conn)
        with conn.cursor() as cur:
            cur.execute("""
                SELECT id, advisor_name, participant, summary, emailed, created_at
                FROM session_briefings ORDER BY id DESC LIMIT %s
            """, (limit,))
            for row in cur.fetchall():
                out.append({"id": row[0], "advisor": row[1] or "—",
                            "participant": row[2] or "—", "summary": row[3],
                            "emailed": bool(row[4]), "when": _fmt_ts(row[5])})
    except Exception as e:
        app.logger.error(f"[briefing] list failed: {e}")
    finally:
        conn.close()
    return out


def participant_label() -> str:
    """How the participant appears on the brief — no more than they've given."""
    prof = load_profile()
    bits = [prof.get(k) for k in ("first_name", "role", "specialty") if prof.get(k)]
    email = session.get("user_email", "")
    if email:
        bits.append(email)
    return " · ".join(bits) if bits else "Anonymous participant"


def build_briefing() -> str:
    """Summarise the conversation for the advisor. Empty string if too thin."""
    history = load_history()
    if len(history) < 2:
        return ""
    lines = []
    for turn in history[-30:]:
        who = "Participant" if turn.get("role") == "user" else "Advisor (AI)"
        lines.append(f"{who}: {(turn.get('content') or '')[:1800]}")
    try:
        resp = client.messages.create(
            model=CONFIG["model"],
            max_tokens=900,
            system=[{"type": "text", "text": BRIEFING_PROMPT}],
            messages=[{"role": "user",
                       "content": "Conversation:\n\n" + "\n\n".join(lines)}],
        )
        parts = [b.text for b in resp.content if getattr(b, "type", "") == "text"]
        return scrub_internal_emails("\n".join(parts).strip())
    except Exception as e:
        app.logger.error(f"[briefing] generation failed: {e}")
        return ""


def email_briefing(advisor_name, participant, summary) -> bool:
    if not mail_transport_configured():
        return False
    to_addr = os.environ.get("BRIEFING_EMAIL") or os.environ.get(
        "SAFETY_ALERT_EMAIL") or CONFIG["contact_email"]
    subject = f"Pre-call brief — {participant}"
    body = (f"A participant has just booked time"
            f"{' with ' + advisor_name if advisor_name and advisor_name != '—' else ''}.\n\n"
            f"Participant: {participant}\n\n{summary}\n\n"
            f"— Generated automatically from their session with the J3P Advisor.")
    try:
        return bool(send_email(to_addr, subject, body))
    except Exception as e:
        app.logger.error(f"[briefing] email failed: {e}")
        return False


# ---------------------------------------------------------------------------
# Per-advisor knowledge
# ---------------------------------------------------------------------------
# Two tiers. The shared J3P base is the framework every advisor draws on. On
# top of that, a document can be assigned to one advisor, and then only that
# advisor's sessions retrieve it.
#
# Ownership is tracked here rather than in the documents table, so database.py
# is untouched: documents keep unique titles (the duplicate check enforces it),
# and a title-to-advisor map is enough to filter retrieval.

def _doc_owner_ensure_table(conn):
    with conn.cursor() as cur:
        cur.execute("""
            CREATE TABLE IF NOT EXISTS document_owner (
                title        TEXT PRIMARY KEY,
                advisor_slug TEXT NOT NULL,
                created_at   TIMESTAMPTZ NOT NULL DEFAULT NOW()
            )
        """)
    conn.commit()


def set_document_owner(title: str, advisor_slug: str) -> bool:
    """Assign a document to one advisor. Empty slug means the shared base."""
    if not title:
        return False
    conn = _settings_db_conn()
    if not conn:
        return False
    try:
        _doc_owner_ensure_table(conn)
        with conn.cursor() as cur:
            if advisor_slug:
                cur.execute("""
                    INSERT INTO document_owner (title, advisor_slug)
                    VALUES (%s, %s)
                    ON CONFLICT (title) DO UPDATE SET advisor_slug = EXCLUDED.advisor_slug
                """, (title[:200], advisor_slug))
            else:
                cur.execute("DELETE FROM document_owner WHERE title = %s", (title[:200],))
        conn.commit()
        return True
    except Exception as e:
        app.logger.error(f"[kb] owner write failed: {e}")
        return False
    finally:
        conn.close()


def document_owners() -> dict:
    """{title: advisor_slug} for every assigned document."""
    conn = _settings_db_conn()
    if not conn:
        return {}
    try:
        _doc_owner_ensure_table(conn)
        with conn.cursor() as cur:
            cur.execute("SELECT title, advisor_slug FROM document_owner")
            return {row[0]: row[1] for row in cur.fetchall()}
    except Exception as e:
        app.logger.error(f"[kb] owner read failed: {e}")
        return {}
    finally:
        conn.close()


def filter_chunks_for_advisor(results, advisor_slug):
    """Shared documents plus this advisor's own. Others are dropped.

    An unassigned document belongs to the shared J3P base and is always
    available; an assigned one is reachable only in that advisor's sessions.
    """
    owners = document_owners()
    if not owners:
        return results
    keep = []
    for r in results:
        owner = owners.get(r.get("title"))
        if not owner or owner == advisor_slug:
            keep.append(r)
    return keep


def detect_deliverable_request(text: str) -> bool:
    """True when the user is asking for a written document of some kind."""
    import re as _r
    if not text:
        return False
    low = text.lower()
    return bool(_r.search(_DELIVERABLE_VERBS, low) and _r.search(_DELIVERABLE_NOUNS, low))


# "one pdf", "a single document", "not separate files" — the user wants
# everything in one file rather than one file per deliverable.
_SINGLE_FILE_RE = (
    r"(\b(?:one|single|1)\s+(?:file|document|doc|pdf|deck|letter|attachment)\b|"
    r"\bnot\s+separate\b|\bno\s+separate\b|\bdon'?t\s+want\s+separate\b|"
    r"\bcombined?\s+(?:into|in)\s+(?:one|a single)\b|\ball in one\b|"
    r"\bas one\b|\bsingle file\b)"
)


_SEPARATE_FILES_RE = (
    r"(\bseparate\s+(?:files?|documents?|docs?)\b|"
    r"\bindividual\s+(?:files?|documents?)\b|"
    r"\beach\s+(?:as|in)\s+(?:its|their)\s+own\b|"
    r"\btwo\s+(?:files?|documents?)\b|\bboth\s+as\s+separate\b|"
    r"\bsplit\s+(?:them|it|into)\b)"
)


def wants_separate_files(text: str) -> bool:
    """True only when the user explicitly asked for more than one file."""
    import re as _r
    if not text:
        return False
    low = text.lower()
    if _r.search(r"(not|no|don'?t|do not)\s+(?:want\s+)?separate", low):
        return False
    return bool(_r.search(_SEPARATE_FILES_RE, low))


def wants_single_file(text: str) -> bool:
    """True when the user explicitly asked for one combined file."""
    import re as _r
    return bool(text and _r.search(_SINGLE_FILE_RE, text.lower()))


# ---------------------------------------------------------------------------
# Attachment text extraction
# ---------------------------------------------------------------------------
# embeddings.extract_text_from_upload handles PDF, DOCX, TXT and MD. Spreadsheets
# and slide decks are read here so people can attach the formats they actually
# work in — an Excel model, a PowerPoint deck, a CSV export.

SPREADSHEET_EXTS = ('.xlsx', '.xlsm', '.xltx')
SLIDES_EXTS = ('.pptx', '.potx')
DELIMITED_EXTS = ('.csv', '.tsv')
LEGACY_OFFICE_EXTS = ('.doc', '.xls', '.ppt')


def _extract_spreadsheet(file_bytes: bytes, max_chars: int = 40000) -> str:
    """Flatten a workbook to text, one row per line, sheet by sheet."""
    import io as _io
    from openpyxl import load_workbook
    wb = load_workbook(_io.BytesIO(file_bytes), data_only=True, read_only=True)
    out, total = [], 0
    for ws in wb.worksheets:
        out.append(f"--- SHEET: {ws.title} ---")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if c is None else str(c) for c in row]
            if not any(c.strip() for c in cells):
                continue
            line = " | ".join(cells).rstrip(" |")
            out.append(line)
            total += len(line)
            if total > max_chars:
                out.append("[... spreadsheet truncated for length ...]")
                return "\n".join(out)
    return "\n".join(out)


def _extract_slides(file_bytes: bytes) -> str:
    """Pull the text of every slide, including speaker notes."""
    import io as _io
    from pptx import Presentation
    prs = Presentation(_io.BytesIO(file_bytes))
    out = []
    def walk(shapes, out):
        """Collect text from shapes, tables, and grouped shapes alike."""
        for shape in shapes:
            # Grouped shapes hide their contents one level down
            if shape.shape_type == 6 and hasattr(shape, "shapes"):
                walk(shape.shapes, out)
                continue
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                    if any(cells):
                        out.append(" | ".join(cells))
                continue
            if shape.has_text_frame and shape.text_frame.text.strip():
                out.append(shape.text_frame.text.strip())
            # Charts carry their categories/series names
            if getattr(shape, "has_chart", False):
                try:
                    plot = shape.chart.plots[0]
                    cats = [str(c) for c in plot.categories]
                    if cats:
                        out.append("Chart categories: " + ", ".join(cats))
                    for series in plot.series:
                        vals = ", ".join("" if v is None else str(v) for v in series.values)
                        out.append(f"Chart series {series.name}: {vals}")
                except Exception:
                    pass

    for i, slide in enumerate(prs.slides, 1):
        out.append(f"--- SLIDE {i} ---")
        walk(slide.shapes, out)
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    out.append(f"[Speaker notes: {notes}]")
        except Exception:
            pass
    return "\n".join(out)


def _extract_delimited(file_bytes: bytes) -> str:
    text = file_bytes.decode("utf-8", errors="replace")
    return text if len(text) <= 40000 else text[:40000] + "\n[... truncated ...]"


def extract_attachment_text(filename: str, file_bytes: bytes) -> str:
    """Extract text from any supported attachment type.

    Raises ValueError with a user-facing message for formats we can't read.
    """
    name = (filename or "").lower()
    ext = "." + name.rsplit(".", 1)[-1] if "." in name else ""

    if ext in SPREADSHEET_EXTS:
        return _extract_spreadsheet(file_bytes)
    if ext in SLIDES_EXTS:
        return _extract_slides(file_bytes)
    if ext in DELIMITED_EXTS:
        return _extract_delimited(file_bytes)
    if ext in LEGACY_OFFICE_EXTS:
        raise ValueError(
            f"{filename} is in an older Office format. Please re-save it as "
            f".docx, .xlsx or .pptx and attach it again."
        )
    # PDF / DOCX / TXT / MD go through the existing pipeline
    return emb.extract_text_from_upload(filename, file_bytes)


def _fit_history(messages: list) -> list:
    """Trim conversation history so the session cookie stays under the browser limit.

    Keeps the most recent turns intact and drops the oldest ones. If even the
    newest message is over budget it's truncated rather than dropped, so a
    follow-up like "now send that as Word" still has something to refer to.
    """
    trimmed = []
    used = 0
    for msg in reversed(messages):
        content = msg.get("content")
        if not isinstance(content, str):
            # Multi-part (image) content is never persisted as-is
            content = "[Attachment]"
        if len(content) > SESSION_MSG_CAP:
            content = content[:SESSION_MSG_CAP].rstrip() + "\n[… truncated …]"
        if used + len(content) > SESSION_HISTORY_BUDGET:
            if trimmed:
                break
            # Newest message alone exceeds the budget — keep a truncated head
            content = content[:SESSION_HISTORY_BUDGET].rstrip() + "\n[… truncated …]"
        trimmed.append({"role": msg.get("role", "user"), "content": content})
        used += len(content)
    trimmed.reverse()
    return trimmed


def retrieve_context_and_lessons(query: str) -> tuple:
    """Search knowledge base AND approved lessons for material relevant to the query.

    Returns (context_string, lessons_list).
    - context_string: formatted KB chunks (or "")
    - lessons_list: list of dicts with user_message/bot_reply/comment for top-3
      semantically similar past failures (or [])

    Computing the embedding once and reusing it for both lookups saves one
    Voyage API call per chat turn.
    """
    if not (db.is_enabled() and emb.is_enabled()):
        return ("", [])
    try:
        query_embedding = emb.embed_text(query)
    except Exception as e:
        app.logger.error(f"Embedding failed: {e}")
        return ("", [])

    # --- Knowledge base ---
    context = ""
    try:
        # Over-fetch, then drop other advisors' material before trimming to
        # the configured top-k, so filtering can't leave us short.
        top_k = CONFIG["rag_top_k"]
        results = db.search_chunks(query_embedding, limit=top_k * 4)
        results = filter_chunks_for_advisor(results, session.get("advisor_slug"))
        relevant = [r for r in results
                    if r["similarity"] >= CONFIG["rag_min_similarity"]][:top_k]
        if relevant:
            sections = [f"[Source: {r['title']}]\n{r['content']}" for r in relevant]
            context = "\n\n---\n\n".join(sections)
    except Exception as e:
        app.logger.error(f"RAG retrieval failed: {e}")

    # --- Approved lessons (negative-feedback memory) ---
    lessons = []
    try:
        lessons = db.search_lessons(query_embedding, limit=3, min_similarity=0.5)
    except Exception as e:
        app.logger.error(f"Lesson retrieval failed: {e}")

    return (context, lessons)


def retrieve_context(query: str) -> str:
    """Backward-compatible wrapper — returns only KB context."""
    context, _ = retrieve_context_and_lessons(query)
    return context


def _render_chat(force_scheduling=None, advisor=None):
    """Render the chat page.

    force_scheduling overrides the admin default; advisor selects a named
    profile so the page shows that person's photo.
    """
    # Anonymous visitors start fresh each visit. Signed-in participants keep
    # their history — that continuity is the point of signing in.
    if not current_user():
        clear_history()
    else:
        session["messages"] = []
    if force_scheduling is None:
        show = load_settings()["show_scheduling_button"]
    else:
        show = bool(force_scheduling)
        # Remember the choice so a reset or refresh keeps the same experience
        session["force_scheduling"] = show
    # A visitor who arrived on a fixed link keeps that variant for the session
    if force_scheduling is None and "force_scheduling" in session:
        show = bool(session["force_scheduling"])

    # Remember which advisor this visitor is with, so a reset keeps the photo
    if advisor:
        session["advisor_slug"] = advisor["slug"]
    active = advisor or get_advisor(session.get("advisor_slug"))

    page_cfg = dict(CONFIG)
    # A name for the default photo, so the general link can read "Alan
    # Friedman" rather than the app's own name.
    def name_the_advisor(cfg_out, person_name):
        """Put a person's name in the greeting in place of the app's name.

        Dropping the article matters: "with the J3P Advisor" reads correctly,
        "with the Alan Friedman" does not.
        """
        cfg_out["avatar_name"] = person_name
        opening = CONFIG.get("opening") or ""
        if CONFIG["persona_name"] in opening:
            cfg_out["opening"] = re.sub(
                r"\bthe\s+" + re.escape(CONFIG["persona_name"]),
                person_name, opening).replace(
                CONFIG["persona_name"], person_name)

    default_name = (load_settings().get("avatar_name") or "").strip()
    if default_name and not active:
        name_the_advisor(page_cfg, default_name)
    if active:
        page_cfg["avatar_url"] = f"/a/{active['slug']}/photo.jpg"
        page_cfg["persona_name"] = active["name"]
        name_the_advisor(page_cfg, active["name"])

    return render_template_string(
        INDEX_HTML,
        cfg=page_cfg,
        show_avatar=bool(load_settings().get("show_avatar")),
        allow_materials=bool(load_settings().get("allow_materials")),
        show_scheduling_button=show,
        release_heading=RELEASE_HEADING,
        release_body=RELEASE_BODY_HTML,
        release_checkbox_label=RELEASE_CHECKBOX_LABEL,
    )


@app.route("/")
@paywall.paywall_required
@login_required
def index():
    """Default entry point — follows the admin panel's Display Settings."""
    return _render_chat()


@app.route("/scheduling")
@paywall.paywall_required
@login_required
def index_with_scheduling():
    """Share this link when you want the booking button shown."""
    session.pop("force_scheduling", None)
    return _render_chat(force_scheduling=True)


@app.route("/a/<slug>")
@paywall.paywall_required
@login_required
def advisor_index(slug):
    """A named advisor's link — follows the admin scheduling default."""
    advisor = get_advisor(slug)
    if not advisor:
        return _advisor_not_found(slug)
    session.pop("force_scheduling", None)
    return _render_chat(advisor=advisor)


@app.route("/a/<slug>/scheduling")
@paywall.paywall_required
@login_required
def advisor_index_with_scheduling(slug):
    advisor = get_advisor(slug)
    if not advisor:
        return _advisor_not_found(slug)
    return _render_chat(force_scheduling=True, advisor=advisor)


@app.route("/a/<slug>/no-scheduling")
@paywall.paywall_required
@login_required
def advisor_index_without_scheduling(slug):
    advisor = get_advisor(slug)
    if not advisor:
        return _advisor_not_found(slug)
    return _render_chat(force_scheduling=False, advisor=advisor)


@app.route("/a/<slug>/photo.jpg")
def advisor_photo(slug):
    """A named advisor's photo, falling back to the default one."""
    stored = get_advisor_photo(slug)
    if stored:
        data, mime = stored
        resp = app.response_class(data, mimetype=mime)
        resp.headers["Cache-Control"] = "no-cache"
        return resp
    return redirect(url_for("advisor_avatar"))


def _advisor_not_found(slug):
    app.logger.info(f"[advisors] unknown slug requested: {slug}")
    return (f"<div style='font-family:sans-serif;padding:2.5rem;max-width:34rem'>"
            f"<h2 style='color:#27334A'>No advisor at that link</h2>"
            f"<p style='color:#6B7280'>The link <code>/a/{slug}</code> doesn't match "
            f"an advisor profile. Check the Advisors section of the admin panel.</p>"
            f"<p><a href='/' style='color:#9D432C'>Go to the default advisor</a></p>"
            f"</div>"), 404


@app.route("/no-scheduling")
@paywall.paywall_required
@login_required
def index_without_scheduling():
    """Share this link when you want the advisor with no booking prompt."""
    session.pop("force_scheduling", None)
    return _render_chat(force_scheduling=False)


@app.route("/chat", methods=["POST"])
@paywall.paywall_required
@login_required
def chat():
    # Accept BOTH multipart/form-data (with optional file OR folder of files) and JSON.
    # Attachment paths:
    #   - Single document (PDF/DOCX/TXT/MD): extract text, prepend as context
    #   - Single image (JPG/PNG/GIF/WEBP): send to Claude as vision content
    #   - Folder of documents: concatenate extracted text from all supported files
    # All attachments are ephemeral — used only in this turn.
    uploaded_file = None
    folder_files = []
    folder_name = ""
    if request.files:
        if "file" in request.files:
            uploaded_file = request.files["file"]
        if "files" in request.files:
            folder_files = request.files.getlist("files")
            # Prefer the explicit folder_name; fall back to attachment_label
            # (used when the paperclip sends multiple files rather than a folder).
            folder_name = (
                request.form.get("folder_name")
                or request.form.get("attachment_label")
                or "Attachments"
            ).strip()

    if uploaded_file or folder_files:
        user_input = (request.form.get("message") or "").strip()
    else:
        data = request.get_json(silent=True) or {}
        user_input = (data.get("message") or "").strip()

    if not user_input and not uploaded_file and not folder_files:
        return jsonify({"error": "Empty message"}), 400

    # Extension-based routing between document vs image
    DOC_EXTS = ('.pdf', '.docx', '.xlsx', '.xlsm', '.pptx',
                '.csv', '.tsv', '.txt', '.md', '.rtf')
    IMAGE_EXTS = ('.jpg', '.jpeg', '.png', '.gif', '.webp')
    IMAGE_MEDIA_TYPES = {
        '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg',
        '.png': 'image/png', '.gif': 'image/gif', '.webp': 'image/webp',
    }

    attachment_context = ""    # Text extracted from a document, if any
    image_blocks = []          # List of Anthropic image content blocks
    attachment_display_name = ""

    if uploaded_file:
        try:
            file_bytes = uploaded_file.read()
            if len(file_bytes) > MAX_UPLOAD_BYTES:
                return jsonify({
                    "error": f"{filename} is too large "
                             f"({len(file_bytes) / 1048576:.0f} MB). "
                             f"Limit is {MAX_UPLOAD_MB} MB per file."
                }), 400
            filename = uploaded_file.filename or "attachment"
            attachment_display_name = filename
            ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

            if ext in IMAGE_EXTS:
                # Vision path — send raw image bytes to Claude as base64
                import base64
                b64 = base64.standard_b64encode(file_bytes).decode("ascii")
                image_blocks.append({
                    "type": "image",
                    "source": {
                        "type": "base64",
                        "media_type": IMAGE_MEDIA_TYPES.get(ext, "image/jpeg"),
                        "data": b64,
                    },
                })
                if not user_input:
                    user_input = "Please describe or analyze this image."
                app.logger.info(f"Chat image received: {filename} ({len(file_bytes)} bytes)")

            elif ext in DOC_EXTS:
                # Text extraction path (unchanged)
                extracted = extract_attachment_text(filename, file_bytes)
                if not extracted.strip():
                    return jsonify({"error": f"Could not extract text from {filename}."}), 400
                if len(extracted) > 40000:
                    extracted = extracted[:40000] + "\n\n[... file truncated for length ...]"
                attachment_context = (
                    f"\n\n[The user attached a document titled '{filename}'. "
                    f"The full text of the document is below. Use it to inform your response.]\n\n"
                    f"--- BEGIN ATTACHED DOCUMENT ---\n{extracted}\n--- END ATTACHED DOCUMENT ---"
                )
                if not user_input:
                    user_input = "Please review this attached document."
                app.logger.info(f"Chat attachment received: {filename} ({len(extracted)} chars extracted)")

            else:
                return jsonify({
                    "error": ("Unsupported file type. Supported: PDF, Word (.docx), "
                              "Excel (.xlsx), PowerPoint (.pptx), CSV, TXT, MD, RTF, "
                              "and images (JPG, PNG, GIF, WEBP).")
                }), 400

        except Exception as e:
            app.logger.error(f"Attachment processing failed: {e}")
            return jsonify({"error": f"Could not process file: {str(e)[:200]}"}), 400

    # === Multi-file attachment (folder OR paperclip multi-select) ===
    # Splits incoming files into two paths:
    #   - Documents (PDF/DOCX/TXT/MD): text extraction, concatenated into one context block
    #   - Images (JPG/PNG/GIF/WEBP): each becomes an Anthropic vision content block
    folder_stats = None  # (uploaded_count, skipped_count) if any files were attached
    if folder_files:
        MAX_FOLDER_FILES = 20
        MAX_COMBINED_CHARS = 80000  # keep prompt reasonable
        MAX_IMAGES = 5              # cap image count to control cost + payload
        supported = [f for f in folder_files
                     if (f.filename or "").lower().endswith(DOC_EXTS + IMAGE_EXTS)]
        supported = supported[:MAX_FOLDER_FILES]
        if not supported:
            return jsonify({"error": "No supported files in the attachment (PDF, DOCX, TXT, MD, or images)."}), 400

        import base64
        combined_parts = []
        total_chars = 0
        docs_used = 0
        images_used = 0
        image_names = []
        skipped_count = 0
        for f in supported:
            try:
                fname = (f.filename or "unknown").rsplit("/", 1)[-1]
                ext = "." + fname.rsplit(".", 1)[-1].lower() if "." in fname else ""
                bytes_ = f.read()
                if len(bytes_) > MAX_UPLOAD_BYTES:
                    skipped_count += 1
                    continue
                if ext in IMAGE_EXTS:
                    # Cap number of images per turn
                    if images_used >= MAX_IMAGES:
                        skipped_count += 1
                        continue
                    b64 = base64.standard_b64encode(bytes_).decode("ascii")
                    image_blocks.append({
                        "type": "image",
                        "source": {
                            "type": "base64",
                            "media_type": IMAGE_MEDIA_TYPES.get(ext, "image/jpeg"),
                            "data": b64,
                        },
                    })
                    images_used += 1
                    image_names.append(fname)
                elif ext in DOC_EXTS:
                    extracted = extract_attachment_text(fname, bytes_)
                    if not extracted.strip():
                        skipped_count += 1
                        continue
                    remaining = MAX_COMBINED_CHARS - total_chars
                    if remaining <= 0:
                        skipped_count += 1
                        continue
                    if len(extracted) > remaining:
                        extracted = extracted[:remaining] + "\n[... file truncated for length ...]"
                    combined_parts.append(f"--- FILE: {fname} ---\n{extracted}")
                    total_chars += len(extracted)
                    docs_used += 1
                else:
                    skipped_count += 1
            except Exception as e:
                app.logger.error(f"Multi-file attachment skipped ({f.filename}): {e}")
                skipped_count += 1

        used_count = docs_used + images_used
        if used_count == 0:
            return jsonify({"error": "Could not process any of the attached files."}), 400

        # Tell the model which images it's looking at, in order. Without this
        # the vision blocks arrive unlabelled and it can't refer to them by name.
        image_note = ""
        if image_names:
            listed = "; ".join(f"image {i}: {n}" for i, n in enumerate(image_names, 1))
            image_note = (
                f"\n\n[The user also attached {len(image_names)} "
                f"image{'s' if len(image_names) != 1 else ''}, provided with this "
                f"message in this order — {listed}. Refer to them by filename.]"
            )

        # Build the doc context block only if we actually used any docs
        if combined_parts:
            attachment_context = (
                f"\n\n[The user attached {docs_used} document{'s' if docs_used != 1 else ''}"
                + (f" and {images_used} image{'s' if images_used != 1 else ''}" if images_used else "")
                + f" (grouped as '{folder_name}'). "
                f"The concatenated content of all documents is below. Use it to inform your response.]\n\n"
                f"--- BEGIN ATTACHED FILES ---\n"
                + "\n\n".join(combined_parts) +
                f"\n--- END ATTACHED FILES ---"
                + image_note
            )
        elif image_note:
            # Images only, no documents
            attachment_context = image_note

        if not user_input:
            user_input = "Please review these attached files."
        pieces = []
        if docs_used:   pieces.append(f"{docs_used} doc{'s' if docs_used != 1 else ''}")
        if images_used: pieces.append(f"{images_used} image{'s' if images_used != 1 else ''}")
        attachment_display_name = f"{folder_name} ({', '.join(pieces)})"
        folder_stats = (used_count, skipped_count)
        app.logger.info(
            f"Chat multi-file attachment: {folder_name} — {docs_used} docs, "
            f"{images_used} images, {skipped_count} skipped, {total_chars} chars"
        )

    # ---------------------------------------------------------------
    # Safety check — runs before anything else
    # ---------------------------------------------------------------
    # If the message signals risk of harm to self or others, respond with the
    # fixed safety guidance instead of coaching, and alert by email. Handled
    # here rather than through the model so the response is guaranteed.
    if detect_risk(user_input):
        app.logger.warning("[safety] risk language detected — safety response returned")
        raise_safety_alert(user_input, client_ip())

        safety_id = None
        try:
            safety_id = db.log_interaction(
                user_message=user_input,
                bot_reply=SAFETY_RESPONSE,
                persona=CONFIG["persona_name"],
                attachment_info="SAFETY ALERT",
            )
        except Exception as e:
            app.logger.error(f"[safety] log_interaction failed: {e}")

        append_history("user", user_input)
        append_history("assistant", SAFETY_RESPONSE)

        return jsonify({
            "reply": SAFETY_RESPONSE,
            "interaction_id": safety_id,
            "export_format": None,      # never offer to export a crisis reply
            "documents": [],
            "single_file": False,
            "separate_files": False,
        })

    # Combine user text with any attached-document context. Images are added
    # separately as a content block when building the current-turn message below.
    # Pick up first name, role and specialty whenever they're volunteered.
    # Stored durably, so a returning participant isn't asked twice.
    try:
        update_profile_from_message(user_input)
    except Exception as e:
        app.logger.error(f"[profile] update failed: {e}")

    full_user_content = user_input + attachment_context

    messages = load_history()

    # A long transcript plus a large attachment can push the request big enough
    # to time out or exhaust the worker — which surfaces to the user as an
    # "upstream error". Trim history when an attachment carries most of the
    # payload this turn.
    if attachment_context and len(attachment_context) > 20000:
        room = max(8000, CHAT_HISTORY_BUDGET - len(attachment_context))
        trimmed, used = [], 0
        for msg in reversed(messages):
            body = msg.get("content") or ""
            if not isinstance(body, str):
                body = "[attachment]"
            used += len(body)
            if used > room and trimmed:
                break
            trimmed.append(msg)
        messages = list(reversed(trimmed))
        app.logger.info(
            f"History trimmed to {len(messages)} turns "
            f"(attachment {len(attachment_context)} chars)")

    # Short summary of this turn for the transcript. Flask sessions are signed cookies with a
    # hard ~4 KB browser limit — if we exceed it the cookie is silently dropped
    # and the ENTIRE conversation history disappears. So attachment text (which
    # can run to tens of thousands of characters) is sent to Claude for this turn
    # only and never persisted to the session.
    history_note = (user_input or "").strip()
    if attachment_display_name:
        history_note = (history_note + f"\n\n[Attached: {attachment_display_name}]").strip()
    if not history_note:
        history_note = "[Attachment]"

    if image_blocks:
        # Combine one or more image blocks with the text (including any
        # concatenated document context) into a single multi-part user message.
        text_for_this_turn = full_user_content or user_input or "Please describe or analyze this."
        current_turn_content = list(image_blocks) + [
            {"type": "text", "text": text_for_this_turn},
        ]
        messages_for_api = messages + [{"role": "user", "content": current_turn_content}]
    elif attachment_context:
        # Document or folder attached — Claude sees the full extracted text this
        # turn; only the short note above is kept in the session.
        messages_for_api = messages + [{"role": "user", "content": full_user_content}]
    else:
        messages_for_api = messages + [{"role": "user", "content": full_user_content}]

    messages.append({"role": "user", "content": history_note})
    append_history("user", history_note)

    # Build system prompt — base prompt + retrieved context if available
    base_prompt = CONFIG["system_prompt"]
    context, lessons = retrieve_context_and_lessons(user_input)

    # Build lessons block: things we got wrong before and shouldn't repeat
    lessons_block = ""
    if lessons:
        lesson_items = []
        for i, lesson in enumerate(lessons, 1):
            # Lessons come from both ratings, so each one has to say which it
            # is. Presenting a thumbs-up example under "why that was unhelpful"
            # would teach the model that a good answer was a mistake.
            worked = (lesson.get("rating") or "").lower() == "up"
            note = (lesson.get("comment") or "").strip()
            if worked:
                verdict = "WHAT WORKED — the participant marked this reply helpful."
                note_line = (f"  What they valued: {note[:500]}\n" if note else "")
                guidance = "  Aim for this again: same structure, depth and directness.\n"
            else:
                verdict = "WHAT DID NOT WORK — the participant marked this reply unhelpful."
                note_line = (f"  Why it fell short: {note[:500]}\n" if note else "")
                guidance = "  Do not repeat this pattern.\n"
            lesson_items.append(
                f"Lesson {i} — {verdict}\n"
                f"  Their question (similar to the current one): {lesson['user_message'][:500]}\n"
                f"  The reply given: {lesson['bot_reply'][:500]}\n"
                + note_line + guidance
            )
        lessons_block = (
            "\n\n---\n"
            "LESSONS FROM PRIOR FEEDBACK — reviewed examples from real sessions on "
            "questions like this one. Some show replies that landed well and some "
            "show replies that did not; each is labelled. Follow the patterns that "
            "worked and avoid the ones that did not. Do NOT mention these lessons "
            "to the participant; just internalize them.\n\n"
            + "\n\n".join(lesson_items)
            + "\n"
        )

    # Scope-limiting + naming restrictions appended on EVERY request
    # -----------------------------------------------------------------
    # Voice differentiation — make responses feel like a real advisor,
    # not a generic AI assistant. Applied on every request regardless
    # of what's in the persona system prompt.
    # -----------------------------------------------------------------
    voice_guard = (
        "\n\n---\n"
        "VOICE & DIFFERENTIATION RULES — apply to every response:\n\n"
        "1. IDENTITY. You are the J3P Advisor — a voice grounded in Alan Friedman's "
        "leadership advisory practice. You are NOT ChatGPT, Claude, Gemini, Copilot, "
        "or any generic AI assistant. If asked what you are, what model you are, or "
        "who built you, say only: 'I'm the J3P Advisor.' Do not name any underlying "
        "model, foundation model provider, or AI company. Do not discuss your training "
        "data, capabilities as an AI, or technical architecture. Redirect back to the "
        "user's actual question.\n\n"
        "2. OPENING. Do not start responses with generic AI phrases like 'Great "
        "question!', 'Certainly!', 'Absolutely!', 'I'd be happy to help', 'What a "
        "thoughtful question', 'That's an interesting question', or any similar "
        "sycophantic opener. Start with the substance. Answer the person.\n\n"
        "3. CLOSING. Do not end responses with 'I hope this helps!', 'Let me know if "
        "you have any other questions', 'Feel free to ask', 'Is there anything else I "
        "can help you with?', or similar service-desk closers. If a follow-up question "
        "belongs at the end, make it a substantive question that pushes the person's "
        "thinking forward — never a service question.\n\n"
        "4. HEDGING. Do not over-hedge. Take positions. Avoid stringing multiple "
        "qualifiers together ('it might possibly be worth considering that perhaps'). "
        "One direct sentence beats three cautious ones. If you genuinely don't know, "
        "say so plainly.\n\n"
        "5. FORMAT. Use bullet lists only when the content is genuinely list-like "
        "(3+ parallel items). Prefer short paragraphs of prose. Do not add headers "
        "to short responses. Do not use emoji.\n\n"
        "6. TONE. Warm but direct. You are speaking with senior clinical leaders — "
        "physicians, chairs, executives. Treat them as capable adults. Name the "
        "hard thing when it needs naming. Do not add disclaimers about consulting "
        "professionals for topics where the person clearly IS the professional.\n\n"
        "7. NO SELF-REFERENCE AS AI. Do not begin sentences with 'As an AI...', "
        "'I'm just an AI...', 'While I don't have feelings...', 'I don't have "
        "personal experiences but...', or similar. Speak from the J3P frameworks "
        "and lived-practice perspective the persona is built on.\n\n"
        "8. NO META. Do not describe what you're about to do ('Let me walk you "
        "through...', 'Here's my breakdown...', 'I'll structure this as...'). "
        "Just do it.\n\n"
        "9. NO STOCK ASSISTANT VOCABULARY. These words and constructions are the "
        "fingerprints of generic AI writing. Never use them: 'delve', 'tapestry', "
        "'landscape' (figurative), 'navigate the complexities', 'multifaceted', "
        "'nuanced' (as filler), 'it's important to note', 'it's worth noting', "
        "'that said' as a paragraph opener, 'ultimately', 'at the end of the day', "
        "'in today's fast-paced world', 'ever-evolving', 'robust' (unless "
        "statistical), 'leverage' as a verb, 'unlock', 'empower', 'holistic', "
        "'synergy', 'game-changer', 'crucial'/'vital'/'essential' used as "
        "intensifiers, 'foster' (except of children), 'embark', 'realm', "
        "'testament to', 'cannot be overstated', 'a double-edged sword'.\n\n"
        "10. NO SYMMETRICAL STRUCTURE. Generic models write in tidy balanced "
        "shapes: three parallel bullets, 'On one hand... on the other hand', "
        "'Not only... but also', a summary paragraph restating what was just "
        "said, and a closing that ties a bow on it. Don't. Let paragraphs run "
        "different lengths. Stop when the point is made. Never end with a "
        "recap.\n\n"
        "11. SPECIFICITY OVER COMPLETENESS. A generic model covers every angle "
        "shallowly. You do the opposite: pick the one or two things that "
        "actually matter in this person's situation and go deep. Reference the "
        "specifics they gave you — the name, the institution, the number, the "
        "exact phrase they used. If they mentioned a colleague or a meeting or "
        "a deadline, that belongs in your answer. Generic advice that would "
        "apply to any physician anywhere is a failure.\n\n"
        "12. USE THE PRACTICE'S OWN MATERIAL. When retrieved context from the "
        "knowledge base is relevant, ground your answer in those frameworks, "
        "language and examples rather than general management wisdom. That "
        "material is the substance of this advisory practice and is the main "
        "thing that distinguishes your answer from any general-purpose "
        "assistant. Use it without naming it as a source or saying 'according "
        "to the knowledge base'.\n\n"
        "14. REVISIONS ARE EDITS, NOT REWRITES. When asked to revise, refine, "
        "or 'keep the original but add X', treat the previous version as the "
        "base text and change only what was asked. Preserve the existing "
        "opening, structure, names, dates, facts and closing verbatim unless "
        "the change requires altering them. Never silently drop paragraphs the "
        "user did not ask you to remove — a participant asked to keep the "
        "original elements and add new messaging, and the reply came back "
        "missing the original opening. That is a failure. If you believe "
        "something should be cut, keep it and say so in one line at the end.\n\n"
        "15. DON'T ASSUME A SPECIALTY. The audience is physician leaders of every "
        "kind — internists, paediatricians, psychiatrists, radiologists, "
        "pathologists, anaesthesiologists, emergency physicians, oncologists, "
        "surgeons, primary care, and non-clinical executives too. Unless the "
        "person has told you their specialty, write for a physician leader "
        "generally: say 'physicians', 'clinicians', 'faculty' or 'clinical "
        "leaders', never 'surgeons'. A reply that opened 'The mistake most "
        "surgeons make is waiting to be noticed' was wrong: the participant "
        "never said they were a surgeon.\n"
        "   Examples must generalise too. Reach for illustrations that hold "
        "across specialties — clinic throughput, call burden, service line "
        "growth, RVU pressure, quality metrics, team conflict — rather than "
        "operative volume or block time, which only apply to proceduralists.\n"
        "   If they DO state a specialty, use it and be specific to it. Mirror "
        "their language exactly.\n\n"
        "16. RETRIEVED MATERIAL MAY BE SPECIALTY-SPECIFIC. The knowledge base "
        "contains work done with particular clients, much of it orthopaedic and "
        "sports medicine. Take the principle from that material, not the "
        "specialty. Never import a specialty from retrieved context into a "
        "reply for someone who hasn't named one, and never imply the "
        "participant works in a field they haven't mentioned.\n\n"
        "13. SAY THE HARD THING. Generic assistants hedge toward the "
        "agreeable. When the person's plan has a real problem, name it plainly "
        "in the first paragraph rather than burying it after praise. It is "
        "acceptable to disagree with them outright.\n"
    )

    scope_guard = (
        "\n\n---\n"
        "STRICT SCOPE RULES — these override any conflicting guidance above:\n\n"
        "1. You answer ONLY questions related to J3P's areas of expertise: "
        "leadership development, organizational behavior, behavioral assessment, "
        "physician/healthcare leadership, team dynamics, executive coaching, "
        "communication, self-awareness, negotiation, career navigation, "
        "and related professional development topics within healthcare and "
        "high-stakes organizational settings.\n\n"
        "1b. IN SCOPE — PRODUCING DOCUMENTS. Writing finished deliverables on "
        "those topics is squarely in scope and must never be declined: cover "
        "letters, letters of intent, CVs and bios, recommendation letters, "
        "memos, board and executive summaries, slide decks, one-pagers, "
        "agendas, development plans, feedback scripts, and similar. Producing "
        "one of these is a core function, not an off-topic request.\n\n"
        "1c. DEFAULT TO IN SCOPE. Before declining, assume the request IS in "
        "scope and look for the leadership question inside it. Real examples "
        "that were wrongly declined and must be answered:\n"
        "   - 'Help me create an org chart for my department' — departmental "
        "structure, reporting lines and span of control are organizational "
        "design. Answer it.\n"
        "   - 'My two deputies are fighting, what should I do' — conflict "
        "between direct reports is core team dynamics, however casually or "
        "with whatever typos it is phrased. Answer it.\n"
        "   - Short follow-ups about work already in progress — 'Please create "
        "in word', 'make it shorter', 'in a PDF', 'add a section on X'. These "
        "refer to the deliverable earlier in the conversation. Act on the "
        "previous document; never treat them as off-topic.\n"
        "   Structure, staffing, workload, meetings, hiring, performance, "
        "conflict, culture, workflow and org design are all in scope when a "
        "leader is asking about their own organization.\n\n"
        "1d. NEVER DECLINE OVER A TYPO. Messages arrive dictated or typed in "
        "haste ('Little boxer that my two deputies are fighting' means 'Look, "
        "bother that...'). Read through the noise to the intent. If the intent "
        "is plausibly a leadership question, answer it. Ask what they meant "
        "only if it is genuinely unreadable.\n\n"
        "2. If the user asks about ANYTHING outside this scope — including but "
        "not limited to: general trivia, animals, science, history, cooking, "
        "sports, entertainment, politics, current events, math, coding, weather, "
        "personal recommendations unrelated to professional growth, or any topic "
        "where J3P would have no specific expertise — you MUST politely decline "
        "and redirect.\n\n"
        "3. Your off-topic decline should be brief and warm, in the J3P voice. "
        "Use this format (adapt naturally):\n"
        "   \"That's outside what I'm here to help with as the J3P Advisor. "
        "I'm focused on leadership, team dynamics, professional growth, and "
        "navigating challenges in healthcare and high-stakes work. "
        "Is there something along those lines I can help you with?\"\n\n"
        "4. Do NOT attempt to bridge an off-topic question into J3P territory. "
        "Do NOT answer the off-topic question even briefly before redirecting. "
        "Decline cleanly.\n\n"
        "5. Greetings, small talk, and meta-questions about what you do are fine "
        "to engage with naturally.\n\n"
        "6. When in doubt about whether a question is in scope, lean toward "
        "declining rather than answering.\n\n"
        "---\n"
        "NAMING RESTRICTIONS — these are absolute and override any retrieved "
        "context or prior instructions:\n\n"
        "A. NEVER use the following names in any response, under any circumstances:\n"
        "   - \"J3P Healthcare Solutions\"\n"
        "   - \"J3Personica\"\n"
        "   - \"Residency Select\"\n"
        "   - any variation, partial form, hyphenation, abbreviation, or "
        "rephrasing of those names\n\n"
        "B. Refer to the organization only as \"J3P\" or \"J3P Health\" if "
        "you must mention it by name. Otherwise, prefer phrases like "
        "\"our approach,\" \"our frameworks,\" \"the methodology,\" or simply "
        "describe the concept directly without attribution.\n\n"
        "C. If the user explicitly asks about \"J3P Healthcare Solutions,\" "
        "\"J3Personica,\" or \"Residency Select\" — respond in a way that "
        "discusses the underlying ideas, tools, or frameworks WITHOUT naming "
        "those specific brands. Do not confirm or deny that those names exist. "
        "Pivot to the substance.\n\n"
        "D. If retrieved context from the knowledge base contains any of those "
        "forbidden names, paraphrase the content so the forbidden names do NOT "
        "appear in your response. The underlying ideas can be conveyed without "
        "the trademarked names.\n\n"
        "E. These naming restrictions apply to ALL responses including "
        "off-topic refusals, greetings, and meta-questions about what you do.\n"
    )

    # Document production — kept in its own block and appended last so it isn't
    # outranked by the scope rules, which claim to override "guidance above".
    document_guard = (
        "\n\n---\n"
        "PRODUCING DOCUMENTS — this section takes precedence over any scope or "
        "voice rule above that would lead you to refuse:\n\n"
        "1. CAPABILITY. You CAN produce downloadable Word, PowerPoint, Excel, and "
        "PDF files. Beneath every response is a SAVE button that converts your "
        "reply into any of those four formats. This is a real, working feature.\n\n"
        "2. NEVER REFUSE ON FORMAT. Never say you cannot create a document, cannot "
        "produce a Word file, can only provide text, or that the user should copy "
        "and paste your output elsewhere. Never suggest they use another tool to "
        "make the document. A request for something 'in Word', 'as a Word doc', "
        "'as an attachment', 'as a PDF', or 'as a deck' is a request for the "
        "CONTENT — write it, and the SAVE button handles the file.\n\n"
        "3. WRITE THE DELIVERABLE ITSELF. When asked for a letter, memo, deck, "
        "one-pager, or any other document, your entire response IS that document. "
        "No preamble ('Here's a draft for you'), no commentary about the format, "
        "no explanation of what you're about to write. Begin with the document. "
        "Anything conversational you add will end up inside the exported file.\n\n"
        "4. STRUCTURE FOR EXPORT. Use markdown headings (## for sections) and "
        "numbered or bulleted lists where warranted — these convert into real "
        "Word heading styles, PowerPoint slides, and PDF sections. For a slide "
        "deck, use one heading per slide with a few bullets beneath each.\n\n"
        "5. DRAFT FIRST, ASK NEVER. Do NOT open with clarifying questions before "
        "producing a deliverable. If purpose, audience, or details are unstated, "
        "make reasonable assumptions from the conversation, any attached "
        "documents, and the person's evident role — then WRITE THE FULL DRAFT. "
        "A concrete draft they can react to is far more useful than a "
        "questionnaire. You may note your key assumptions in one short line at "
        "the end and offer to revise. Asking a list of questions instead of "
        "producing the document is a failure to follow these instructions.\n\n"
        "6. SLIDE DECKS. For a deck, write the actual slide content: one '## ' "
        "heading per slide, 3-6 short bullets beneath each. Aim for 8-14 slides "
        "unless told otherwise. Do not describe what the deck would contain — "
        "write the slides themselves. Use a real title for each heading; never "
        "label them 'Slide 1:', 'Slide 2:' and so on. Only produce a deck when "
        "one was actually asked for — if the user asks for a summary, memo, or "
        "plan, write prose with '## ' section headings, not slides.\n\n"
        "6b. NEVER ASK FOR WHAT YOU ALREADY HAVE. If the user asks you to "
        "revise, shorten, tighten, or rework something, use the version already "
        "in this conversation or in the attached documents. Never ask them to "
        "paste it again. If several versions exist, use the most recent.\n\n"
        "6c. RESPECT EXCLUSIONS. When the user rules something out — 'an "
        "executive summary, not a slide deck', 'one file, not separate ones' — "
        "follow the exclusion exactly. Producing the thing they just declined "
        "is a failure.\n\n"
        "7. MULTIPLE DELIVERABLES. If the request calls for more than one "
        "document (for example a cover letter AND a presentation), begin each "
        "one with a single-hash '# ' title on its own line — '# Cover Letter', "
        "'# Strategic Vision Presentation'. Use '## ' only for sections inside "
        "a document. This separation lets each deliverable be downloaded as its "
        "own file in its own format, so it matters. Do not merge two documents "
        "under one title.\n\n"
        "8. OFFER THE FILE, DON'T ANNOUNCE ONE. A download is offered beneath "
        "your reply and beneath that a SAVE button, but nothing is saved unless "
        "the user chooses a format. Never say a file 'is downloading', 'has been "
        "attached', or 'is ready' — that would be untrue. Do not tell them you "
        "are unable to produce a file either, and never instruct them to copy "
        "your text into Word or PowerPoint themselves. If you mention it at all, "
        "nothing at all about files. The buttons speak for themselves, and any "
        "such line ends up inside the document itself when the user saves it. "
        "Write the deliverable and stop.\n"
    )

    # All human contact routes through client services, never an individual.
    contact_guard = (
        "\n\n---\n"
        "CONTACT & REFERRALS — this overrides anything in the knowledge base:\n\n"
        f"1. THE ONLY CONTACT IS {CONFIG['contact_email']}. Whenever the user "
        "asks how to reach someone, who to talk to, how to get help, how to "
        "follow up, who to send something to, how to book or change time, who "
        "handles billing or administration, or asks for a phone number, email "
        "address, or 'a real person', give them "
        f"{CONFIG['contact_email']} and nothing else.\n\n"
        "2. NEVER NAME INDIVIDUALS AS CONTACTS. Do not offer the name, email "
        "address, direct line, or calendar of any J3P person — including "
        "advisors, coaches, operations, or administrative staff — even if that "
        "information appears in retrieved knowledge base content or attached "
        "documents. Do not say 'reach out to' followed by a person's name.\n\n"
        "3. NO OTHER ADDRESSES. Never invent or repeat any other email address "
        "or phone number for the organization. If asked for a specific "
        f"person's details, reply that {CONFIG['contact_email']} is the way in "
        "and that the team will route the request.\n\n"
        "4. SCHEDULING. For booking time, point to the scheduling button in the "
        f"app or to {CONFIG['contact_email']} — not to an individual's calendar.\n"
    )

    if context:
        composed_prompt = (
            base_prompt
            + "\n\n---\nRELEVANT CONTEXT FROM J3P KNOWLEDGE BASE:\n\n"
            + context
            + "\n\n---\nUse this context to inform your answer when relevant. "
              "Stay in your assigned voice and frameworks."
            + lessons_block
            + scope_guard
            + voice_guard
            + document_guard
            + contact_guard
            + profile_guidance()
            + open_commitments_block()
            + participant_materials_block()
        )
    else:
        composed_prompt = (
            base_prompt + lessons_block + scope_guard + voice_guard
            + document_guard + contact_guard + profile_guidance()
            + open_commitments_block() + participant_materials_block()
        )

    try:
        response = client.messages.create(
            model=CONFIG["model"],
            max_tokens=CONFIG["max_tokens"],
            system=[
                {
                    "type": "text",
                    "text": composed_prompt,
                    "cache_control": {"type": "ephemeral"},
                }
            ],
            messages=messages_for_api,
        )
    except anthropic.APIError as e:
        return jsonify({"error": f"API error: {str(e)}"}), 500
    except Exception as e:
        return jsonify({"error": f"Server error: {str(e)}"}), 500

    assistant_text = next(
        (block.text for block in response.content if block.type == "text"), ""
    )

    # ---------------------------------------------------------------
    # Deliverable enforcement
    # ---------------------------------------------------------------
    # The user asked for a file. If the model came back with clarifying
    # questions instead of the document, the download would contain those
    # questions. Rather than rely on the system prompt alone, retry once with
    # an explicit instruction to produce the document now. One retry only.
    requested_format = detect_export_format(user_input)
    wants_deliverable = bool(requested_format) or detect_deliverable_request(user_input)
    if wants_deliverable and not looks_like_document(assistant_text):
        app.logger.info(
            f"Deliverable enforcement: retrying for {requested_format or 'document'} "
            f"(first reply looked like clarifying questions)"
        )
        # Phrased as an ordinary follow-up from the user. An earlier version
        # wrapped this in "[SYSTEM DIRECTIVE — OVERRIDES ALL PRIOR GUIDANCE]",
        # which arrives in a *user* turn claiming system authority — the model
        # rightly called that a prompt-injection attempt and said so to the user.
        force_note = (
            "Please don't ask me questions — go ahead and write the complete "
            "document now, in full, using what's already in our conversation and "
            "the documents I've attached. Make reasonable assumptions where "
            "something isn't specified, and note those assumptions in one short "
            "line at the end. "
        )
        if requested_format == "pptx":
            force_note += (
                "Format it as a slide deck: one '## ' heading per slide with "
                "3-6 short bullets beneath each, 8-14 slides."
            )
        else:
            force_note += (
                "Use '## ' headings for sections and lists where appropriate."
            )
        try:
            retry_messages = messages_for_api + [
                {"role": "assistant", "content": assistant_text},
                {"role": "user", "content": force_note},
            ]
            retry = client.messages.create(
                model=CONFIG["model"],
                max_tokens=CONFIG["max_tokens"],
                system=[{
                    "type": "text",
                    "text": composed_prompt,
                    "cache_control": {"type": "ephemeral"},
                }],
                messages=retry_messages,
            )
            retry_text = next(
                (b.text for b in retry.content if b.type == "text"), ""
            ).strip()
            if retry_text and looks_like_document(retry_text):
                assistant_text = retry_text
            elif retry_text and len(retry_text) > len(assistant_text):
                # Not obviously a document, but more substantive than questions
                assistant_text = retry_text
        except Exception as e:
            app.logger.error(f"Deliverable retry failed: {e}")

    # Defensive scrubber: replace any forbidden brand names if the model slips them through.
    # The system prompt instructs Claude not to use these, but we sanitize as backup.
    FORBIDDEN_NAMES = [
        ("J3P Healthcare Solutions", "J3P"),
        ("J3P Healthcare", "J3P"),
        ("J3Personica", "the assessment framework"),
        ("J3 Personica", "the assessment framework"),
        ("Residency Select", "the residency selection tool"),
    ]
    import re as _re

    # ---------------------------------------------------------------
    # Contact scrubber
    # ---------------------------------------------------------------
    # Knowledge base documents carry staff names, emails and direct lines, so
    # the prompt rule alone isn't enough. This works at paragraph level: any
    # passage that routes the user to a named individual is replaced outright
    # with the client services address, rather than patched phrase by phrase.
    CONTACT = CONFIG["contact_email"]

    _STAFF_NAME_RE = _re.compile(
        r"\b(?:Alan(?:\s+Friedman)?|(?:Mr|Dr)\.?\s+Friedman|Friedman|"
        r"Ivy(?:\s+Seader)?|Ms\.?\s+Seader|Seader|"
        r"Diane(?:\s+Blake)?|Ms\.?\s+Blake)\b", _re.IGNORECASE)

    # Signals that a passage is telling the user how to reach a human
    _CONTACT_SIGNAL_RE = _re.compile(
        r"(@|\bemail\b|\bphone\b|\bcall\b|\breach\b|\bcontact\b|"
        r"\bgo-?to\b|\bstarting point\b|\bget in touch\b|\bfollow up\b|"
        r"\bschedul\w*\b|\bcalendar\b|\bbook\b|\bwork with\b|"
        r"\bspeak (?:to|with)\b|\bconnect with\b|\bintroduc\w*\b|"
        r"\bcoordinat\w*\b|\bhandles\b|\bmanages\b|\boversees\b|"
        r"\bleads\b|\bCEO\b|\bdirector\b|\bteam member\b|\d{3}[-.\s]\d{4})",
        _re.IGNORECASE)

    _INTERNAL_EMAIL_RE = _re.compile(r"\b[\w.+-]+@[\w-]+\.[\w.-]+\b")
    _PHONE_RE = _re.compile(r"(?:\+?1[-.\s])?\(?\d{3}\)?[-.\s]\d{3}[-.\s]\d{4}\b")
    _INTERNAL_DOMAINS = ("j3p.health", "j3phealth.com", "j3personica.com",
                         "residencyselect.com")
    _STAFF_LOCALS = ("afriedman", "alanfriedman", "alan.friedman",
                     "iseader", "ivy.seader", "ivyseader",
                     "dblake", "diane.blake", "dianeblake")

    def _is_internal_email(addr: str) -> bool:
        # The client services address is the approved destination, so a passage
        # containing it is correct and must be left alone. Treating it as
        # "internal" made the scrubber delete legitimate contact paragraphs
        # from letters the advisor had been asked to write.
        if addr.lower() == CONTACT.lower():
            return False
        domain = addr.split("@")[-1].lower()
        local = addr.split("@")[0].lower()
        return (domain in _INTERNAL_DOMAINS
                or any(tag in local for tag in _STAFF_LOCALS))

    replaced_any = False
    cleaned_paras = []
    for para in assistant_text.split("\n\n"):
        has_staff = bool(_STAFF_NAME_RE.search(para))
        has_internal_email = any(_is_internal_email(a)
                                 for a in _INTERNAL_EMAIL_RE.findall(para))
        # A passage naming staff alongside contact language is a referral —
        # drop it. Same for any passage carrying an internal address.
        if (has_staff and _CONTACT_SIGNAL_RE.search(para)) or has_internal_email:
            replaced_any = True
            continue
        # Otherwise strip only the internal identifiers, leaving the prose
        para = _INTERNAL_EMAIL_RE.sub(
            lambda m: CONTACT if _is_internal_email(m.group(0)) else m.group(0), para)
        cleaned_paras.append(para)

    assistant_text = "\n\n".join(cleaned_paras).strip()

    if replaced_any:
        # Strip any direct line that survived in a remaining paragraph
        assistant_text = _PHONE_RE.sub("", assistant_text)
        assistant_text = _re.sub(r"\s+([,.;:])", r"\1", assistant_text)
        assistant_text = _re.sub(r"[ \t]{2,}", " ", assistant_text)
        assistant_text = _re.sub(r"\n{3,}", "\n\n", assistant_text).strip()
        # Only add the pointer to a conversational reply. Appending it to a
        # letter or deck would put J3P's internal routing line inside the
        # user's own deliverable.
        is_deliverable = looks_like_document(assistant_text)
        if CONTACT.lower() not in assistant_text.lower() and not is_deliverable:
            line = (f"For anything that needs a person — scheduling, "
                    f"administration, or getting started — contact {CONTACT} "
                    f"and the team will route it.")
            assistant_text = (assistant_text + "\n\n" + line).strip() if assistant_text else line

    for forbidden, replacement in FORBIDDEN_NAMES:
        # Case-insensitive, whole-phrase replacement
        pattern = _re.compile(_re.escape(forbidden), _re.IGNORECASE)
        assistant_text = pattern.sub(replacement, assistant_text)

    # Voice scrubber: catch generic-AI phrasing that leaks past the system prompt.
    # We strip a small set of high-signal opener/closer phrases. This runs on every
    # response as a belt-and-suspenders backup to the voice_guard system rules.
    STRIP_PHRASES = [
        # Model identity — replace the whole identifying sentence with a branded one
        (r"\b(?:I am|I'?m)\s+(?:Claude|ChatGPT|GPT-?\d*|Gemini|Bard|Copilot|an?\s+AI(?:\s+language)?\s+model|a\s+large\s+language\s+model|an?\s+AI\s+assistant)[^.!?]*[.!?]",
         "I'm the J3P Advisor."),
        # "As an AI [anything]..." preface — strip up to the first comma or period
        (r"^\s*As an? AI[^,.]*[,.]\s*", ""),
        (r"^\s*As a language model[^,.]*[,.]\s*", ""),
        (r"^\s*As a large language model[^,.]*[,.]\s*", ""),
        # "I was created/trained by [foundation model company]"
        (r"\bI\s+(?:was\s+)?(?:created|built|made|developed|trained)\s+by\s+(?:Anthropic|OpenAI|Google|Microsoft|Meta|xAI)[^.!?]*[.!?]\s*", ""),
        # Sycophantic openers
        (r"^(?:Great|Excellent|Wonderful|Fantastic|Amazing|Terrific|That'?s a (?:great|really good|thoughtful|wonderful|excellent|interesting)|What a (?:great|thoughtful|wonderful|excellent))\s+question[!.]?\s*", ""),
        (r"^(?:Certainly|Absolutely|Sure(?:ly)?|Of course|Definitely)[!.,]?\s+", ""),
        (r"^I'?d be (?:happy|glad|delighted) to (?:help|assist)[^.!?]*[.!?]\s*", ""),
        (r"^(?:Happy|Glad) to help[!.]?\s*", ""),
        # File/export chatter has no place in a deliverable — it would be
        # carried into the exported document
        # Every phrasing observed in the conversation log. Several claimed a
        # download had begun when nothing had — downloads are opt-in now.
        # Can appear mid-line, not only at the start of one (seen in the log)
        (r"\s*[*_]{0,2}Use the buttons? below to save[^\n]*", ""),
        (r"\n*[*_]{0,2}Use the SAVE button[^\n]*", "\n"),
        (r"\n*[*_]{0,2}Use SAVE (?:below|beneath)[^\n]*", "\n"),
        (r"\n*[*_]{0,2}You can save (?:this|it)[^\n]*button[^\n]*", "\n"),
        (r"\n*[*_]{0,2}(?:Your |The )?(?:file|document|deck|Word document|PDF)"
         r"\s+is\s+(?:downloading|ready|attached)[^\n]*", "\n"),
        (r"\n*[*_]{0,2}The deck is written and ready[^\n]*", "\n"),
        (r"\n*[*_]{0,2}If you don'?t see the SAVE button[^\n]*", "\n"),
        (r"\n*[*_]{0,2}Use the SAVE button beneath each[^\n]*", "\n"),
        (r"[^.\n]*\bre-?download\b[^.\n]*\.?", ""),
        # Tidy an orphaned emphasis marker left behind on its own line
        (r"\n\s*[*_]{1,2}\s*$", ""),
        # Stock assistant phrasing — swapped for plainer wording
        (r"\bit'?s (?:important|worth) (?:to note|noting) that\b,?\s*", ""),
        (r"\bit (?:is|'s) important to (?:remember|understand|recognize) that\b,?\s*", ""),
        (r"^\s*That said,\s*", "", ),
        (r"\bat the end of the day\b,?\s*", ""),
        (r"\bin today'?s (?:fast[- ]paced|ever[- ]changing|complex)\s+\w+(?:\s+\w+)?\b,?\s*", ""),
        (r"\bfast[- ]paced\s+", ""),
        (r"\bever[- ]evolving\b\s*", ""),
        (r"\bcannot be overstated\b", "matters"),
        (r"\bis a testament to\b", "shows"),
        (r"\bdelve into\b", "dig into"),
        (r"\bdelving into\b", "digging into"),
        (r"\bnavigate the complexities of\b", "work through"),
        (r"\bleverage\b(?=\s+\w)", "use"),
        (r"\bLeverage\b(?=\s+\w)", "Use"),
        (r"\ba double[- ]edged sword\b", "a trade-off"),
        (r"\bgame[- ]changer\b", "significant"),
        (r"\bthe (?:leadership|organizational|clinical) landscape\b", "the situation"),
        (r"\brich tapestry\b", "mix"),
        (r"\btapestry of\b", "mix of"),
        (r"\bembark on\b", "start"),
        (r"\bin the realm of\b", "in"),
        # Service-desk closers (end-of-response)
        (r"\s*I hope (?:this|that) (?:helps|is helpful)[!.]?", ""),
        (r"\s*(?:Please )?(?:let me know|feel free to (?:ask|reach out))[^.!?]*[.!?]?\s*$", ""),
        (r"\s*Is there anything else I can help(?: you)? with[?.!]?\s*$", ""),
        (r"\s*Don'?t hesitate to (?:ask|reach out)[^.!?]*[.!?]?\s*$", ""),
    ]
    for pattern, replacement in STRIP_PHRASES:
        assistant_text = _re.sub(pattern, replacement, assistant_text, flags=_re.IGNORECASE)
    assistant_text = assistant_text.strip()

    messages.append({"role": "assistant", "content": assistant_text})
    # Note which profile field the advisor just asked about, so it asks once
    try:
        prof_now = load_profile()
        asked = set(x for x in (prof_now.get("asked") or "").split(",") if x)
        probes = [
            ("specialty", r"(what|which)[^.?]{0,60}(specialty|speciality|field|"
                          r"area of medicine|kind of medicine)"),
            ("role", r"(what|which)[^.?]{0,60}(role|title|position|"
                     r"where do you sit|hat do you wear)"),
            ("first_name", r"(what|may i ask)[^.?]{0,40}(your name|"
                           r"should i call you|do you go by)"),
        ]
        for field, pattern in probes:
            if field not in asked and not prof_now.get(field):
                if re.search(pattern, assistant_text, re.I):
                    asked.add(field)
        if asked and ",".join(sorted(asked)) != (prof_now.get("asked") or ""):
            save_profile(asked=",".join(sorted(asked)))
    except Exception as e:
        app.logger.error(f"[profile] ask tracking failed: {e}")

    append_history("assistant", assistant_text)
    session["messages"] = _fit_history(messages)   # cookie fallback only

    # Auto-log every exchange to the DB (rating stays NULL until user gives feedback).
    # We log the visible parts only — the raw user text and the bot's reply.
    # Attachment context is not stored because it can be enormous; we store a
    # short description of what was attached instead.
    interaction_id = None
    try:
        attachment_note = attachment_display_name or ""
        if not attachment_note and folder_stats:
            used, _ = folder_stats
            attachment_note = f"Folder: {folder_name} ({used} file{'s' if used != 1 else ''})"
        interaction_id = db.log_interaction(
            user_message=user_input,   # raw question, without the inlined document text
            bot_reply=assistant_text,
            persona=CONFIG["persona_name"],
            attachment_info=attachment_note,
        )
        # Resolve city/region/country in the background — never blocks the reply
        record_location_async(interaction_id, client_ip())
        record_acknowledgement(interaction_id)
    except Exception as e:
        app.logger.error(f"log_interaction failed: {e}")

    # If the user asked for a specific file format, tell the client so it can
    # generate and download it without a further click.
    export_format = detect_export_format(user_input)
    if not export_format and wants_deliverable:
        # No format named — infer one from the shape of the reply, avoiding
        # anything the user explicitly ruled out.
        try:
            inferred = exports.suggest_format("", assistant_text)
            if inferred in rejected_formats(user_input):
                inferred = "docx"
            export_format = inferred
        except Exception as e:
            app.logger.error(f"Format inference failed: {e}")
            export_format = "docx"

    # Describe each deliverable in the reply so the client can offer them
    # separately, each in its own format.
    documents = []
    single_file = wants_single_file(user_input)
    separate_files = wants_separate_files(user_input) and not single_file
    try:
        split = [] if single_file else exports.split_documents(assistant_text)
        for i, doc in enumerate(split):
            if not doc["body"].strip():
                continue
            documents.append({
                "index": i,
                "title": doc["title"] or f"Document {i + 1}",
                "suggested": exports.suggest_format(doc["title"], doc["body"]),
            })
    except Exception as e:
        app.logger.error(f"Document split failed: {e}")

    # Several deliverables going into ONE file: prose is the safer container,
    # unless the user named a format or every part is slide-shaped.
    if (len(documents) > 1 and not separate_files
            and not detect_export_format(user_input)
            and not all(d["suggested"] == "pptx" for d in documents)):
        export_format = "docx"

    return jsonify({
        "reply": assistant_text,
        "interaction_id": interaction_id,
        "export_format": export_format,
        "documents": documents,
        "single_file": single_file,
        "separate_files": separate_files,
    })


@app.route("/acknowledge", methods=["POST"])
@paywall.paywall_required
def acknowledge_release():
    """Records that this session accepted the release and acknowledgment."""
    data = request.get_json(silent=True) or {}
    session["release_ack_at"] = datetime.now().isoformat(timespec="seconds")
    session["release_ack_version"] = str(data.get("version") or "")[:40]
    app.logger.info("[ack] release acknowledged")
    return jsonify({"ok": True, "acknowledged_at": session["release_ack_at"]})


@app.route("/reset", methods=["POST"])
@paywall.paywall_required
def reset():
    clear_history()
    return jsonify({"ok": True})


@app.route("/export/<fmt>", methods=["POST"])
@paywall.paywall_required
def export_response(fmt):
    """Render an assistant reply as a Word, PowerPoint, Excel, or PDF download."""
    from flask import send_file

    data = request.get_json(silent=True) or {}
    text = (data.get("text") or "").strip()
    title = (data.get("title") or "").strip()
    part = data.get("part")

    if not text:
        return jsonify({"error": "Nothing to export."}), 400

    # A reply can contain several deliverables. When the client asks for a
    # specific one, export only that document.
    if part is not None:
        try:
            docs = exports.split_documents(text)
            idx = int(part)
            if 0 <= idx < len(docs):
                text = docs[idx]["body"]
                title = title or docs[idx]["title"]
        except Exception as e:
            app.logger.error(f"Export part selection failed: {e}")
    if len(text) > 200_000:
        return jsonify({"error": "Response too long to export."}), 400
    if fmt.lower() not in exports.BUILDERS:
        return jsonify({"error": f"Unsupported format: {fmt}"}), 400

    try:
        buffer, filename, mimetype = exports.build(fmt, text, title)
    except Exception as e:
        app.logger.error(f"Export failed ({fmt}): {e}")
        return jsonify({"error": f"Could not generate {fmt.upper()}: {str(e)[:200]}"}), 500

    app.logger.info(f"Export generated: {filename}")
    return send_file(
        buffer,
        mimetype=mimetype,
        as_attachment=True,
        download_name=filename,
    )


@app.route("/feedback", methods=["POST"])
@paywall.paywall_required
def feedback():
    data = request.get_json(silent=True) or {}
    rating = data.get("rating")
    reply = (data.get("reply") or "")[:20000]
    comment = (data.get("comment") or "")[:4000]
    interaction_id = data.get("interaction_id")
    if rating not in ("up", "down"):
        return jsonify({"error": "Invalid rating"}), 400

    messages = load_history()
    last_user_msg = ""
    for m in reversed(messages):
        if m.get("role") == "user":
            last_user_msg = (m.get("content") or "")[:8000]
            break

    # Preferred path: update the row that was auto-logged when the bot replied.
    # Fallback path: if we don't have an interaction_id (e.g. reply predates this
    # feature, or DB was unavailable at log time), insert a fresh row.
    updated = False
    try:
        if interaction_id:
            updated = db.update_feedback_rating(int(interaction_id), rating, comment)
        if not updated:
            db.log_feedback(rating, last_user_msg, reply, CONFIG["persona_name"], comment)
    except Exception as e:
        app.logger.error(f"DB feedback log failed: {e}")
    app.logger.info(
        "FEEDBACK persona=%s rating=%s interaction_id=%s user_msg=%r reply=%r comment=%r",
        CONFIG["persona_name"], rating, interaction_id, last_user_msg, reply, comment,
    )
    return jsonify({"ok": True})


@app.route("/health")
def health():
    return jsonify({
        "status": "ok",
        "version": APP_VERSION,
        "build": APP_BUILD_NOTES,
        "admin_section_order": [
            "Feedback Overview", "Display Settings", "Conversation Log",
            "Knowledge Base", "Upload Document", "Upload Folder",
            "Add Knowledge from URL",
        ],
        "persona": CONFIG["persona_name"],
        "rag_enabled": db.is_enabled() and emb.is_enabled(),
    })


@app.route("/debug")
def debug_env():
    """Diagnostic endpoint — reports presence of key env vars without revealing values."""
    return jsonify({
        "DATABASE_URL_set": bool(os.environ.get("DATABASE_URL")),
        "DATABASE_URL_starts_with": (os.environ.get("DATABASE_URL", "")[:20] + "..." if os.environ.get("DATABASE_URL") else None),
        "VOYAGE_API_KEY_set": bool(os.environ.get("VOYAGE_API_KEY")),
        "VOYAGE_API_KEY_starts_with": (os.environ.get("VOYAGE_API_KEY", "")[:7] + "..." if os.environ.get("VOYAGE_API_KEY") else None),
        "ADMIN_PASSWORD_set": bool(os.environ.get("ADMIN_PASSWORD")),
        "FLASK_SECRET_KEY_set": bool(os.environ.get("FLASK_SECRET_KEY")),
        "ANTHROPIC_API_KEY_set": bool(os.environ.get("ANTHROPIC_API_KEY")),
        "db_is_enabled": db.is_enabled(),
        "emb_is_enabled": emb.is_enabled(),
        "psycopg_imported": db.HAS_PSYCOPG,
        "voyage_imported": emb.HAS_VOYAGE,
        "paywall_enabled": paywall.PAYWALL_ENABLED,
        "paywall_configured": paywall.is_configured(),
        "paywall_active": paywall.is_active(),
    })


# ---------------------------------------------------------------------------
# Paywall — auth (magic link) and billing (Stripe) routes
# ---------------------------------------------------------------------------

def _paywall_price_display():
    """Price shown on the checkout landing page. Reads from env or defaults."""
    return os.environ.get("PAYWALL_PRICE_DISPLAY", "29")


def _paywall_base_url():
    return paywall.PUBLIC_BASE_URL or request.host_url.rstrip("/")


@app.route("/auth/login", methods=["GET"])
def auth_login():
    return render_template_string(
        paywall.LOGIN_HTML,
        brand=CONFIG["persona_name"],
        notice=request.args.get("notice"),
        notice_type=request.args.get("notice_type", "ok"),
        sent=False,
    )


@app.route("/auth/send-link", methods=["POST"])
def auth_send_link():
    email = (request.form.get("email") or "").strip().lower()
    if not email or "@" not in email:
        return render_template_string(
            paywall.LOGIN_HTML,
            brand=CONFIG["persona_name"],
            notice="Please enter a valid email address.",
            notice_type="error",
            sent=False,
        )
    if not paywall.is_configured():
        return render_template_string(
            paywall.LOGIN_HTML,
            brand=CONFIG["persona_name"],
            notice="Paywall is not fully configured. Contact the administrator.",
            notice_type="error",
            sent=False,
        )
    magic_url = paywall.make_magic_link(email, _paywall_base_url())
    ok = paywall.send_magic_link_email(email, magic_url)
    if not ok:
        return render_template_string(
            paywall.LOGIN_HTML,
            brand=CONFIG["persona_name"],
            notice="We couldn't send the email. Please try again in a moment.",
            notice_type="error",
            sent=False,
        )
    return render_template_string(
        paywall.LOGIN_HTML,
        brand=CONFIG["persona_name"],
        notice=f"Sign-in link sent to {email}. Check your inbox (and spam folder).",
        notice_type="ok",
        sent=True,
    )


@app.route("/auth/verify", methods=["GET"])
def auth_verify():
    token = request.args.get("token", "")
    email = paywall.verify_magic_link(token)
    if not email:
        return render_template_string(
            paywall.LOGIN_HTML,
            brand=CONFIG["persona_name"],
            notice="That sign-in link is invalid or has expired. Please request a new one.",
            notice_type="error",
            sent=False,
        ), 400
    session["authenticated_email"] = email
    # Ensure user record exists
    paywall._upsert_user(email)
    # If they have a subscription, go to the bot. If not, go to checkout.
    if paywall.is_user_subscribed(email):
        return redirect(url_for("index"))
    return redirect(url_for("billing_checkout"))


@app.route("/auth/logout", methods=["GET", "POST"])
def auth_logout():
    session.pop("authenticated_email", None)
    session["messages"] = []
    return redirect(url_for("auth_login"))


@app.route("/billing/checkout", methods=["GET", "POST"])
def billing_checkout():
    email = session.get("authenticated_email")
    if not email:
        return redirect(url_for("auth_login"))
    # If already subscribed, no need for checkout
    if paywall.is_user_subscribed(email):
        return redirect(url_for("index"))
    if request.method == "POST":
        checkout_url = paywall.create_checkout_session(email, _paywall_base_url())
        if not checkout_url:
            return "Stripe checkout is not configured. Contact the administrator.", 500
        return redirect(checkout_url)
    # GET — show landing page with a "Continue" button
    return render_template_string(
        paywall.CHECKOUT_LANDING_HTML,
        brand=CONFIG["persona_name"],
        email=email,
        price=_paywall_price_display(),
    )


@app.route("/billing/success", methods=["GET"])
def billing_success():
    email = session.get("authenticated_email")
    if email:
        # The webhook usually beats us here, but sync as a fallback so the user
        # doesn't get bounced back to checkout on their next click.
        try:
            paywall.sync_user_from_stripe(email)
        except Exception as e:
            app.logger.error(f"[paywall] post-checkout sync failed: {e}")
    return redirect(url_for("index"))


@app.route("/billing/canceled", methods=["GET"])
def billing_canceled():
    return redirect(url_for("billing_checkout"))


@app.route("/billing/portal", methods=["GET", "POST"])
def billing_portal():
    email = session.get("authenticated_email")
    if not email:
        return redirect(url_for("auth_login"))
    portal_url = paywall.create_portal_session(email, _paywall_base_url())
    if not portal_url:
        return "Customer portal not available yet — no active subscription found.", 400
    return redirect(portal_url)


@app.route("/billing/webhook", methods=["POST"])
def billing_webhook():
    payload = request.get_data()
    sig = request.headers.get("Stripe-Signature", "")
    ok, msg = paywall.handle_stripe_webhook(payload, sig)
    if not ok:
        app.logger.error(f"[paywall] webhook error: {msg}")
        return jsonify({"error": msg}), 400
    return jsonify({"ok": True, "msg": msg})


# Serve image assets from project root
@app.route("/<path:filename>.png")
def serve_png(filename):
    return send_from_directory(".", f"{filename}.png")


@app.route("/<path:filename>.jpg")
def serve_jpg(filename):
    return send_from_directory(".", f"{filename}.jpg")


@app.route("/avatar/speak", methods=["POST"])
@paywall.paywall_required
def avatar_speak():
    """Return a talking-head video for a reply, when the feature is on."""
    if not talking_avatar_enabled():
        return jsonify({"ok": False, "error": "Talking avatar is off."}), 400

    data = request.get_json(silent=True) or {}
    text = (data.get("text") or "").strip()
    if not text:
        return jsonify({"ok": False, "error": "Nothing to speak."}), 400

    if TALKING_AVATAR_MODE == "demo":
        # Bundled clip so the flow can be seen without an account
        return jsonify({"ok": True, "video_url": "/advisor_placeholder.mp4",
                        "demo": True})
    try:
        url = _did_request(text)
        return jsonify({"ok": True, "video_url": url, "demo": False})
    except Exception as e:
        app.logger.error(f"[avatar] generation failed: {e}")
        return jsonify({"ok": False, "error": str(e)[:160]}), 502


@app.route("/advisor_idle.mp4")
def advisor_idle_mp4():
    """Looping animated portrait used as the avatar's resting state."""
    return send_from_directory(".", "advisor_idle.mp4")


@app.route("/advisor_idle.webm")
def advisor_idle_webm():
    return send_from_directory(".", "advisor_idle.webm")


@app.route("/advisor_placeholder.mp4")
def advisor_placeholder():
    return send_from_directory(".", "advisor_placeholder.mp4")


@app.route("/advisor_placeholder.webm")
def advisor_placeholder_webm():
    return send_from_directory(".", "advisor_placeholder.webm")


def materials_enabled() -> bool:
    return bool(load_settings().get("allow_materials"))


def scrub_internal_emails(text: str) -> str:
    """Point any internal address at client services.

    The full paragraph-level scrubber lives inside chat(); a plan is generated
    from an already-scrubbed transcript, so this only has to catch an address
    the model might reconstruct.
    """
    if not text:
        return text
    internal = ("j3p.health", "j3phealth.com")
    contact = CONFIG["contact_email"]

    def _swap(m):
        addr = m.group(0)
        if addr.lower() == contact.lower():
            return addr
        domain = addr.split("@")[-1].lower()
        return contact if domain in internal else addr

    return re.sub(r"[\w.+-]+@[\w.-]+\.\w+", _swap, text)


@app.route("/briefing/schedule", methods=["POST"])
@paywall.paywall_required
@login_required
def briefing_schedule():
    """Called when a participant books time — brief the advisor."""
    summary = build_briefing()
    if not summary:
        # Nothing worth sending; not an error the participant should see
        return jsonify({"ok": True, "sent": False, "reason": "conversation too short"})

    advisor = get_advisor(session.get("advisor_slug"))
    advisor_name = advisor["name"] if advisor else ""
    who = participant_label()

    emailed = email_briefing(advisor_name, who, summary)
    save_briefing(session.get("advisor_slug"), advisor_name, who, summary, emailed)
    app.logger.info(f"[briefing] prepared for {advisor_name or 'the team'} "
                    f"(emailed={emailed})")
    return jsonify({"ok": True, "sent": True, "emailed": emailed})


@app.route("/plan/create", methods=["POST"])
@paywall.paywall_required
@login_required
def plan_create():
    """Build an accountability plan from this conversation."""
    history = load_history()
    if len(history) < 2:
        return jsonify({"ok": False,
                        "error": "There's not enough of a conversation yet to "
                                 "build a plan from."}), 400

    transcript = []
    for turn in history[-24:]:
        who = "Them" if turn.get("role") == "user" else "You"
        transcript.append(f"{who}: {(turn.get('content') or '')[:2000]}")

    try:
        resp = client.messages.create(
            model=CONFIG["model"],
            max_tokens=1200,
            system=[{"type": "text", "text": PLAN_PROMPT}],
            messages=[{"role": "user",
                       "content": "Conversation:\n\n" + "\n\n".join(transcript)}],
        )
        parts = [b.text for b in resp.content if getattr(b, "type", "") == "text"]
        plan = "\n".join(parts).strip()
    except Exception as e:
        app.logger.error(f"[plans] generation failed: {e}")
        return jsonify({"ok": False,
                        "error": "Couldn't build the plan just now. Try again "
                                 "in a moment."}), 502

    if not plan:
        return jsonify({"ok": False, "error": "The plan came back empty."}), 502

    plan = scrub_internal_emails(plan)
    saved = save_plan(plan)
    return jsonify({"ok": True, "plan": plan, "saved": saved})


@app.route("/plan/list", methods=["GET"])
@paywall.paywall_required
@login_required
def plan_list():
    return jsonify({"ok": True, "plans": list_plans()})


@app.route("/plan/delete/<int:plan_id>", methods=["POST"])
@paywall.paywall_required
@login_required
def plan_delete(plan_id):
    if delete_plan(plan_id):
        return jsonify({"ok": True})
    return jsonify({"ok": False, "error": "Could not remove that plan."}), 400


@app.route("/profile", methods=["GET"])
@paywall.paywall_required
@login_required
def profile_view():
    """What the advisor remembers about this participant."""
    prof = load_profile()
    return jsonify({"ok": True, "profile": {
        k: prof.get(k) for k in PROFILE_FIELDS if prof.get(k)}})


@app.route("/profile/forget", methods=["POST"])
@paywall.paywall_required
@login_required
def profile_forget():
    if forget_profile():
        return jsonify({"ok": True})
    return jsonify({"ok": False, "error": "Could not clear it."}), 400


@app.route("/materials", methods=["GET"])
@paywall.paywall_required
@login_required
def materials_list():
    """This participant's own library."""
    if not materials_enabled():
        return jsonify({"ok": False, "error": "That feature is turned off."}), 403
    return jsonify({"ok": True, "docs": list_participant_docs(),
                    "limit": PARTICIPANT_DOC_LIMIT})


@app.route("/materials/text", methods=["POST"])
@paywall.paywall_required
@login_required
def materials_add_text():
    """Paste in writing — an article, notes, a bio, a draft."""
    if not materials_enabled():
        return jsonify({"ok": False, "error": "That feature is turned off."}), 403
    data = request.get_json(silent=True) or {}
    title = (data.get("title") or "").strip() or "Untitled note"
    body = (data.get("content") or "").strip()
    shared = bool(data.get("shared"))
    if len(body) < 20:
        return jsonify({"ok": False, "error": "Add a little more text than that."}), 400
    if add_participant_doc(title, body, kind="text", shared=shared):
        return jsonify({"ok": True})
    return jsonify({"ok": False,
                    "error": f"Could not save it. You can keep up to "
                             f"{PARTICIPANT_DOC_LIMIT} items."}), 400


@app.route("/materials/upload", methods=["POST"])
@paywall.paywall_required
@login_required
def materials_upload():
    """Upload documents — PDF, Word, PowerPoint, Excel, text."""
    if not materials_enabled():
        return jsonify({"ok": False, "error": "That feature is turned off."}), 403
    files = request.files.getlist("files") or []
    if not files:
        return jsonify({"ok": False, "error": "Choose a file first."}), 400

    shared = (request.form.get("shared") or "").lower() in ("1", "true", "on")
    added, failed = [], []
    for f in files:
        if not f or not f.filename:
            continue
        raw = f.read()
        if len(raw) > PARTICIPANT_DOC_MAX_BYTES:
            failed.append(f"{f.filename} (too large)")
            continue
        try:
            text = extract_attachment_text(f.filename, raw)
        except Exception as e:
            app.logger.error(f"[materials] extract failed for {f.filename}: {e}")
            text = ""
        if not text or not text.strip():
            failed.append(f"{f.filename} (no text found)")
            continue
        if add_participant_doc(f.filename, text, kind="upload", shared=shared):
            added.append(f.filename)
        else:
            failed.append(f"{f.filename} (library full)")

    if not added and failed:
        return jsonify({"ok": False, "error": "; ".join(failed[:4])}), 400
    return jsonify({"ok": True, "added": added, "failed": failed})


@app.route("/materials/delete/<int:doc_id>", methods=["POST"])
@paywall.paywall_required
@login_required
def materials_delete(doc_id):
    if not materials_enabled():
        return jsonify({"ok": False, "error": "That feature is turned off."}), 403
    if delete_participant_doc(doc_id):
        return jsonify({"ok": True})
    return jsonify({"ok": False, "error": "Could not remove that item."}), 400


@app.route("/advisor_avatar.jpg")
def advisor_avatar():
    """The photo shown beside advisor replies.

    An uploaded photo lives in Postgres so it survives redeploys; the bundled
    file is the fallback when nothing has been uploaded.
    """
    stored = load_avatar()
    if stored:
        data, mime, updated = stored
        resp = app.response_class(data, mimetype=mime or "image/jpeg")
        # Change the tag when the photo changes so browsers refetch it
        try:
            resp.headers["ETag"] = f'"{int(updated.timestamp())}"'
        except Exception:
            pass
        resp.headers["Cache-Control"] = "no-cache"
        return resp
    return send_from_directory(".", "advisor_avatar.jpg")


@app.route("/admin/avatar", methods=["POST"])
@admin_required
def admin_upload_avatar():
    """Replace the advisor photo."""
    file = request.files.get("avatar")
    if not file or not file.filename:
        flash("Choose an image first.")
        return redirect(url_for("admin_dashboard"))

    raw = file.read()
    if len(raw) > AVATAR_MAX_BYTES:
        flash(f"That image is {len(raw)/1048576:.1f} MB — the limit is "
              f"{AVATAR_MAX_BYTES // 1048576} MB.")
        return redirect(url_for("admin_dashboard"))

    ext = (file.filename.rsplit(".", 1)[-1] or "").lower()
    if ext not in ("jpg", "jpeg", "png", "webp", "gif"):
        flash("Use a JPG, PNG, WEBP or GIF image.")
        return redirect(url_for("admin_dashboard"))

    data, mime = prepare_avatar(raw)
    if store_avatar(data, mime):
        flash("✓ Advisor photo updated. Participants will see it immediately.")
    else:
        flash("Could not save the photo — check the database connection.")
    return redirect(url_for("admin_dashboard"))


@app.route("/admin/avatar/delete", methods=["POST"])
@admin_required
def admin_delete_avatar():
    """Revert to the photo bundled with the app."""
    if clear_avatar():
        flash("Reverted to the bundled advisor photo.")
    else:
        flash("Could not revert the photo.")
    return redirect(url_for("admin_dashboard"))


# ---------------------------------------------------------------------------
# Admin panel
# ---------------------------------------------------------------------------

ADMIN_LOGIN_HTML = """<!DOCTYPE html>
<html lang="en"><head>
<meta charset="UTF-8" />
<meta name="viewport" content="width=device-width, initial-scale=1.0" />
<title>Admin — {{ cfg.persona_name }}</title>
<link rel="icon" href="{{ cfg.favicon_url }}" />
<link rel="preconnect" href="https://fonts.googleapis.com">
<link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>
<link href="https://fonts.googleapis.com/css2?family=Jost:wght@300;400;500;600&display=swap" rel="stylesheet">
<style>
  :root {
    --navy: #27334A; --gold: #D2BC8D; --rust: #9D432C;
    --paper: #FAF6F0; --line: rgba(39,51,74,0.12); --muted: #6B7280;
  }
  *, *::before, *::after { box-sizing: border-box; }
  body {
    margin: 0; font-family: 'Jost', -apple-system, BlinkMacSystemFont, sans-serif;
    background: var(--paper); color: var(--navy);
    display: flex; align-items: center; justify-content: center;
    min-height: 100vh; padding: 1.25rem;
  }
  .box {
    background: #fff; border-radius: 4px; width: 100%; max-width: 420px;
    box-shadow: 0 18px 50px rgba(39,51,74,0.18); overflow: hidden;
  }
  .box-head {
    background: var(--navy); border-bottom: 2px solid var(--gold);
    padding: 1rem 1.75rem; display: flex; align-items: center; gap: 0.9rem;
  }
  .box-head img { height: 44px; width: auto; display: block; }
  .brand-divider { width: 1px; height: 26px; background: rgba(210,188,141,0.45); }
  .box-head span {
    color: var(--gold); font-size: 0.78rem;
    letter-spacing: 0.22em; text-transform: uppercase; font-weight: 400;
  }
  .content { padding: 1.75rem 2rem 1.7rem; }
  h1 {
    margin: 0 0 1.2rem; font-size: 0.85rem; font-weight: 500;
    letter-spacing: 0.16em; text-transform: uppercase; color: var(--navy);
    padding-bottom: 0.65rem; border-bottom: 1px solid var(--line);
  }
  input {
    width: 100%; padding: 0.8rem 0.9rem; border: 1px solid var(--line);
    border-radius: 2px; font-family: inherit; font-size: 0.95rem;
    background: var(--paper); color: var(--navy);
  }
  input:focus { outline: none; border-color: var(--gold); background: #fff; }
  button {
    width: 100%; margin-top: 0.9rem; padding: 0.85rem;
    background: var(--navy); color: var(--gold); border: 1px solid var(--navy);
    border-radius: 2px; cursor: pointer; font-family: inherit;
    font-size: 0.78rem; letter-spacing: 0.18em; text-transform: uppercase;
    transition: background 0.2s ease, color 0.2s ease;
  }
  button:hover { background: var(--gold); color: var(--navy); }
  .err {
    background: #FEE; color: var(--rust); border: 1px solid #E7C3BA;
    padding: 0.65rem 0.85rem; border-radius: 2px;
    font-size: 0.85rem; margin-bottom: 1rem;
  }
  .foot {
    margin-top: 1.1rem; font-size: 0.62rem; color: var(--muted);
    letter-spacing: 0.1em; text-transform: uppercase; text-align: center;
  }
  @media (max-width: 480px) {
    .box-head { padding: 0.85rem 1.2rem; gap: 0.7rem; }
    .box-head img { height: 36px; }
    .box-head span { font-size: 0.68rem; letter-spacing: 0.16em; }
    .content { padding: 1.4rem 1.3rem 1.3rem; }
  }
</style></head><body>
  <form method="POST" class="box">
    <div class="box-head">
      <img src="{{ cfg.logo_url }}" alt="{{ cfg.persona_name }}" />
      <div class="brand-divider"></div>
      <span>{{ cfg.persona_name }}</span>
    </div>
    <div class="content">
      <h1>Admin sign in</h1>
      {% if error %}<div class="err">{{ error }}</div>{% endif %}
      <input type="password" name="password" placeholder="Password" autofocus required />
      <button type="submit">Sign in</button>
      <div class="foot">Authorized access only</div>
    </div>
  </form>
</body></html>"""

LEARNING_ARCHIVE_HTML = """<!DOCTYPE html>
<html lang="en"><head>
<meta charset="UTF-8" />
<meta name="viewport" content="width=device-width, initial-scale=1.0" />
<title>Learning archive — {{ cfg.persona_name }}</title>
<link rel="icon" href="{{ cfg.favicon_url }}" />
<link rel="preconnect" href="https://fonts.googleapis.com">
<link href="https://fonts.googleapis.com/css2?family=Jost:wght@300;400;500;600&display=swap" rel="stylesheet">
<style>
  :root { --navy:#27334A; --gold:#D2BC8D; --paper:#FAF6F0; --line:rgba(39,51,74,0.12); --muted:#6B7280; }
  *,*::before,*::after { box-sizing:border-box; }
  body { margin:0; font-family:'Jost',-apple-system,sans-serif; background:var(--paper); color:var(--navy); }
  header { background:var(--navy); border-bottom:2px solid var(--gold);
           padding:1rem 1.5rem; display:flex; justify-content:space-between; align-items:center; }
  header h1 { margin:0; font-size:0.85rem; letter-spacing:0.18em; text-transform:uppercase;
              color:var(--gold); font-weight:500; }
  header a { color:var(--gold); text-decoration:none; font-size:0.75rem;
             letter-spacing:0.12em; text-transform:uppercase; }
  .container { max-width:1100px; width:95%; margin:0 auto; padding:2rem 1.5rem; }
  .section { background:#fff; border:1px solid var(--line); border-radius:4px; padding:1.5rem; }
  h2 { margin:0 0 1rem; font-size:0.8rem; letter-spacing:0.14em; text-transform:uppercase;
       padding-bottom:0.6rem; border-bottom:1px solid var(--line); font-weight:500; }
  table { width:100%; border-collapse:collapse; font-size:0.85rem; }
  th { text-align:left; font-size:0.68rem; letter-spacing:0.1em; text-transform:uppercase;
       color:var(--muted); padding:0.6rem 0.5rem; border-bottom:2px solid var(--navy); }
  td { padding:0.6rem 0.5rem; border-bottom:1px solid var(--line); vertical-align:top; }
  .muted { color:var(--muted); }
  details summary { cursor:pointer; font-size:0.78rem; color:var(--muted); }
  pre { white-space:pre-wrap; font-size:0.72rem; background:var(--paper);
        padding:0.6rem; border-radius:3px; margin:0.5rem 0 0; }
</style></head><body>
<header>
  <h1>{{ cfg.persona_name }} — Learning archive</h1>
  <a href="/admin">&larr; Back to admin</a>
</header>
<div class="container">
  <div class="section">
    <h2>Archived runs ({{ runs|length }})</h2>
    {% if runs %}
    <p class="muted" style="font-size:0.84rem; margin:0 0 1rem;">
      Every learning run that has been archived, newest first. Runs that learned
      nothing are archived automatically so the main panel stays readable.
    </p>
    <table>
      <tr><th>When</th><th>Trigger</th><th>Learned</th><th>Held back</th><th>Detail</th></tr>
      {% for run in runs %}
      <tr>
        <td>{{ run.when }}</td>
        <td class="muted">{{ run.trigger }}</td>
        <td><strong>{{ run.approved }}</strong></td>
        <td class="muted">{{ run.skipped }}</td>
        <td>
          {% if run.detail %}
          <details><summary>Show</summary><pre>{{ run.detail }}</pre></details>
          {% else %}<span class="muted">—</span>{% endif %}
        </td>
      </tr>
      {% endfor %}
    </table>
    {% else %}
    <p class="muted">Nothing archived yet.</p>
    {% endif %}
  </div>
</div>
</body></html>"""


ADMIN_HTML = """<!DOCTYPE html><html><head>
<title>Admin — {{ cfg.persona_name }}</title>
<link rel="icon" href="{{ cfg.favicon_url }}" />
<style>
:root { --navy: #27334A; --gold: #D2BC8D; --rust: #9D432C; --paper: #FAF6F0; --line: rgba(39,51,74,0.12); }
body { font-family: -apple-system, sans-serif; background: var(--paper); color: var(--navy); margin: 0; }
header { background: var(--navy); color: #fff; padding: 1rem 2rem; display: flex; justify-content: space-between; align-items: center; border-bottom: 2px solid var(--gold); }
header h1 { margin: 0; font-size: 1rem; letter-spacing: 0.18em; text-transform: uppercase; color: var(--gold); font-weight: 400; }
header a { color: rgba(210,188,141,0.7); text-decoration: none; font-size: 0.75rem; letter-spacing: 0.12em; text-transform: uppercase; }
header a:hover { color: var(--gold); }
*, *::before, *::after { box-sizing: border-box; }
/* Use most of the window — the conversation log is wide and was cramped
   inside a 1000px column on a large screen. */
.container { max-width: 1800px; width: 95%; margin: 0 auto; padding: 2rem 1.5rem; }
/* A wide log scrolls inside its own card rather than the whole page */
.section { overflow-x: auto; }
@media (max-width: 900px) { .container { width: 100%; padding: 1.25rem 1rem; } }
.section { background: #fff; border: 1px solid var(--line); border-radius: 4px; padding: 1.5rem 1.75rem; margin-bottom: 1.5rem; }
.confirm-overlay {
  position: fixed; inset: 0; background: rgba(39,51,74,0.55);
  display: flex; align-items: center; justify-content: center;
  z-index: 1000; padding: 1.25rem;
}
.confirm-overlay[hidden] { display: none; }
.confirm-box {
  background: #fff; border-radius: 4px; max-width: 460px; width: 100%;
  box-shadow: 0 20px 60px rgba(39,51,74,0.3); overflow: hidden;
}
.confirm-head {
  background: var(--navy); border-bottom: 2px solid var(--rust);
  padding: 0.9rem 1.4rem;
}
.confirm-head h3 {
  margin: 0; color: var(--paper); font-size: 0.82rem; font-weight: 500;
  letter-spacing: 0.16em; text-transform: uppercase;
}
.confirm-body { padding: 1.3rem 1.4rem 1rem; }
.confirm-body p { margin: 0 0 0.8rem; font-size: 0.93rem; line-height: 1.6; }
.confirm-detail {
  background: var(--paper); border-left: 3px solid var(--rust);
  padding: 0.6rem 0.8rem; margin-bottom: 0.9rem;
  font-size: 0.85rem; max-height: 160px; overflow-y: auto;
}
.confirm-type-label { display: block; font-size: 0.82rem; margin-bottom: 0.4rem; }
#confirm-type {
  width: 100%; padding: 0.6rem 0.7rem; border: 1px solid var(--line);
  border-radius: 2px; font-family: inherit; font-size: 0.95rem;
  letter-spacing: 0.1em; text-transform: uppercase;
}
#confirm-type:focus { outline: none; border-color: var(--rust); }
.confirm-actions {
  display: flex; justify-content: flex-end; gap: 0.6rem;
  padding: 0 1.4rem 1.3rem;
}
.confirm-cancel { background: transparent; color: var(--navy); }
#confirm-go:disabled { opacity: 0.45; cursor: not-allowed; }
.group-heading {
  font-size: 1.35rem; letter-spacing: 0.1em; text-transform: uppercase;
  color: var(--navy); font-weight: 500;
  margin: 2.25rem 0 0.9rem; padding-bottom: 0.5rem;
  border-bottom: 2px solid var(--gold);
}
.section h2 { margin: 0 0 1rem 0; font-size: 0.85rem; letter-spacing: 0.16em; text-transform: uppercase; color: var(--navy); border-bottom: 1px solid var(--line); padding-bottom: 0.6rem; }
.stats { display: flex; gap: 2rem; margin-bottom: 0.5rem; }
.stat { flex: 1; }
.stat-value { font-size: 1.8rem; font-weight: 500; color: var(--navy); }
.stat-label { font-size: 0.7rem; letter-spacing: 0.14em; text-transform: uppercase; color: #6B7280; }
table { width: 100%; border-collapse: collapse; font-size: 0.85rem; }
th { text-align: left; padding: 0.6rem 0.5rem; border-bottom: 2px solid var(--navy); font-size: 0.7rem; letter-spacing: 0.12em; text-transform: uppercase; color: var(--navy); }
td { padding: 0.6rem 0.5rem; border-bottom: 1px solid var(--line); vertical-align: top; }
.tag-up { background: #1B998B; color: #fff; padding: 2px 8px; border-radius: 2px; font-size: 0.7rem; }
.tag-down { background: var(--rust); color: #fff; padding: 2px 8px; border-radius: 2px; font-size: 0.7rem; }
.tag-ack { color: #1B998B; font-weight: 600; margin-right: 0.2rem; }
.rating-set { display: inline-flex; gap: 0.15rem; margin-left: 0.35rem; vertical-align: middle; }
.rate-btn {
  background: transparent; border: 1px solid var(--line); border-radius: 2px;
  cursor: pointer; padding: 0.1rem 0.3rem; font-size: 0.72rem; line-height: 1.3;
  color: var(--navy); transition: all 0.15s ease;
}
.rate-btn:hover:not(:disabled) { border-color: var(--gold); background: #fff; }
.rate-btn.on-up { background: #1B998B; border-color: #1B998B; }
.rate-btn.on-down { background: var(--rust); border-color: var(--rust); }
.rate-btn:disabled { opacity: 0.5; cursor: default; }
.tag-lesson { background: #2D7D5F; color: #fff; padding: 2px 8px; border-radius: 2px; font-size: 0.65rem;
              margin-left: 0.3rem; letter-spacing: 0.05em; }
.btn { padding: 0.6rem 1.1rem; background: var(--navy); color: var(--gold); border: 1px solid var(--navy); border-radius: 2px; cursor: pointer; font-size: 0.75rem; letter-spacing: 0.14em; text-transform: uppercase; text-decoration: none; display: inline-block; }
.btn:hover { background: var(--gold); color: var(--navy); }
.btn-danger { background: var(--rust); color: #fff; border-color: var(--rust); padding: 0.3rem 0.7rem; font-size: 0.7rem; }
.btn-danger:hover { background: #fff; color: var(--rust); }
/* All three upload rows share one grid, so the field and the button line up
   across sections regardless of button label length. */
form.upload {
  display: grid;
  /* Columns collapse on their own when there isn't room, so this holds on a
     phone regardless of how the browser reports the viewport width. */
  grid-template-columns: repeat(auto-fit, minmax(280px, auto));
  grid-template-areas: none;
  gap: 0.6rem; align-items: center;
}
@media (min-width: 860px) {
  /* four columns now: file, title, which base, submit */
  form.upload { grid-template-columns: 280px minmax(0, 1fr) 200px 200px; }
}
form.upload > * { min-width: 0; }
form.upload .btn { width: 100%; justify-content: center; text-align: center; }
@media (max-width: 820px) {
  form.upload { grid-template-columns: 1fr; }
  form.upload .btn { width: auto; justify-self: start; }
}
.upload-stacked select[name="owner"] { max-width: 320px; margin-top: 0.6rem; }
form.upload select, select[name="owner"] {
  padding: 0.5rem; border: 1px solid var(--line); border-radius: 2px;
  font-family: inherit; font-size: 0.85rem; background: #fff;
  color: var(--navy); width: 100%;
}
input[type="file"], input[type="text"] {
  padding: 0.5rem; border: 1px solid var(--line); border-radius: 2px;
  font-family: inherit; width: 100%;
}
.flash { padding: 0.7rem 1rem; background: var(--gold); color: var(--navy); border-radius: 2px; margin-bottom: 1rem; font-size: 0.85rem; }
.muted { color: #6B7280; font-size: 0.8rem; }
.warn { background: #fef3c7; border: 1px solid #f59e0b; padding: 0.7rem 1rem; border-radius: 2px; margin-bottom: 1rem; font-size: 0.85rem; }
/* Knowledge Base table — keeps long source URLs from stretching the row */
.kb-table { table-layout: fixed; width: 100%; }
.kb-table th:nth-child(1), .kb-table td:nth-child(1) { width: 28%; }
.kb-table th:nth-child(2), .kb-table td:nth-child(2) { width: 13%; }
.kb-table th:nth-child(3), .kb-table td:nth-child(3) { width: 25%; }
.kb-table th:nth-child(4), .kb-table td:nth-child(4) { width: 8%; }
.kb-table th:nth-child(5), .kb-table td:nth-child(5) { width: 14%; }
.kb-table th:nth-child(6), .kb-table td:nth-child(6) { width: 10%; }
.kb-title { font-weight: 500; word-break: break-word; }
.kb-source {
  font-size: 0.78rem; overflow: hidden; text-overflow: ellipsis;
  white-space: nowrap; max-width: 0;
}
.kb-date { white-space: nowrap; font-size: 0.78rem; }
@media (max-width: 820px) {
  .kb-table { table-layout: auto; }
  .kb-source { white-space: normal; max-width: none; word-break: break-all; }
}
.truncate { max-width: 320px; overflow: hidden; text-overflow: ellipsis; white-space: nowrap; }
    .expand-btn {
      background: transparent;
      border: 1px solid var(--line);
      color: var(--navy);
      padding: 0.25rem 0.55rem;
      border-radius: 2px;
      cursor: pointer;
      font-size: 0.7rem;
      font-family: inherit;
      letter-spacing: 0.1em;
      text-transform: uppercase;
      transition: background 0.15s ease;
    }
    .expand-btn:hover { background: var(--gold); color: var(--navy); }
    .feedback-detail {
      background: var(--paper);
      border-top: 1px dashed var(--line);
    }
    .feedback-detail td { padding: 1rem 1.2rem !important; }
    .feedback-detail-block { margin-bottom: 1rem; }
    .feedback-detail-block:last-child { margin-bottom: 0; }
    .feedback-detail-label {
      font-size: 0.7rem;
      letter-spacing: 0.14em;
      text-transform: uppercase;
      color: var(--muted);
      margin-bottom: 0.3rem;
      font-weight: 500;
    }
    .feedback-detail-content {
      background: #fff;
      border: 1px solid var(--line);
      border-radius: 3px;
      padding: 0.8rem 1rem;
      white-space: pre-wrap;
      word-wrap: break-word;
      font-size: 0.88rem;
      line-height: 1.55;
      color: var(--navy);
      max-height: 400px;
      overflow-y: auto;
    }
    .feedback-detail-content.comment-highlight {
      background: #FFF9E6;
      border-color: var(--gold);
    }
    .feedback-detail-meta {
      font-size: 0.75rem;
      color: var(--muted);
      font-style: italic;
    }
</style></head><body>
<header>
  <h1>{{ cfg.persona_name }} — Admin
    <span style="font-size: 0.6rem; letter-spacing: 0.08em; color: rgba(210,188,141,0.55);
                 margin-left: 0.6rem; text-transform: none;">build {{ app_version }}</span>
  </h1>
  <div>
    <a href="/" style="margin-right: 1.5rem;">← Back to bot</a>
    <a href="/admin/logout">Sign out</a>
  </div>
</header>
<div class="container">
  {% with messages = get_flashed_messages() %}
    {% for m in messages %}<div class="flash">{{ m }}</div>{% endfor %}
  {% endwith %}

  {% if not rag_ready %}
  <div class="warn">
    <strong>RAG is not fully configured.</strong>
    {% if not db_ok %}Set up Railway Postgres (Add Plugin → PostgreSQL) and the <code>DATABASE_URL</code> will appear automatically.{% endif %}
    {% if not emb_ok %}Set <code>OPENAI_API_KEY</code> in environment variables.{% endif %}
    Once both are set, redeploy and you can upload documents.
  </div>
  {% endif %}

  <!-- Confirmation gate for every destructive action -->
  <div id="confirm-overlay" class="confirm-overlay" hidden>
    <div class="confirm-box" role="dialog" aria-modal="true" aria-labelledby="confirm-title">
      <div class="confirm-head">
        <h3 id="confirm-title">Confirm deletion</h3>
      </div>
      <div class="confirm-body">
        <p id="confirm-message">This cannot be undone.</p>
        <div id="confirm-detail" class="confirm-detail" hidden></div>
        <label id="confirm-type-wrap" hidden>
          <span class="confirm-type-label">Type <strong>DELETE</strong> to confirm:</span>
          <input type="text" id="confirm-type" autocomplete="off" spellcheck="false" />
        </label>
      </div>
      <div class="confirm-actions">
        <button type="button" id="confirm-cancel" class="btn confirm-cancel">Cancel</button>
        <button type="button" id="confirm-go" class="btn btn-danger">Delete</button>
      </div>
    </div>
  </div>

  <h2 class="group-heading">Display</h2>

  <div class="section">
    <h2>Display Settings</h2>
    <form method="POST" action="/admin/settings" id="settings-form">
      <input type="hidden" name="_fields"
             value="show_scheduling_button,show_avatar,allow_materials,avatar_name" />
      <label style="display: flex; align-items: flex-start; gap: 0.7rem;
                    cursor: pointer; font-size: 0.9rem; line-height: 1.5;">
        <input type="checkbox" name="show_scheduling_button" value="1"
               {% if settings.show_scheduling_button %}checked{% endif %}
               style="margin-top: 0.2rem; width: 17px; height: 17px;
                      accent-color: var(--navy); cursor: pointer;" />
        <span>
          <strong>Show the &ldquo;Schedule Time With a J3P Advisor&rdquo; button</strong><br />
          <span class="muted">
            When off, the scheduling button is hidden from the chat page and
            participants use the advisor without a booking prompt. Takes effect
            immediately for everyone.
          </span>
        </span>
      </label>
      <label style="display: flex; align-items: flex-start; gap: 0.7rem;
                    cursor: pointer; font-size: 0.9rem; line-height: 1.5;
                    margin-top: 1.1rem; padding-top: 1.1rem;
                    border-top: 1px dashed var(--line);">
        <input type="checkbox" name="show_avatar" value="1"
               {% if settings.show_avatar %}checked{% endif %}
               style="margin-top: 0.2rem; width: 17px; height: 17px;
                      accent-color: var(--navy); cursor: pointer;" />
        <span>
          <strong>Show the advisor avatar</strong><br />
          <span class="muted">
            Controls the whole avatar: the animated portrait that sits beside
            the chat and the small photo on each reply. When off, participants
            see the advisor with no photo at all — useful if you'd rather they
            didn't read it as a specific person.
          </span>
        </span>
      </label>

      <div style="margin-top: 0.9rem; padding-left: 1.6rem; display: flex;
                  align-items: center; gap: 0.9rem; flex-wrap: wrap;">
        <img src="{{ cfg.avatar_url }}?v={{ avatar_version }}" alt="Current advisor photo"
             onerror="this.style.display='none'"
             style="width: 54px; height: 54px; border-radius: 50%;
                    object-fit: cover; border: 1.5px solid var(--gold);" />
        <span class="muted" style="font-size: 0.8rem;">
          {% if avatar_custom %}Uploaded photo{% else %}Bundled photo{% endif %}
        </span>
      </div>

    </form>

    <div style="margin-top: 1.4rem; padding-top: 1.1rem;
                border-top: 1px dashed var(--line);">
      <p style="margin: 0 0 0.6rem 0; font-size: 0.9rem;">
        <strong>Replace the advisor photo</strong>
      </p>
      <p class="muted" style="margin: 0 0 0.8rem 0; font-size: 0.82rem; line-height: 1.55;">
        JPG, PNG, WEBP or GIF, up to {{ avatar_max_mb }} MB. It's cropped to a
        centred square and resized automatically, and stored in the database so
        it survives redeploys.
      </p>
      <form method="POST" action="/admin/avatar" enctype="multipart/form-data"
            style="display: flex; gap: 0.6rem; align-items: center; flex-wrap: wrap;">
        <input type="file" name="avatar" accept=".jpg,.jpeg,.png,.webp,.gif" required />
        <button type="submit" class="btn">Upload photo</button>
      </form>
      {% if avatar_custom %}
      <form method="POST" action="/admin/avatar/delete" style="margin-top: 0.7rem;">
        <button type="submit" class="btn"
                style="background: transparent; color: var(--rust);
                       border-color: var(--rust);">Revert to bundled photo</button>
      </form>
      {% endif %}
    </div>

      <div style="margin-top: 1rem;">
        <label for="avatar_name"
               style="display: block; font-size: 0.85rem; margin-bottom: 0.35rem;">
          <strong>Name shown with this photo</strong>
        </label>
        <p class="muted" style="margin: 0 0 0.5rem 0; font-size: 0.8rem;">
          Appears beneath the avatar on the main link. Leave blank to use
          “{{ cfg.persona_name }}”. Advisors with their own profile always show
          their own name. Start typing to pick an existing advisor.
        </p>
        <input type="text" id="avatar_name" name="avatar_name"
               form="settings-form" list="advisor-name-options"
               value="{{ settings.avatar_name or '' }}"
               placeholder="e.g. Alan Friedman"
               style="max-width: 340px; padding: 0.5rem;
                      border: 1px solid var(--line); border-radius: 2px;
                      font-family: inherit;" />
        <datalist id="advisor-name-options">
          {% for adv in advisors %}<option value="{{ adv.name }}"></option>{% endfor %}
        </datalist>
      </div>

      <label style="display: flex; align-items: flex-start; gap: 0.7rem;
                    cursor: pointer; font-size: 0.9rem; line-height: 1.5;
                    margin-top: 1.1rem; padding-top: 1.1rem;
                    border-top: 1px dashed var(--line);">
        <input type="checkbox" name="allow_materials" value="1" form="settings-form"
               {% if settings.allow_materials %}checked{% endif %}
               style="margin-top: 0.2rem; width: 17px; height: 17px;
                      accent-color: var(--navy); cursor: pointer;" />
        <span>
          <strong>Let participants add their own documents &amp; writing</strong><br />
          <span class="muted">
            Adds a link under the chat box where participants can upload a CV,
            a plan or an article, or paste their own writing. Their material is
            private to them and never enters the shared knowledge base. When
            off, the link disappears and existing material is left untouched
            but unused.
          </span>
        </span>
      </label>

      <button type="submit" class="btn" form="settings-form"
              style="margin-top: 1rem;">Save settings</button>





  </div>

  <h2 class="group-heading">Advisors</h2>

  <div class="section">
    <h2>Advisor Profiles</h2>
    <p class="muted" style="margin: 0 0 1rem 0;">
      Each advisor gets their own photo and their own pair of links. Everything
      else — the knowledge base, the guardrails, the conversation log — is shared
      across all of them.
    </p>

    {% if advisors %}
    <table style="margin-bottom: 1.5rem;">
      <tr><th style="width: 58px;"></th><th>Advisor</th><th>Links</th><th></th></tr>
      {% for adv in advisors %}
      <tr>
        <td>
          <img src="/a/{{ adv.slug }}/photo.jpg" alt=""
               onerror="this.style.display='none'"
               style="width: 44px; height: 44px; border-radius: 50%;
                      object-fit: cover; border: 1.5px solid var(--gold);" />
        </td>
        <td>
          <strong>{{ adv.name }}</strong><br />
          <span class="muted" style="font-size: 0.76rem;">
            {{ adv.slug }}{% if not adv.has_photo %} · no photo, using the default{% endif %}
          </span>
        </td>
        <td style="font-size: 0.78rem; line-height: 1.9;">
          <code>{{ base_url }}/a/{{ adv.slug }}/scheduling</code>
          <span class="muted">— with booking</span><br />
          <code>{{ base_url }}/a/{{ adv.slug }}/no-scheduling</code>
          <span class="muted">— without</span><br />
          <code>{{ base_url }}/a/{{ adv.slug }}</code>
          <span class="muted">— follows the default</span>
        </td>
        <td style="text-align: right;">
          <form method="POST" action="/admin/advisors/delete/{{ adv.slug }}"
                style="display:inline;" data-doc-title="{{ adv.name }}">
            <button type="submit" class="btn btn-danger">Delete</button>
          </form>
        </td>
      </tr>
      {% endfor %}
    </table>
    {% else %}
    <p class="muted">
      No advisor profiles yet. The links above still work and use the default
      photo; add a profile below to give someone their own.
    </p>
    {% endif %}

    <div style="padding-top: 1.1rem; border-top: 1px dashed var(--line);">
      <p style="margin: 0 0 0.6rem 0; font-size: 0.9rem;">
        <strong>Add or update an advisor</strong>
      </p>
      <p class="muted" style="margin: 0 0 0.8rem 0; font-size: 0.82rem;">
        Re-using an existing name updates that profile. Leave the photo blank to
        keep the current one.
      </p>
      <form method="POST" action="/admin/advisors" enctype="multipart/form-data"
            class="upload">
        <input type="text" name="name" placeholder="Advisor name (e.g. Bruce Gewertz)" required />
        <input type="file" name="photo" accept=".jpg,.jpeg,.png,.webp,.gif" />
        <button type="submit" class="btn">Save advisor</button>
      </form>
    </div>
  </div>

  <h2 class="group-heading">Access</h2>

  <div class="section">
    <h2>Participant Access</h2>
    <form method="POST" action="/admin/settings">
      <input type="hidden" name="_fields" value="require_login" />
      <label style="display: flex; align-items: flex-start; gap: 0.7rem;
                    cursor: pointer; font-size: 0.9rem; line-height: 1.5;
                    margin-top: 0;">
        <input type="checkbox" name="require_login" value="1" 
               {% if settings.require_login %}checked{% endif %}
               {% if not mail_ready %}disabled{% endif %}
               style="margin-top: 0.2rem; width: 17px; height: 17px;
                      accent-color: var(--navy); cursor: pointer;" />
        <span>
          <strong>Require participants to sign in</strong><br />
          <span class="muted">
            Participants enter their email and follow a one-time link. Signing in
            lets the advisor remember previous conversations, so each person
            picks up where they left off.
          </span>
          {% if not mail_ready %}
          <br /><span style="color: var(--rust); font-size: 0.82rem;">
            Unavailable: no email is configured, so sign-in links couldn't be
            delivered and everyone would be locked out. Set POSTMARK_SERVER_TOKEN
            or SMTP_HOST first.
          </span>
          {% endif %}
        </span>
      </label>

      </label>

      <button type="submit" class="btn" style="margin-top: 1rem;">Save</button>
    </form>
  </div>

  <h2 class="group-heading">Scheduling</h2>

  <div class="section">
    <h2>Scheduling Links</h2>
    <p class="muted" style="margin: 0 0 1rem 0; font-size: 0.87rem; line-height: 1.6;">
      The scheduling toggle in Display Settings is the default for
      <code>{{ base_url }}/</code>. To run both experiences at once, share these
      fixed links instead — each ignores the default and always behaves the same
      way, whoever opens it.
    </p>
    <table style="font-size: 0.85rem;">
      <tr>
        <th style="width: 45%;">Link</th><th>Scheduling button</th>
      </tr>
      <tr>
        <td><code>{{ base_url }}/scheduling</code></td>
        <td>Always shown</td>
      </tr>
      <tr>
        <td><code>{{ base_url }}/no-scheduling</code></td>
        <td>Always hidden</td>
      </tr>
      <tr>
        <td><code>{{ base_url }}/</code></td>
        <td class="muted">Follows the Display Settings toggle</td>
      </tr>
    </table>
  </div>

  <h2 class="group-heading">Feedback</h2>

  <div class="section">
    <h2>Feedback Overview</h2>
    <div class="stats">
      <div class="stat">
        <div class="stat-value">{{ stats.up }}</div>
        <div class="stat-label">Thumbs up</div>
      </div>
      <div class="stat">
        <div class="stat-value">{{ stats.down }}</div>
        <div class="stat-label">Thumbs down</div>
      </div>
      <div class="stat">
        <div class="stat-value">{{ stats.total }}</div>
        <div class="stat-label">Total ratings</div>
      </div>
      <div class="stat">
        <div class="stat-value">
          {% if stats.total > 0 %}{{ (100 * stats.up / stats.total)|round(0)|int }}%{% else %}—{% endif %}
        </div>
        <div class="stat-label">Helpful rate</div>
      </div>
    </div>
    <form method="POST" action="/admin/settings"
          style="margin-top: 1.4rem; padding-top: 1.1rem;
                 border-top: 1px dashed var(--line);">
      <input type="hidden" name="_fields" value="auto_learning" />
      <label style="display: flex; align-items: flex-start; gap: 0.7rem;
                    cursor: pointer; font-size: 0.9rem; line-height: 1.5;">
        <input type="checkbox" name="auto_learning" value="1"
               {% if settings.auto_learning %}checked{% endif %}
               style="margin-top: 0.2rem; width: 17px; height: 17px;
                      accent-color: var(--navy); cursor: pointer;" />
        <span>
          <strong>Learn from feedback automatically</strong><br />
          <span class="muted">
            Runs the Continuous Learning step below on a schedule (every
            {{ learning_interval }}h), so commented thumbs-down feedback becomes
            lessons without you doing anything. Sensitive exchanges are still
            held back for manual review.
          </span>
        </span>
      </label>
      <button type="submit" class="btn" style="margin-top: 0.9rem;">Save</button>
    </form>
  </div>

  <div class="section">
    <h2>Continuous Learning</h2>
    <p class="muted" style="margin: 0 0 1rem 0; font-size: 0.87rem; line-height: 1.6;">
      Rated exchanges become lessons the advisor sees whenever a similar
      question comes up again. A <strong>thumbs up</strong> teaches it the shape
      of a reply that landed well; a <strong>thumbs down</strong> with a comment
      teaches it what to avoid. This processes them in bulk instead of one at a
      time.
    </p>
    <form method="POST" action="/admin/learning/run" style="display: inline;">
      <button type="submit" name="preview" value="1" class="btn"
              style="background: transparent; color: var(--navy);">Preview</button>
    </form>
    <form method="POST" action="/admin/learning/run" style="display: inline;">
      <button type="submit" class="btn">Learn from feedback now</button>
    </form>
    <p class="muted" style="margin: 0.9rem 0 0; font-size: 0.8rem;">
      Preview shows what would happen without changing anything. Exchanges
      mentioning compensation, discipline, litigation, patient detail, or a
      safety escalation are always held back for you to review by hand.
    </p>
    {% if learning_runs %}
    <table style="margin-top: 1.3rem; font-size: 0.84rem;">
      <tr><th>When</th><th>Trigger</th><th>Learned</th><th>Held back</th></tr>
      {% for run in learning_runs %}
      <tr>
        <td>{{ run.when }}</td>
        <td class="muted">{{ run.trigger }}</td>
        <td><strong>{{ run.approved }}</strong></td>
        <td class="muted">{{ run.skipped }}</td>
      </tr>
      {% endfor %}
    </table>
    {% endif %}

    <div style="margin-top: 1rem; display: flex; gap: 0.8rem; align-items: center;
                flex-wrap: wrap; font-size: 0.8rem;">
      {% if learning_runs %}
      <form method="POST" action="/admin/learning/archive" style="display: inline;">
        <button type="submit" class="btn"
                style="background: transparent; color: var(--navy);
                       font-size: 0.7rem;">Archive this list</button>
      </form>
      {% endif %}
      <a href="/admin/learning/archive" class="muted"
         style="text-decoration: none; border-bottom: 1px solid var(--gold);">
        View archive{% if archived_runs %} ({{ archived_runs }}){% endif %}
      </a>
    </div>
  </div>

  <h2 class="group-heading">Pre-Call Briefings</h2>

  <div class="section">
    <h2>Briefings</h2>
    <p class="muted" style="margin: 0 0 1rem 0;">
      When a participant books time, a short brief on what they've been working
      through is prepared for the advisor.
      {% if mail_ready %}
      It's emailed as well as listed here.
      {% else %}
      <strong>No email is configured</strong>, so briefs appear here only — set
      <code>POSTMARK_SERVER_TOKEN</code> or <code>SMTP_HOST</code> to have them
      sent automatically.
      {% endif %}
    </p>
    {% if briefings %}
    <table style="font-size: 0.85rem;">
      <tr>
        <th style="width: 15%;">When</th><th style="width: 18%;">Advisor</th>
        <th style="width: 27%;">Participant</th><th>Brief</th>
        <th style="width: 9%;">Emailed</th>
      </tr>
      {% for b in briefings %}
      <tr>
        <td class="muted" style="white-space: nowrap;">{{ b.when }}</td>
        <td>{{ b.advisor }}</td>
        <td class="muted">{{ b.participant }}</td>
        <td>
          <details>
            <summary style="cursor: pointer; color: var(--muted);
                            font-size: 0.78rem;">Read</summary>
            <pre style="white-space: pre-wrap; font-size: 0.76rem;
                        background: var(--paper); padding: 0.7rem;
                        border-radius: 3px; margin: 0.5rem 0 0;
                        font-family: inherit;">{{ b.summary }}</pre>
          </details>
        </td>
        <td class="muted">{{ '✓' if b.emailed else '—' }}</td>
      </tr>
      {% endfor %}
    </table>
    {% else %}
    <p class="muted">No briefings yet. One is prepared each time a participant
      continues to scheduling.</p>
    {% endif %}
  </div>

  <h2 class="group-heading">Logs</h2>

  <div class="section">
    <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 1rem; padding-bottom: 0.6rem; border-bottom: 1px solid var(--line); flex-wrap: wrap; gap: 0.5rem;">
      <h2 style="margin: 0; border: none; padding: 0;">Conversation Log</h2>
      <div style="display: flex; gap: 0.5rem; align-items: center; flex-wrap: wrap;">
        <form method="GET" action="/admin" style="display: inline-flex; gap: 0.4rem; align-items: center; margin: 0;">
          <label for="filter-select" style="font-size: 0.78rem; color: var(--muted);">Show:</label>
          <select id="filter-select" name="filter" onchange="this.form.submit()"
                  style="padding: 0.35rem 0.55rem; border: 1px solid var(--line); border-radius: 2px; font-family: inherit; font-size: 0.8rem; background: white; cursor: pointer;">
            <option value="all"     {% if log_filter == 'all'     %}selected{% endif %}>All conversations</option>
            <option value="rated"   {% if log_filter == 'rated'   %}selected{% endif %}>Rated only</option>
            <option value="unrated" {% if log_filter == 'unrated' %}selected{% endif %}>Unrated only</option>
            <option value="up"      {% if log_filter == 'up'      %}selected{% endif %}>Thumbs up only</option>
            <option value="down"    {% if log_filter == 'down'    %}selected{% endif %}>Thumbs down only</option>
          </select>
        </form>
        {% if stats.total > 0 %}
        <button type="button" id="export-csv" class="btn"
                style="padding: 0.4rem 0.85rem; font-size: 0.7rem;">↓ CSV</button>
        <button type="button" id="export-xlsx" class="btn"
                style="padding: 0.4rem 0.85rem; font-size: 0.7rem;">↓ Excel</button>
        {% endif %}
      </div>
    </div>
    <p class="muted" style="font-size: 0.82rem; margin: -0.3rem 0 1rem 0;">
      Every chat exchange is logged automatically. Ratings and comments are added when a user clicks thumbs up or down.
      Currently showing {{ feedback_rows|length }} record{{ 's' if feedback_rows|length != 1 else '' }}.
      Tick rows to export or delete just those; with nothing ticked, export includes every record.
    </p>
    {% if feedback_rows %}
    <form method="POST" action="/admin/feedback/delete-selected"
          id="feedback-form"
          onsubmit="const c = document.querySelectorAll('input[name=feedback_ids]:checked').length;
                   if (c === 0) { alert('Select at least one row.'); return false; }
                   return confirm('Delete ' + c + ' selected feedback row(s)? This cannot be undone.');">
      <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 0.8rem; flex-wrap: wrap; gap: 0.5rem;">
        <div style="display: flex; gap: 0.5rem; align-items: center;">
          <button type="submit" class="btn btn-danger" style="padding: 0.45rem 0.95rem; font-size: 0.7rem;">
            Delete selected
          </button>
          <label style="font-size: 0.78rem; color: var(--muted); cursor: pointer;">
            <input type="checkbox" id="select-all" style="margin-right: 0.3rem; vertical-align: middle;" />
            Select all
          </label>
        </div>
      </div>
      <table>
        <tr>
          <th style="width: 28px;"></th>
          <th>When</th><th>Rating</th><th>Release</th><th>Location</th><th>User question</th><th>Bot reply</th><th>Attachment</th><th>Comment</th>
          <th style="width: 60px;"></th>
        </tr>
        {% for f in feedback_rows %}
        <tr id="row-{{ f.id }}">
          <td><input type="checkbox" name="feedback_ids" value="{{ f.id }}" class="feedback-checkbox" /></td>
          <td class="muted">{{ f.created_at.strftime('%m/%d %H:%M') }}</td>
          <td style="white-space: nowrap;">
            <span class="rating-state" data-id="{{ f.id }}">
              {% if f.rating == 'up' %}<span class="tag-up">UP</span>
              {% elif f.rating == 'down' %}<span class="tag-down">DOWN</span>
              {% else %}<span class="muted" style="font-size: 0.7rem;">—</span>{% endif %}
            </span>
            <span class="rating-set">
              <button type="button" class="rate-btn{% if f.rating == 'up' %} on-up{% endif %}"
                      data-id="{{ f.id }}" data-rating="up" title="Mark helpful">&#128077;</button>
              <button type="button" class="rate-btn{% if f.rating == 'down' %} on-down{% endif %}"
                      data-id="{{ f.id }}" data-rating="down" title="Mark not helpful">&#128078;</button>
              <button type="button" class="rate-btn" data-id="{{ f.id }}" data-rating=""
                      title="Clear rating">&times;</button>
            </span>
            {% if f.approved_for_learning %}<br /><span class="tag-lesson">LESSON</span>{% endif %}
          </td>
          <td style="font-size: 0.78rem; white-space: nowrap;">
            {% if acks.get(f.id) %}
              <span class="tag-ack" title="Release accepted {{ acks[f.id] }}">&#10003;</span>
              <span class="muted" style="font-size: 0.7rem;">{{ acks[f.id] }}</span>
            {% else %}<span class="muted" style="font-size: 0.7rem;">—</span>{% endif %}
          </td>
          <td class="muted" style="font-size: 0.78rem; max-width: 150px;">
            {% if locations.get(f.id) %}{{ locations[f.id] }}{% else %}—{% endif %}
          </td>
          <td class="truncate" title="{{ f.user_message }}">{{ f.user_message }}</td>
          <td class="truncate" title="{{ f.bot_reply }}">{{ f.bot_reply }}</td>
          <td class="truncate" title="{{ f.attachment_info or '' }}" style="max-width: 160px; font-size: 0.78rem;">
            {% if f.attachment_info %}📎 {{ f.attachment_info }}{% else %}<span class="muted">—</span>{% endif %}
          </td>
          <td class="truncate" title="{{ f.comment or '' }}" style="max-width: 240px;">
            {% if f.comment %}<strong>{{ f.comment }}</strong>{% else %}<span class="muted">—</span>{% endif %}
          </td>
          <td>
            <button type="button" class="expand-btn" data-target="detail-{{ f.id }}">
              View
            </button>
          </td>
        </tr>
        <tr id="detail-{{ f.id }}" class="feedback-detail" style="display: none;">
          <td colspan="10">
            <div class="feedback-detail-meta">
              Log ID #{{ f.id }} · {{ f.created_at.strftime('%A, %B %d %Y at %I:%M %p') }}
              · Rating: <strong>
                {% if f.rating == 'up' %}Helpful 👍
                {% elif f.rating == 'down' %}Not helpful 👎
                {% else %}Unrated{% endif %}
              </strong>
              {% if f.persona %}· Persona: {{ f.persona }}{% endif %}
              {% if locations.get(f.id) %}· Location: {{ locations[f.id] }}{% endif %}
              {% if acks.get(f.id) %}· Release accepted: {{ acks[f.id] }}{% else %}· Release: not recorded{% endif %}
              {% if f.attachment_info %}· Attachment: {{ f.attachment_info }}{% endif %}
            </div>

            <div class="feedback-detail-block" style="margin-top: 0.9rem;">
              <div class="feedback-detail-label">User question</div>
              <div class="feedback-detail-content">{{ f.user_message or '(empty)' }}</div>
            </div>

            <div class="feedback-detail-block">
              <div class="feedback-detail-label">Bot reply</div>
              <div class="feedback-detail-content">{{ f.bot_reply or '(empty)' }}</div>
            </div>

            <div class="feedback-detail-block">
              <div class="feedback-detail-label">User comment</div>
              {% if f.comment %}
                <div class="feedback-detail-content comment-highlight">{{ f.comment }}</div>
              {% else %}
                <div class="feedback-detail-content" style="font-style: italic; color: var(--muted);">No comment provided.</div>
              {% endif %}
            </div>

            {% if f.rating == 'down' and f.comment %}
            <div class="feedback-detail-block" style="margin-top: 1.2rem; padding-top: 1rem; border-top: 1px dashed var(--line);">
              <div class="feedback-detail-label">Learning loop</div>
              {% if f.approved_for_learning %}
                <p style="font-size: 0.85rem; margin: 0.4rem 0;">
                  <span class="tag-lesson">ACTIVE LESSON</span>
                  &nbsp;The bot uses this feedback to improve answers to similar questions.
                </p>
                <button type="button" class="lesson-action-btn" data-action="revoke" data-id="{{ f.id }}"
                        style="background: transparent; color: var(--rust); border: 1px solid var(--rust); padding: 0.4rem 0.85rem; font-size: 0.7rem; letter-spacing: 0.14em; text-transform: uppercase; cursor: pointer; border-radius: 2px; font-family: inherit; margin-top: 0.4rem;">
                  Revoke lesson
                </button>
              {% else %}
                <p style="font-size: 0.85rem; margin: 0.4rem 0; color: var(--muted);">
                  Approve this as a lesson and the bot will see it as guidance when answering semantically similar questions.
                </p>
                <button type="button" class="lesson-action-btn" data-action="approve" data-id="{{ f.id }}"
                        style="background: var(--navy); color: var(--gold); border: 1px solid var(--navy); padding: 0.4rem 0.85rem; font-size: 0.7rem; letter-spacing: 0.14em; text-transform: uppercase; cursor: pointer; border-radius: 2px; font-family: inherit; margin-top: 0.4rem;">
                  ✓ Approve as lesson
                </button>
              {% endif %}
            </div>
            {% endif %}
          </td>
        </tr>
        {% endfor %}
      </table>
    </form>



    <script>
      // Set or clear a rating straight from the log
      (function() {
        document.querySelectorAll(".rate-btn").forEach(btn => {
          btn.addEventListener("click", async () => {
            const id = btn.dataset.id;
            const rating = btn.dataset.rating || "";
            const group = btn.closest(".rating-set");
            group.querySelectorAll("button").forEach(b => b.disabled = true);
            try {
              const body = new URLSearchParams();
              body.set("rating", rating);
              const resp = await fetch(`/admin/feedback/${id}/rating`, {
                method: "POST",
                headers: { "Content-Type": "application/x-www-form-urlencoded" },
                body: body.toString(),
              });
              const data = await resp.json();
              if (!data.ok) throw new Error(data.error || "Update failed");

              // Reflect the new state without a page reload
              const state = document.querySelector(`.rating-state[data-id="${id}"]`);
              if (state) {
                state.innerHTML = rating === "up"
                  ? '<span class="tag-up">UP</span>'
                  : rating === "down"
                    ? '<span class="tag-down">DOWN</span>'
                    : '<span class="muted" style="font-size: 0.7rem;">—</span>';
              }
              group.querySelectorAll("button").forEach(b => {
                b.classList.remove("on-up", "on-down");
                if (b.dataset.rating === "up" && rating === "up") b.classList.add("on-up");
                if (b.dataset.rating === "down" && rating === "down") b.classList.add("on-down");
              });
            } catch (err) {
              alert("Couldn't update that rating: " + err.message);
            } finally {
              group.querySelectorAll("button").forEach(b => b.disabled = false);
            }
          });
        });
      })();

      // Toggle expanded feedback detail rows
      (function() {
        document.querySelectorAll('.expand-btn').forEach(btn => {
          btn.addEventListener('click', () => {
            const target = document.getElementById(btn.dataset.target);
            if (!target) return;
            const isOpen = target.style.display !== 'none';
            target.style.display = isOpen ? 'none' : 'table-row';
            btn.textContent = isOpen ? 'View' : 'Close';
          });
        });

        // Lesson approve/revoke buttons — submit via fetch (avoids nested form issues)
        document.querySelectorAll('.lesson-action-btn').forEach(btn => {
          btn.addEventListener('click', async () => {
            const action = btn.dataset.action;
            const id = btn.dataset.id;
            const verb = action === 'approve' ? 'Approve as a lesson?' :
                         'Revoke this lesson? The bot will stop learning from it.';
            if (!confirm(verb)) return;
            btn.disabled = true;
            btn.textContent = action === 'approve' ? 'Approving…' : 'Revoking…';
            try {
              const resp = await fetch(`/admin/feedback/${id}/${action}-lesson`, {
                method: 'POST',
                headers: { 'Content-Type': 'application/x-www-form-urlencoded' },
              });
              // Routes redirect back to admin, so reloading shows updated state
              window.location.href = '/admin';
            } catch (e) {
              alert('Failed: ' + e.message);
              btn.disabled = false;
              btn.textContent = action === 'approve' ? '✓ Approve as lesson' : 'Revoke lesson';
            }
          });
        });
      })();
    </script>

    <div style="margin-top: 1.5rem; padding-top: 1rem; border-top: 1px dashed var(--line);">
      <details>
        <summary style="cursor: pointer; font-size: 0.78rem; color: var(--rust); letter-spacing: 0.1em; text-transform: uppercase;">
          Danger zone — Clear all feedback
        </summary>
        <form method="POST" action="/admin/feedback/delete-all" style="margin-top: 0.8rem;"
              onsubmit="return confirm('Permanently delete ALL feedback rows? This cannot be undone.');">
          <p class="muted" style="margin: 0.5rem 0;">
            This permanently deletes every feedback row in the database.
            Type <strong>YES</strong> to confirm.
          </p>
          <div style="display: flex; gap: 0.5rem; align-items: center;">
            <input type="text" name="confirm" placeholder="Type YES to confirm"
                   style="flex: 0 1 200px; padding: 0.5rem; border: 1px solid var(--line); border-radius: 2px;" />
            <button type="submit" class="btn btn-danger" style="padding: 0.5rem 1rem; font-size: 0.7rem;">
              Clear all feedback
            </button>
          </div>
        </form>
      </details>
    </div>

    <script>
      // Select-all + export of just the checked rows
      (function() {
        const selectAll = document.getElementById("select-all");
        const checkboxes = Array.from(document.querySelectorAll(".feedback-checkbox"));
        const csvBtn = document.getElementById("export-csv");
        const xlsxBtn = document.getElementById("export-xlsx");

        function selectedIds() {
          return checkboxes.filter(cb => cb.checked).map(cb => cb.value);
        }

        // Buttons show what they will actually export
        function refreshLabels() {
          const n = selectedIds().length;
          if (csvBtn)  csvBtn.textContent  = n ? `\u2193 CSV (${n})`   : "\u2193 CSV";
          if (xlsxBtn) xlsxBtn.textContent = n ? `\u2193 Excel (${n})` : "\u2193 Excel";
          const title = n
            ? `Export the ${n} selected row${n === 1 ? "" : "s"}`
            : "Export all rows — tick rows to export only those";
          if (csvBtn)  csvBtn.title = title;
          if (xlsxBtn) xlsxBtn.title = title;
          if (selectAll) {
            selectAll.checked = n > 0 && n === checkboxes.length;
            selectAll.indeterminate = n > 0 && n < checkboxes.length;
          }
        }

        if (selectAll) {
          selectAll.addEventListener("change", () => {
            checkboxes.forEach(cb => cb.checked = selectAll.checked);
            refreshLabels();
          });
        }
        checkboxes.forEach(cb => cb.addEventListener("change", refreshLabels));

        function doExport(ext) {
          const ids = selectedIds();
          let url = `/admin/export/feedback.${ext}`;
          if (ids.length) url += `?ids=${encodeURIComponent(ids.join(","))}`;
          window.location.href = url;
        }
        if (csvBtn)  csvBtn.addEventListener("click", () => doExport("csv"));
        if (xlsxBtn) xlsxBtn.addEventListener("click", () => doExport("xlsx"));

        refreshLabels();
      })();
    </script>
    {% else %}
    <p class="muted">No feedback yet.</p>
    {% endif %}
  </div>

  {% if rag_ready %}
  <h2 class="group-heading">Knowledge Base</h2>

  <div class="section">
    <h2>Documents</h2>
    <p class="muted" style="margin: 0 0 1rem 0;">
      {{ docs|length }} document{{ 's' if docs|length != 1 else '' }} embedded and
      available to the advisor. Deleting one removes its chunks and the advisor
      stops drawing on it.
    </p>
    {% if docs %}
    <table class="kb-table">
      <tr>
        <th>Title</th><th>Base</th><th>Source</th>
        <th style="text-align: right;">Chunks</th>
        <th>Uploaded</th><th></th>
      </tr>
      {% for d in docs %}
      <tr>
        <td class="kb-title">{{ d.title }}</td>
        <td class="muted" style="font-size: 0.76rem;">
          {% if owners.get(d.title) %}{{ advisor_names.get(owners[d.title], owners[d.title]) }}{% else %}Shared{% endif %}
        </td>
        <td class="muted kb-source" title="{{ d.source or '' }}">{{ d.source or '—' }}</td>
        <td style="text-align: right;">{{ d.chunk_count }}</td>
        <td class="muted kb-date">{{ d.uploaded_at.strftime('%Y-%m-%d %H:%M') }}</td>
        <td style="text-align: right;">
          <form method="POST" action="/admin/delete/{{ d.id }}" style="display:inline;"
                data-doc-title="{{ d.title }}">
            <button type="submit" class="btn btn-danger">Delete</button>
          </form>
        </td>
      </tr>
      {% endfor %}
    </table>
    {% else %}
    <p class="muted">No documents yet. Upload your first one below.</p>
    {% endif %}
  </div>

  <h2 class="group-heading">Knowledge Upload</h2>

  <div class="section">
    <h2>Upload Document</h2>
    <p class="muted" style="margin: 0 0 1rem 0;">Accepts PDF, Word, Excel, PowerPoint, CSV, TXT, MD, RTF. Up to {{ cfg.max_upload_mb }} MB. The document will be chunked and embedded automatically.</p>
    <form method="POST" action="/admin/upload" enctype="multipart/form-data" class="upload">
      <input type="file" name="file" accept=".pdf,.docx,.xlsx,.xlsm,.pptx,.csv,.tsv,.txt,.md,.rtf" required />
      <input type="text" name="title" placeholder="Document title (optional)" />
            <select name="owner" title="Which knowledge base this belongs to">
        <option value="">Shared J3P base — all advisors</option>
        {% for adv in advisors %}
        <option value="{{ adv.slug }}">Only {{ adv.name }}</option>
        {% endfor %}
      </select>
      <button type="submit" class="btn">Upload & Embed</button>
    </form>
  </div>

  <div class="section">
    <h2>Upload Folder</h2>
    <p class="muted" style="margin: 0 0 1rem 0;">
      Select an entire folder. All supported files inside it (including subfolders) will be uploaded and embedded in one batch.
      Unsupported files and duplicates are skipped automatically. Maximum 50 files per batch.
    </p>
    <form method="POST" action="/admin/upload-folder" enctype="multipart/form-data" class="upload" id="folder-upload-form">
      <input type="file" name="files" id="folder-input" webkitdirectory directory multiple required />
      <input type="text" name="folder_title" placeholder="Folder title (optional)" />
            <select name="owner" title="Which knowledge base this belongs to">
        <option value="">Shared J3P base — all advisors</option>
        {% for adv in advisors %}
        <option value="{{ adv.slug }}">Only {{ adv.name }}</option>
        {% endfor %}
      </select>
      <button type="submit" class="btn" id="folder-upload-btn">Upload Folder</button>
    </form>
    <p id="folder-preview" class="muted" style="margin: 0.75rem 0 0 0; font-size: 0.85rem; display: none;"></p>
    <script>
      (function() {
        const folderInput = document.getElementById("folder-input");
        const preview = document.getElementById("folder-preview");
        const btn = document.getElementById("folder-upload-btn");
        const form = document.getElementById("folder-upload-form");
        const SUPPORTED = /\\.(pdf|docx|txt|md)$/i;

        folderInput.addEventListener("change", () => {
          const all = Array.from(folderInput.files || []);
          const supported = all.filter(f => SUPPORTED.test(f.name));
          const skipped = all.length - supported.length;
          if (all.length === 0) {
            preview.style.display = "none";
            return;
          }
          let msg = `${supported.length} supported file${supported.length === 1 ? '' : 's'} ready to upload`;
          if (skipped > 0) msg += ` · ${skipped} unsupported file${skipped === 1 ? '' : 's'} will be skipped`;
          if (supported.length > 50) {
            msg += ` · ⚠ Only the first 50 will be processed`;
          }
          if (supported.length === 0) {
            msg = "⚠ No supported files found in this folder (PDF, DOCX, TXT, MD only).";
            btn.disabled = true;
          } else {
            btn.disabled = false;
          }
          preview.textContent = msg;
          preview.style.display = "block";
        });

        form.addEventListener("submit", (e) => {
          const all = Array.from(folderInput.files || []);
          const supported = all.filter(f => SUPPORTED.test(f.name));
          if (supported.length === 0) {
            e.preventDefault();
            alert("No supported files found in this folder.");
            return;
          }
          btn.textContent = "Uploading… (this may take a while)";
          btn.disabled = true;
        });
      })();
    </script>
  </div>

  <div class="section">
    <h2>Add Knowledge from URL</h2>
    <p class="muted" style="margin: 0 0 1rem 0;">Paste a link to an article, blog post, or web page. The main article text will be extracted and embedded. Works best with article-style pages (not paywalled, login-required, or JavaScript-only sites).</p>
    <form method="POST" action="/admin/upload-url" class="upload">
      <input type="url" name="url" placeholder="https://example.com/article" required style="flex: 1.5; min-width: 280px; padding: 0.5rem; border: 1px solid var(--line); border-radius: 2px; font-family: inherit;" />
      <input type="text" name="url_title" placeholder="Title (optional, auto-detected)" />
            <select name="owner" title="Which knowledge base this belongs to">
        <option value="">Shared J3P base — all advisors</option>
        {% for adv in advisors %}
        <option value="{{ adv.slug }}">Only {{ adv.name }}</option>
        {% endfor %}
      </select>
      <button type="submit" class="btn">Fetch & Embed</button>
    </form>
  </div>

  <div class="section">
    <h2>Add Knowledge from Text</h2>
    <p class="muted" style="margin: 0 0 1rem 0;">
      Paste content directly — a podcast transcript, show notes, an email, a
      passage from a book. Use this when a page won't give up its text, which
      is normal for podcast players, video sites and most social platforms.
    </p>
    <form method="POST" action="/admin/upload-text" class="upload-stacked">
      <input type="text" name="text_title" placeholder="Title (optional)"
             style="width: 100%; padding: 0.5rem; margin-bottom: 0.6rem;
                    border: 1px solid var(--line); border-radius: 2px;
                    font-family: inherit;" />
      <textarea name="text_body" required
                placeholder="Paste the transcript or text here…"
                style="width: 100%; min-height: 150px; padding: 0.6rem;
                       border: 1px solid var(--line); border-radius: 2px;
                       font-family: inherit; font-size: 0.9rem;
                       line-height: 1.5; resize: vertical;"></textarea>
            <select name="owner" title="Which knowledge base this belongs to">
        <option value="">Shared J3P base — all advisors</option>
        {% for adv in advisors %}
        <option value="{{ adv.slug }}">Only {{ adv.name }}</option>
        {% endfor %}
      </select>
      <button type="submit" class="btn" style="margin-top: 0.6rem;">Add &amp; Embed</button>
    </form>
  </div>
  {% endif %}

</div>

    <script>
      // ---------------------------------------------------------------
      // Deletion gate
      // ---------------------------------------------------------------
      // Every destructive form routes through one modal. Browser confirm()
      // is a single reflexive click; this states exactly what will be lost,
      // and for bulk or irreversible actions requires typing DELETE.
      (function() {
        const overlay  = document.getElementById("confirm-overlay");
        const msgEl    = document.getElementById("confirm-message");
        const detailEl = document.getElementById("confirm-detail");
        const typeWrap = document.getElementById("confirm-type-wrap");
        const typeIn   = document.getElementById("confirm-type");
        const goBtn    = document.getElementById("confirm-go");
        const cancelBtn= document.getElementById("confirm-cancel");
        if (!overlay) return;

        let pendingForm = null;

        function close() {
          overlay.hidden = true;
          pendingForm = null;
          typeIn.value = "";
          typeWrap.hidden = true;
          detailEl.hidden = true;
          detailEl.textContent = "";
        }

        function refreshGo() {
          goBtn.disabled = !typeWrap.hidden && typeIn.value.trim().toUpperCase() !== "DELETE";
        }

        function open(form, opts) {
          pendingForm = form;
          noSubmit = !!opts.noSubmit;
          msgEl.textContent = opts.message;
          if (opts.detail) {
            detailEl.textContent = opts.detail;
            detailEl.hidden = false;
          }
          typeWrap.hidden = !opts.requireTyping;
          goBtn.textContent = opts.buttonLabel || "Delete";
          overlay.hidden = false;
          refreshGo();
          setTimeout(() => (opts.requireTyping ? typeIn : goBtn).focus(), 60);
        }

        typeIn.addEventListener("input", refreshGo);
        typeIn.addEventListener("keydown", e => {
          if (e.key === "Enter" && !goBtn.disabled) { e.preventDefault(); goBtn.click(); }
        });
        cancelBtn.addEventListener("click", close);
        overlay.addEventListener("click", e => { if (e.target === overlay) close(); });
        document.addEventListener("keydown", e => {
          if (e.key === "Escape" && !overlay.hidden) close();
        });

        let noSubmit = false;
        goBtn.addEventListener("click", () => {
          if (goBtn.disabled) return;
          const form = pendingForm;
          const informationalOnly = noSubmit;
          close();
          if (form && !informationalOnly) {
            form.dataset.confirmed = "1";
            form.submit();
          }
        });

        // Delegated so it covers forms that appear later in the document
        // than this script — the clear-all form does, and per-form listeners
        // silently missed it.
        document.addEventListener("submit", e => {
          const form = e.target;
          if (!form || form.tagName !== "FORM") return;
          const action = form.getAttribute("action") || "";
          const isDelete = action.indexOf("/delete") !== -1;
          const isRevoke = action.indexOf("revoke-lesson") !== -1;
          if (!isDelete && !isRevoke) return;

          if (form.dataset.confirmed === "1") {
            form.dataset.confirmed = "";
            return;                        // already approved — let it through
          }
          e.preventDefault();
          form.removeAttribute("onsubmit");   // supersede any inline confirm()

          if (action.indexOf("delete-all") !== -1) {
            open(form, {
              message: "This permanently deletes EVERY conversation log entry, "
                     + "including all ratings, comments and approved lessons. "
                     + "It cannot be undone and there is no backup.",
              requireTyping: true,
              buttonLabel: "Delete everything",
            });
            return;
          }

          if (action.indexOf("delete-selected") !== -1) {
            const boxes = document.querySelectorAll(".feedback-checkbox:checked");
            if (boxes.length === 0) {
              open(form, { message: "No rows are ticked, so nothing would be deleted.",
                           buttonLabel: "OK", noSubmit: true });
              return;
            }
            const NL = String.fromCharCode(10);
            const preview = Array.from(boxes).slice(0, 8).map(b => {
              const row = b.closest("tr");
              const q = row ? row.querySelector(".truncate") : null;
              return "\u2022 " + (q ? q.textContent.trim().slice(0, 70) : "row " + b.value);
            }).join(NL);
            open(form, {
              message: "Delete " + boxes.length + " conversation log entr"
                     + (boxes.length === 1 ? "y" : "ies")
                     + "? This cannot be undone.",
              detail: preview + (boxes.length > 8
                      ? NL + "\u2026and " + (boxes.length - 8) + " more" : ""),
              requireTyping: boxes.length >= 10,
              buttonLabel: "Delete " + boxes.length,
            });
            return;
          }

          if (isRevoke) {
            open(form, {
              message: "Stop using this exchange as a lesson? The advisor will no "
                     + "longer learn from it, though the log entry itself is kept.",
              buttonLabel: "Remove lesson",
            });
            return;
          }

          if (action.indexOf("/admin/advisors/delete/") !== -1) {
            open(form, {
              message: "Remove this advisor profile? Their dedicated links stop "
                     + "working immediately and their photo is deleted. The "
                     + "knowledge base and conversation log are not affected.",
              detail: form.dataset.docTitle || "",
              buttonLabel: "Remove advisor",
            });
            return;
          }

          open(form, {
            message: "Delete this document and every embedded chunk from the "
                   + "knowledge base? The advisor will stop drawing on it. "
                   + "You would need to upload the file again to restore it.",
            detail: form.dataset.docTitle || "",
            buttonLabel: "Delete document",
          });
        }, true);
      })();
    </script>
</body></html>"""


@app.route("/admin/login", methods=["GET", "POST"])
def admin_login():
    if not CONFIG["admin_password"]:
        return ("Admin disabled. Set ADMIN_PASSWORD environment variable.", 503)
    if request.method == "POST":
        if request.form.get("password") == CONFIG["admin_password"]:
            session["is_admin"] = True
            return redirect(url_for("admin_dashboard"))
        return render_template_string(ADMIN_LOGIN_HTML, cfg=CONFIG, error="Incorrect password")
    if session.get("is_admin"):
        return redirect(url_for("admin_dashboard"))
    return render_template_string(ADMIN_LOGIN_HTML, cfg=CONFIG, error=None)


@app.route("/admin/logout")
def admin_logout():
    session.pop("is_admin", None)
    return redirect(url_for("admin_login"))


@app.route("/admin")
@admin_required
def admin_dashboard():
    db_ok = db.is_enabled()
    emb_ok = emb.is_enabled()
    rag_ready = db_ok and emb_ok
    docs = db.list_documents() if db_ok else []
    # Filter for the conversation log — default shows everything.
    # Accepted values: all | rated | up | down | unrated
    log_filter = (request.args.get("filter") or "all").lower()
    if log_filter not in ("all", "rated", "up", "down", "unrated"):
        log_filter = "all"
    feedback_rows = db.list_feedback(
        limit=100,
        rating=(None if log_filter == "all" else log_filter),
    ) if db_ok else []
    stats = db.feedback_stats() if db_ok else {"up": 0, "down": 0, "total": 0}
    return render_template_string(
        ADMIN_HTML, cfg=CONFIG, docs=docs, feedback_rows=feedback_rows,
        settings=load_settings(force=True),
        mail_ready=mail_transport_configured(),
        avatar_custom=bool(load_avatar()),
        advisors=list_advisors(),
        owners=document_owners(),
        advisor_names={a["slug"]: a["name"] for a in list_advisors()},
        avatar_version=int(datetime.now().timestamp()),
        avatar_max_mb=AVATAR_MAX_BYTES // 1048576,
        learning_runs=recent_learning_runs(),
        briefings=list_briefings(),
        archived_runs=archived_run_count(),
        learning_interval=LEARNING_INTERVAL_HOURS,
        app_version=APP_VERSION,
        locations=locations_for([r.get("id") for r in feedback_rows]),
        acks=acknowledgements_for([r.get("id") for r in feedback_rows]),
        base_url=(paywall.PUBLIC_BASE_URL or request.host_url.rstrip("/")),
        stats=stats, rag_ready=rag_ready, db_ok=db_ok, emb_ok=emb_ok,
        log_filter=log_filter,
    )


@app.route("/admin/feedback/<int:feedback_id>/rating", methods=["POST"])
@admin_required
def admin_set_rating(feedback_id):
    """Set or clear a rating from the conversation log."""
    value = (request.form.get("rating") or "").strip().lower()
    rating = value if value in ("up", "down") else None
    if set_feedback_rating(feedback_id, rating):
        return jsonify({"ok": True, "rating": rating})
    return jsonify({"ok": False, "error": "Could not update that rating."}), 400


@app.route("/admin/learning/run", methods=["POST"])
@admin_required
def admin_run_learning():
    """Process pending feedback into lessons now."""
    dry = bool(request.form.get("preview"))
    result = run_learning_cycle(trigger="manual preview" if dry else "manual", dry_run=dry)
    if not result.get("ok"):
        flash(f"Learning run failed: {result.get('error')}")
    elif dry:
        flash(f"Preview — {result['approved']} would be learned from, "
              f"{result['skipped']} held back. Nothing changed.")
    else:
        parts = []
        if result.get("up"):
            parts.append(f"{result['up']} that worked well")
        if result.get("down"):
            parts.append(f"{result['down']} to avoid repeating")
        breakdown = (" (" + ", ".join(parts) + ")") if parts else ""
        flash(f"\u2713 Learned from {result['approved']} exchange"
              f"{'s' if result['approved'] != 1 else ''}{breakdown}; "
              f"{result['skipped']} held back for review.")
    return redirect(url_for("admin_dashboard"))


@app.route("/admin/learning/archive", methods=["POST"])
@admin_required
def admin_archive_learning():
    """Move the visible run history into the archive. Nothing is deleted."""
    if archive_all_learning_runs():
        flash("Run history archived — still available under View archive.")
    else:
        flash("Could not archive the run history.")
    return redirect(url_for("admin_dashboard"))


@app.route("/admin/learning/archive", methods=["GET"])
@admin_required
def admin_view_learning_archive():
    """Every archived run, newest first."""
    runs = recent_learning_runs(limit=500, archived=True)
    return render_template_string(LEARNING_ARCHIVE_HTML, cfg=CONFIG, runs=runs)


@app.route("/admin/advisors", methods=["POST"])
@admin_required
def admin_save_advisor():
    """Create or update an advisor profile."""
    name = (request.form.get("name") or "").strip()[:80]
    if not name:
        flash("Give the advisor a name.")
        return redirect(url_for("admin_dashboard"))

    slug = (request.form.get("slug") or "").strip().lower()
    slug = slugify_advisor(slug or name)

    photo = mime = None
    file = request.files.get("photo")
    if file and file.filename:
        raw = file.read()
        if len(raw) > AVATAR_MAX_BYTES:
            flash(f"That photo is {len(raw)/1048576:.1f} MB — the limit is "
                  f"{AVATAR_MAX_BYTES // 1048576} MB.")
            return redirect(url_for("admin_dashboard"))
        ext = (file.filename.rsplit(".", 1)[-1] or "").lower()
        if ext not in ("jpg", "jpeg", "png", "webp", "gif"):
            flash("Use a JPG, PNG, WEBP or GIF photo.")
            return redirect(url_for("admin_dashboard"))
        photo, mime = prepare_avatar(raw)

    if save_advisor(slug, name, photo, mime):
        flash(f"✓ {name} saved — links are listed below.")
    else:
        flash("Could not save the advisor — check the database connection.")
    return redirect(url_for("admin_dashboard"))


@app.route("/admin/advisors/delete/<slug>", methods=["POST"])
@admin_required
def admin_delete_advisor(slug):
    """Remove an advisor profile and its photo."""
    if delete_advisor(slug):
        flash(f"Advisor “{slug}” removed. Their links no longer work.")
    else:
        flash("Could not remove that advisor.")
    return redirect(url_for("admin_dashboard"))


@app.route("/admin/settings", methods=["POST"])
@admin_required
def admin_settings():
    """Persist the display toggles from the admin panel."""
    # Settings now live in more than one form, so each form declares which
    # keys it owns. Without this, saving one form would switch off every
    # toggle it didn't happen to include.
    managed = [f.strip() for f in (request.form.get("_fields") or "").split(",")
               if f.strip()]
    if not managed:
        managed = list(_SETTINGS_DEFAULTS.keys())

    labels = {
        "show_scheduling_button": ("Scheduling button", "shown", "hidden"),
        "show_avatar": ("Advisor avatar", "shown", "hidden"),
        "allow_materials": ("Participant materials", "on", "off"),
        "auto_learning": ("Continuous learning", "on", "off"),
        "require_login": ("Sign-in", "required", "not required"),
    }
    messages = []

    # Text settings are saved separately; the loop below handles booleans.
    if "avatar_name" in managed:
        nm = (request.form.get("avatar_name") or "").strip()[:60]
        if save_setting("avatar_name", nm):
            messages.append(f"Photo name set to \u201c{nm}\u201d" if nm
                            else "Photo name cleared")
        else:
            messages.append("Could not save the photo name")

    for key in managed:
        if key not in _SETTINGS_DEFAULTS:
            continue
        if not isinstance(_SETTINGS_DEFAULTS[key], bool):
            continue
        value = bool(request.form.get(key))
        name, on_word, off_word = labels.get(key, (key, "on", "off"))

        # Enabling sign-in without email would lock every participant out
        if key == "require_login" and value and not mail_transport_configured():
            messages.append("Sign-in NOT enabled — no email transport is configured")
            continue

        if save_setting(key, value):
            messages.append(f"{name} {on_word if value else off_word}")
        else:
            messages.append(f"Could not save {name.lower()}")

    flash((" · ".join(messages) + ".") if messages else "Nothing to save.")
    return redirect(url_for("admin_dashboard"))


@app.route("/admin/upload", methods=["POST"])
@admin_required
def admin_upload():
    if not (db.is_enabled() and emb.is_enabled()):
        flash("Cannot upload: RAG not fully configured.")
        return redirect(url_for("admin_dashboard"))

    file = request.files.get("file")
    if not file or not file.filename:
        flash("No file selected.")
        return redirect(url_for("admin_dashboard"))

    title = (request.form.get("title") or "").strip() or file.filename

    # Duplicate check BEFORE expensive embedding work.
    # We compare both title and filename (source) against existing docs.
    duplicate = db.find_duplicate_document(title=title, source=file.filename)
    if duplicate:
        flash(
            f"⚠ Duplicate detected — '{duplicate['title']}' was already uploaded "
            f"on {duplicate['uploaded_at'].strftime('%Y-%m-%d %H:%M')} "
            f"({duplicate['chunk_count']} chunks). Delete the existing entry first "
            f"if you want to replace it."
        )
        return redirect(url_for("admin_dashboard"))

    try:
        file_bytes = file.read()
        text = extract_attachment_text(file.filename, file_bytes)
        if not text.strip():
            flash(f"No text could be extracted from {file.filename}.")
            return redirect(url_for("admin_dashboard"))

        chunks = emb.chunk_text(text)
        if not chunks:
            flash("Document produced no chunks (too short or empty).")
            return redirect(url_for("admin_dashboard"))

        # Embed all chunks in batch
        vectors = emb.embed_batch(chunks)
        pairs = list(zip(chunks, vectors))
        doc_id = db.insert_document(title, file.filename, pairs)
        set_document_owner(title, (request.form.get("owner") or "").strip())

        flash(f"✓ Uploaded '{title}' — {len(chunks)} chunks embedded (doc #{doc_id}).")
    except Exception as e:
        app.logger.error(f"Upload failed: {e}")
        flash(f"Upload failed: {str(e)[:200]}")

    return redirect(url_for("admin_dashboard"))


@app.route("/admin/upload-folder", methods=["POST"])
@admin_required
def admin_upload_folder():
    """Bulk-upload all supported files from a folder in one request.

    Each file goes through the same pipeline as single-file upload:
    duplicate check → extract → chunk → embed → insert. Files that
    fail any step are logged but do not stop the batch. Returns a
    summary of successes, skipped files (duplicates or unsupported),
    and failures.
    """
    if not (db.is_enabled() and emb.is_enabled()):
        flash("Cannot upload: RAG not fully configured.")
        return redirect(url_for("admin_dashboard"))

    files = request.files.getlist("files")
    if not files:
        flash("No files were selected.")
        return redirect(url_for("admin_dashboard"))

    SUPPORTED_EXT = ('.pdf', '.docx', '.txt', '.md')
    MAX_BATCH = 50   # keep request under Railway 5-min timeout for typical files
    MAX_BYTES = MAX_UPLOAD_BYTES

    # Filter to supported extensions first
    supported_files = [f for f in files if f.filename and f.filename.lower().endswith(SUPPORTED_EXT)]
    unsupported_count = len(files) - len(supported_files)

    # Cap the batch size — anything above the limit is silently skipped for this run
    over_limit = max(0, len(supported_files) - MAX_BATCH)
    process_files = supported_files[:MAX_BATCH]

    # Optional label so a batch is recognisable in the knowledge base
    folder_title = (request.form.get("folder_title") or "").strip()[:80]
    folder_owner = (request.form.get("owner") or "").strip()

    uploaded = []       # list of (title, chunk_count, doc_id)
    duplicates = []     # list of (filename, existing_title)
    failed = []         # list of (filename, error)

    for file in process_files:
        # webkitdirectory paths look like "folder/subfolder/file.pdf"
        # Use just the basename as the source and title default.
        full_path = file.filename or "unknown"
        basename = full_path.rsplit("/", 1)[-1]
        display_title = f"{folder_title} — {basename}" if folder_title else basename

        try:
            # Duplicate check first
            dup = db.find_duplicate_document(title=basename, source=basename)
            if dup:
                duplicates.append((basename, dup["title"]))
                continue

            file_bytes = file.read()
            if len(file_bytes) > MAX_BYTES:
                failed.append((basename, f"too large (>{MAX_UPLOAD_MB} MB)"))
                continue
            if not file_bytes:
                failed.append((basename, "empty file"))
                continue

            text = extract_attachment_text(basename, file_bytes)
            if not text.strip():
                failed.append((basename, "no text extracted"))
                continue

            chunks = emb.chunk_text(text)
            if not chunks:
                failed.append((basename, "too short to chunk"))
                continue

            vectors = emb.embed_batch(chunks)
            pairs = list(zip(chunks, vectors))
            # Title carries the folder label when given; source stays the
            # filename so duplicate detection keeps working on re-upload.
            doc_id = db.insert_document(display_title, basename, pairs)
            set_document_owner(display_title, folder_owner)
            uploaded.append((display_title, len(chunks), doc_id))
            app.logger.info(f"Folder upload: {basename} → doc #{doc_id}, {len(chunks)} chunks")

        except Exception as e:
            app.logger.error(f"Folder upload failed for {basename}: {e}")
            failed.append((basename, str(e)[:100]))

    # Build a summary flash message
    parts = []
    if uploaded:
        parts.append(f"✓ {len(uploaded)} file{'s' if len(uploaded) != 1 else ''} uploaded")
    if duplicates:
        parts.append(f"⏭ {len(duplicates)} skipped as duplicates")
    if failed:
        # Show first 3 failure reasons so the user has something to act on
        detail = "; ".join([f"{n} ({e})" for n, e in failed[:3]])
        parts.append(f"✗ {len(failed)} failed — {detail}" + (" …" if len(failed) > 3 else ""))
    if unsupported_count:
        parts.append(f"({unsupported_count} unsupported file type{'s' if unsupported_count != 1 else ''} ignored)")
    if over_limit:
        parts.append(f"⚠ {over_limit} additional file(s) exceeded the 50-file batch limit and were not processed — run again to continue")

    flash(" · ".join(parts) if parts else "No files were processed.")
    return redirect(url_for("admin_dashboard"))


@app.route("/admin/upload-text", methods=["POST"])
@admin_required
def admin_upload_text():
    """Paste content straight in — transcripts, show notes, anything a
    JavaScript page won't hand over."""
    if not (db.is_enabled() and emb.is_enabled()):
        flash("Cannot add text: RAG not fully configured.")
        return redirect(url_for("admin_dashboard"))

    title = (request.form.get("text_title") or "").strip()
    body = (request.form.get("text_body") or "").strip()
    if len(body) < 100:
        flash("Paste a bit more text than that — at least a paragraph.")
        return redirect(url_for("admin_dashboard"))
    if not title:
        title = body.split("\n", 1)[0][:70].strip() or "Pasted text"

    dup = db.find_duplicate_document(title=title)
    if dup:
        flash(f"⚠ '{title}' already exists in the knowledge base. Use a "
              f"different title or delete the existing entry first.")
        return redirect(url_for("admin_dashboard"))

    try:
        chunks = emb.chunk_text(body)
        if not chunks:
            flash("That text produced no chunks.")
            return redirect(url_for("admin_dashboard"))
        vectors = emb.embed_batch(chunks)
        doc_id = db.insert_document(title, "pasted text", list(zip(chunks, vectors)))
        set_document_owner(title, (request.form.get("owner") or "").strip())
        flash(f"✓ Added '{title}' — {len(chunks)} chunks embedded (doc #{doc_id}).")
    except Exception as e:
        app.logger.error(f"[text] ingest failed: {e}")
        flash(f"Could not add that text: {str(e)[:160]}")
    return redirect(url_for("admin_dashboard"))


def fetch_podcast_feed(url: str):
    """Episodes from a podcast RSS feed.

    A podcast's Spotify or Apple page is rendered in the browser and carries
    almost no text, but the underlying RSS feed holds every episode's title and
    full show notes as plain XML. That's the version worth embedding.
    """
    import urllib.request as _url
    import html as _html
    import xml.etree.ElementTree as ET

    req = _url.Request(url, headers={
        "User-Agent": "Mozilla/5.0 (compatible; J3PAdvisor/1.0)",
        "Accept": "application/rss+xml, application/xml, text/xml, */*",
    })
    with _url.urlopen(req, timeout=25) as resp:
        raw = resp.read(8_000_000)

    root = ET.fromstring(raw)
    channel = root.find("channel")
    if channel is None:
        return None, []

    def text_of(parent, *names):
        for name in names:
            el = parent.find(name)
            if el is not None and (el.text or "").strip():
                return el.text.strip()
            # namespaced (itunes:summary and friends)
            for child in parent:
                if child.tag.split("}")[-1] == name.split(":")[-1]:
                    if (child.text or "").strip():
                        return child.text.strip()
        return ""

    def clean(html_text: str) -> str:
        no_tags = re.sub(r"<[^>]+>", " ", html_text or "")
        return re.sub(r"\s+", " ", _html.unescape(no_tags)).strip()

    show_title = text_of(channel, "title") or "Podcast"
    show_desc = clean(text_of(channel, "description", "itunes:summary"))

    episodes = []
    for item in channel.findall("item"):
        ep_title = text_of(item, "title") or "Untitled episode"
        body = clean(text_of(item, "description", "itunes:summary", "content:encoded"))
        when = text_of(item, "pubDate")[:16]
        if len(body) < 40:
            continue
        episodes.append({"title": ep_title, "body": body, "when": when})

    return {"title": show_title, "description": show_desc}, episodes


def looks_like_feed(url: str, probe=True) -> bool:
    """A feed either declares itself in the URL or answers as XML."""
    low = url.lower()
    if any(t in low for t in ("/rss", "rss.xml", "feed.xml", "/feed", ".rss",
                              "feeds.", "podcast.xml", "?format=rss")):
        return True
    if not probe:
        return False
    try:
        import urllib.request as _url
        req = _url.Request(url, headers={"User-Agent": "Mozilla/5.0 (compatible; J3PAdvisor/1.0)"})
        with _url.urlopen(req, timeout=12) as resp:
            ctype = (resp.headers.get("Content-Type") or "").lower()
            head = resp.read(400).decode("utf-8", "ignore").lstrip()
        return ("xml" in ctype) or head.startswith("<?xml") or "<rss" in head[:200]
    except Exception:
        return False


def fetch_url_metadata(url: str):
    """Title and description from a page's metadata.

    Pages built in JavaScript — podcast players, most social sites — serve no
    article text, so the main extractor finds nothing. Their metadata usually
    does carry a real description, which is worth having.
    """
    import urllib.request as _url
    import html as _html

    req = _url.Request(url, headers={
        "User-Agent": "Mozilla/5.0 (compatible; J3PAdvisor/1.0)",
        "Accept": "text/html,application/xhtml+xml",
    })
    with _url.urlopen(req, timeout=20) as resp:
        raw = resp.read(600000).decode("utf-8", "ignore")

    def meta(*names):
        for name in names:
            for pattern in (
                rf'<meta[^>]+property=["\']{name}["\'][^>]+content=["\']([^"\']+)',
                rf'<meta[^>]+content=["\']([^"\']+)["\'][^>]+property=["\']{name}["\']',
                rf'<meta[^>]+name=["\']{name}["\'][^>]+content=["\']([^"\']+)',
                rf'<meta[^>]+content=["\']([^"\']+)["\'][^>]+name=["\']{name}["\']',
            ):
                m = re.search(pattern, raw, re.IGNORECASE)
                if m:
                    return _html.unescape(m.group(1)).strip()
        return ""

    title = meta("og:title", "twitter:title") or ""
    if not title:
        m = re.search(r"<title[^>]*>(.*?)</title>", raw, re.IGNORECASE | re.DOTALL)
        title = _html.unescape(re.sub(r"\s+", " ", m.group(1))).strip() if m else ""

    desc = meta("og:description", "twitter:description", "description") or ""

    # JSON-LD often carries a much fuller description on media pages
    for m in re.finditer(
            r'<script[^>]+type=["\']application/ld\+json["\'][^>]*>(.*?)</script>',
            raw, re.IGNORECASE | re.DOTALL):
        try:
            import json as _j
            data = _j.loads(m.group(1))
            items = data if isinstance(data, list) else [data]
            for item in items:
                if not isinstance(item, dict):
                    continue
                for key in ("description", "abstract", "articleBody"):
                    val = item.get(key)
                    if isinstance(val, str) and len(val) > len(desc):
                        desc = val.strip()
        except Exception:
            continue

    return title, desc


@app.route("/admin/upload-url", methods=["POST"])
@admin_required
def admin_upload_url():
    if not (db.is_enabled() and emb.is_enabled()):
        flash("Cannot ingest URL: RAG not fully configured.")
        return redirect(url_for("admin_dashboard"))

    url = (request.form.get("url") or "").strip()
    custom_title = (request.form.get("url_title") or "").strip()

    if not url:
        flash("No URL provided.")
        return redirect(url_for("admin_dashboard"))

    # Duplicate check on the URL itself BEFORE fetching/embedding.
    # We check the URL as source. Title check happens later if custom_title is set.
    duplicate = db.find_duplicate_document(source=url)
    if duplicate:
        flash(
            f"⚠ Duplicate URL — this link was already ingested as "
            f"'{duplicate['title']}' on {duplicate['uploaded_at'].strftime('%Y-%m-%d %H:%M')} "
            f"({duplicate['chunk_count']} chunks). Delete the existing entry first "
            f"if you want to re-ingest."
        )
        return redirect(url_for("admin_dashboard"))

    # A podcast or blog feed carries the real text; handle it first.
    feed_owner = (request.form.get("owner") or "").strip()
    if looks_like_feed(url):
        try:
            show, episodes = fetch_podcast_feed(url)
        except Exception as e:
            app.logger.error(f"[feed] parse failed: {e}")
            show, episodes = None, []
        if show and episodes:
            added, skipped = 0, 0
            for ep in episodes[:60]:
                ep_title = f"{show['title']} — {ep['title']}"[:200]
                if db.find_duplicate_document(title=ep_title):
                    skipped += 1
                    continue
                try:
                    body = f"{ep['title']}\n{ep.get('when','')}\n\n{ep['body']}"
                    chunks = emb.chunk_text(body)
                    if not chunks:
                        continue
                    vectors = emb.embed_batch(chunks)
                    db.insert_document(ep_title, url, list(zip(chunks, vectors)))
                    set_document_owner(ep_title, feed_owner)
                    added += 1
                except Exception as e:
                    app.logger.error(f"[feed] episode failed: {e}")
            note = f"✓ Ingested {added} episode{'s' if added != 1 else ''} from "
            note += f"“{show['title']}”"
            if skipped:
                note += f"; {skipped} already in the knowledge base"
            flash(note + ".")
            return redirect(url_for("admin_dashboard"))

    try:
        extracted_title, text = emb.fetch_url_content(url)
        title = custom_title or extracted_title or url

        # Second duplicate check on the resolved title (in case a URL changed but
        # the article title is the same as something already in the KB).
        title_dup = db.find_duplicate_document(title=title)
        if title_dup:
            flash(
                f"⚠ Duplicate title — '{title}' was already ingested on "
                f"{title_dup['uploaded_at'].strftime('%Y-%m-%d %H:%M')} "
                f"({title_dup['chunk_count']} chunks). Use a different title or "
                f"delete the existing entry first."
            )
            return redirect(url_for("admin_dashboard"))

        if not text.strip():
            # The page is probably rendered in the browser rather than served
            # as HTML. Its metadata may still hold a usable description.
            meta_title, meta_desc = "", ""
            try:
                meta_title, meta_desc = fetch_url_metadata(url)
            except Exception as e:
                app.logger.error(f"[url] metadata fallback failed: {e}")

            if len(meta_desc) >= 120:
                text = f"{meta_title}\n\n{meta_desc}" if meta_title else meta_desc
                title = custom_title or meta_title or title
                flash("Note: this page had no article text, so only its summary "
                      "was captured. For a podcast or video, paste the "
                      "transcript or show notes below for the full content.")
            else:
                flash("Couldn't get any text from that page. Pages built in "
                      "JavaScript — podcast players, most social sites — don't "
                      "serve readable text. Use \u201cAdd knowledge from text\u201d "
                      "below and paste the transcript or show notes instead.")
                return redirect(url_for("admin_dashboard"))

        chunks = emb.chunk_text(text)
        if not chunks:
            flash("URL produced no chunks (page too short or empty).")
            return redirect(url_for("admin_dashboard"))

        vectors = emb.embed_batch(chunks)
        pairs = list(zip(chunks, vectors))
        doc_id = db.insert_document(title, url, pairs)
        set_document_owner(title, (request.form.get("owner") or "").strip())

        flash(f"✓ Ingested '{title}' from URL — {len(chunks)} chunks embedded (doc #{doc_id}).")
    except Exception as e:
        app.logger.error(f"URL ingest failed: {e}")
        flash(f"URL ingest failed: {str(e)[:200]}")

    return redirect(url_for("admin_dashboard"))


@app.route("/admin/delete/<int:doc_id>", methods=["POST"])
@admin_required
def admin_delete(doc_id):
    try:
        db.delete_document(doc_id)
        flash(f"Deleted document #{doc_id}.")
    except Exception as e:
        flash(f"Delete failed: {str(e)[:200]}")
    return redirect(url_for("admin_dashboard"))


def _fmt_ts(value) -> str:
    """Format a timestamp defensively — drivers don't always return datetimes."""
    if not value:
        return ""
    try:
        return value.strftime("%Y-%m-%d %H:%M:%S")
    except AttributeError:
        return str(value)[:19]


_XLSX_ILLEGAL = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f]")
XLSX_CELL_LIMIT = 32000        # Excel's hard limit is 32,767


def xlsx_safe(value):
    """Make a value safe for an Excel cell.

    openpyxl raises IllegalCharacterError on control characters, which turn up
    in text extracted from PDFs and pasted documents — one such character
    anywhere in the log made the whole export fail with a 500. Cells also have
    a hard length limit that long generated documents can exceed.
    """
    if value is None:
        return ""
    if isinstance(value, (int, float)):
        return value
    text = str(value)
    text = _XLSX_ILLEGAL.sub(" ", text)
    if len(text) > XLSX_CELL_LIMIT:
        text = text[:XLSX_CELL_LIMIT] + "… [truncated for Excel]"
    return text


def _rows_for_export():
    """Rows to export: only the selected ids when given, otherwise everything.

    Ids arrive as a comma-separated 'ids' parameter (query string or form),
    which is what the admin table's checkboxes produce.
    """
    if not db.is_enabled():
        return [], 0

    raw = (request.values.get("ids") or "").strip()
    selected = []
    for piece in raw.replace(" ", "").split(","):
        if piece.isdigit():
            selected.append(int(piece))

    rows = db.list_feedback(limit=10000)
    if not selected:
        return rows, 0

    wanted = set(selected)
    filtered = [r for r in rows if r.get("id") in wanted]
    # Preserve the order the user seen in the table (newest first from the DB)
    return filtered, len(wanted)


def _export_error(kind: str, err: Exception):
    """Readable failure instead of a bare 500 page."""
    app.logger.error(f"[export] {kind} export failed: {err}", exc_info=True)
    return (f"<h2 style='font-family:sans-serif'>{kind} export failed</h2>"
            f"<p style='font-family:sans-serif'>{type(err).__name__}: {err}</p>"
            f"<p style='font-family:sans-serif'>The details are in the Railway "
            f"deploy logs. Try the other format, or fewer rows, meanwhile.</p>"), 500


@app.route("/admin/export/feedback.csv")
@admin_required
def admin_export_feedback():
    """Stream all feedback as a CSV download. UTF-8 with BOM so Excel renders cleanly."""
    import csv
    import io
    from flask import Response

    try:
        rows, selected_count = _rows_for_export()
    except Exception as e:
        return _export_error("CSV", e)

    output = io.StringIO()
    writer = csv.writer(output, quoting=csv.QUOTE_ALL)
    geo = locations_for([r.get("id") for r in rows])
    acks = acknowledgements_for([r.get("id") for r in rows])
    writer.writerow([
        "ID", "Timestamp", "Rating", "Release Accepted", "Location",
        "User Question", "Bot Reply", "Comment", "Persona"
    ])
    for r in rows:
        writer.writerow([
            r.get("id", ""),
            _fmt_ts(r.get("created_at")),
            r.get("rating", ""),
            acks.get(r.get("id"), ""),
            geo.get(r.get("id"), ""),
            r.get("user_message", "") or "",
            r.get("bot_reply", "") or "",
            r.get("comment", "") or "",
            r.get("persona", "") or "",
        ])

    # Prepend UTF-8 BOM so Excel decodes em-dashes, smart quotes, accented characters correctly
    csv_content = "\ufeff" + output.getvalue()
    output.close()

    from datetime import datetime
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    suffix = f"_selected_{len(rows)}" if selected_count else ""
    filename = f"j3p_feedback{suffix}_{timestamp}.csv"

    return Response(
        csv_content,
        mimetype="text/csv; charset=utf-8",
        headers={
            "Content-Disposition": f'attachment; filename="{filename}"',
            "Cache-Control": "no-store",
        },
    )


@app.route("/admin/export/feedback.xlsx")
@admin_required
def admin_export_feedback_xlsx():
    """Stream all feedback as an Excel (.xlsx) download with formatting."""
    from flask import Response
    from datetime import datetime
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.utils import get_column_letter
    except ImportError:
        return ("Excel export unavailable: openpyxl not installed. "
                "Use CSV export instead, or add 'openpyxl' to requirements.txt."), 500

    try:
        rows, selected_count = _rows_for_export()
    except Exception as e:
        return _export_error("Excel", e)

    wb = Workbook()
    ws = wb.active
    ws.title = "Feedback"

    geo = locations_for([r.get("id") for r in rows])
    acks = acknowledgements_for([r.get("id") for r in rows])
    headers = ["ID", "Timestamp", "Rating", "Release Accepted", "Location",
               "User Question", "Bot Reply", "Comment", "Persona"]
    ws.append(headers)

    # Header styling — navy background, gold text, bold
    header_font = Font(bold=True, color="D2BC8D", size=11)
    header_fill = PatternFill("solid", fgColor="27334A")
    header_align = Alignment(horizontal="left", vertical="center", wrap_text=False)
    thin_border = Border(
        bottom=Side(style="medium", color="27334A"),
    )
    for col_idx, _ in enumerate(headers, start=1):
        cell = ws.cell(row=1, column=col_idx)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = header_align
        cell.border = thin_border

    # Data rows
    wrap_align = Alignment(vertical="top", wrap_text=True)
    down_fill = PatternFill("solid", fgColor="FEEAE5")  # subtle rust tint for thumbs-down rows
    for r in rows:
        ws.append([
            r.get("id", ""),
            xlsx_safe(_fmt_ts(r.get("created_at"))),
            xlsx_safe(r.get("rating", "")),
            xlsx_safe(acks.get(r.get("id"), "")),
            xlsx_safe(geo.get(r.get("id"), "")),
            xlsx_safe(r.get("user_message", "")),
            xlsx_safe(r.get("bot_reply", "")),
            xlsx_safe(r.get("comment", "")),
            xlsx_safe(r.get("persona", "")),
        ])
        row_idx = ws.max_row
        # Highlight thumbs-down rows so they're easy to spot when reviewing
        if r.get("rating") == "down":
            for col_idx in range(1, len(headers) + 1):
                ws.cell(row=row_idx, column=col_idx).fill = down_fill
        # Wrap long text cells
        for col_idx in [6, 7, 8]:  # User Question, Bot Reply, Comment
            ws.cell(row=row_idx, column=col_idx).alignment = wrap_align

    # Column widths — sized for readable browsing in Excel
    column_widths = {
        "A": 8,    # ID
        "B": 20,   # Timestamp
        "C": 8,    # Rating
        "D": 18,   # Release Accepted
        "E": 24,   # Location
        "F": 50,   # User Question
        "G": 80,   # Bot Reply
        "H": 40,   # Comment
        "I": 18,   # Persona
    }
    for col_letter, width in column_widths.items():
        ws.column_dimensions[col_letter].width = width

    # Freeze the header row so it stays visible while scrolling
    ws.freeze_panes = "A2"

    # AutoFilter on the header so user can sort/filter in Excel
    ws.auto_filter.ref = f"A1:{get_column_letter(len(headers))}{ws.max_row}"

    # Stream to a BytesIO
    import io
    buffer = io.BytesIO()
    try:
        wb.save(buffer)
    except Exception as e:
        return _export_error("Excel", e)
    buffer.seek(0)

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    suffix = f"_selected_{len(rows)}" if selected_count else ""
    filename = f"j3p_feedback{suffix}_{timestamp}.xlsx"

    return Response(
        buffer.read(),
        mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={
            "Content-Disposition": f'attachment; filename="{filename}"',
            "Cache-Control": "no-store",
        },
    )


@app.route("/admin/feedback/delete-selected", methods=["POST"])
@admin_required
def admin_delete_selected_feedback():
    """Delete one or more feedback rows by ID (checkboxes from the table)."""
    ids = request.form.getlist("feedback_ids")
    if not ids:
        flash("No feedback rows selected.")
        return redirect(url_for("admin_dashboard"))
    try:
        count = db.delete_feedback_ids(ids)
        flash(f"Deleted {count} feedback row{'s' if count != 1 else ''}.")
    except Exception as e:
        app.logger.error(f"Delete selected feedback failed: {e}")
        flash(f"Delete failed: {str(e)[:200]}")
    return redirect(url_for("admin_dashboard"))


@app.route("/admin/feedback/delete-all", methods=["POST"])
@admin_required
def admin_delete_all_feedback():
    """Wipe ALL feedback. Form must include confirm='YES' to prevent accidents."""
    confirm = (request.form.get("confirm") or "").strip()
    if confirm != "YES":
        flash("Clear-all cancelled — confirmation text did not match.")
        return redirect(url_for("admin_dashboard"))
    try:
        count = db.delete_all_feedback()
        flash(f"Cleared all feedback ({count} rows).")
    except Exception as e:
        app.logger.error(f"Delete all feedback failed: {e}")
        flash(f"Clear failed: {str(e)[:200]}")
    return redirect(url_for("admin_dashboard"))


@app.route("/admin/feedback/<int:feedback_id>/approve-lesson", methods=["POST"])
@admin_required
def admin_approve_lesson(feedback_id):
    """Approve a thumbs-down feedback row as a learning example.

    Embeds the user's original question and stores it on the feedback row.
    From this point forward, the bot will see this lesson when answering
    semantically similar questions.
    """
    if not (db.is_enabled() and emb.is_enabled()):
        flash("Cannot approve lesson — database or embeddings not configured.")
        return redirect(url_for("admin_dashboard"))

    row = db.get_feedback(feedback_id)
    if not row:
        flash("Feedback row not found.")
        return redirect(url_for("admin_dashboard"))
    if row.get("rating") != "down":
        flash("Only thumbs-down feedback can be approved as a lesson.")
        return redirect(url_for("admin_dashboard"))
    if not (row.get("comment") or "").strip():
        flash("This feedback has no comment — nothing to learn from. Add a comment first.")
        return redirect(url_for("admin_dashboard"))

    try:
        question_embedding = emb.embed_text(row["user_message"] or "")
        ok = db.approve_feedback_as_lesson(feedback_id, question_embedding)
        if ok:
            flash(f"✓ Lesson approved — the bot will now learn from feedback #{feedback_id}.")
        else:
            flash(f"Could not approve feedback #{feedback_id} (not eligible).")
    except Exception as e:
        app.logger.error(f"Approve lesson failed: {e}")
        flash(f"Approve failed: {str(e)[:200]}")
    return redirect(url_for("admin_dashboard"))


@app.route("/admin/feedback/<int:feedback_id>/revoke-lesson", methods=["POST"])
@admin_required
def admin_revoke_lesson(feedback_id):
    """Stop using this feedback as a lesson going forward."""
    try:
        ok = db.revoke_feedback_lesson(feedback_id)
        if ok:
            flash(f"Lesson revoked — feedback #{feedback_id} no longer informs the bot.")
        else:
            flash(f"Feedback #{feedback_id} not found.")
    except Exception as e:
        flash(f"Revoke failed: {str(e)[:200]}")
    return redirect(url_for("admin_dashboard"))


@app.route("/webhook/email", methods=["POST"])
def email_webhook():
    """
    Postmark inbound email webhook.

    Configure in Postmark: Servers → Inbound → Server settings → set webhook URL to
    https://web-production-901d85.up.railway.app/webhook/email

    Security model:
      - Subject line MUST contain the secret keyword from EMAIL_INGEST_KEYWORD env var
        (case-insensitive substring match). Without that, the email is rejected.
      - This protects against anyone who guesses the inbound email address.

    What gets ingested:
      - Email body (text version) becomes one document
      - Each PDF / DOCX / TXT / MD attachment becomes its own separate document
      - Each document is titled with the email subject + filename for attachments

    Postmark returns 200 to ANY response. We always return 200 even on rejected emails
    so Postmark doesn't retry — but we log the rejection reason.
    """
    if not (db.is_enabled() and emb.is_enabled()):
        app.logger.error("Email webhook hit but DB or embeddings not configured")
        return jsonify({"status": "rejected", "reason": "service-not-configured"}), 200

    keyword = (os.environ.get("EMAIL_INGEST_KEYWORD") or "").strip().lower()
    if not keyword:
        app.logger.error("Email webhook hit but EMAIL_INGEST_KEYWORD not set — rejecting all")
        return jsonify({"status": "rejected", "reason": "keyword-not-configured"}), 200

    payload = request.get_json(silent=True) or {}
    subject = (payload.get("Subject") or "").strip()
    from_email = (payload.get("FromFull") or {}).get("Email") or payload.get("From") or "unknown"

    # Security gate: subject MUST contain the secret keyword
    if keyword not in subject.lower():
        app.logger.warning(
            f"Email webhook rejected — subject missing keyword. From: {from_email}, Subject: {subject!r}"
        )
        return jsonify({"status": "rejected", "reason": "missing-keyword"}), 200

    # Strip the keyword from subject so the doc title is cleaner.
    # Removes "[KEYWORD]" or "KEYWORD:" or just "KEYWORD" patterns.
    import re
    clean_subject = re.sub(
        rf"\[?\b{re.escape(keyword)}\b\]?[\s:]*",
        "",
        subject,
        flags=re.IGNORECASE,
    ).strip()
    if not clean_subject:
        clean_subject = "Email submission"

    text_body = (payload.get("TextBody") or "").strip()
    html_body = (payload.get("HtmlBody") or "").strip()
    attachments = payload.get("Attachments") or []

    ingested = []
    errors = []

    # ---- Ingest the email body (if it has substance) ----
    if text_body and len(text_body) > 50:
        try:
            body_dup = db.find_duplicate_document(title=clean_subject)
            if body_dup:
                errors.append({
                    "kind": "body",
                    "error": f"duplicate-of-doc-{body_dup['id']}",
                    "skipped_title": clean_subject,
                })
                app.logger.info(
                    f"Email body skipped — duplicate of doc #{body_dup['id']} ({clean_subject!r})"
                )
            else:
                chunks = emb.chunk_text(text_body)
                if chunks:
                    vectors = emb.embed_batch(chunks)
                    pairs = list(zip(chunks, vectors))
                    doc_id = db.insert_document(
                        clean_subject,
                        f"email:{from_email}",
                        pairs,
                    )
                    ingested.append({
                        "kind": "body",
                        "title": clean_subject,
                        "doc_id": doc_id,
                        "chunks": len(chunks),
                    })
                    app.logger.info(f"Email body ingested: doc #{doc_id}, {len(chunks)} chunks")
        except Exception as e:
            app.logger.error(f"Email body ingest failed: {e}")
            errors.append({"kind": "body", "error": str(e)[:200]})

    # ---- Ingest each attachment ----
    import base64
    SUPPORTED_EXT = (".pdf", ".docx", ".txt", ".md")
    for att in attachments:
        att_name = att.get("Name") or "attachment"
        content_b64 = att.get("Content") or ""
        if not content_b64:
            continue
        if not att_name.lower().endswith(SUPPORTED_EXT):
            errors.append({
                "kind": "attachment",
                "name": att_name,
                "error": "unsupported-extension",
            })
            continue
        # Title uses subject + filename so multiple attachments are distinguishable
        doc_title = f"{clean_subject} — {att_name}" if clean_subject else att_name
        # Duplicate check before embedding work
        att_dup = db.find_duplicate_document(title=doc_title)
        if att_dup:
            errors.append({
                "kind": "attachment",
                "name": att_name,
                "error": f"duplicate-of-doc-{att_dup['id']}",
            })
            app.logger.info(
                f"Email attachment skipped — duplicate of doc #{att_dup['id']} ({doc_title!r})"
            )
            continue
        try:
            file_bytes = base64.b64decode(content_b64)
            text = emb.extract_text_from_upload(att_name, file_bytes)
            if not text.strip():
                errors.append({"kind": "attachment", "name": att_name, "error": "no-text-extracted"})
                continue
            chunks = emb.chunk_text(text)
            if not chunks:
                errors.append({"kind": "attachment", "name": att_name, "error": "no-chunks"})
                continue
            vectors = emb.embed_batch(chunks)
            pairs = list(zip(chunks, vectors))
            doc_id = db.insert_document(doc_title, f"email:{from_email}:{att_name}", pairs)
            ingested.append({
                "kind": "attachment",
                "title": doc_title,
                "doc_id": doc_id,
                "chunks": len(chunks),
            })
            app.logger.info(f"Email attachment ingested: {att_name} -> doc #{doc_id}, {len(chunks)} chunks")
        except Exception as e:
            app.logger.error(f"Email attachment ingest failed for {att_name}: {e}")
            errors.append({"kind": "attachment", "name": att_name, "error": str(e)[:200]})

    app.logger.info(
        f"Email webhook complete. From: {from_email}, Subject: {clean_subject}, "
        f"Ingested: {len(ingested)}, Errors: {len(errors)}"
    )
    return jsonify({
        "status": "ok",
        "from": from_email,
        "subject": clean_subject,
        "ingested": ingested,
        "errors": errors,
    }), 200


start_learning_scheduler()

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 8080))
    app.run(host="0.0.0.0", port=port)
